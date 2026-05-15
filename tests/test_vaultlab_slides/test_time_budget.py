"""Tests for :mod:`vaultlab.slides.time_budget` — time-budget audit.

Each test builds a tiny .pptx via python-pptx (kept independent of the
lab template so the test runs even when the template isn't bundled) and
asserts that :func:`vaultlab.slides.time_budget.audit_time_budget` returns
the expected per-slide estimate kind and aggregate budget verdict.
"""

from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL")
pptx_mod = pytest.importorskip("pptx")

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from vaultlab.slides.time_budget import (
    SlideTimeEstimate,
    TimeBudgetReport,
    audit_time_budget,
)


# ---------------------------------------------------------------------------
# Helpers — minimal pptx fixtures (no template required)
# ---------------------------------------------------------------------------


def _make_image(path: Path, color: str = "red") -> Path:
    Image.new("RGB", (200, 150), color).save(str(path))
    return path


def _add_title(slide, text: str, *, pt: int = 32) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.name = "Roboto"


def _add_bullets(slide, bullets: list[str], *, pt: int = 24) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(4))
    tf = box.text_frame
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(pt)
        r.font.name = "Roboto"


def _add_picture(slide, img: Path) -> None:
    slide.shapes.add_picture(str(img), Inches(3), Inches(2), Inches(7), Inches(4))


@pytest.fixture
def small_deck(tmp_path: Path) -> Path:
    """A 4-slide deck: title + bullets + figure+bullets + references."""
    fig = _make_image(tmp_path / "fig.png", "blue")
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) Title slide — short title, no body
    s1 = prs.slides.add_slide(blank)
    _add_title(s1, "Solid Deck Title")

    # 2) Bullets-only — a typical text slide with 3 bullets
    s2 = prs.slides.add_slide(blank)
    _add_title(s2, "Three findings from the spatial analysis")
    _add_bullets(s2, ["Finding one in detail", "Finding two in detail", "Finding three"])

    # 3) Figure + bullets
    s3 = prs.slides.add_slide(blank)
    _add_title(s3, "Spatial transcriptomics maps cell types in tissue")
    _add_picture(s3, fig)
    _add_bullets(s3, ["Captures spatial context", "Resolves cell neighborhoods"])

    # 4) References
    s4 = prs.slides.add_slide(blank)
    _add_title(s4, "References")
    _add_bullets(s4, ["Pentimalli et al., 2025", "Doe et al., 2024"])

    out = tmp_path / "small_deck.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def divider_only_deck(tmp_path: Path) -> Path:
    """A 1-slide deck that's title-only (no body) — should classify as title/divider."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Section Divider")
    out = tmp_path / "divider.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def figure_heavy_deck(tmp_path: Path) -> Path:
    """A deck with 8 figure-heavy slides, each with a picture + bullets."""
    fig = _make_image(tmp_path / "fig.png", "green")
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # title
    s0 = prs.slides.add_slide(blank)
    _add_title(s0, "Heavy Deck Title")

    for i in range(8):
        s = prs.slides.add_slide(blank)
        _add_title(s, f"Detailed analysis of dataset {i + 1}")
        _add_picture(s, fig)
        _add_bullets(s, ["bullet one", "bullet two", "bullet three"])

    out = tmp_path / "heavy.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestAuditTimeBudgetBasics:
    def test_returns_report_type(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=10)
        assert isinstance(report, TimeBudgetReport)
        assert report.deck_path == small_deck
        assert report.budget_minutes == 10
        assert report.qa_reserve_minutes == 5

    def test_per_slide_one_estimate_per_slide(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=10)
        assert len(report.per_slide) == 4
        assert all(isinstance(s, SlideTimeEstimate) for s in report.per_slide)
        assert [s.slide_index for s in report.per_slide] == [0, 1, 2, 3]

    def test_each_estimate_has_kind_and_rationale(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=10)
        for est in report.per_slide:
            assert est.kind, f"Slide {est.slide_index} missing kind"
            assert est.estimate_seconds > 0
            assert est.rationale, f"Slide {est.slide_index} missing rationale"


class TestTimeBudgetVerdict:
    def test_generous_budget_is_under(self, small_deck: Path) -> None:
        """Budget = 10 minutes (5 effective after Q&A reserve) > 4 small slides."""
        report = audit_time_budget(small_deck, budget_minutes=10)
        assert report.over_budget() is False
        assert report.estimated_total_minutes > 0

    def test_tight_budget_flags_over(self, small_deck: Path) -> None:
        """Budget = 2 minutes (-5 reserve = -3 effective) → instantly over."""
        report = audit_time_budget(small_deck, budget_minutes=2)
        assert report.over_budget() is True

    def test_figure_heavy_blows_short_budget(self, figure_heavy_deck: Path) -> None:
        """9 slides, mostly figure-heavy, should bust a 5-minute budget."""
        report = audit_time_budget(figure_heavy_deck, budget_minutes=5)
        assert report.over_budget() is True


class TestPerSlideKinds:
    def test_title_slide_short_estimate(self, small_deck: Path) -> None:
        """Title slide should be the shortest (~15s)."""
        report = audit_time_budget(small_deck, budget_minutes=10)
        title_est = report.per_slide[0]
        assert title_est.kind == "title"
        assert title_est.estimate_seconds <= 30

    def test_figure_slide_kind(self, small_deck: Path) -> None:
        """Figure + bullets slide should be classified as 'figure'."""
        report = audit_time_budget(small_deck, budget_minutes=10)
        # Slide index 2 in the small_deck fixture is the figure slide
        fig_est = report.per_slide[2]
        assert fig_est.kind == "figure"
        # Figure slides should take longer than title slides
        assert fig_est.estimate_seconds >= report.per_slide[0].estimate_seconds

    def test_bullets_slide_kind(self, small_deck: Path) -> None:
        """Bullets-only slide should be classified as 'bullets'."""
        report = audit_time_budget(small_deck, budget_minutes=10)
        bullets_est = report.per_slide[1]
        assert bullets_est.kind == "bullets"
        # Bullets slides are mid-range
        assert 20 <= bullets_est.estimate_seconds <= 120

    def test_references_slide_kind(self, small_deck: Path) -> None:
        """A slide titled 'References' should be classified as 'references'."""
        report = audit_time_budget(small_deck, budget_minutes=10)
        ref_est = report.per_slide[3]
        assert ref_est.kind == "references"


class TestBudgetMath:
    def test_budget_seconds_reflects_reserve(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=10, qa_reserve_minutes=5)
        assert report.budget_seconds == 5 * 60

    def test_custom_qa_reserve(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=15, qa_reserve_minutes=10)
        assert report.budget_seconds == 5 * 60
        assert report.qa_reserve_minutes == 10

    def test_estimated_total_matches_per_slide_sum(self, small_deck: Path) -> None:
        report = audit_time_budget(small_deck, budget_minutes=10)
        assert report.estimated_total_seconds == sum(
            s.estimate_seconds for s in report.per_slide
        )


class TestErrors:
    def test_missing_pptx_raises(self, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            audit_time_budget(tmp_path / "nope.pptx", budget_minutes=10)
