"""Tests for the lifted slide-layout primitives.

Ported from ``bobby-tools/tests/test_bobby_slides/test_layout.py``. Adapts
imports to the vaultlab namespace.
"""
from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL", reason="Pillow required for layout tests")
pptx = pytest.importorskip("pptx", reason="python-pptx required for layout tests")

from PIL import Image
from pptx import Presentation as PptxPresentation

from vaultlab.slides.layouts import (
    add_figure_slide,
    add_multi_figure_slide,
    add_references_slide,
    add_section_divider,
    add_text_slide,
    add_title_slide,
)
from vaultlab.slides.template import default_font, lab_template_path, load_template


pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — skipping lifted-layout tests",
)


def make_test_image(path: Path, color: str = "red", size: tuple = (200, 150)):
    img = Image.new("RGB", size, color)
    img.save(str(path))
    return path


def _save_and_reload(pres, tmp_path: Path, name: str):
    """Save to tmp_path and reload — verifies the .pptx is valid."""
    out = tmp_path / name
    pres.save(str(out))
    return PptxPresentation(str(out))


def _all_runs_use_roboto(pres) -> bool:
    """Check that every text run in the presentation uses Roboto."""
    found_any = False
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            found_any = True
                            if run.font.name != default_font():
                                return False
    return found_any


class TestAddTitleSlide:
    def test_basic(self, tmp_path):
        pres = load_template()
        slide = add_title_slide(pres, "Test Title")
        assert slide is not None
        reloaded = _save_and_reload(pres, tmp_path, "title.pptx")
        assert len(reloaded.slides) > 0

    def test_with_subtitle_and_author(self, tmp_path):
        pres = load_template()
        add_title_slide(pres, "Main", subtitle="Sub", author="Bobby Ni")
        reloaded = _save_and_reload(pres, tmp_path, "title2.pptx")
        assert len(reloaded.slides) > 0
        assert _all_runs_use_roboto(reloaded)


class TestAddFigureSlide:
    def test_with_image(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        add_figure_slide(pres, fig, title="Figure 1", caption="A test figure")
        reloaded = _save_and_reload(pres, tmp_path, "fig.pptx")
        assert len(reloaded.slides) > 0

    def test_missing_image_does_not_crash(self, tmp_path):
        pres = load_template()
        add_figure_slide(pres, "/nonexistent/x.png", title="Missing")
        reloaded = _save_and_reload(pres, tmp_path, "missing.pptx")
        assert len(reloaded.slides) > 0

    def test_with_bullets(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "b.png"
        make_test_image(fig)
        add_figure_slide(
            pres, fig,
            title="With Bullets",
            bullets=["n=100", "p<0.001", "Effect size: 0.8"],
        )
        reloaded = _save_and_reload(pres, tmp_path, "bullets.pptx")
        assert len(reloaded.slides) > 0

    def test_with_citation(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "c.png"
        make_test_image(fig)
        add_figure_slide(
            pres, fig,
            title="Cited",
            citation_source="Smith et al., 2024, Nature",
        )
        reloaded = _save_and_reload(pres, tmp_path, "cite.pptx")
        all_text = []
        for s in reloaded.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        assert any("Smith et al" in t for t in all_text)


class TestAddMultiFigureSlide:
    def test_two_figures(self, tmp_path):
        pres = load_template()
        figs = []
        for i, color in enumerate(["red", "blue"]):
            p = tmp_path / f"m{i}.png"
            make_test_image(p, color)
            figs.append({"path": str(p), "label": chr(ord("A") + i), "caption": f"Panel {i}"})
        add_multi_figure_slide(pres, figs, title="Two Figs")
        reloaded = _save_and_reload(pres, tmp_path, "two.pptx")
        assert len(reloaded.slides) > 0

    def test_four_figures_grid(self, tmp_path):
        pres = load_template()
        figs = []
        for i, color in enumerate(["red", "blue", "green", "yellow"]):
            p = tmp_path / f"g{i}.png"
            make_test_image(p, color)
            figs.append({"path": str(p), "label": chr(ord("A") + i)})
        add_multi_figure_slide(pres, figs, title="2x2 grid")
        reloaded = _save_and_reload(pres, tmp_path, "grid.pptx")
        assert len(reloaded.slides) > 0

    def test_zero_figures_does_not_crash(self, tmp_path):
        pres = load_template()
        add_multi_figure_slide(pres, [], title="Empty")
        reloaded = _save_and_reload(pres, tmp_path, "empty_multi.pptx")
        assert len(reloaded.slides) > 0


class TestAddTextSlide:
    def test_with_bullets(self, tmp_path):
        pres = load_template()
        add_text_slide(pres, "Bullets", ["First", "Second", "Third"])
        reloaded = _save_and_reload(pres, tmp_path, "text.pptx")
        all_text = []
        for s in reloaded.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        assert any("First" in t for t in all_text)

    def test_no_bullets(self, tmp_path):
        pres = load_template()
        add_text_slide(pres, "Just title", [])
        reloaded = _save_and_reload(pres, tmp_path, "title_only.pptx")
        assert len(reloaded.slides) > 0


class TestAddSectionDivider:
    def test_renders(self, tmp_path):
        pres = load_template()
        add_section_divider(pres, "Methods")
        reloaded = _save_and_reload(pres, tmp_path, "div.pptx")
        all_text = []
        for s in reloaded.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        assert any("Methods" in t for t in all_text)


class TestAddReferencesSlide:
    def test_with_refs(self, tmp_path):
        pres = load_template()
        add_references_slide(pres, [
            "Smith et al., 2024, Nature",
            "Jones et al., 2023, Cell",
        ])
        reloaded = _save_and_reload(pres, tmp_path, "refs.pptx")
        all_text = []
        for s in reloaded.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        assert any("Smith" in t for t in all_text)
        assert any("Jones" in t for t in all_text)

    def test_empty_refs(self, tmp_path):
        pres = load_template()
        add_references_slide(pres, [])
        reloaded = _save_and_reload(pres, tmp_path, "no_refs.pptx")
        assert len(reloaded.slides) > 0


class TestFontConsistency:
    def test_all_layouts_use_roboto(self, tmp_path):
        """Build a deck with every layout type and verify Roboto everywhere."""
        pres = load_template()
        fig = tmp_path / "all.png"
        make_test_image(fig)

        add_title_slide(pres, "Title", subtitle="Sub", author="Bobby")
        add_section_divider(pres, "Section 1")
        add_figure_slide(pres, fig, title="Fig", caption="Cap", bullets=["a", "b"])
        add_multi_figure_slide(pres, [{"path": str(fig), "label": "A"}], title="Multi")
        add_text_slide(pres, "Text", ["bullet 1"])
        add_references_slide(pres, ["Ref 1"])

        reloaded = _save_and_reload(pres, tmp_path, "full.pptx")
        assert _all_runs_use_roboto(reloaded)
