"""Tests for the lifted speaker-notes formatter (vaultlab.slides.notes).

Ported from ``bobby-tools/tests/test_bobby_slides/test_speaker.py``.
"""
from __future__ import annotations

import pytest

from vaultlab.slides.notes import (
    attach_to_slide,
    format_speaker_notes,
    parse_speaker_notes,
)


class TestFormat:
    def test_empty_returns_empty_string(self):
        assert format_speaker_notes(None) == ""
        assert format_speaker_notes({}) == ""

    def test_full_note(self):
        notes = {
            "hook": "Lipids tell a story",
            "key_claim": "Three programs distinguish IBD",
            "evidence": "MALDI-IMS heatmap, n=47",
            "key_terms": ["IBD", "phospholipid"],
            "click": "Reveal panel B",
            "transition": "Methods next",
        }
        result = format_speaker_notes(notes)
        assert "HOOK: Lipids tell a story" in result
        assert "KEY CLAIM: Three programs distinguish IBD" in result
        assert "KEY TERMS: IBD, phospholipid" in result
        assert "TRANSITION: Methods next" in result

    def test_partial_note_only_renders_present_keys(self):
        notes = {"hook": "H", "transition": "T"}
        result = format_speaker_notes(notes)
        assert "HOOK: H" in result
        assert "TRANSITION: T" in result
        assert "KEY CLAIM" not in result
        assert "EVIDENCE" not in result

    def test_order_is_consistent(self):
        notes = {
            "transition": "T",
            "hook": "H",
            "evidence": "E",
            "key_claim": "K",
        }
        result = format_speaker_notes(notes)
        assert result.index("HOOK") < result.index("KEY CLAIM")
        assert result.index("KEY CLAIM") < result.index("EVIDENCE")
        assert result.index("EVIDENCE") < result.index("TRANSITION")

    def test_string_key_terms_works(self):
        notes = {"key_terms": "single-term"}
        result = format_speaker_notes(notes)
        assert "KEY TERMS: single-term" in result


class TestParse:
    def test_empty_returns_empty_dict(self):
        assert parse_speaker_notes("") == {}
        assert parse_speaker_notes("   ") == {}

    def test_parses_full_note(self):
        text = (
            "- HOOK: Lipids tell a story\n"
            "- KEY CLAIM: Three programs distinguish IBD\n"
            "- KEY TERMS: IBD, phospholipid, sphingolipid\n"
            "- TRANSITION: Methods next\n"
        )
        result = parse_speaker_notes(text)
        assert result["hook"] == "Lipids tell a story"
        assert result["key_claim"] == "Three programs distinguish IBD"
        assert result["key_terms"] == ["IBD", "phospholipid", "sphingolipid"]
        assert result["transition"] == "Methods next"

    def test_round_trip(self):
        original = {
            "hook": "H",
            "key_claim": "K",
            "evidence": "E",
            "key_terms": ["A", "B"],
            "click": "C",
            "transition": "T",
        }
        formatted = format_speaker_notes(original)
        parsed = parse_speaker_notes(formatted)
        assert parsed == original

    def test_ignores_non_dash_lines(self):
        text = "Some intro\n- HOOK: H\nMore prose\n"
        result = parse_speaker_notes(text)
        assert result == {"hook": "H"}


class TestAttachToSlide:
    def test_attaches_to_pptx_slide(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation as PptxPresentation

        from vaultlab.slides.layouts import add_text_slide
        from vaultlab.slides.template import lab_template_path, load_template

        if lab_template_path() is None:
            pytest.skip("Hickey lab template not bundled")

        pres = load_template()
        slide = add_text_slide(pres, "Test", ["A bullet"])
        attach_to_slide(slide, {"hook": "H", "transition": "T"})

        out = tmp_path / "with_notes.pptx"
        pres.save(str(out))

        reloaded = PptxPresentation(str(out))
        last_slide = reloaded.slides[-1]
        notes_text = last_slide.notes_slide.notes_text_frame.text
        assert "HOOK: H" in notes_text
        assert "TRANSITION: T" in notes_text

    def test_empty_notes_are_noop(self, tmp_path):
        pytest.importorskip("pptx")
        from vaultlab.slides.layouts import add_text_slide
        from vaultlab.slides.template import lab_template_path, load_template

        if lab_template_path() is None:
            pytest.skip("Hickey lab template not bundled")

        pres = load_template()
        slide = add_text_slide(pres, "Test", ["A"])
        attach_to_slide(slide, None)
        attach_to_slide(slide, {})
        # No exception = pass
