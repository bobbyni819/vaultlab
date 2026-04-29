"""End-to-end render tests — require python-pptx; marked slow.

These verify the renderer produces a valid .pptx file by re-opening the
output through python-pptx and inspecting the shape tree.
"""

from __future__ import annotations

from pathlib import Path

import pytest

pytestmark = pytest.mark.slow


@pytest.fixture
def pptx() -> object:
    """Skip the entire module when python-pptx isn't installed."""
    try:
        import pptx  # type: ignore[import-not-found]

        return pptx
    except ImportError:
        pytest.skip("python-pptx not installed")


class TestRenderRoundTrip:
    def test_render_simple_deck(self, tmp_path: Path, pptx) -> None:
        from vaultlab.slides import Deck, Slide, render_pptx

        deck = Deck(
            title="Test deck",
            slides=[
                Slide(layout="title", title="Test deck", subtitle="A test"),
                Slide(
                    layout="content_with_bullets",
                    title="Outline",
                    bullets=["Background", "Methods", "Results", "Discussion"],
                ),
            ],
        )

        out = render_pptx(deck, tmp_path / "deck.pptx")
        assert out.exists()
        assert out.stat().st_size > 0

        # Re-open and verify structure
        from pptx import Presentation

        p = Presentation(str(out))
        assert len(p.slides) == 2

    def test_speaker_notes_persisted(self, tmp_path: Path, pptx) -> None:
        from pptx import Presentation

        from vaultlab.slides import Deck, Slide, render_pptx

        deck = Deck(
            title="Test",
            slides=[
                Slide(
                    layout="content_with_bullets",
                    title="With notes",
                    bullets=["bullet"],
                    speaker_notes="The presenter says these words.",
                )
            ],
        )
        out = render_pptx(deck, tmp_path / "x.pptx")
        p = Presentation(str(out))
        notes_text = p.slides[0].notes_slide.notes_text_frame.text
        assert "presenter says these words" in notes_text

    def test_missing_figure_raises(self, tmp_path: Path, pptx) -> None:
        from vaultlab.slides import Deck, Slide, render_pptx
        from vaultlab.slides.render import RenderError

        deck = Deck(
            title="Test",
            slides=[
                Slide(
                    layout="figure_with_caption",
                    title="Fig",
                    figure_path="absolutely-does-not-exist.png",
                    caption="caption",
                )
            ],
        )
        with pytest.raises(RenderError, match="Figure not found"):
            render_pptx(deck, tmp_path / "x.pptx")
