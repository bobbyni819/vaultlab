"""Tests for vaultlab.slides.annotate — figure annotation overlays.

Ported from ``bobby-tools/tests/test_bobby_slides/test_annotations.py``.
"""
from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL")
pptx = pytest.importorskip("pptx")

from PIL import Image
from pptx import Presentation as PptxPresentation

from vaultlab.slides import (
    add_annotations,
    add_figure_slide,
    build_from_plan,
    load_template,
)
from vaultlab.slides.template import lab_template_path


pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — annotation tests need it",
)


def make_test_image(path: Path, color: str = "red"):
    Image.new("RGB", (400, 300), color).save(str(path))
    return path


def _save_and_reload(pres, tmp_path: Path, name: str):
    out = tmp_path / name
    pres.save(str(out))
    return PptxPresentation(str(out))


class TestAnnotations:
    def test_circle_annotation(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) == 1
        anns = add_annotations(slide, pictures[0], [
            {"type": "circle", "x": 0.5, "y": 0.5, "r": 0.1, "color": "FF0000"},
        ], with_animations=False)
        assert len(anns) == 1

    def test_rect_annotation(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [
            {"type": "rect", "bbox": [0.1, 0.1, 0.4, 0.4], "color": "00FF00"},
        ], with_animations=False)
        assert len(anns) == 1

    def test_arrow_annotation(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [
            {"type": "arrow", "from": [0.2, 0.2], "to": [0.6, 0.7], "color": "FFEB3B"},
        ], with_animations=False)
        assert len(anns) == 1

    def test_label_annotation(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [
            {"type": "label", "x": 0.5, "y": 0.9, "text": "Key area"},
        ], with_animations=False)
        assert len(anns) == 1

    def test_multiple_annotations(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [
            {"type": "circle", "x": 0.3, "y": 0.3, "r": 0.05},
            {"type": "rect", "bbox": [0.5, 0.5, 0.7, 0.7]},
            {"type": "arrow", "from": [0.1, 0.1], "to": [0.9, 0.9]},
        ], with_animations=False)
        assert len(anns) == 3

    def test_pptx_reloads_after_annotations(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        add_annotations(slide, pictures[0], [
            {"type": "circle", "x": 0.5, "y": 0.5, "r": 0.1},
        ], with_animations=False)
        reloaded = _save_and_reload(pres, tmp_path, "ann.pptx")
        assert len(reloaded.slides) > 0

    def test_unknown_annotation_type_skipped(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [
            {"type": "circle", "x": 0.5, "y": 0.5, "r": 0.1},
            {"type": "weird_unknown_type"},
            {"type": "rect", "bbox": [0.1, 0.1, 0.3, 0.3]},
        ], with_animations=False)
        assert len(anns) == 2  # circle + rect, weird type skipped

    def test_empty_annotations_returns_empty(self, tmp_path):
        pres = load_template()
        fig = tmp_path / "f.png"
        make_test_image(fig)
        slide = add_figure_slide(pres, fig, title="Test")
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        anns = add_annotations(slide, pictures[0], [], with_animations=False)
        assert anns == []


class TestAnnotationsViaBuildFromPlan:
    def test_annotations_in_plan(self, tmp_path):
        fig = tmp_path / "f.png"
        make_test_image(fig)
        plan = {
            "slides": [
                {"type": "figure", "title": "T", "image_path": str(fig),
                 "annotations": [
                     {"type": "circle", "x": 0.5, "y": 0.5, "r": 0.1, "color": "FF0000"},
                     {"type": "label", "x": 0.5, "y": 0.85, "text": "Important"},
                 ]},
            ],
        }
        out = tmp_path / "deck.pptx"
        result = build_from_plan(plan, out, write_marp=False)
        prs = PptxPresentation(str(result["pptx"]))
        last = prs.slides[-1]
        # Should have: title text + picture + circle + label = 4 shapes (at least)
        assert len(last.shapes) >= 4

    def test_annotations_with_click_index(self, tmp_path):
        fig = tmp_path / "f.png"
        make_test_image(fig)
        plan = {
            "slides": [
                {"type": "figure", "title": "T", "image_path": str(fig),
                 "annotations": [
                     {"type": "circle", "x": 0.3, "y": 0.3, "r": 0.05, "click_index": 0},
                     {"type": "circle", "x": 0.7, "y": 0.7, "r": 0.05, "click_index": 1},
                 ]},
            ],
        }
        out = tmp_path / "deck.pptx"
        result = build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert result["pptx"].exists()

    def test_annotations_failure_doesnt_break_deck(self, tmp_path):
        """Even if annotations spec is malformed, deck should still render."""
        fig = tmp_path / "f.png"
        make_test_image(fig)
        plan = {
            "slides": [
                {"type": "figure", "title": "T", "image_path": str(fig),
                 "annotations": [
                     {"type": "circle"},  # missing coords — uses defaults
                     {},  # totally empty — should be skipped
                 ]},
            ],
        }
        out = tmp_path / "deck.pptx"
        result = build_from_plan(plan, out, write_marp=False)
        assert result["pptx"].exists()
