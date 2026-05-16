"""Tests for :mod:`vaultlab.slides.self_review` — composite deck review pass.

Builds tiny test decks via :func:`vaultlab.slides.build_from_plan` and the
typed-Deck path, then runs :func:`vaultlab.slides.self_review.review_deck`
to assert the audits flag (or don't flag) the expected issues.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

PIL = pytest.importorskip("PIL")
pptx_mod = pytest.importorskip("pptx")

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from vaultlab.slides import build_from_plan
from vaultlab.slides.self_review import (
    ReviewReport,
    SlideReview,
    render_review_html,
    review_deck,
    write_review_report,
)
from vaultlab.slides.template import lab_template_path


pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — self_review tests need build_from_plan",
)


def _make_image(path: Path, color: str = "red") -> Path:
    Image.new("RGB", (200, 150), color).save(str(path))
    return path


# ---------------------------------------------------------------------------
# Helpers — build a known-good deck and a known-bad deck
# ---------------------------------------------------------------------------


@pytest.fixture
def good_deck(tmp_path: Path) -> Path:
    """A 5-slide deck that should pass every audit.

    Body slides (figure + text) carry a dual-format speaker-notes block so
    the speaker-notes structural audit (added 2026-05-15) passes too.
    """
    fig = _make_image(tmp_path / "good_fig.png", "green")
    # 200-400 word script — meets the speaker-notes-short / -long bounds.
    sample_script = (
        "The slide before us shows the key data point I want to land in the "
        "room. We start from the working hypothesis stated in the section "
        "divider and then walk through the three supporting observations one "
        "at a time so the audience can track the chain of reasoning. The "
        "first observation is anchored in the leftmost panel where the "
        "raw signal is rendered without any background subtraction so the "
        "viewer can judge the raw effect size; we then move rightward to "
        "the background subtracted panel and finally to the quantification "
        "panel which collapses the comparison into a single boxplot. Each "
        "of these three views answers a different question and together "
        "they give us a triangulated read on the underlying biology. We "
        "deliberately separated the views rather than overlaying them "
        "because overlay charts hide the kind of low frequency texture "
        "that often signals a methodological artefact, and we want that "
        "texture visible to anyone who asks during Q and A about whether "
        "the trend is robust or merely an artefact of the colour scale. "
        "Before we close the slide I want to flag the single most important "
        "limit on this measurement which is the donor count for the rightmost "
        "panel; we will return to that limit on the next slide where the "
        "full per donor breakdown is shown so the audience knows we are not "
        "hiding the variability that the boxplot collapses."
    )
    # ``build_from_plan`` consumes the ``speaker_notes`` key as a dict and
    # routes it through ``attach_to_slide`` which renders the dual-format
    # mental-map + script string onto the slide's notes pane.
    speaker_notes = {
        "hook": "One slide, three views of the same effect.",
        "key_claim": "The effect survives background subtraction.",
        "evidence": "Three-panel figure with quantification on the right.",
        "key_terms": ["signal-to-noise", "background subtraction"],
        "click": "Click reveals the quantification panel last.",
        "transition": "Next slide shows the per-donor breakdown.",
        "script": sample_script,
    }
    plan = {
        "title": "Solid Deck",
        "author": "Test Author",
        "slides": [
            {
                "type": "title",
                "title": "Solid Deck",
                "subtitle": "All audits should pass",
                "author": "Test Author",
            },
            {"type": "section_divider", "title": "Background"},
            {
                "type": "figure",
                "title": "Spatial transcriptomics maps cell types in tissue",
                "image_path": str(fig),
                "caption": "Composite immunofluorescence",
                "bullets": ["Captures spatial context", "Resolves cell neighborhoods"],
                "citation_source": "Pentimalli et al., 2025",
                "speaker_notes": speaker_notes,
            },
            {
                "type": "text",
                "title": "Three findings emerge from the spatial analysis",
                "bullets": ["Finding one with evidence", "Finding two with evidence"],
                "speaker_notes": speaker_notes,
            },
            {
                "type": "references",
                "title": "References",
                "references": ["Pentimalli et al., 2025"],
            },
        ],
    }
    out = tmp_path / "good_deck.pptx"
    build_from_plan(plan, out, write_marp=False)
    return out


@pytest.fixture
def tiny_title_deck(tmp_path: Path) -> Path:
    """Synthetic deck with an intentionally too-small title font (8pt).

    Built directly via python-pptx (not build_from_plan, which enforces
    the hard rules) so we can introduce the violation deterministically.
    """
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Some descriptive but tiny title"
    run.font.size = Pt(8)  # << violation — well under 28pt floor
    run.font.name = "Roboto"

    # Add a small body so the slide isn't classified as a section divider.
    bbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(3))
    btf = bbox.text_frame
    bp = btf.paragraphs[0]
    brun = bp.add_run()
    brun.text = "Body text large enough to read"
    brun.font.size = Pt(24)
    brun.font.name = "Roboto"

    out = tmp_path / "tiny_title.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def overlap_deck(tmp_path: Path) -> Path:
    """Synthetic deck with two heavily overlapping shapes (critical issue)."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Title at the top — large and Roboto so we don't trip other audits.
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Slide title that is sufficiently descriptive"
    r.font.size = Pt(32)
    r.font.name = "Roboto"

    # Two big boxes that overlap by ~5 inches.
    for left in (Inches(1), Inches(2)):
        b = slide.shapes.add_textbox(left, Inches(2), Inches(8), Inches(4))
        bf = b.text_frame
        bp = bf.paragraphs[0]
        br = bp.add_run()
        br.text = "Lots of body content here"
        br.font.size = Pt(24)
        br.font.name = "Roboto"

    out = tmp_path / "overlap.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestReviewDeckGoodCase:
    def test_known_good_deck_has_no_criticals(self, good_deck: Path) -> None:
        report = review_deck(good_deck)
        # Diagnostics surface as the failure message so a regression is loud.
        criticals = [i for i in report.all_issues() if i.get("severity") == "critical"]
        assert report.n_critical == 0, f"Expected zero criticals, got: {criticals}"
        assert report.ok() is True

    def test_review_report_dataclass_shape(self, good_deck: Path) -> None:
        report = review_deck(good_deck)
        assert isinstance(report, ReviewReport)
        assert report.pptx_path == good_deck
        assert report.n_slides == 5
        assert all(isinstance(s, SlideReview) for s in report.per_slide)

    def test_summary_lines_non_empty(self, good_deck: Path) -> None:
        report = review_deck(good_deck)
        lines = report.summary_lines()
        assert lines, "summary_lines should return at least one line"
        assert any("5" in line for line in lines), "Summary should mention slide count"


class TestReviewDeckCriticalCases:
    def test_tiny_title_font_is_critical(self, tiny_title_deck: Path) -> None:
        report = review_deck(tiny_title_deck)
        rules = {i["rule"] for i in report.all_issues() if i.get("severity") == "critical"}
        assert "min-title-font" in rules, (
            f"Expected min-title-font critical issue. Got: {list(report.all_issues())}"
        )
        assert report.ok() is False
        assert report.n_critical >= 1

    def test_shape_overlap_is_critical(self, overlap_deck: Path) -> None:
        report = review_deck(overlap_deck)
        rules = {i["rule"] for i in report.all_issues() if i.get("severity") == "critical"}
        assert "no-shape-overlap" in rules, (
            f"Expected overlap critical issue. Got: {list(report.all_issues())}"
        )
        assert report.ok() is False


class TestReviewDeckWarningCases:
    def test_short_title_flags_warning(self, tmp_path: Path) -> None:
        # A normal-looking deck except one slide has a one-word title.
        fig = _make_image(tmp_path / "fig.png")
        plan = {
            "title": "Deck",
            "author": "Test",
            "slides": [
                {"type": "title", "title": "Deck", "subtitle": "Subtitle", "author": "T"},
                {"type": "section_divider", "title": "Section"},
                {
                    "type": "text",
                    "title": "Findings",  # one-word → warning
                    "bullets": ["Bullet one", "Bullet two"],
                },
                {
                    "type": "figure",
                    "title": "Bigger descriptive title for figure slide",
                    "image_path": str(fig),
                    "caption": "Cap",
                    "bullets": ["Bullet"],
                },
            ],
        }
        out = tmp_path / "short_title.pptx"
        build_from_plan(plan, out, write_marp=False)
        report = review_deck(out)

        slide2 = report.per_slide[2]
        rules = {i["rule"] for i in slide2.issues}
        assert "descriptive-title" in rules, (
            f"Expected descriptive-title warning on slide 2. Issues: {slide2.issues}"
        )


class TestReviewDeckStoryArc:
    def test_empty_deck_flags_story_arc(self, tmp_path: Path) -> None:
        prs = PptxPresentation()
        out = tmp_path / "empty.pptx"
        prs.save(str(out))
        report = review_deck(out)
        rules = {i["rule"] for i in report.story_arc_issues}
        assert "story-arc-empty" in rules
        assert report.ok() is False


class TestReviewHtml:
    def test_render_review_html_is_non_empty(self, good_deck: Path) -> None:
        report = review_deck(good_deck)
        html = render_review_html(report)
        assert isinstance(html, str)
        assert len(html) > 500, "HTML report should not be a stub"
        # Should reference the deck name + the audit framing.
        assert good_deck.stem in html

    def test_render_review_html_contains_issue_text(self, tiny_title_deck: Path) -> None:
        report = review_deck(tiny_title_deck)
        html = render_review_html(report)
        # The audit_html builder uppercases the issue kind for display.
        assert "MIN-TITLE-FONT" in html, "Issue rule name should appear in the HTML"
        # The detail text should surface in the per-slide card.
        assert "28pt minimum" in html or "28pt" in html

    def test_write_review_report_writes_html_and_sidecars(
        self, good_deck: Path, tmp_path: Path
    ) -> None:
        out = tmp_path / "review.html"
        report = review_deck(good_deck)
        written = write_review_report(report, out)
        assert written.exists()
        assert written.stat().st_size > 500
        # Red Line #2 — sidecar pair next to the HTML.
        assert (out.parent / "review.html.provenance.json").exists()
        assert (out.parent / "review.html.method.md").exists()


class TestReviewDeckErrors:
    def test_missing_pptx_raises(self, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            review_deck(tmp_path / "nope.pptx")


# ---------------------------------------------------------------------------
# WCAG color-contrast check (deferred-followups bundle, 2026-05-15)
# ---------------------------------------------------------------------------


class TestColorContrastCheck:
    """Per-text-shape WCAG contrast ratio.

    AA threshold = 4.5 for normal text; <3.0 = critical, 3.0-4.5 = warning.
    Theme/scheme colors must NOT trip the check (no false positives).
    """

    def _build_contrast_deck(
        self,
        tmp_path: Path,
        *,
        fg: tuple[int, int, int],
        bg: tuple[int, int, int] | None = None,
        name: str = "contrast.pptx",
    ) -> Path:
        """Build a single-slide pptx with a known fg/bg color pair."""
        from pptx.dml.color import RGBColor

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Title at top (large + Roboto, satisfies other audits).
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tp = tbox.text_frame.paragraphs[0]
        tr = tp.add_run()
        tr.text = "Slide title that is sufficiently descriptive"
        tr.font.size = Pt(32)
        tr.font.name = "Roboto"

        # Body shape — apply explicit fill if bg given; else leave default.
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        if bg is not None:
            body.fill.solid()
            body.fill.fore_color.rgb = RGBColor(*bg)
        bp = body.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = "Body text that should be legible at projector distance"
        br.font.size = Pt(24)
        br.font.name = "Roboto"
        br.font.color.rgb = RGBColor(*fg)

        out = tmp_path / name
        prs.save(str(out))
        return out

    def test_low_contrast_text_is_critical(self, tmp_path: Path) -> None:
        # Light-gray text on white background — ratio ≈ 1.6 (well below 3.0).
        deck = self._build_contrast_deck(
            tmp_path, fg=(0xDD, 0xDD, 0xDD), bg=(0xFF, 0xFF, 0xFF)
        )
        report = review_deck(deck)
        rules = {(i["severity"], i["rule"]) for i in report.all_issues()}
        assert ("critical", "color-contrast") in rules, (
            f"Expected color-contrast critical. Got: {list(report.all_issues())}"
        )

    def test_borderline_contrast_is_warning(self, tmp_path: Path) -> None:
        # Mid-gray text on white — ratio in the 3.0-4.5 band (warning, not critical).
        # #888888 on white has ratio ≈ 3.54.
        deck = self._build_contrast_deck(
            tmp_path, fg=(0x88, 0x88, 0x88), bg=(0xFF, 0xFF, 0xFF)
        )
        report = review_deck(deck)
        contrast_issues = [
            i for i in report.all_issues() if i.get("rule") == "color-contrast"
        ]
        assert any(i.get("severity") == "warning" for i in contrast_issues), (
            f"Expected color-contrast warning at borderline ratio. Got: {contrast_issues}"
        )
        # And not critical, since we're above 3.0.
        assert not any(i.get("severity") == "critical" for i in contrast_issues)

    def test_high_contrast_is_silent(self, tmp_path: Path) -> None:
        # Black on white — ratio = 21:1 (max). No contrast issue should fire.
        deck = self._build_contrast_deck(
            tmp_path, fg=(0x00, 0x00, 0x00), bg=(0xFF, 0xFF, 0xFF)
        )
        report = review_deck(deck)
        contrast_issues = [
            i for i in report.all_issues() if i.get("rule") == "color-contrast"
        ]
        assert not contrast_issues, (
            f"High-contrast deck should not raise contrast issues. Got: {contrast_issues}"
        )

    def test_unset_run_color_does_not_flag(self, tmp_path: Path) -> None:
        """No fg color set → theme inherited → must skip silently (no false positive)."""
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tp = tbox.text_frame.paragraphs[0]
        tr = tp.add_run()
        tr.text = "Slide title that is sufficiently descriptive"
        tr.font.size = Pt(32)
        tr.font.name = "Roboto"

        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        bp = body.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = "Body text with no explicit color"
        br.font.size = Pt(24)
        br.font.name = "Roboto"
        # NOTE: deliberately not touching br.font.color

        out = tmp_path / "theme.pptx"
        prs.save(str(out))
        report = review_deck(out)
        contrast_issues = [
            i for i in report.all_issues() if i.get("rule") == "color-contrast"
        ]
        assert not contrast_issues, (
            f"Themed colors should not be flagged. Got: {contrast_issues}"
        )

    def test_contrast_helpers_match_wcag(self) -> None:
        from vaultlab.slides.self_review import _contrast_ratio

        # Black on white is exactly 21:1 (the WCAG-defined max).
        assert abs(_contrast_ratio((0, 0, 0), (255, 255, 255)) - 21.0) < 0.01
        # Same color → ratio = 1.0.
        assert abs(_contrast_ratio((128, 128, 128), (128, 128, 128)) - 1.0) < 1e-9
        # Symmetry: order of args must not matter.
        a = _contrast_ratio((0x11, 0x55, 0x99), (0xEE, 0xEE, 0xEE))
        b = _contrast_ratio((0xEE, 0xEE, 0xEE), (0x11, 0x55, 0x99))
        assert abs(a - b) < 1e-9


# ---------------------------------------------------------------------------
# Phase 7.3 close-out — time-budget + Q&A anticipator wiring
# ---------------------------------------------------------------------------


class TestReviewDeckAddOns:
    """``review_deck`` should expose optional time-budget + Q&A audits."""

    def test_time_budget_not_run_by_default(self, good_deck: Path) -> None:
        report = review_deck(good_deck)
        assert report.time_budget is None
        assert report.anticipated_questions == []

    def test_time_budget_runs_when_requested(self, good_deck: Path) -> None:
        report = review_deck(good_deck, budget_minutes=10)
        assert report.time_budget is not None
        assert report.time_budget.budget_minutes == 10
        assert report.time_budget.per_slide

    def test_time_budget_over_flags_arc_warning(self, good_deck: Path) -> None:
        # 1-min slot is impossible to meet for any non-trivial deck.
        report = review_deck(good_deck, budget_minutes=1, qa_reserve_minutes=0)
        rules = {i["rule"] for i in report.story_arc_issues}
        assert "time-budget-over" in rules, (
            f"Expected time-budget-over warning. Got: {report.story_arc_issues}"
        )

    def test_anticipate_questions_runs_when_requested(self, good_deck: Path) -> None:
        report = review_deck(good_deck, anticipate_questions=True, qa_n_questions=5)
        # Good deck may not trigger heuristic questions; just confirm
        # the list shape is correct (no exception, list-typed).
        assert isinstance(report.anticipated_questions, list)
        # Length is bounded by the cap
        assert len(report.anticipated_questions) <= 5


# ---------------------------------------------------------------------------
# Speaker-notes structure audit (deferred sub-goal 5.4)
# ---------------------------------------------------------------------------
#
# Bobby's hard rule from ``feedback_slide_hard_rules``: every body slide
# needs two-tier speaker notes (mental_map heading + 200-400 word script).
# review_deck now audits the rendered notes for that structure. Title /
# divider slides are exempt from the body-word-count floor.


def _attach_notes(slide: Any, notes_text: str) -> None:
    """Drop notes text onto a python-pptx slide's notes slide."""
    slide.notes_slide.notes_text_frame.text = notes_text


def _build_minimal_slide(prs: PptxPresentation, *, title: str, body: str | None = None) -> Any:
    """Make a single Roboto-clean slide so other audits stay silent."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tr = tbox.text_frame.paragraphs[0].add_run()
    tr.text = title
    tr.font.size = Pt(32)
    tr.font.name = "Roboto"
    if body is not None:
        bbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4))
        br = bbox.text_frame.paragraphs[0].add_run()
        br.text = body
        br.font.size = Pt(24)
        br.font.name = "Roboto"
    return slide


def _add_figure_picture(slide: Any, fig_path: Path) -> None:
    """Attach a real PICTURE shape so the auditor classifies as 'figure'."""
    slide.shapes.add_picture(str(fig_path), Inches(1.5), Inches(2.5), Inches(8), Inches(4))


class TestSpeakerNotesAudit:
    """Audit speaker-notes structure per ``feedback_slide_hard_rules``."""

    @staticmethod
    def _well_structured_notes() -> str:
        """Mental map + ~250-word script, both required."""
        mental_map = (
            "- HOOK: We engineer T cells in three flavors.\n"
            "- KEY CLAIM: TCR, TAA, and CAR therapies differ in MHC-dependence.\n"
            "- EVIDENCE: Three-panel comparison in this figure.\n"
            "- KEY TERMS: scFv, ITAM, zeta-chain\n"
            "- CLICK: First click reveals annotations 1-4.\n"
            "- TRANSITION: Next slide shows the CAR construct."
        )
        script_body = (
            "These three panels show the three ways an engineered T cell recognises "
            "a tumour antigen. In panel a the introduced TCR competes with the "
            "endogenous TCR for an MHC presented peptide; recognition is therefore "
            "MHC restricted and requires the patient's HLA haplotype to match the "
            "engineered TCR's specificity. Panel b shows a TAA targeted approach "
            "where a soluble or membrane bound antigen is recognised independent "
            "of MHC. Panel c illustrates the CAR architecture: a single chain "
            "variable fragment fused to a transmembrane region and a zeta chain "
            "based intracellular signalling module that delivers signal one when "
            "antigen is engaged. The construct is independent of MHC presentation "
            "which broadens applicability but constrains target choice to surface "
            "antigens. We will return to these three architectures throughout the "
            "talk because each implies a different manufacturing footprint, a "
            "different safety profile, and a different set of failure modes that "
            "we will need to control for in clinical translation work going "
            "forward in this thesis programme. Across all three architectures "
            "the question of antigen escape remains a dominant clinical risk."
        )
        return f"{mental_map}\n\n--- DETAILED SCRIPT ---\n{script_body}"

    def test_figure_slide_with_empty_notes_is_critical(self, tmp_path: Path) -> None:
        """Empty notes on a figure / data slide → critical."""
        fig = _make_image(tmp_path / "fig.png", "blue")
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide first so deck has the expected opener.
        _build_minimal_slide(prs, title="Solid Title Deck Opener")
        _attach_notes(prs.slides[0], self._well_structured_notes())

        # Figure slide with EMPTY notes — the violation.
        figure_slide = _build_minimal_slide(
            prs, title="A descriptive title for this figure slide"
        )
        _add_figure_picture(figure_slide, fig)
        _attach_notes(figure_slide, "")  # critical

        out = tmp_path / "empty_notes_figure.pptx"
        prs.save(str(out))
        report = review_deck(out)

        # Find the figure-slide review
        fig_review = next(s for s in report.per_slide if s.slide_type == "figure")
        rules = [(i["severity"], i["rule"]) for i in fig_review.issues]
        assert ("critical", "speaker-notes-empty") in rules, (
            f"Expected critical speaker-notes-empty on figure slide. Got: {fig_review.issues}"
        )

    def test_notes_without_mental_map_is_warning(self, tmp_path: Path) -> None:
        """Body slide whose notes have script-only (no HOOK/CLAIM markers) → warning."""
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        _build_minimal_slide(prs, title="Solid Title Deck Opener")
        _attach_notes(prs.slides[0], self._well_structured_notes())

        body = _build_minimal_slide(
            prs,
            title="A descriptive body title here for testing",
            body="Body content for the slide that is long enough to read",
        )
        # 150-ish word free-form script with NO mental-map markers
        script_only = (
            "We open with a brief overview of the methodology. The protocol uses "
            "a standard flow cytometry pipeline with twelve color panels and a "
            "fluorescence minus one control for every channel. Compensation is "
            "computed from single stain controls run alongside the samples each "
            "day; we do not rely on stored compensation matrices for clinical "
            "samples because day to day laser intensity varies enough to affect "
            "the diagonal terms in the matrix. After acquisition we gate first "
            "on singlets then on live cells using a viability dye then on the "
            "lineage marker of interest and finally on the activation marker. "
            "Manual gating is reviewed by a second analyst blinded to the "
            "treatment arm to limit gating drift across the cohort, which we "
            "monitor as a daily quality control measure tracked over time."
        )
        _attach_notes(body, script_only)

        out = tmp_path / "notes_no_map.pptx"
        prs.save(str(out))
        report = review_deck(out)

        body_review = report.per_slide[1]
        rules = [(i["severity"], i["rule"]) for i in body_review.issues]
        assert ("warning", "speaker-notes-mental-map") in rules, (
            f"Expected warning speaker-notes-mental-map. Got: {body_review.issues}"
        )

    def test_short_script_on_body_slide_is_warning(self, tmp_path: Path) -> None:
        """< 100 word script body on a non-title/divider slide → warning."""
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        _build_minimal_slide(prs, title="Solid Title Deck Opener")
        _attach_notes(prs.slides[0], self._well_structured_notes())

        body = _build_minimal_slide(
            prs,
            title="A descriptive body title for the audit",
            body="Body that is long enough to project",
        )
        short_notes = (
            "- HOOK: opener\n"
            "- KEY CLAIM: claim\n\n"
            "--- DETAILED SCRIPT ---\n"
            "This is a very thin script with way too few words to actually present from."
        )
        _attach_notes(body, short_notes)

        out = tmp_path / "short_notes.pptx"
        prs.save(str(out))
        report = review_deck(out)

        body_review = report.per_slide[1]
        rules = [(i["severity"], i["rule"]) for i in body_review.issues]
        assert ("warning", "speaker-notes-short") in rules, (
            f"Expected warning speaker-notes-short. Got: {body_review.issues}"
        )

    def test_title_slide_empty_notes_is_exempt(self, tmp_path: Path) -> None:
        """Title slide with empty notes must not trip the audit."""
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        # Lone slide is the title — empty notes is fine for the opener.
        _build_minimal_slide(prs, title="Just A Title")
        _attach_notes(prs.slides[0], "")

        out = tmp_path / "title_empty_notes.pptx"
        prs.save(str(out))
        report = review_deck(out)

        title_review = report.per_slide[0]
        rules = {i["rule"] for i in title_review.issues}
        # No speaker-notes rules should fire on the title slide
        assert not any(r.startswith("speaker-notes") for r in rules), (
            f"Title slide should be exempt; got speaker-notes issues: {title_review.issues}"
        )

    def test_well_structured_notes_pass_clean(self, tmp_path: Path) -> None:
        """Mental map + 200-400 word script → no speaker-notes issues."""
        fig = _make_image(tmp_path / "fig.png", "green")
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        _build_minimal_slide(prs, title="Solid Title Deck Opener")
        _attach_notes(prs.slides[0], self._well_structured_notes())

        figure_slide = _build_minimal_slide(
            prs, title="A descriptive title for the figure slide"
        )
        _add_figure_picture(figure_slide, fig)
        _attach_notes(figure_slide, self._well_structured_notes())

        out = tmp_path / "clean_notes.pptx"
        prs.save(str(out))
        report = review_deck(out)

        all_rules = {i["rule"] for s in report.per_slide for i in s.issues}
        assert not any(r.startswith("speaker-notes") for r in all_rules), (
            "Well-structured notes should pass clean. "
            f"Issues: {[s.issues for s in report.per_slide]}"
        )

    def test_long_script_is_warning(self, tmp_path: Path) -> None:
        """> 500 word script on a non-deep-think slide → warning."""
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        _build_minimal_slide(prs, title="Solid Title Deck Opener")
        _attach_notes(prs.slides[0], self._well_structured_notes())

        body = _build_minimal_slide(
            prs,
            title="A descriptive body title for length audit",
            body="Body that is long enough to project",
        )
        long_script = " ".join(["word"] * 520)
        long_notes = (
            "- HOOK: opener\n"
            "- KEY CLAIM: claim\n"
            "- EVIDENCE: figure\n\n"
            f"--- DETAILED SCRIPT ---\n{long_script}"
        )
        _attach_notes(body, long_notes)

        out = tmp_path / "long_notes.pptx"
        prs.save(str(out))
        report = review_deck(out)

        body_review = report.per_slide[1]
        rules = [(i["severity"], i["rule"]) for i in body_review.issues]
        assert ("warning", "speaker-notes-long") in rules, (
            f"Expected warning speaker-notes-long. Got: {body_review.issues}"
        )
