from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL")
pptx_mod = pytest.importorskip("pptx")

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


def test_panel_layout_contract_validates_dense_non_overlapping_layout() -> None:
    from vaultlab.slides.panel_contract import (
        PanelLayoutContract,
        PanelSlot,
        audit_panel_layout_contract,
    )

    contract = PanelLayoutContract(
        figure_id="FigSynthetic",
        slide_width_in=13.333,
        slide_height_in=7.5,
        min_gutter_in=0.04,
        panel_letter_font_pt=14.0,
        min_effective_text_pt=5.5,
        panels=[
            PanelSlot(letter="A", image_path="a.png", slot_in=[0.5, 0.8, 3.0, 2.0]),
            PanelSlot(letter="B", image_path="b.png", slot_in=[3.58, 0.8, 3.0, 2.0]),
            PanelSlot(letter="C", image_path="c.png", slot_in=[0.5, 2.9, 3.0, 2.0]),
            PanelSlot(letter="D", image_path="d.png", slot_in=[3.58, 2.9, 3.0, 2.0]),
        ],
    )

    audit = audit_panel_layout_contract(contract)

    assert audit.ok() is True
    assert audit.overall_severity == "pass"
    assert audit.n_fail == 0


def test_panel_contract_helpers_are_exported_from_slides_package() -> None:
    from vaultlab.slides import (
        PanelLayoutContract,
        PanelSlot,
        audit_panel_layout_contract,
        extract_pptx_slide_geometry,
    )

    assert PanelLayoutContract.__name__ == "PanelLayoutContract"
    assert PanelSlot.__name__ == "PanelSlot"
    assert callable(audit_panel_layout_contract)
    assert callable(extract_pptx_slide_geometry)


def test_panel_layout_contract_flags_overlap_bounds_and_letter_font() -> None:
    from vaultlab.slides.panel_contract import (
        PanelLayoutContract,
        PanelSlot,
        audit_panel_layout_contract,
    )

    contract = PanelLayoutContract(
        figure_id="FigBad",
        slide_width_in=13.333,
        slide_height_in=7.5,
        min_gutter_in=0.1,
        panel_letter_font_pt=8.0,
        min_effective_text_pt=5.5,
        panels=[
            PanelSlot(letter="A", image_path="a.png", slot_in=[0.5, 0.8, 4.0, 2.0]),
            PanelSlot(letter="B", image_path="b.png", slot_in=[4.45, 0.8, 4.0, 2.0]),
            PanelSlot(letter="C", image_path="c.png", slot_in=[12.9, 6.8, 1.0, 1.0]),
        ],
    )

    audit = audit_panel_layout_contract(contract)
    failures = {(issue.rule, issue.severity) for issue in audit.issues}

    assert audit.ok() is False
    assert audit.overall_severity == "fail"
    assert ("panel-slot-overlap", "fail") in failures
    assert ("panel-slot-bounds", "fail") in failures
    assert ("panel-letter-font", "fail") in failures


def test_extract_pptx_slide_geometry_reads_shape_text_coordinates_and_media(
    tmp_path: Path,
) -> None:
    from vaultlab.slides.panel_contract import extract_pptx_slide_geometry

    image_path = tmp_path / "panel.png"
    Image.new("RGB", (320, 240), "navy").save(image_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.0), Inches(0.5))
    title_run = title.text_frame.paragraphs[0].add_run()
    title_run.text = "Synthetic panel slide"
    title_run.font.size = Pt(24)

    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(1.0),
        width=Inches(4.0),
        height=Inches(3.0),
    )
    picture.name = "panel-A-image"

    label = slide.shapes.add_textbox(Inches(0.28), Inches(0.85), Inches(0.25), Inches(0.25))
    label_run = label.text_frame.paragraphs[0].add_run()
    label_run.text = "A"
    label_run.font.size = Pt(14)

    pptx_path = tmp_path / "synthetic_panel.pptx"
    prs.save(pptx_path)

    geometry = extract_pptx_slide_geometry(pptx_path, slide_number=1)

    assert geometry.slide_width_in == pytest.approx(13.333, abs=0.001)
    assert geometry.slide_height_in == pytest.approx(7.5)
    assert len(geometry.shapes) == 3
    assert any(shape.text == "Synthetic panel slide" for shape in geometry.shapes)
    picture_shapes = [shape for shape in geometry.shapes if shape.kind == "picture"]
    assert len(picture_shapes) == 1
    assert picture_shapes[0].name == "panel-A-image"
    assert picture_shapes[0].media_width_px == 320
    assert picture_shapes[0].media_height_px == 240
    assert picture_shapes[0].x_in == pytest.approx(0.5)
    assert picture_shapes[0].w_in == pytest.approx(4.0)
