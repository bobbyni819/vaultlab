"""Tests for :mod:`vaultlab.slides.version_diff` — slide-level pptx diff."""

from __future__ import annotations

import copy
from pathlib import Path

import pytest

pptx_mod = pytest.importorskip("pptx")

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from vaultlab.slides.version_diff import DeckDiff, SlideDiff, diff_decks


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(32)
    r.font.name = "Roboto"


def _add_body(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24)
    r.font.name = "Roboto"


def _make_deck(tmp_path: Path, slides: list[tuple[str, str]], name: str = "deck.pptx") -> Path:
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for title, body in slides:
        s = prs.slides.add_slide(blank)
        _add_title(s, title)
        if body:
            _add_body(s, body)
    out = tmp_path / name
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestUnchangedDeck:
    def test_identical_decks_report_no_changes(self, tmp_path: Path) -> None:
        spec = [("First slide", "Body one"), ("Second slide", "Body two")]
        a = _make_deck(tmp_path, spec, name="a.pptx")
        b = _make_deck(tmp_path, spec, name="b.pptx")
        diff = diff_decks(a, b)
        assert isinstance(diff, DeckDiff)
        assert diff.n_added == 0
        assert diff.n_removed == 0
        assert diff.n_modified == 0
        # Every slide should map and be unchanged
        assert all(s.change_kind == "unchanged" for s in diff.slides)


class TestAddRemove:
    def test_added_slide_detected(self, tmp_path: Path) -> None:
        a = _make_deck(tmp_path, [("First", "Body")], name="a.pptx")
        b = _make_deck(
            tmp_path,
            [("First", "Body"), ("Second", "Brand new slide")],
            name="b.pptx",
        )
        diff = diff_decks(a, b)
        assert diff.n_added == 1
        added = [s for s in diff.slides if s.change_kind == "added"]
        assert len(added) == 1
        assert added[0].slide_index_a is None
        assert added[0].slide_index_b == 1

    def test_removed_slide_detected(self, tmp_path: Path) -> None:
        a = _make_deck(
            tmp_path,
            [("First", "Body"), ("Doomed", "Will be removed")],
            name="a.pptx",
        )
        b = _make_deck(tmp_path, [("First", "Body")], name="b.pptx")
        diff = diff_decks(a, b)
        assert diff.n_removed == 1
        removed = [s for s in diff.slides if s.change_kind == "removed"]
        assert len(removed) == 1
        assert removed[0].slide_index_a == 1
        assert removed[0].slide_index_b is None


class TestModification:
    def test_text_modification_detected(self, tmp_path: Path) -> None:
        """Changing a body text on one slide → modified with text field change."""
        a = _make_deck(
            tmp_path,
            [("Title slide", "Original body text"), ("Other", "Same body")],
            name="a.pptx",
        )
        b = _make_deck(
            tmp_path,
            [("Title slide", "Modified body text"), ("Other", "Same body")],
            name="b.pptx",
        )
        diff = diff_decks(a, b)
        modified = [s for s in diff.slides if s.change_kind == "modified"]
        assert len(modified) == 1
        m = modified[0]
        assert m.slide_index_a == 0
        assert m.slide_index_b == 0
        assert m.field_changes, "Expected at least one field-level change record"
        # At least one of the field changes should mention the text changing
        changes_blob = " ".join(
            f"{shape_id} {before} {after}" for shape_id, before, after in m.field_changes
        )
        assert "Original" in changes_blob or "Modified" in changes_blob

    def test_title_modification_detected(self, tmp_path: Path) -> None:
        a = _make_deck(tmp_path, [("Title A", "Body")], name="a.pptx")
        b = _make_deck(tmp_path, [("Title B", "Body")], name="b.pptx")
        diff = diff_decks(a, b)
        modified = [s for s in diff.slides if s.change_kind == "modified"]
        assert modified, f"Expected modified slide. Diff slides: {diff.slides}"


class TestAggregate:
    def test_combined_add_remove_modify(self, tmp_path: Path) -> None:
        a = _make_deck(
            tmp_path,
            [
                ("Stable", "Same"),
                ("Doomed", "Removed in B"),
                ("Will change", "Old body"),
            ],
            name="a.pptx",
        )
        b = _make_deck(
            tmp_path,
            [
                ("Stable", "Same"),
                ("Will change", "New body"),
                ("Brand new", "Added in B"),
            ],
            name="b.pptx",
        )
        diff = diff_decks(a, b)
        # Counts should be sensible — at least one of each kind
        assert diff.n_added >= 1
        assert diff.n_removed >= 1
        assert diff.n_modified >= 1


class TestErrors:
    def test_missing_pptx_a_raises(self, tmp_path: Path) -> None:
        b = _make_deck(tmp_path, [("First", "Body")], name="b.pptx")
        with pytest.raises(FileNotFoundError):
            diff_decks(tmp_path / "nope.pptx", b)

    def test_missing_pptx_b_raises(self, tmp_path: Path) -> None:
        a = _make_deck(tmp_path, [("First", "Body")], name="a.pptx")
        with pytest.raises(FileNotFoundError):
            diff_decks(a, tmp_path / "nope.pptx")


class TestSlideDiffShape:
    def test_slide_diff_fields(self, tmp_path: Path) -> None:
        a = _make_deck(tmp_path, [("Title", "Body")], name="a.pptx")
        b = _make_deck(tmp_path, [("Title", "Body")], name="b.pptx")
        diff = diff_decks(a, b)
        for s in diff.slides:
            assert isinstance(s, SlideDiff)
            assert s.change_kind in {"added", "removed", "modified", "unchanged"}
