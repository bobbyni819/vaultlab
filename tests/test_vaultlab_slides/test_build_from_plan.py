"""Tests for vaultlab.slides.deck.build_from_plan — dict-plan deck builder.

Ported from ``bobby-tools/tests/test_bobby_slides/test_build_from_plan.py``.
The mandatory ``test_synthetic_plan_all_types`` test covers all six slide
types (title / section_divider / figure / multi_figure / text /
references) called out in the lift spec.
"""
from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL")
pptx = pytest.importorskip("pptx")

from PIL import Image
from pptx import Presentation as PptxPresentation

from vaultlab.slides import KBReader, build_from_plan
from vaultlab.slides.template import lab_template_path


pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — build_from_plan tests need it",
)


def make_test_image(path: Path, color: str = "red"):
    Image.new("RGB", (200, 150), color).save(str(path))
    return path


@pytest.fixture
def kb_root(tmp_path):
    """Minimal KB so we can test KB logging."""
    (tmp_path / "Output" / "Reports").mkdir(parents=True)
    return tmp_path


@pytest.fixture
def sample_plan(tmp_path):
    fig1 = tmp_path / "fig1.png"
    fig2 = tmp_path / "fig2.png"
    make_test_image(fig1, "red")
    make_test_image(fig2, "blue")

    return {
        "title": "Test Talk",
        "author": "Bobby Ni",
        "subtitle": "BME Retreat",
        "topic": "Phospholipid Programs",
        "kb": "metabolism",
        "slides": [
            {
                "type": "title",
                "title": "Phospholipid Programs",
                "subtitle": "BME Retreat 2026",
                "author": "Bobby Ni",
                "speaker_notes": {
                    "hook": "Lipids tell a story",
                    "key_terms": ["IBD", "phospholipid"],
                },
            },
            {"type": "section_divider", "title": "Background"},
            {
                "type": "figure",
                "title": "MALDI-IMS workflow",
                "image_path": str(fig1),
                "caption": "Sample prep through analysis",
                "bullets": ["n=47 patients", "Spatial resolution: 10um"],
                "citation_source": "Smith et al., 2024",
                "speaker_notes": {
                    "hook": "How do we measure lipids in tissue?",
                    "key_claim": "MALDI-IMS gives spatial lipid maps",
                },
            },
            {
                "type": "multi_figure",
                "title": "Three lipid programs",
                "figures": [
                    {"path": str(fig1), "label": "A", "caption": "Phospholipids"},
                    {"path": str(fig2), "label": "B", "caption": "Sphingolipids"},
                ],
            },
            {"type": "text", "title": "Conclusions", "bullets": ["Three programs", "IBD-specific"]},
            {"type": "references", "references": ["Smith et al., 2024, Nature"]},
        ],
    }


class TestBuildFromPlanBasics:
    def test_creates_pptx(self, sample_plan, tmp_path):
        out = tmp_path / "deck.pptx"
        result = build_from_plan(sample_plan, out, write_marp=False)
        assert result["pptx"] == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_pptx_has_all_slides(self, sample_plan, tmp_path):
        out = tmp_path / "deck.pptx"
        build_from_plan(sample_plan, out, write_marp=False)
        prs = PptxPresentation(str(out))
        assert len(prs.slides) >= 6

    def test_speaker_notes_attached(self, sample_plan, tmp_path):
        out = tmp_path / "deck.pptx"
        build_from_plan(sample_plan, out, write_marp=False)
        prs = PptxPresentation(str(out))
        all_notes = []
        for s in prs.slides:
            try:
                all_notes.append(s.notes_slide.notes_text_frame.text)
            except Exception:
                continue
        joined = "\n".join(all_notes)
        assert "HOOK: Lipids tell a story" in joined


class TestSyntheticPlanAllTypes:
    """Mandatory test from the lift spec: 5+ slide plan covering all six types.

    Verifies build_from_plan accepts a typed plan and produces a real .pptx
    with the expected slide count + shape names per slide type.
    """

    def test_synthetic_plan_all_types(self, tmp_path):
        fig = tmp_path / "synthetic.png"
        make_test_image(fig, "red")

        plan = {
            "title": "Synthetic L4 Deck",
            "author": "Test Author",
            "topic": "synthetic-coverage",
            "slides": [
                {"type": "title", "title": "Synthetic L4 Deck",
                 "subtitle": "All six types", "author": "Test Author"},
                {"type": "section_divider", "title": "Section 1"},
                {"type": "figure", "title": "Single Figure",
                 "image_path": str(fig), "caption": "Solo",
                 "bullets": ["Bullet A", "Bullet B"],
                 "citation_source": "Cite 2024"},
                {"type": "multi_figure", "title": "Two Panels",
                 "figures": [
                     {"path": str(fig), "label": "A", "caption": "Panel A"},
                     {"path": str(fig), "label": "B", "caption": "Panel B"},
                 ]},
                {"type": "text", "title": "Conclusions",
                 "bullets": ["Conclusion one", "Conclusion two"]},
                {"type": "references", "title": "References",
                 "references": ["Smith 2024", "Jones 2023"]},
            ],
        }

        out = tmp_path / "synthetic.pptx"
        result = build_from_plan(plan, out, write_marp=False)
        assert result["pptx"].exists()

        prs = PptxPresentation(str(out))
        # Plan has 6 slides; lab template adds none (strip_starter_slides).
        assert len(prs.slides) == 6, f"Expected 6 slides, got {len(prs.slides)}"

        # Each slide has at least one shape (even an "empty" multi_figure
        # slide gets a title shape).
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) >= 1, f"slide {i} has no shapes"

        # Title slide (index 0) — first text frame contains the title text
        title_text = next(
            (sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame),
            "",
        )
        assert "Synthetic L4 Deck" in title_text

        # Section divider (index 1) — should contain the section name
        divider_texts = [
            sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame
        ]
        assert any("Section 1" in t for t in divider_texts)

        # Figure slide (index 2) — has a picture (shape_type 13) and text
        figure_shapes = prs.slides[2].shapes
        has_picture = any(sh.shape_type == 13 for sh in figure_shapes)
        assert has_picture, "Figure slide should contain a picture"

        # Multi-figure slide (index 3) — has 2 pictures
        multi_pictures = [
            sh for sh in prs.slides[3].shapes if sh.shape_type == 13
        ]
        assert len(multi_pictures) == 2, (
            f"multi_figure slide should have 2 pictures, got {len(multi_pictures)}"
        )

        # Text slide (index 4) — title + bullets
        text_texts = [
            sh.text_frame.text for sh in prs.slides[4].shapes if sh.has_text_frame
        ]
        joined = "\n".join(text_texts)
        assert "Conclusions" in joined
        assert "Conclusion one" in joined

        # References slide (index 5)
        refs_texts = [
            sh.text_frame.text for sh in prs.slides[5].shapes if sh.has_text_frame
        ]
        joined_refs = "\n".join(refs_texts)
        assert "References" in joined_refs
        assert "Smith 2024" in joined_refs


class TestMarpMirror:
    def test_writes_marp_by_default(self, sample_plan, tmp_path):
        out = tmp_path / "deck.pptx"
        result = build_from_plan(sample_plan, out)
        assert "marp" in result
        assert result["marp"].suffix == ".md"
        assert result["marp"].exists()
        assert "marp: true" in result["marp"].read_text(encoding="utf-8")

    def test_skips_marp_when_disabled(self, sample_plan, tmp_path):
        out = tmp_path / "deck.pptx"
        result = build_from_plan(sample_plan, out, write_marp=False)
        assert "marp" not in result


class TestKBLogging:
    def test_writes_report_when_kb_provided(self, sample_plan, tmp_path, kb_root):
        out = tmp_path / "deck.pptx"
        reader = KBReader(kb_root)
        result = build_from_plan(sample_plan, out, write_marp=False, kb_log=reader)
        assert "report" in result
        assert result["report"].exists()
        assert result["report"].parent == kb_root / "Output" / "Reports"

    def test_appends_to_log(self, sample_plan, tmp_path, kb_root):
        out = tmp_path / "deck.pptx"
        reader = KBReader(kb_root)
        build_from_plan(sample_plan, out, write_marp=False, kb_log=reader)
        log_path = kb_root / "_Log.md"
        assert log_path.exists()
        log_text = log_path.read_text(encoding="utf-8")
        assert "compile" in log_text
        assert "Phospholipid Programs" in log_text

    def test_report_filename_uses_topic(self, sample_plan, tmp_path, kb_root):
        out = tmp_path / "deck.pptx"
        reader = KBReader(kb_root)
        result = build_from_plan(sample_plan, out, write_marp=False, kb_log=reader)
        assert "phospholipid-programs" in result["report"].name
        assert result["report"].name.endswith("-deck.md")

    def test_report_includes_sources(self, sample_plan, tmp_path, kb_root):
        out = tmp_path / "deck.pptx"
        reader = KBReader(kb_root)
        result = build_from_plan(sample_plan, out, write_marp=False, kb_log=reader)
        report_text = result["report"].read_text(encoding="utf-8")
        assert "Smith et al., 2024" in report_text


class TestEmptyPlan:
    def test_empty_slides_does_not_crash(self, tmp_path):
        out = tmp_path / "empty.pptx"
        plan = {"title": "Empty", "slides": []}
        result = build_from_plan(plan, out, write_marp=False)
        assert result["pptx"].exists()


class TestWithAnimations:
    def _count_click_effects_in_deck(self, pptx_path: Path) -> int:
        prs = PptxPresentation(str(pptx_path))
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        count = 0
        for slide in prs.slides:
            timing = slide._element.find(f"{ns}timing")
            if timing is None:
                continue
            for ctn in timing.iter(f"{ns}cTn"):
                if ctn.get("nodeType") == "clickEffect":
                    count += 1
        return count

    def test_animations_off_by_default(self, sample_plan, tmp_path):
        out = tmp_path / "noanim.pptx"
        build_from_plan(sample_plan, out, write_marp=False)
        assert self._count_click_effects_in_deck(out) == 0

    def test_animations_on_text_slides(self, tmp_path):
        out = tmp_path / "anim_text.pptx"
        plan = {
            "title": "Animated",
            "slides": [
                {"type": "text", "title": "T", "bullets": ["one", "two", "three"]},
            ],
        }
        build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert self._count_click_effects_in_deck(out) >= 3

    def test_animations_on_multi_figure(self, tmp_path):
        out = tmp_path / "anim_multi.pptx"
        fig1, fig2 = tmp_path / "1.png", tmp_path / "2.png"
        make_test_image(fig1, "red")
        make_test_image(fig2, "blue")
        plan = {
            "title": "Animated multi",
            "slides": [
                {
                    "type": "multi_figure",
                    "title": "Two panels",
                    "figures": [
                        {"path": str(fig1), "label": "A"},
                        {"path": str(fig2), "label": "B"},
                    ],
                },
            ],
        }
        build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert self._count_click_effects_in_deck(out) >= 2

    def test_title_slide_not_animated(self, tmp_path):
        out = tmp_path / "anim_title.pptx"
        plan = {
            "title": "T",
            "slides": [
                {"type": "title", "title": "T", "subtitle": "S", "author": "A"},
            ],
        }
        build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert self._count_click_effects_in_deck(out) == 0

    def test_section_divider_not_animated(self, tmp_path):
        out = tmp_path / "anim_div.pptx"
        plan = {
            "title": "T",
            "slides": [
                {"type": "section_divider", "title": "Methods"},
            ],
        }
        build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert self._count_click_effects_in_deck(out) == 0

    def test_single_bullet_no_reveal(self, tmp_path):
        out = tmp_path / "anim_single.pptx"
        plan = {
            "title": "T",
            "slides": [
                {"type": "text", "title": "T", "bullets": ["only one"]},
            ],
        }
        build_from_plan(plan, out, write_marp=False, with_animations=True)
        assert self._count_click_effects_in_deck(out) == 0
