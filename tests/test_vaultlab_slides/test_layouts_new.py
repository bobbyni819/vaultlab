"""Tests for the 4 new slide-layout primitives (sub-goal 5.3).

Layouts under test:
    - add_equation_slide
    - add_table_slide
    - add_comparison_table_slide
    - add_acknowledgments_grid_slide

Each test verifies:
1. The layout function returns a slide object.
2. Round-tripping the .pptx through python-pptx works.
3. The hard slide rules are honored (Roboto font, min sizes 28/24/18,
   no overlapping shapes).
4. Expected content actually lands on the slide.
"""

from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL", reason="Pillow required for layout tests")
pptx = pytest.importorskip("pptx", reason="python-pptx required for layout tests")

from pptx import Presentation as PptxPresentation

from vaultlab.slides.layouts import (
    add_acknowledgments_grid_slide,
    add_comparison_table_slide,
    add_equation_slide,
    add_table_slide,
)
from vaultlab.slides.template import default_font, lab_template_path, load_template

pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — skipping new-layout tests",
)


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #


def _save_and_reload(pres, tmp_path: Path, name: str):
    """Save to tmp_path and reload — verifies the .pptx is valid."""
    out = tmp_path / name
    pres.save(str(out))
    assert out.exists()
    return PptxPresentation(str(out))


def _all_text(pres) -> list[str]:
    out: list[str] = []
    for s in pres.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return out


def _all_runs_use_roboto(pres) -> bool:
    found_any = False
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            found_any = True
                            if run.font.name != default_font():
                                return False
    return found_any


def _all_runs_respect_min_sizes(pres) -> bool:
    """Every run should be >= 18pt (the caption minimum).

    The sized helpers may emit 9pt citation-source footers and 12pt
    italic captions on existing layouts; the new layouts must respect
    the 18-pt caption floor for their primary content.
    """
    min_pt = 18
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is None:
                            continue
                        size_pt = run.font.size.pt
                        if size_pt < min_pt:
                            return False
    return True


def _shape_bbox(shape):
    """Return (left, top, right, bottom) in EMU for any shape."""
    return (
        shape.left,
        shape.top,
        shape.left + shape.width,
        shape.top + shape.height,
    )


def _rects_overlap(a, b, tol: int = 0) -> bool:
    al, at, ar, ab = a
    bl, bt, br, bb = b
    return not (ar <= bl + tol or br <= al + tol or ab <= bt + tol or bb <= at + tol)


def _no_significant_overlap(pres, tol_emu: int = 9144) -> bool:
    """No two non-title text/picture shapes should overlap by more than ~0.01".

    tol_emu = 9144 EMU ~ 0.01 inch slack to allow for sub-pixel rounding.
    """
    for slide in pres.slides:
        rects = []
        for shape in slide.shapes:
            try:
                rects.append(_shape_bbox(shape))
            except Exception:
                continue
        n = len(rects)
        for i in range(n):
            for j in range(i + 1, n):
                if _rects_overlap(rects[i], rects[j], tol=-tol_emu):
                    # Negative tolerance means the overlap must be >tol_emu
                    # in each axis to count as a violation.
                    return False
    return True


# --------------------------------------------------------------------------- #
# add_equation_slide
# --------------------------------------------------------------------------- #


class TestAddEquationSlide:
    def test_returns_slide(self, tmp_path):
        pres = load_template()
        slide = add_equation_slide(
            pres,
            equation="E = mc^2",
            title="Mass-energy equivalence",
        )
        assert slide is not None
        reloaded = _save_and_reload(pres, tmp_path, "eq.pptx")
        assert len(reloaded.slides) == 1

    def test_title_and_equation_rendered(self, tmp_path):
        pres = load_template()
        add_equation_slide(
            pres,
            equation="dC/dt = k * (C_ss - C)",
            title="First-order kinetics",
            caption="Where k is the rate constant",
        )
        reloaded = _save_and_reload(pres, tmp_path, "eq2.pptx")
        texts = _all_text(reloaded)
        assert any("First-order kinetics" in t for t in texts)
        assert any("dC/dt" in t for t in texts)
        assert any("rate constant" in t for t in texts)

    def test_uses_roboto(self, tmp_path):
        pres = load_template()
        add_equation_slide(pres, equation="y = mx + b", title="Linear model")
        reloaded = _save_and_reload(pres, tmp_path, "eq3.pptx")
        assert _all_runs_use_roboto(reloaded)

    def test_no_overlap(self, tmp_path):
        pres = load_template()
        add_equation_slide(
            pres,
            equation="alpha + beta = gamma",
            title="A descriptive sentence-title for the equation",
            caption="Caption goes here",
        )
        reloaded = _save_and_reload(pres, tmp_path, "eq4.pptx")
        assert _no_significant_overlap(reloaded)

    def test_no_caption_is_ok(self, tmp_path):
        pres = load_template()
        add_equation_slide(pres, equation="x = 1", title="Trivial")
        reloaded = _save_and_reload(pres, tmp_path, "eq5.pptx")
        assert len(reloaded.slides) == 1


# --------------------------------------------------------------------------- #
# add_table_slide
# --------------------------------------------------------------------------- #


class TestAddTableSlide:
    def test_basic_rows(self, tmp_path):
        pres = load_template()
        rows = [
            ["Condition", "n", "p-value"],
            ["Control", "12", "—"],
            ["Treatment A", "12", "0.04"],
            ["Treatment B", "12", "<0.001"],
        ]
        slide = add_table_slide(
            pres,
            rows,
            title="Experimental conditions",
        )
        assert slide is not None
        reloaded = _save_and_reload(pres, tmp_path, "tbl.pptx")
        texts = _all_text(reloaded)
        assert any("Experimental conditions" in t for t in texts)
        # At least one cell text should land somewhere in the slide
        assert any("Treatment A" in t for t in texts) or any(
            "Treatment A" in cell.text
            for s in reloaded.slides
            for shape in s.shapes
            if shape.has_table
            for row in shape.table.rows
            for cell in row.cells
        )

    def test_header_styled_bold(self, tmp_path):
        pres = load_template()
        rows = [
            ["Header A", "Header B"],
            ["v1", "v2"],
        ]
        add_table_slide(pres, rows, title="Header check")
        reloaded = _save_and_reload(pres, tmp_path, "tblhdr.pptx")
        # Find the table and check first row has bold runs.
        found_bold = False
        for s in reloaded.slides:
            for shape in s.shapes:
                if shape.has_table:
                    first_row = shape.table.rows[0]
                    for cell in first_row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.bold:
                                    found_bold = True
        assert found_bold

    def test_fallback_when_too_many_rows(self, tmp_path):
        pres = load_template()
        # Header + 11 body rows = 12 total → exceeds the 10-body-row cap.
        rows = [["col1", "col2"]] + [[f"r{i}", f"v{i}"] for i in range(11)]
        add_table_slide(
            pres,
            rows,
            title="Big table",
            appendix_message="See appendix for full table",
        )
        reloaded = _save_and_reload(pres, tmp_path, "tblbig.pptx")
        texts = _all_text(reloaded)
        assert any("appendix" in t.lower() for t in texts)

    def test_uses_roboto_and_min_sizes(self, tmp_path):
        pres = load_template()
        rows = [["A", "B"], ["1", "2"]]
        add_table_slide(pres, rows, title="Min-size table")
        reloaded = _save_and_reload(pres, tmp_path, "tblsizes.pptx")
        assert _all_runs_use_roboto(reloaded)
        assert _all_runs_respect_min_sizes(reloaded)

    def test_empty_rows_does_not_crash(self, tmp_path):
        pres = load_template()
        add_table_slide(pres, [], title="Empty")
        reloaded = _save_and_reload(pres, tmp_path, "tblempty.pptx")
        assert len(reloaded.slides) == 1


# --------------------------------------------------------------------------- #
# add_comparison_table_slide
# --------------------------------------------------------------------------- #


class TestAddComparisonTableSlide:
    def test_basic_two_columns(self, tmp_path):
        pres = load_template()
        slide = add_comparison_table_slide(
            pres,
            left_header="Approach A",
            right_header="Approach B",
            left_bullets=["Fast to run", "Lower fidelity"],
            right_bullets=["Slow to run", "Publication-grade fidelity"],
            title="Approach A vs Approach B trade-off",
        )
        assert slide is not None
        reloaded = _save_and_reload(pres, tmp_path, "cmp.pptx")
        texts = _all_text(reloaded)
        assert any("Approach A" in t for t in texts)
        assert any("Approach B" in t for t in texts)
        assert any("Fast to run" in t for t in texts)
        assert any("Publication-grade fidelity" in t for t in texts)

    def test_with_key_insight_callout(self, tmp_path):
        pres = load_template()
        add_comparison_table_slide(
            pres,
            left_header="Before",
            right_header="After",
            left_bullets=["Slow", "Error-prone"],
            right_bullets=["Fast", "Reproducible"],
            title="Before vs After the refactor",
            key_insight="The refactor unlocked an order-of-magnitude speedup",
        )
        reloaded = _save_and_reload(pres, tmp_path, "cmpinsight.pptx")
        texts = _all_text(reloaded)
        assert any("order-of-magnitude" in t for t in texts)

    def test_uses_roboto_min_sizes(self, tmp_path):
        pres = load_template()
        add_comparison_table_slide(
            pres,
            left_header="A",
            right_header="B",
            left_bullets=["one", "two"],
            right_bullets=["three"],
            title="Compare",
        )
        reloaded = _save_and_reload(pres, tmp_path, "cmpsizes.pptx")
        assert _all_runs_use_roboto(reloaded)
        assert _all_runs_respect_min_sizes(reloaded)

    def test_no_overlap(self, tmp_path):
        pres = load_template()
        add_comparison_table_slide(
            pres,
            left_header="LeftLeftLeft",
            right_header="RightRightRight",
            left_bullets=["one bullet", "two bullets", "three"],
            right_bullets=["a", "b", "c"],
            title="No overlap please",
            key_insight="Centered key insight callout",
        )
        reloaded = _save_and_reload(pres, tmp_path, "cmpoverlap.pptx")
        assert _no_significant_overlap(reloaded)


# --------------------------------------------------------------------------- #
# add_acknowledgments_grid_slide
# --------------------------------------------------------------------------- #


class TestAddAcknowledgmentsGridSlide:
    def test_basic(self, tmp_path):
        pres = load_template()
        people = [
            ("John Hickey", "PI", "Duke BME"),
            ("Bobby Ni", "Student", "Duke BME"),
            ("Collaborator A", "Postdoc", "Stanford"),
        ]
        slide = add_acknowledgments_grid_slide(
            pres,
            people,
            title="Acknowledgments",
        )
        assert slide is not None
        reloaded = _save_and_reload(pres, tmp_path, "ack.pptx")
        texts = _all_text(reloaded)
        assert any("John Hickey" in t for t in texts)
        assert any("Bobby Ni" in t for t in texts)
        assert any("Stanford" in t for t in texts)

    def test_role_optional_affiliation(self, tmp_path):
        pres = load_template()
        people = [
            ("Alice", "Reviewer", ""),
            ("Bob", "", "Some University"),
            ("Carol", "Mentor", "ACME"),
        ]
        add_acknowledgments_grid_slide(pres, people)
        reloaded = _save_and_reload(pres, tmp_path, "ack2.pptx")
        texts = _all_text(reloaded)
        assert any("Alice" in t for t in texts)
        assert any("Bob" in t for t in texts)
        assert any("Carol" in t for t in texts)

    def test_uses_roboto_min_sizes(self, tmp_path):
        pres = load_template()
        people = [("X", "Role", "Place")]
        add_acknowledgments_grid_slide(pres, people, title="Sized ack")
        reloaded = _save_and_reload(pres, tmp_path, "acksizes.pptx")
        assert _all_runs_use_roboto(reloaded)
        assert _all_runs_respect_min_sizes(reloaded)

    def test_empty_people_does_not_crash(self, tmp_path):
        pres = load_template()
        add_acknowledgments_grid_slide(pres, [], title="No one")
        reloaded = _save_and_reload(pres, tmp_path, "ackempty.pptx")
        assert len(reloaded.slides) == 1

    def test_large_grid(self, tmp_path):
        pres = load_template()
        people = [
            (f"Person {i}", f"Role {i}", f"Aff {i}") for i in range(12)
        ]
        add_acknowledgments_grid_slide(pres, people, title="Big grid")
        reloaded = _save_and_reload(pres, tmp_path, "ackbig.pptx")
        texts = _all_text(reloaded)
        # Make sure at least the first and last names land
        assert any("Person 0" in t for t in texts)
        assert any("Person 11" in t for t in texts)

    def test_no_overlap(self, tmp_path):
        pres = load_template()
        people = [
            ("Person A", "Role A", "Affil A"),
            ("Person B", "Role B", "Affil B"),
            ("Person C", "Role C", "Affil C"),
            ("Person D", "Role D", "Affil D"),
        ]
        add_acknowledgments_grid_slide(pres, people, title="No overlap ack")
        reloaded = _save_and_reload(pres, tmp_path, "ackoverlap.pptx")
        assert _no_significant_overlap(reloaded)


# --------------------------------------------------------------------------- #
# Integration: all four layouts can be composed into one deck
# --------------------------------------------------------------------------- #


class TestFullDeck:
    def test_compose_all_four(self, tmp_path):
        pres = load_template()
        add_equation_slide(pres, equation="y = f(x)", title="Equation")
        add_table_slide(
            pres,
            [["A", "B"], ["1", "2"], ["3", "4"]],
            title="Table",
        )
        add_comparison_table_slide(
            pres,
            left_header="Old",
            right_header="New",
            left_bullets=["slow"],
            right_bullets=["fast"],
            title="Comparison",
            key_insight="It is faster now",
        )
        add_acknowledgments_grid_slide(
            pres,
            [("Bobby", "Student", "Duke")],
            title="Thanks",
        )
        reloaded = _save_and_reload(pres, tmp_path, "all_four.pptx")
        assert len(reloaded.slides) == 4
        assert _all_runs_use_roboto(reloaded)
