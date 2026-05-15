"""Tests for :mod:`vaultlab.slides.qa_anticipator` — Q&A anticipator.

Heuristic-mode tests build tiny pptx fixtures with text that should
trigger specific question templates (statistical claims, comparisons,
future work, limitations). LLM-mode is exercised with a fake
runner_callback.
"""

from __future__ import annotations

from pathlib import Path

import pytest

pptx_mod = pytest.importorskip("pptx")

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from vaultlab.slides.qa_anticipator import (
    AnticipatedQuestion,
    anticipate_qa,
)


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------


def _add_text(slide, text: str, *, top_in: float = 0.3, pt: int = 24) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top_in), Inches(12), Inches(2))
    tf = box.text_frame
    lines = text.splitlines() or [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(pt)
        r.font.name = "Roboto"


def _make_deck(tmp_path: Path, slides_spec: list[tuple[str, str]]) -> Path:
    """Build a deck where each slide gets a title (top) and a body text shape."""
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for title, body in slides_spec:
        s = prs.slides.add_slide(blank)
        _add_text(s, title, top_in=0.3, pt=32)
        _add_text(s, body, top_in=2.0, pt=24)
    out = tmp_path / "qa.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Heuristic-mode tests
# ---------------------------------------------------------------------------


class TestHeuristicMode:
    def test_statistical_claim_triggers_question(self, tmp_path: Path) -> None:
        """A slide with a p-value claim → at least one question anchored to that slide."""
        out = _make_deck(
            tmp_path,
            [
                ("Title slide", "Welcome to the talk"),
                ("Survival analysis", "Cohort A vs cohort B: p<0.001, n=42"),
                ("Future work", "Plan to extend in 6 months"),
            ],
        )
        questions = anticipate_qa(out, n_questions=10)
        # At least one question anchored to slide 1 (stats claim)
        stats_qs = [q for q in questions if q.anchor_slide_index == 1]
        assert stats_qs, f"Expected a question anchored to slide 1 (stats). Got: {questions}"

    def test_comparison_triggers_question(self, tmp_path: Path) -> None:
        """Slide with 'compared to' / 'vs' → comparison-flavored question."""
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Method A vs method B", "Our approach compared to the baseline shows lift"),
            ],
        )
        questions = anticipate_qa(out, n_questions=10)
        cmp_qs = [q for q in questions if q.anchor_slide_index == 1]
        assert cmp_qs, f"Expected a comparison question. Got: {questions}"
        # Confirm at least one of the comparison-anchored questions mentions
        # 'why' / 'comparison' / 'baseline' (template language).
        joined = " ".join(q.question.lower() for q in cmp_qs)
        assert any(token in joined for token in ("why", "compar", "baseline", "alternative"))

    def test_future_work_triggers_timeline_question(self, tmp_path: Path) -> None:
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Future work", "Next steps include scaling and extension"),
            ],
        )
        questions = anticipate_qa(out, n_questions=10)
        # At least one question anchored to the future-work slide.
        fw_qs = [q for q in questions if q.anchor_slide_index == 1]
        assert fw_qs
        joined = " ".join(q.question.lower() for q in fw_qs)
        assert any(token in joined for token in ("timeline", "when", "next", "schedule"))

    def test_limitations_triggers_plan_question(self, tmp_path: Path) -> None:
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Limitations", "Sample size is small; replicates limited to 3 donors"),
            ],
        )
        questions = anticipate_qa(out, n_questions=10)
        lim_qs = [q for q in questions if q.anchor_slide_index == 1]
        assert lim_qs
        joined = " ".join(q.question.lower() for q in lim_qs)
        assert any(token in joined for token in ("plan", "address", "mitigate"))


class TestQuestionShape:
    def test_returns_anticipated_question_dataclass(self, tmp_path: Path) -> None:
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Stats slide", "p<0.05 n=10"),
            ],
        )
        questions = anticipate_qa(out, n_questions=5)
        assert questions, "Expected at least one question"
        q = questions[0]
        assert isinstance(q, AnticipatedQuestion)
        assert q.question
        assert q.why_likely
        assert 0.0 <= q.confidence <= 1.0
        assert q.anchor_slide_index is not None

    def test_n_questions_caps_output(self, tmp_path: Path) -> None:
        """Even when many triggers fire, output respects n_questions."""
        # Build a deck with multiple statistical / comparison / future / limitation slides
        spec = [("Title", "Body")]
        for i in range(10):
            spec.append((f"Stats slide {i}", "p<0.05 n=42 compared to baseline"))
        out = _make_deck(tmp_path, spec)
        questions = anticipate_qa(out, n_questions=3)
        assert len(questions) <= 3


class TestLLMMode:
    def test_runner_callback_is_called(self, tmp_path: Path) -> None:
        """When a runner_callback is supplied, it's invoked and its return used."""
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Some slide", "Some claims"),
            ],
        )
        calls = []

        def fake_runner(prompt: str) -> str:
            calls.append(prompt)
            # JSON-shaped response (the parser should accept JSON-like output)
            return (
                '[{"question": "How did you validate this?", '
                '"anchor_slide_index": 1, "why_likely": "stats not verified", '
                '"confidence": 0.8}]'
            )

        questions = anticipate_qa(out, runner_callback=fake_runner, n_questions=5)
        assert calls, "Expected the runner_callback to be invoked"
        assert questions, "Expected questions parsed from runner output"
        assert questions[0].question == "How did you validate this?"
        assert questions[0].anchor_slide_index == 1

    def test_runner_failure_falls_back_to_heuristics(self, tmp_path: Path) -> None:
        """If the callback raises, heuristic mode still produces output."""
        out = _make_deck(
            tmp_path,
            [
                ("Title", "Body"),
                ("Stats slide", "p<0.001 n=42"),
            ],
        )

        def bad_runner(prompt: str) -> str:
            raise RuntimeError("LLM down")

        questions = anticipate_qa(out, runner_callback=bad_runner, n_questions=10)
        # Heuristic fallback still surfaces the stats question
        assert questions
        assert any(q.anchor_slide_index == 1 for q in questions)


class TestErrors:
    def test_missing_pptx_raises(self, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            anticipate_qa(tmp_path / "nope.pptx")
