"""Tests for the lifted vaultlab.slides.animations OOXML engine.

Ported from ``bobby-tools/tests/test_bobby_slides/test_animation.py``.
"""
from __future__ import annotations

from pathlib import Path

import pytest

PIL = pytest.importorskip("PIL")
pptx = pytest.importorskip("pptx")
lxml = pytest.importorskip("lxml")

from PIL import Image
from pptx import Presentation as PptxPresentation

from vaultlab.slides import (
    add_figure_slide,
    add_text_slide,
    load_template,
)
from vaultlab.slides.animations import (
    appear_on_click,
    bullet_reveal,
    fade_on_click,
    panel_buildup,
)
from vaultlab.slides.template import lab_template_path

_P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


pytestmark = pytest.mark.skipif(
    lab_template_path() is None,
    reason="Hickey Lab template not bundled — animation tests need it",
)


def _save_and_reload(pres, tmp_path: Path, name: str):
    out = tmp_path / name
    pres.save(str(out))
    return PptxPresentation(str(out))


def _count_click_effects(slide) -> int:
    """Count the number of clickEffect nodes in a slide's timing tree."""
    timing = slide._element.find(f"{_P_NS}timing")
    if timing is None:
        return 0
    count = 0
    for ctn in timing.iter(f"{_P_NS}cTn"):
        if ctn.get("nodeType") == "clickEffect":
            count += 1
    return count


def _has_timing(slide) -> bool:
    return slide._element.find(f"{_P_NS}timing") is not None


class TestAppearOnClick:
    def test_adds_timing_to_slide(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["A"])
        shape = slide.shapes[-1]
        appear_on_click(slide, shape, click_index=0)
        assert _has_timing(slide)
        assert _count_click_effects(slide) >= 1

    def test_multiple_clicks_one_slide(self, tmp_path):
        pres = load_template()
        img = tmp_path / "i.png"
        Image.new("RGB", (100, 100), "red").save(img)
        slide = add_figure_slide(pres, img, title="Multi", bullets=["a", "b"])
        appear_on_click(slide, slide.shapes[0], click_index=0)
        appear_on_click(slide, slide.shapes[1], click_index=1)
        assert _count_click_effects(slide) == 2

    def test_pptx_reloads_without_error(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["A"])
        appear_on_click(slide, slide.shapes[-1], click_index=0)
        reloaded = _save_and_reload(pres, tmp_path, "appear.pptx")
        assert len(reloaded.slides) > 0


class TestFadeOnClick:
    def test_adds_fade_effect(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["A"])
        fade_on_click(slide, slide.shapes[-1], click_index=0)
        assert _count_click_effects(slide) == 1

    def test_uses_fade_preset(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["A"])
        fade_on_click(slide, slide.shapes[-1], click_index=0)
        timing = slide._element.find(f"{_P_NS}timing")
        for ctn in timing.iter(f"{_P_NS}cTn"):
            if ctn.get("nodeType") == "clickEffect":
                assert ctn.get("presetID") == "10"
                break

    def test_pptx_reloads_without_error(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["A"])
        fade_on_click(slide, slide.shapes[-1], click_index=0)
        reloaded = _save_and_reload(pres, tmp_path, "fade.pptx")
        assert len(reloaded.slides) > 0


class TestBulletReveal:
    def test_three_bullets_three_clicks(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["one", "two", "three"])
        bullet_reveal(slide, slide.shapes[-1].text_frame)
        assert _count_click_effects(slide) == 3

    def test_pptx_reloads_without_error(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "T", ["one", "two", "three"])
        bullet_reveal(slide, slide.shapes[-1].text_frame)
        reloaded = _save_and_reload(pres, tmp_path, "bullets.pptx")
        assert len(reloaded.slides) > 0

    def test_no_bullets_does_not_crash(self, tmp_path):
        pres = load_template()
        slide = add_text_slide(pres, "Title only", [])
        bullet_reveal(slide, slide.shapes[-1].text_frame)
        _save_and_reload(pres, tmp_path, "nobullets.pptx")


class TestPanelBuildup:
    def test_three_panels_three_clicks(self, tmp_path):
        pres = load_template()
        img = tmp_path / "p.png"
        Image.new("RGB", (100, 100), "blue").save(img)
        slide = add_figure_slide(pres, img, title="Panel A")
        from pptx.util import Inches
        s2 = slide.shapes.add_picture(str(img), Inches(2), Inches(2), Inches(1), Inches(1))
        s3 = slide.shapes.add_picture(str(img), Inches(4), Inches(2), Inches(1), Inches(1))
        panel_buildup(slide, [slide.shapes[1], s2, s3])
        assert _count_click_effects(slide) == 3

    def test_pptx_reloads_without_error(self, tmp_path):
        pres = load_template()
        img = tmp_path / "p.png"
        Image.new("RGB", (100, 100), "blue").save(img)
        slide = add_figure_slide(pres, img, title="P")
        from pptx.util import Inches
        s2 = slide.shapes.add_picture(str(img), Inches(2), Inches(2), Inches(1), Inches(1))
        panel_buildup(slide, [slide.shapes[1], s2])
        reloaded = _save_and_reload(pres, tmp_path, "buildup.pptx")
        assert len(reloaded.slides) > 0


class TestTimingIdempotence:
    def test_calling_twice_doesnt_break(self, tmp_path):
        pres = load_template()
        img = tmp_path / "i.png"
        Image.new("RGB", (100, 100), "red").save(img)
        slide = add_figure_slide(pres, img, title="Title")
        appear_on_click(slide, slide.shapes[0], 0)
        appear_on_click(slide, slide.shapes[1], 1)
        reloaded = _save_and_reload(pres, tmp_path, "idem.pptx")
        assert len(reloaded.slides) > 0
        assert _count_click_effects(reloaded.slides[-1]) == 2
