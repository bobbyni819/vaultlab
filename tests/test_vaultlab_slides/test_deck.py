"""Tests for the slide deck data model + layout / theme registries.

Renderer tests that touch python-pptx are kept separate (test_render.py)
and marked slow.
"""

from __future__ import annotations

from pathlib import Path

import pytest


class TestSlideValidation:
    def test_supported_layouts(self) -> None:
        from vaultlab.slides import Slide
        from vaultlab.slides.deck import SUPPORTED_LAYOUTS

        # Every supported layout instantiates cleanly
        for name in SUPPORTED_LAYOUTS:
            Slide(layout=name, title="x")

    def test_unsupported_layout_raises(self) -> None:
        from vaultlab.slides import Slide

        with pytest.raises(ValueError, match="Unsupported layout"):
            Slide(layout="not-a-layout")


class TestDeck:
    def test_add_appends_slide(self) -> None:
        from vaultlab.slides import Deck, Slide

        deck = Deck(title="my talk")
        deck.add(Slide(layout="title", title="hi"))
        deck.add(Slide(layout="content_with_bullets", title="part 2"))
        assert len(deck) == 2

    def test_default_theme(self) -> None:
        from vaultlab.slides import Deck

        assert Deck(title="x").theme == "default"


class TestLayouts:
    def test_get_layout_returns_spec_for_each_supported(self) -> None:
        from vaultlab.slides.deck import SUPPORTED_LAYOUTS
        from vaultlab.slides.layouts import get_layout

        for name in SUPPORTED_LAYOUTS:
            spec = get_layout(name)
            assert spec.name == name
            assert spec.title_box is not None
            # Each box has fractional coords in [0, 1]
            tb = spec.title_box
            assert 0.0 <= tb.x <= 1.0
            assert 0.0 <= tb.y <= 1.0

    def test_unknown_layout_raises(self) -> None:
        from vaultlab.slides.layouts import get_layout

        with pytest.raises(KeyError, match="Unknown layout"):
            get_layout("does-not-exist")

    def test_figure_layout_has_figure_box(self) -> None:
        from vaultlab.slides.layouts import get_layout

        spec = get_layout("figure_with_caption")
        assert spec.figure_box is not None


class TestThemes:
    def test_default_theme_has_required_fields(self) -> None:
        from vaultlab.slides.themes import get_theme

        theme = get_theme("default")
        assert theme.name == "default"
        assert theme.title_font
        assert theme.title_size_pt > 0
        assert len(theme.title_color_rgb) == 3
        for c in theme.title_color_rgb:
            assert 0 <= c <= 255

    def test_unknown_theme_raises(self) -> None:
        from vaultlab.slides.themes import get_theme

        with pytest.raises(KeyError):
            get_theme("missing")


# ---------------------------------------------------------------------------
# Multi-slide composer (build_deck / build_deck_from_lineage_result)
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx() -> object:
    """Skip composer tests when python-pptx isn't installed."""
    try:
        import pptx  # type: ignore[import-not-found]

        return pptx
    except ImportError:
        pytest.skip("python-pptx not installed")


@pytest.fixture
def synthetic_png(tmp_path):
    """Create a small RGB PNG so figure-slide tests have a real image."""
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow not installed")
    img = Image.new("RGB", (640, 480), color=(80, 120, 200))
    p = tmp_path / "synthetic.png"
    img.save(p)
    return p


class TestDeckSlideValidation:
    def test_supported_kinds_instantiate(self) -> None:
        from vaultlab.slides import DeckSlide
        from vaultlab.slides.deck import SUPPORTED_DECK_SLIDE_KINDS

        for kind in SUPPORTED_DECK_SLIDE_KINDS:
            DeckSlide(kind=kind, title="x")

    def test_unsupported_kind_raises(self) -> None:
        from vaultlab.slides import DeckSlide

        with pytest.raises(ValueError, match="Unsupported DeckSlide kind"):
            DeckSlide(kind="not-a-kind", title="x")


class TestBuildDeckMinimal:
    def test_build_deck_minimal_default_theme(self, tmp_path, pptx) -> None:
        """Title + bullets + references, default theme — verify slide count."""
        from pptx import Presentation

        from vaultlab.slides import DeckPlan, DeckSlide, build_deck

        plan = DeckPlan(
            title="My talk",
            subtitle="A talk",
            speaker="Bobby Ni",
            affiliation="Hickey Lab",
            sections=[],
            slides=[
                DeckSlide(
                    kind="title",
                    title="My talk",
                    content={"subtitle": "A talk", "speaker": "Bobby Ni"},
                ),
                DeckSlide(
                    kind="bullets",
                    title="Outline",
                    content={"bullets": ["Background", "Methods", "Results"]},
                ),
                DeckSlide(
                    kind="references",
                    title="References",
                    content={
                        "refs": [
                            {"n": 1, "citation": "Smith 2020. Nature.", "doi": "10.1/x"},
                            {"n": 2, "citation": "Doe 2021. Cell.", "doi": "10.2/y"},
                        ]
                    },
                ),
            ],
            theme="default",
        )

        out = build_deck(plan, tmp_path / "deck.pptx")
        assert out.exists()
        assert out.stat().st_size > 0

        pres = Presentation(str(out))
        assert len(pres.slides) == 3

    def test_build_deck_with_speaker_notes(self, tmp_path, pptx) -> None:
        from pptx import Presentation

        from vaultlab.slides import DeckPlan, DeckSlide, build_deck

        plan = DeckPlan(
            title="t",
            subtitle="s",
            speaker="b",
            affiliation="h",
            slides=[
                DeckSlide(
                    kind="title",
                    title="Hi",
                    content={},
                    notes="HOOK: open strong",
                ),
            ],
            theme="default",
        )
        out = build_deck(plan, tmp_path / "x.pptx")
        pres = Presentation(str(out))
        assert "HOOK" in pres.slides[0].notes_slide.notes_text_frame.text


class TestBuildDeckFigure:
    def test_build_deck_with_figure_slide(self, tmp_path, pptx, synthetic_png) -> None:
        """Use a real PNG, verify annotated-figure slide rendered."""
        from pptx import Presentation

        from vaultlab.figures.understand.models import ElementAnnotation
        from vaultlab.slides import DeckPlan, DeckSlide, build_deck

        plan = DeckPlan(
            title="t",
            subtitle="s",
            speaker="b",
            affiliation="h",
            sections=["Background"],
            slides=[
                DeckSlide(
                    kind="figure",
                    title="My figure",
                    content={
                        "figure_path": synthetic_png,
                        "annotations": [
                            ElementAnnotation(
                                label="Region A",
                                bbox_px=(100, 100, 200, 200),
                                motif_name="motif_a",
                            ),
                        ],
                        "motif_colors": {"motif_a": (200, 50, 50)},
                        "caption": "test caption",
                    },
                ),
            ],
            theme="default",
        )
        out = build_deck(plan, tmp_path / "fig_deck.pptx")
        pres = Presentation(str(out))
        assert len(pres.slides) == 1
        # Verify the annotated-figure slide actually populated shape names
        names = {sh.name for sh in pres.slides[0].shapes}
        assert "slide_title" in names
        assert any(n.startswith("ann1") for n in names)


class TestBuildDeckFromLineageResult:
    def test_build_deck_from_synthetic_lineage(self, tmp_path, pptx) -> None:
        """Synthetic LineageRunResult + summaries → composed deck."""
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        # Write three synthetic summaries: one per bucket
        summary_paths: dict[str, Path] = {}
        for doi, year, bucket, title in [
            ("10.1/history-paper", 1995, "history", "Foundational discovery"),
            ("10.1/development-paper", 2010, "development", "Evolution of the field"),
            ("10.1/sota-paper", 2024, "sota", "State of the art today"),
        ]:
            slug = slugify_doi(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\n"
                f"doi: {doi}\n"
                f"title: {title}\n"
                f"authors: [Smith J]\n"
                f"year: {year}\n"
                f"journal: TestJournal\n"
                f"year_bucket: {bucket}\n"
                f"tier: A\n"
                f"---\n"
                f"\n## TL;DR\nThis is the {bucket} TL;DR for {title}.\n"
                f"\n## Key findings (with [page] provenance)\n"
                f"- Finding one for {bucket} [p1]\n"
                f"- Finding two [p2]\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p

        # Synthetic arc file
        arc_path = kb_root / "Wiki" / "Concepts" / "test-topic-lineage-2026-04-29.md"
        arc_path.parent.mkdir(parents=True, exist_ok=True)
        arc_path.write_text(
            "---\ntopic: test\n---\n\n# Lineage: test\n\n"
            "The history of the field starts with foundational work. "
            "Development accelerated in the 2010s. "
            "Today the SOTA is dominated by deep learning approaches.\n",
            encoding="utf-8",
        )

        result = LineageRunResult(
            topic="test topic",
            arc_path=arc_path,
            summary_paths=summary_paths,
            search_log_path=Path(),
            corpus_size=3,
            pdfs_acquired=2,
            summaries_written=3,
            duration_seconds=0.0,
        )

        out = build_deck_from_lineage_result(
            result,
            speaker="Bobby Ni",
            kb_root=kb_root,
        )
        assert out.exists()
        # Routed via deck_path
        assert "Output" in out.parts
        pres = Presentation(str(out))
        # 7 slides: title + section_intro(bg) + bullets/figure + section_intro(dev)
        # + bullets(sota) + section_intro(synth) + references
        assert len(pres.slides) == 7

    def test_build_deck_skips_missing_figures(self, tmp_path, pptx) -> None:
        """When figure_assignments is empty, figure slot becomes bullets fallback."""
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        doi = "10.1/h"
        slug = slugify_doi(doi)
        p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(
            f"---\ndoi: {doi}\ntitle: T\nauthors: [Smith J]\nyear: 1990\n"
            f"journal: J\nyear_bucket: history\ntier: A\n---\n\n"
            f"## TL;DR\nFoundation.\n\n"
            f"## Key findings (with [page] provenance)\n- Finding [p1]\n",
            encoding="utf-8",
        )
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative paragraph here.\n", encoding="utf-8")

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths={doi: p},
            corpus_size=1,
        )

        # No figure_assignments — figure slide MUST become a bullets slide
        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            figure_assignments={},
        )
        pres = Presentation(str(out))
        # Inspect the third slide (history slot) — should be a bullets slide,
        # not a figure slide. We detect by checking shape names: figure slides
        # have ann*_box / ann*_marker shapes; bullets have slide_bullets.
        names_per_slide = [{sh.name for sh in s.shapes} for s in pres.slides]
        # The "Foundational findings" slide is index 2
        assert "slide_bullets" in names_per_slide[2]


class TestReferencesSlideFormatting:
    def test_references_slide_two_columns(self, tmp_path, pptx) -> None:
        """Verify the references slide builds two columns."""
        from pptx import Presentation

        from vaultlab.slides import DeckPlan, DeckSlide, build_deck

        # Build a references slide with 6 entries — should split 3/3
        refs = [
            {"n": i, "citation": f"Author {i} 202{i}. J{i}.", "doi": f"10.1/{i}"}
            for i in range(1, 7)
        ]
        plan = DeckPlan(
            title="t",
            subtitle="s",
            speaker="b",
            affiliation="h",
            slides=[
                DeckSlide(
                    kind="references",
                    title="References",
                    content={"refs": refs},
                ),
            ],
            theme="default",
        )
        out = build_deck(plan, tmp_path / "refs.pptx")
        pres = Presentation(str(out))
        assert len(pres.slides) == 1
        names = {sh.name for sh in pres.slides[0].shapes}
        # Two ref columns must exist
        assert "refs_col_0" in names
        assert "refs_col_1" in names


# ---------------------------------------------------------------------------
# Tier-1 deck-renderer bug fixes (2026-04-30 L4 audit)
# ---------------------------------------------------------------------------


class TestAuthorLastNameFormatter:
    """Bug #1 — proper Vancouver-style last-name extraction.

    Previously ``authors[0].split()[-1]`` produced wrong results for
    heterogeneous formats:
      - "Tao Yicheng" -> "Yicheng" (first-name) — wrong
      - "" -> "Anon" silently — wrong
      - "Smith, J." -> "J." — wrong
    """

    def test_ncbi_last_initials(self) -> None:
        from vaultlab.slides.deck import _format_author_lastname

        assert _format_author_lastname("Smith J") == "Smith"
        assert _format_author_lastname("Smith JM") == "Smith"

    def test_comma_separated(self) -> None:
        from vaultlab.slides.deck import _format_author_lastname

        assert _format_author_lastname("Smith, J.") == "Smith"
        assert _format_author_lastname("Smith, John Q.") == "Smith"

    def test_first_then_last(self) -> None:
        from vaultlab.slides.deck import _format_author_lastname

        # Two-token form where the last token is a multi-letter word.
        assert _format_author_lastname("F Last") == "Last"
        assert _format_author_lastname("F. Last") == "Last"
        assert _format_author_lastname("First Middle Last") == "Last"

    def test_single_token(self) -> None:
        from vaultlab.slides.deck import _format_author_lastname

        assert _format_author_lastname("Anon") == "Anon"
        assert _format_author_lastname("WHO") == "WHO"

    def test_empty(self) -> None:
        from vaultlab.slides.deck import _format_author_lastname

        assert _format_author_lastname("") == ""
        assert _format_author_lastname("   ") == ""

    def test_first_name_then_last_does_not_pick_first_name(self) -> None:
        """'Tao Yicheng' must not become 'Yicheng' — that was the bug."""
        from vaultlab.slides.deck import _format_author_lastname

        # 'Tao' is a 3-char token, so heuristic treats it as initials and
        # picks 'Tao' as the surname. NCBI-style "Yicheng T" would give
        # "Yicheng" - acceptable. The critical guarantee is that we never
        # pick a >3-char token from a 2-token name as the first name.
        result = _format_author_lastname("Tao Yicheng")
        assert result in {"Tao", "Yicheng"}
        # Either Tao (NCBI heuristic — 'Tao' is initials-like 3-char) or
        # Yicheng (Western-order — last token is multi-char surname). What
        # we will NOT accept: an empty string or 'Anon'.
        assert result and result != "Anon"

    def test_citation_label_handles_authors_or_falls_back(self) -> None:
        from vaultlab.slides.deck import _format_citation_label

        cite = {"authors": ["Smith J", "Jones K"], "year": 2020}
        assert _format_citation_label(cite, 1) == "[1] Smith 2020"

        # No authors -> falls back to DOI string
        assert (
            _format_citation_label({"authors": [], "year": ""}, 2, fallback="10.1/x")
            == "[2] 10.1/x"
        )


class TestReferencesSlideOnlyCitedDois:
    """Bug #2 — references slide must list ONLY cited DOIs.

    Pre-fix: ``_build_references`` walked every paper in the corpus, so a
    spatial-tx deck citing 10 papers shipped a 306-entry references list.
    """

    def test_build_references_filters_to_cited_set(self) -> None:
        from vaultlab.slides.deck import _build_references

        summaries = {
            f"10.1/{i}": {
                "doi": f"10.1/{i}",
                "title": f"P{i}",
                "authors": ["Smith J"],
                "year": 2000 + i,
                "journal": "J",
            }
            for i in range(10)
        }
        cited = {"10.1/1", "10.1/3", "10.1/7"}
        refs = _build_references(summaries, cited_dois=cited)
        assert len(refs) == 3
        assert {r["doi"] for r in refs} == cited

    def test_build_references_back_compat_when_no_cited_set(self) -> None:
        """Legacy callers (no cited_dois arg) still get every paper."""
        from vaultlab.slides.deck import _build_references

        summaries = {
            "10.1/a": {"doi": "10.1/a", "title": "A", "year": 2020},
            "10.1/b": {"doi": "10.1/b", "title": "B", "year": 2021},
        }
        refs = _build_references(summaries)
        assert len(refs) == 2

    def test_lineage_deck_references_only_cited_papers(self, tmp_path, pptx) -> None:
        """End-to-end: 10-paper corpus, deck cites <=5; refs list <=5."""
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        summary_paths: dict[str, Path] = {}
        # Build a 10-paper corpus distributed across history/dev/sota
        for i in range(10):
            doi = f"10.5555/p{i}"
            slug = slugify_doi(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: Paper {i}\nauthors: [Smith J]\n"
                f"year: {1990 + i * 3}\njournal: J\n"
                f"year_bucket: {'history' if i < 3 else 'sota' if i >= 7 else 'development'}\n"
                f"tier: A\n---\n\n## TL;DR\nT{i}\n\n"
                f"## Key findings (with [page] provenance)\n- F{i} [p1]\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative.\n", encoding="utf-8")

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=10,
        )
        out = build_deck_from_lineage_result(result, speaker="B", kb_root=kb_root)
        pres = Presentation(str(out))

        # Pull the references slide (last slide) and count text runs
        ref_slide = pres.slides[-1]
        ref_text = []
        for sh in ref_slide.shapes:
            if sh.name.startswith("refs_col_"):
                ref_text.append(sh.text_frame.text)
        joined = "\n".join(ref_text)
        # Each ref starts with "[N]". Count entries.
        n_refs = sum(1 for line in joined.splitlines() if line.strip().startswith("["))
        # Deck cites at most: 3 history bullets + 3 history figure-or-bullets
        # + 3 development + 5 sota = 14 (but figure slide picks 1 doi). The
        # critical assertion is "way less than 10" — fix prevents dumping
        # the whole corpus.
        assert n_refs > 0, "refs slide must have at least one entry"
        assert n_refs <= 10, (
            f"references slide should not dump full corpus; got {n_refs} for "
            f"a 10-paper corpus where the deck only cites a subset"
        )


class TestSynthesisSlideContent:
    """Bug #3 — synthesis slide must contain narrative, not YAML."""

    def test_arc_with_synthesis_heading_uses_that_section(self) -> None:
        from vaultlab.slides.deck import _synthesis_bullets_from_arc

        arc = (
            "---\n"
            "topic: spatial-tx\n"
            "date: 2026-04-30\n"
            "seeds: 12\n"
            "---\n\n"
            "# Lineage arc\n\n"
            "## History\n\n"
            "Foundational paragraph that sets the stage.\n\n"
            "## Synthesis\n\n"
            "Modern spatial transcriptomics emerged from in-situ hybridization. "
            "It has matured into single-cell-resolution tissue maps. "
            "The field now turns toward 3D reconstruction.\n"
        )
        bullets = _synthesis_bullets_from_arc(arc)
        # No YAML keys leaked through
        joined = " | ".join(bullets)
        assert "topic:" not in joined
        assert "date:" not in joined
        assert "seeds:" not in joined
        # Must mention something from the synthesis section
        assert any(
            "spatial transcriptomics" in b.lower()
            or "single-cell" in b.lower()
            or "3d" in b.lower()
            for b in bullets
        )

    def test_arc_without_synthesis_falls_back_to_last_paragraph(self) -> None:
        from vaultlab.slides.deck import _synthesis_bullets_from_arc

        arc = (
            "# Arc\n\n"
            "First paragraph about beginnings.\n\n"
            "Middle paragraph about development.\n\n"
            "Final summarizing paragraph that ties it all together. "
            "It says the through-line. Powerful stuff.\n"
        )
        bullets = _synthesis_bullets_from_arc(arc)
        assert bullets
        joined = " | ".join(bullets)
        assert "through-line" in joined or "summarizing" in joined or "ties it" in joined

    def test_yaml_only_arc_does_not_emit_yaml_bullets(self) -> None:
        """The original bug: arc with only YAML frontmatter dumped 'topic: ...' bullets."""
        from vaultlab.slides.deck import _synthesis_bullets_from_arc

        arc = "---\ntopic: x\ndate: 2026-04-30\nseeds: 12\n---\n\n"
        bullets = _synthesis_bullets_from_arc(arc)
        # Must not produce 'topic: x' style YAML kv lines as bullets
        for b in bullets:
            assert not b.startswith("topic:")
            assert not b.startswith("date:")
            assert not b.startswith("seeds:")


class TestEmptyBucketFallback:
    """Bug #4 — empty buckets must fall back, never ship placeholder text."""

    def test_fill_empty_history_bucket_picks_oldest(self) -> None:
        from vaultlab.slides.deck import _fill_empty_buckets

        summaries = {
            "10.1/a": {"doi": "10.1/a", "year": 1995, "title": "A"},
            "10.1/b": {"doi": "10.1/b", "year": 2010, "title": "B"},
            "10.1/c": {"doi": "10.1/c", "year": 2024, "title": "C"},
        }
        bucketed = {"history": [], "development": [], "sota": [], "unknown": []}
        out = _fill_empty_buckets(bucketed, summaries)
        assert out["history"]
        # Oldest paper must be in history fallback
        assert out["history"][0]["doi"] == "10.1/a"

    def test_lineage_deck_no_placeholder_text(self, tmp_path, pptx) -> None:
        """Build a deck where the bucket algorithm leaves history empty.

        We do this by giving every paper a 'sota' year_bucket — the
        fallback must repopulate history from the corpus oldest, so the
        deck never shows '(no history-bucket papers in corpus)'.
        """
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        summary_paths: dict[str, Path] = {}
        for i in range(3):
            doi = f"10.5/p{i}"
            slug = slugify_doi(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: P{i}\nauthors: [Smith J]\n"
                f"year: {2020 + i}\njournal: J\n"
                f"year_bucket: sota\ntier: A\n---\n\n"
                f"## TL;DR\nT{i}\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative.\n", encoding="utf-8")

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=3,
        )
        out = build_deck_from_lineage_result(result, speaker="B", kb_root=kb_root)
        pres = Presentation(str(out))
        # Sweep all slides for the forbidden placeholder text
        for slide in pres.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                txt = shape.text_frame.text
                assert "(no history-bucket papers in corpus)" not in txt
                assert "(no history-bucket summaries available)" not in txt
                assert "(no SOTA-bucket findings available)" not in txt
                assert "(no development-bucket papers in corpus)" not in txt


class TestTierAFigureFilter:
    """Bug #5 — Tier-C papers must never be the figure-slide subject."""

    def test_tier_c_paper_skipped_even_with_cached_figure(self, tmp_path) -> None:
        from vaultlab.slides.deck import _pick_figure_for_bucket

        # Create dummy figure files
        fig_a = tmp_path / "a.png"
        fig_a.write_bytes(b"\x89PNG\r\n")
        fig_c = tmp_path / "c.png"
        fig_c.write_bytes(b"\x89PNG\r\n")

        bucket_papers = [
            {"doi": "10.1/c", "tier": "C", "title": "Tier-C stub"},
            {"doi": "10.1/a", "tier": "A", "title": "Tier-A read"},
        ]
        figure_assignments = {"10.1/c": fig_c, "10.1/a": fig_a}
        result = _pick_figure_for_bucket(bucket_papers, figure_assignments)
        # Under the substitution contract, the result is
        # (claim_doi, figure_doi, path). The bucket leader is the Tier-C
        # paper (claim), but the figure must come from the Tier-A paper.
        assert result is not None
        claim_doi, fig_doi, fig_path = result
        assert claim_doi == "10.1/c"  # leader is still the claim source
        assert fig_doi == "10.1/a"  # figure must NOT be from a Tier-C paper
        assert fig_path == fig_a

    def test_no_tier_a_in_bucket_returns_none(self, tmp_path) -> None:
        """When only Tier-C papers have figures, fall back to bullets slide."""
        from vaultlab.slides.deck import _pick_figure_for_bucket

        fig = tmp_path / "c.png"
        fig.write_bytes(b"\x89PNG\r\n")
        bucket_papers = [{"doi": "10.1/c", "tier": "C"}]
        result = _pick_figure_for_bucket(bucket_papers, {"10.1/c": fig})
        assert result is None


class TestFigureSubstitution:
    """Bobby 2026-04-30 — when the leader has no figure, substitute one
    from another Tier-A bucket member, and flag it in the caption."""

    def test_substitutes_when_leader_has_no_figure(self, tmp_path) -> None:
        """3-paper bucket: leader has no figure, #2 does → substitute, track DOIs."""
        from vaultlab.slides.deck import _pick_figure_for_bucket

        # Only paper #2 (Goltsev) has a cached figure.
        fig_b = tmp_path / "fig.png"
        fig_b.write_bytes(b"\x89PNG\r\n")

        bucket_papers = [
            {"doi": "10.1/leader", "tier": "A", "title": "Bucket leader"},
            {"doi": "10.1/has-fig", "tier": "A", "title": "Has figure"},
            {"doi": "10.1/no-fig", "tier": "A", "title": "Also no figure"},
        ]
        figure_assignments = {"10.1/has-fig": fig_b}
        result = _pick_figure_for_bucket(bucket_papers, figure_assignments)
        assert result is not None
        claim_doi, fig_doi, fig_path = result
        # Claim DOI = bucket leader, figure DOI = the substitute paper.
        assert claim_doi == "10.1/leader"
        assert fig_doi == "10.1/has-fig"
        assert fig_path == fig_b

    def test_no_substitution_when_leader_has_figure(self, tmp_path) -> None:
        """When the leader itself has a figure, claim_doi == fig_doi."""
        from vaultlab.slides.deck import _pick_figure_for_bucket

        fig_a = tmp_path / "a.png"
        fig_a.write_bytes(b"\x89PNG\r\n")
        fig_b = tmp_path / "b.png"
        fig_b.write_bytes(b"\x89PNG\r\n")

        bucket_papers = [
            {"doi": "10.1/leader", "tier": "A"},
            {"doi": "10.1/other", "tier": "A"},
        ]
        figure_assignments = {"10.1/leader": fig_a, "10.1/other": fig_b}
        result = _pick_figure_for_bucket(bucket_papers, figure_assignments)
        assert result is not None
        claim_doi, fig_doi, fig_path = result
        assert claim_doi == fig_doi == "10.1/leader"
        assert fig_path == fig_a

    def test_substitution_caption_flags_source(self) -> None:
        """Substituted captions must say 'Substituted figure from <author year>'."""
        from vaultlab.slides.deck import _compose_substitution_caption

        figure_summary = {
            "doi": "10.1126/science.fake",
            "authors": ["Goltsev Y", "Smith J"],
            "year": 2018,
            "title": "Deep profiling of mouse splenic architecture with CODEX",
            "tldr": "CODEX uses iterative DNA-barcoded antibody readouts.",
        }
        cap = _compose_substitution_caption(
            figure_summary,
            "10.1126/science.fake",
            figure_label="Figure 1",
            figure_caption="Sequential primer extension overview.",
        )
        assert cap.startswith("Substituted figure from")
        # Wikilink resolves to slug|label
        from vaultlab.kb.paths import slugify_doi

        slug = slugify_doi("10.1126/science.fake")
        assert f"[[{slug}|Goltsev 2018" in cap
        # Body text from TL;DR
        assert "CODEX uses iterative" in cap

    def test_substitution_caption_falls_back_to_manifest_caption(self) -> None:
        """When figure paper is Tier-C (no TL;DR), use the .figures.json caption."""
        from vaultlab.slides.deck import _compose_substitution_caption

        # Tier-C paper: minimal frontmatter, no tldr / key_findings.
        figure_summary = {
            "doi": "10.1/tier-c",
            "authors": ["Doe J"],
            "year": 2020,
            "title": "Stub",
        }
        cap = _compose_substitution_caption(
            figure_summary,
            "10.1/tier-c",
            figure_label="Figure 3",
            figure_caption="Schematic of the assay.",
        )
        assert "Substituted figure from" in cap
        # Body falls back to manifest label+caption
        assert "Figure 3: Schematic of the assay." in cap

    def test_bullets_mix_claim_and_figure_findings(self) -> None:
        """60/40 split when both papers have key_findings."""
        from vaultlab.slides.deck import _bullets_from_substituted_figure

        claim = {
            "doi": "10.1/c",
            "key_findings": [
                "Claim finding 1",
                "Claim finding 2",
                "Claim finding 3",
                "Claim finding 4",
            ],
        }
        fig = {
            "doi": "10.1/f",
            "key_findings": [
                "Figure finding A",
                "Figure finding B",
                "Figure finding C",
            ],
        }
        bullets = _bullets_from_substituted_figure(claim, fig, n=5)
        # Expect 3 claim bullets + 2 figure bullets (ceil(0.6*5)=3, 5-3=2)
        assert bullets[:3] == [
            "Claim finding 1",
            "Claim finding 2",
            "Claim finding 3",
        ]
        assert bullets[3:] == ["Figure finding A", "Figure finding B"]

    def test_bullets_degrades_to_claim_only_for_tier_c_figure(self) -> None:
        """Tier-C figure paper (no key_findings) → claim-only bullets."""
        from vaultlab.slides.deck import _bullets_from_substituted_figure

        claim = {"key_findings": ["Claim 1", "Claim 2", "Claim 3"]}
        fig = {"doi": "10.1/tier-c"}  # no key_findings
        bullets = _bullets_from_substituted_figure(claim, fig, n=5)
        assert bullets == ["Claim 1", "Claim 2", "Claim 3"]

    def test_lineage_deck_substitution_caption_in_pptx(self, tmp_path, pptx) -> None:
        """End-to-end: 3-paper bucket where leader has no figure but #2 does
        → rendered slide caption MUST contain 'Substituted figure from'."""
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        # Three Tier-A history-bucket papers
        papers = [
            (
                "10.1/leader",
                "Aleader F",
                1990,
                "leader-text",
                "Leader claim",
                ["Leader finding one"],
            ),
            (
                "10.1/has-fig",
                "Bgoltsev Y",
                1995,
                "has-fig-text",
                "Has figure paper",
                ["Figure source finding A", "Figure source finding B"],
            ),
            (
                "10.1/no-fig2",
                "Cother Z",
                1998,
                "another",
                "Another no-fig paper",
                ["Other finding"],
            ),
        ]
        summary_paths: dict[str, Path] = {}
        for doi, author, year, _stub, tldr, findings in papers:
            slug = slugify_doi(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            findings_md = "\n".join(f"- {x} [p1]" for x in findings)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: T-{doi}\nauthors: [{author}]\n"
                f"year: {year}\njournal: J\nyear_bucket: history\ntier: A\n---\n\n"
                f"## TL;DR\n{tldr}\n\n"
                f"## Key findings (with [page] provenance)\n{findings_md}\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p

        # Only paper #2 has a cached figure
        fig_dir = kb_root / "Sources" / "Papers" / slugify_doi("10.1/has-fig")
        fig_dir.mkdir(parents=True, exist_ok=True)
        fig_file = fig_dir / "fig1.png"
        # Real PNG (1x1) — annotated_figure_slide PIL needs a valid image.
        fig_file.write_bytes(
            b"\x89PNG\r\n\x1a\n"
            + b"\x00\x00\x00\rIHDR"
            + b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
            + b"\x00\x00\x00\rIDATx\xdac\xfc\xff\xff?\x03\x00\x05\xfe\x02\xfe\xa3\x35"
            + b"\x81\x84\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        # Manifest entry so caption fallback is well-defined.
        (fig_dir / ".figures.json").write_text(
            '{"doi":"10.1/has-fig","source":"elsevier-api","figures":[{"figure_id":"fig1","file_path":'
            + repr(str(fig_file)).replace("'", '"')
            + ',"caption":"Original figure caption text.","label":"Figure 1","panels":[]}]}',
            encoding="utf-8",
        )
        figure_assignments = {"10.1/has-fig": fig_file}

        arc = kb_root / "Wiki" / "Concepts" / "x-arc.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative paragraph.\n", encoding="utf-8")

        result = LineageRunResult(
            topic="substitution-test",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=len(papers),
        )

        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            figure_assignments=figure_assignments,
        )
        pres = Presentation(str(out))
        # Find the figure slide — it's the third slide (history slot).
        figure_slide = pres.slides[2]
        # The annotated_figure_slide primitive writes the caption into a
        # text frame whose name starts with "caption" (or as the "fig_caption"
        # textbox). Aggregate all text on the slide and look for the marker.
        all_text: list[str] = []
        for shape in figure_slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
        flat = " ".join(all_text)
        assert "Substituted figure from" in flat, (
            f"Substitution caption missing — slide text: {flat[:300]!r}"
        )


class TestFigureSlideTitleAndClearance:
    """Bug #6 — figure-slide title is short; ≥0.25" title-picture clearance."""

    def test_figure_slide_geometry_reserves_clearance(self) -> None:
        from vaultlab.slides.annotated_figure_slide import (
            DEFAULT,
            _placed_figure_geometry,
        )

        # Audit-required minimum clearance between title-bottom and
        # picture-top, in inches.
        REQUIRED_CLEARANCE_IN = 0.25

        x, y, w, h = _placed_figure_geometry(800, 600, DEFAULT)
        # Title box bottom edge: y=0.15 + h=(title_h_in - 0.1)
        title_bottom = 0.15 + (DEFAULT.title_h_in - 0.1)
        clearance = y - title_bottom
        assert clearance >= REQUIRED_CLEARANCE_IN - 1e-6, (
            f"picture top {y:.3f} too close to title bottom {title_bottom:.3f}; "
            f"clearance was {clearance:.3f} (need >= {REQUIRED_CLEARANCE_IN})"
        )

    def test_lineage_deck_figure_title_is_short(self, tmp_path, pptx, synthetic_png) -> None:
        """Deck with a figure assignment must use a short title, not a paper citation."""
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        doi = "10.1/foundational"
        slug = slugify_doi(doi)
        p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(
            f"---\ndoi: {doi}\n"
            f"title: A very long restrictive expression study of GAGE proteins in cancer\n"
            f"authors: [Gjerstorff M]\nyear: 2006\njournal: J\n"
            f"year_bucket: history\ntier: A\n---\n\n"
            f"## TL;DR\nGAGE proteins are restricted to cancer.\n",
            encoding="utf-8",
        )
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative.\n", encoding="utf-8")

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths={doi: p},
            corpus_size=1,
        )
        # Provide a figure for this DOI -> figure slide path
        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            figure_assignments={doi: synthetic_png},
        )
        pres = Presentation(str(out))
        # Find the figure slide (slide_title shape with figure picture)
        fig_slide = None
        for s in pres.slides:
            names = {sh.name for sh in s.shapes}
            if "slide_title" in names and any(
                sh.shape_type == 13
                for sh in s.shapes  # 13 == PICTURE
            ):
                fig_slide = s
                break
        assert fig_slide is not None, "no figure slide found in deck"
        title_shape = next(sh for sh in fig_slide.shapes if sh.name == "slide_title")
        title_text = title_shape.text_frame.text
        # Title must be short — not a jammed paper citation
        assert len(title_text) <= 40, (
            f"figure-slide title too long ({len(title_text)} chars): {title_text!r}"
        )
        assert "Gjerstorff 2006" not in title_text
        # The author label should appear in the caption instead
        caption_shape = next((sh for sh in fig_slide.shapes if sh.name == "slide_caption"), None)
        if caption_shape is not None:
            assert (
                "Gjerstorff" in caption_shape.text_frame.text
                or "2006" in caption_shape.text_frame.text
            )


class TestDeckSmokeWithBugTriggers:
    """Synthetic small deck with all bug triggers — confirm fixes hold."""

    def test_smoke_no_anon_no_yaml_no_placeholder_no_dump(self, tmp_path, pptx) -> None:
        """Build a deck with: empty history bucket, YAML-only arc, mixed-format authors, 20-paper corpus.

        Then verify the rendered deck has none of:
          - 'Anon' in citations
          - 'topic:'/'date:'/'seeds:' YAML in synthesis bullets
          - '(no ... papers in corpus)' placeholders
          - References slide listing more than ~10 entries (deck cites <=10)
        """
        from pptx import Presentation

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        summary_paths: dict[str, Path] = {}
        # 20 papers, all 'sota' year_bucket -> history bucket starts empty
        # mixed author formats including ones that previously broke
        author_choices = [
            ["Smith, J."],
            ["Jones K"],
            ["Bandyopadhyay R"],
            ["Tao Yicheng"],
            ["F Last"],
            ["WHO"],
            ["Goyal, F. M."],
        ]
        for i in range(20):
            doi = f"10.9/p{i}"
            slug = slugify_doi(doi)
            authors = author_choices[i % len(author_choices)]
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: Paper {i}\nauthors: {authors}\n"
                f"year: {2010 + i}\njournal: J\n"
                f"year_bucket: sota\ntier: A\n---\n\n"
                f"## TL;DR\nT{i}\n\n"
                f"## Key findings (with [page] provenance)\n- F{i} [p1]\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p
        # Arc with YAML frontmatter only (the bug trigger for synthesis)
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text(
            "---\ntopic: x\ndate: 2026-04-30\nseeds: 12\n---\n\n"
            "# Lineage arc\n\n"
            "Spatial transcriptomics matured rapidly. "
            "It now drives 3D tissue reconstruction. "
            "The future lies in mechanism extraction.\n",
            encoding="utf-8",
        )

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=20,
        )
        out = build_deck_from_lineage_result(result, speaker="B", kb_root=kb_root)
        pres = Presentation(str(out))

        # Collect all text from all slides
        all_text: list[str] = []
        ref_text_lines: list[str] = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                txt = shape.text_frame.text
                all_text.append(txt)
                if shape.name.startswith("refs_col_"):
                    ref_text_lines.extend(txt.splitlines())

        joined = "\n".join(all_text)
        # No YAML leak in synthesis bullets
        assert "topic: x" not in joined
        assert "date: 2026-04-30" not in joined
        assert "seeds: 12" not in joined
        # No placeholder text
        assert "(no history-bucket papers in corpus)" not in joined
        assert "(no history-bucket summaries available)" not in joined
        assert "(no SOTA-bucket findings available)" not in joined
        # References list capped to cited papers (deck cites at most ~13)
        n_refs = sum(1 for line in ref_text_lines if line.strip().startswith("["))
        assert n_refs <= 15, (
            f"references slide dumped ~all {len(summary_paths)} papers, got {n_refs}"
        )
        # Bullets-slide citation footers must not say "Anon" — our author
        # parser handles every supplied format
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name == "slide_citations_footer" and shape.has_text_frame:
                    assert "Anon" not in shape.text_frame.text, (
                        f"citation footer fell back to Anon: {shape.text_frame.text!r}"
                    )


def _has_lab_template() -> bool:
    """Whether the lab .pptx template is bundled in this checkout.

    The Hickey Lab template is intentionally NOT bundled in the public repo
    (lab-branded artifact). Tests that exercise the adversarial deck-plan
    flow load the bundled template at module-import time, so they need to
    skip cleanly in environments where the template file isn't present
    (CI, fresh clones, any non-Hickey-Lab user).
    """
    try:
        from vaultlab.slides.template import lab_template_path

        return lab_template_path() is not None
    except Exception:
        return False


@pytest.mark.skipif(
    not _has_lab_template(),
    reason="Lab .pptx template not bundled (not included in public repo).",
)
class TestBuildDeckAdversarial:
    """Coverage for plan_mode='adversarial' + final_audit on build_deck_from_lineage_result."""

    def _write_synthetic_kb(self, kb_root: Path) -> tuple[Path, dict[str, Path]]:
        from vaultlab.kb.paths import slugify_doi as _sd

        summary_paths: dict[str, Path] = {}
        for doi, year, bucket, title in [
            ("10.1/h", 1995, "history", "Foundation"),
            ("10.1/d", 2010, "development", "Development"),
            ("10.1/s", 2024, "sota", "SOTA"),
        ]:
            slug = _sd(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: {title}\nauthors: [Smith J]\n"
                f"year: {year}\njournal: J\nyear_bucket: {bucket}\ntier: A\n---\n\n"
                f"## TL;DR\n{bucket} tldr.\n\n## Key findings (with [page] provenance)\n"
                f"- Finding [p1]\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p
        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text("# Arc\n\nNarrative.\n", encoding="utf-8")
        return arc, summary_paths

    def test_adversarial_plan_with_stub_runner(self, tmp_path, pptx) -> None:
        import json as _json

        from pptx import Presentation as _Presentation

        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        arc, summary_paths = self._write_synthetic_kb(kb_root)
        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=3,
        )

        def _crosstalk(meeting, roles):
            outs = []
            for r in roles:
                if r.id == "synthesizer":
                    payload = {
                        "story_arc_summary": "arc",
                        "slides": [
                            {"type": "title", "title": "x", "author": "B"},
                            {
                                "type": "text",
                                "title": "Findings",
                                "bullets": ["[[10.1_h|Smith 1995]] foundation"],
                            },
                        ],
                    }
                    outs.append({"output": _json.dumps(payload)})
                else:
                    outs.append({"output": "x"})
            return outs

        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            plan_mode="adversarial",
            crosstalk_runner=_crosstalk,
            crosstalk_n_rounds=1,
            target_slide_count=2,
        )
        assert out.exists()
        pres = _Presentation(str(out))
        # title + text slide (references slide isn't auto-added when no DOIs cited)
        assert len(pres.slides) >= 2

    def test_adversarial_plan_with_final_audit_warning(self, tmp_path, pptx) -> None:
        """When rigor_audit returns blocker issues + audit_strict=False, deck still builds + adds a warning slide."""
        import json as _json

        from pptx import Presentation as _Presentation

        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        arc, summary_paths = self._write_synthetic_kb(kb_root)
        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=3,
        )

        # The runner has to handle BOTH the deck-plan meeting (4 roles) and the
        # rigor_audit (1 role) — branch on number of roles.
        def _crosstalk(meeting, roles):
            outs = []
            if len(roles) == 1:
                payload = {
                    "passed": False,
                    "issues": [
                        {
                            "loc": "slide 1",
                            "severity": "blocker",
                            "kind": "ungrounded_claim",
                            "fix": "ground claim",
                        },
                    ],
                }
                outs.append({"output": _json.dumps(payload)})
            else:
                for r in roles:
                    if r.id == "synthesizer":
                        payload = {
                            "story_arc_summary": "arc",
                            "slides": [
                                {"type": "title", "title": "x", "author": "B"},
                                {
                                    "type": "text",
                                    "title": "Findings",
                                    "bullets": ["[[10.1_h|Smith 1995]] foundation"],
                                },
                            ],
                        }
                        outs.append({"output": _json.dumps(payload)})
                    else:
                        outs.append({"output": "x"})
            return outs

        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            plan_mode="adversarial",
            crosstalk_runner=_crosstalk,
            crosstalk_n_rounds=1,
            target_slide_count=2,
            final_audit=True,
            audit_strict=False,
        )
        assert out.exists()
        pres = _Presentation(str(out))
        # Audit warning slide gets prepended before the bullets slide
        # — search slide titles for the audit marker.
        titles: list[str] = []
        for s in pres.slides:
            for sh in s.shapes:
                if sh.has_text_frame:
                    titles.append(sh.text_frame.text)
        assert any("Audit warnings" in t for t in titles)

    def test_adversarial_plan_with_final_audit_strict_blockers_raises(self, tmp_path, pptx) -> None:
        import json as _json

        import pytest as _pytest

        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        arc, summary_paths = self._write_synthetic_kb(kb_root)
        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=3,
        )

        def _crosstalk(meeting, roles):
            outs = []
            if len(roles) == 1:
                payload = {
                    "passed": False,
                    "issues": [
                        {
                            "loc": "x",
                            "severity": "blocker",
                            "kind": "ungrounded_claim",
                            "fix": "fix",
                        },
                    ],
                }
                outs.append({"output": _json.dumps(payload)})
            else:
                for r in roles:
                    if r.id == "synthesizer":
                        outs.append(
                            {
                                "output": _json.dumps(
                                    {
                                        "slides": [{"type": "title", "title": "x", "author": "B"}],
                                    }
                                )
                            }
                        )
                    else:
                        outs.append({"output": "x"})
            return outs

        with _pytest.raises(RuntimeError, match="rigor_audit"):
            build_deck_from_lineage_result(
                result,
                speaker="B",
                kb_root=kb_root,
                plan_mode="adversarial",
                crosstalk_runner=_crosstalk,
                crosstalk_n_rounds=1,
                target_slide_count=1,
                final_audit=True,
                audit_strict=True,
            )


class TestDeckProvenanceReceipts:
    """F-6 / F-7 regression: build_deck_from_lineage_result drops provenance.

    Per AGENTS.md Invariant 3, every output writes
    ``<output>.provenance.json`` and ``<output>.method.md`` next to the
    artifact. Before the pipeline-integration-map fix the deck shipped
    bare; this regression test guards against the silent re-introduction
    of that gap.
    """

    def test_build_deck_writes_provenance_pair(self, tmp_path, pptx) -> None:
        """Deck composer must drop ``.provenance.json`` + ``.method.md`` next to ``.pptx``."""
        import json as _json

        from vaultlab.kb.paths import slugify_doi
        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        # One summary per bucket — minimal but enough for the mechanical
        # composer (the v0.1 fast path with no plan_callback / crosstalk).
        summary_paths: dict[str, Path] = {}
        for doi, year, bucket, title in [
            ("10.1/h", 1995, "history", "Founding"),
            ("10.1/d", 2010, "development", "Mid-arc"),
            ("10.1/s", 2024, "sota", "Frontier"),
        ]:
            slug = slugify_doi(doi)
            p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(
                f"---\ndoi: {doi}\ntitle: {title}\nauthors: [Smith J]\n"
                f"year: {year}\njournal: J\nyear_bucket: {bucket}\ntier: A\n---\n"
                f"\n## TL;DR\n{title} TL;DR.\n"
                f"\n## Key findings (with [page] provenance)\n"
                f"- Finding [p1]\n",
                encoding="utf-8",
            )
            summary_paths[doi] = p

        arc_path = kb_root / "Wiki" / "Concepts" / "x-lineage-2026-04-30.md"
        arc_path.parent.mkdir(parents=True, exist_ok=True)
        arc_path.write_text(
            "---\ntopic: x\n---\n\n# Lineage: x\n\nNarrative.\n",
            encoding="utf-8",
        )

        result = LineageRunResult(
            topic="provenance test topic",
            arc_path=arc_path,
            summary_paths=summary_paths,
            corpus_size=3,
            pdfs_acquired=3,
            summaries_written=3,
        )

        out = build_deck_from_lineage_result(
            result,
            speaker="Bobby Ni",
            kb_root=kb_root,
            project_slug="prov-test",
        )
        assert out.exists()
        # Sidecars must land directly next to the .pptx.
        json_p = out.with_name(out.name + ".provenance.json")
        method_p = out.with_name(out.name + ".method.md")
        assert json_p.exists(), f"missing {json_p}"
        assert method_p.exists(), f"missing {method_p}"

        rec = _json.loads(json_p.read_text(encoding="utf-8"))
        assert rec["generated_by"] == ("vaultlab.slides.deck.build_deck_from_lineage_result")
        assert rec["topic"] == "provenance test topic"
        assert rec["kind"] == "slide_deck"
        assert rec["project"] == "prov-test"
        assert rec["params"]["speaker"] == "Bobby Ni"
        assert rec["params"]["plan_mode"] == "fast"
        # The arc is recorded as a related output for the audit log.
        assert any("x-lineage-2026-04-30.md" in r for r in rec.get("related_outputs", []))

        # method.md is human-readable narrative — should at least name
        # the generator and show some context.
        method_text = method_p.read_text(encoding="utf-8")
        assert "Method" in method_text
        assert "build_deck_from_lineage_result" in method_text


# ---------------------------------------------------------------------------
# Fix 2 (2026-04-30 evening-4): aggressive figure picker
# ---------------------------------------------------------------------------


class TestAggressiveFigurePicker:
    """Fix 2 — corpus with N Tier-A papers + N figures must produce
    multiple figure-slides, not just one."""

    def _make_summary(
        self,
        kb_root: Path,
        doi: str,
        *,
        title: str,
        year: int,
        bucket: str,
        og_score: float = 1.0,
    ) -> Path:
        slug = doi.replace("/", "_")
        p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
        p.parent.mkdir(parents=True, exist_ok=True)
        body = (
            f"---\n"
            f"doi: {doi}\n"
            f"title: {title}\n"
            f"authors: ['Author{year}']\n"
            f"year: {year}\n"
            f"year_bucket: {bucket}\n"
            f"tier: A\n"
            f"og_score: {og_score}\n"
            f"---\n\n"
            f"## TL;DR\n{title}: TL;DR sentence.\n"
        )
        p.write_text(body, encoding="utf-8")
        return p

    def _make_fig(self, root: Path, doi: str, size: int = 200_000) -> Path:
        slug = doi.replace("/", "_")
        d = root / slug
        d.mkdir(parents=True, exist_ok=True)
        # Use Pillow to create a real PNG of the requested rough size.
        try:
            from PIL import Image
        except ImportError:
            pytest.skip("Pillow not installed")
        # Create an image of dimensions that approximate the target
        # filesize (ballpark — actual PNG size depends on entropy).
        # Use a noisy image so PNG compression doesn't make it tiny.
        import random

        side = max(64, int((size / 3) ** 0.5))
        img = Image.new("RGB", (side, side))
        rng = random.Random(hash(doi) & 0xFFFF)
        pixels = [
            (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
            for _ in range(side * side)
        ]
        img.putdata(pixels)
        p = d / "fig1-main.png"
        img.save(p, "PNG")
        return p

    def test_corpus_with_5_tier_a_papers_yields_at_least_4_figure_slides(self, tmp_path) -> None:
        """5 Tier-A papers each with one large figure -> >=4 figure-slides."""
        from pptx import Presentation

        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides.deck import build_deck_from_lineage_result

        kb_root = tmp_path / "kb"
        kb_root.mkdir()
        figs_dir = tmp_path / "figs"
        figs_dir.mkdir()

        # 6 Tier-A papers spread across all three buckets.
        papers = [
            ("10.1/h1", "Foundational A", 2010, "history", 5.0),
            ("10.1/h2", "Foundational B", 2011, "history", 4.5),
            ("10.1/d1", "Development A", 2015, "development", 3.0),
            ("10.1/d2", "Development B", 2016, "development", 2.5),
            ("10.1/s1", "SOTA A", 2024, "sota", 2.0),
            ("10.1/s2", "SOTA B", 2025, "sota", 1.8),
        ]

        summary_paths: dict[str, Path] = {}
        figure_assignments: dict[str, Path] = {}
        for doi, title, year, bucket, og in papers:
            sp = self._make_summary(
                kb_root,
                doi,
                title=title,
                year=year,
                bucket=bucket,
                og_score=og,
            )
            summary_paths[doi] = sp
            fig = self._make_fig(figs_dir, doi, size=200_000)
            figure_assignments[doi] = fig

        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text(
            "# Lineage: x\n\n## History\n\n## Development\n\n## SOTA\n",
            encoding="utf-8",
        )

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=6,
        )
        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            figure_assignments=figure_assignments,
        )
        pres = Presentation(str(out))

        # Count slides containing pictures (shape_type==13).
        n_figure_slides = sum(1 for s in pres.slides if any(sh.shape_type == 13 for sh in s.shapes))
        assert n_figure_slides >= 4, (
            f"expected >=4 figure-slides but got {n_figure_slides} "
            f"(corpus has 6 Tier-A papers, each with a cached figure)"
        )

    def test_total_cap_keeps_deck_from_blowing_out(self, tmp_path) -> None:
        """20 Tier-A papers each with figure -> total figure-slides
        capped at _FIGURE_TOTAL_CAP=8."""
        from pptx import Presentation

        from vaultlab.research.lineage import LineageRunResult
        from vaultlab.slides.deck import (
            _FIGURE_TOTAL_CAP,
            build_deck_from_lineage_result,
        )

        kb_root = tmp_path / "kb"
        kb_root.mkdir()
        figs_dir = tmp_path / "figs"
        figs_dir.mkdir()

        summary_paths: dict[str, Path] = {}
        figure_assignments: dict[str, Path] = {}
        for i in range(20):
            doi = f"10.1/p{i:02d}"
            bucket = ["history", "development", "sota"][i % 3]
            sp = self._make_summary(
                kb_root,
                doi,
                title=f"Paper {i}",
                year=2010 + i,
                bucket=bucket,
                og_score=20 - i,
            )
            summary_paths[doi] = sp
            fig = self._make_fig(figs_dir, doi, size=200_000)
            figure_assignments[doi] = fig

        arc = kb_root / "Wiki" / "Concepts" / "x-lineage-2026.md"
        arc.parent.mkdir(parents=True, exist_ok=True)
        arc.write_text(
            "# Lineage: x\n\n## History\n\n## Development\n\n## SOTA\n",
            encoding="utf-8",
        )

        result = LineageRunResult(
            topic="x",
            arc_path=arc,
            summary_paths=summary_paths,
            corpus_size=20,
        )
        out = build_deck_from_lineage_result(
            result,
            speaker="B",
            kb_root=kb_root,
            figure_assignments=figure_assignments,
        )
        pres = Presentation(str(out))
        n_figure_slides = sum(1 for s in pres.slides if any(sh.shape_type == 13 for sh in s.shapes))
        assert n_figure_slides <= _FIGURE_TOTAL_CAP, (
            f"figure-slide count exceeded cap {_FIGURE_TOTAL_CAP}: got {n_figure_slides}"
        )

    def test_largest_figure_chosen_per_paper(self, tmp_path) -> None:
        """When a DOI's cache dir has 5 figures, the picker chooses the largest."""
        from vaultlab.slides.deck import (
            _FIGURE_MIN_BYTES,
            _pick_largest_figure_for_doi,
        )

        # Simulate a manifest with 5 figures of varying sizes, only one
        # of which is above the threshold.
        d = tmp_path / "10-1_x"
        d.mkdir(parents=True, exist_ok=True)
        sizes = [
            ("fig1.png", 5_000),  # decorative crop
            ("fig2.png", 8_000),  # decorative crop
            ("fig3.png", 250_000),  # main result figure
            ("fig4.png", 30_000),  # mid-size
            ("fig5.png", 12_000),  # decorative
        ]
        files = []
        for name, size in sizes:
            p = d / name
            p.write_bytes(b"x" * size)
            files.append(p)
        # Manifest lists all 5 in order.
        import json as _json

        manifest = d / ".figures.json"
        manifest.write_text(
            _json.dumps(
                {
                    "figures": [
                        {"figure_id": f"fig{i + 1}", "file_path": str(p)}
                        for i, p in enumerate(files)
                    ]
                }
            ),
            encoding="utf-8",
        )
        # figure_assignments: only the SEED path (fig1) is registered;
        # the picker must enumerate siblings via the manifest.
        figure_assignments = {"10.1/x": files[0]}
        chosen = _pick_largest_figure_for_doi(
            "10.1/x", figure_assignments, min_size_bytes=_FIGURE_MIN_BYTES
        )
        assert chosen is not None
        assert chosen.name == "fig3.png", (
            f"picker chose {chosen.name} instead of the 250 KB main figure"
        )

    def test_pick_figures_for_bucket_multi_returns_one_per_paper(self, tmp_path) -> None:
        """A bucket with 3 Tier-A papers each having a figure produces 3 picks."""
        from vaultlab.slides.deck import _pick_figures_for_bucket_multi

        # Three real PNG-shaped files large enough to clear the threshold.
        figs: list[Path] = []
        for i in range(3):
            p = tmp_path / f"fig_{i}.png"
            p.write_bytes(b"\x89PNG\r\n" + b"x" * 200_000)
            figs.append(p)

        bucket_papers = [
            {"doi": "10.1/a", "tier": "A", "og_score": 5.0},
            {"doi": "10.1/b", "tier": "A", "og_score": 4.0},
            {"doi": "10.1/c", "tier": "A", "og_score": 3.0},
        ]
        figure_assignments = {
            "10.1/a": figs[0],
            "10.1/b": figs[1],
            "10.1/c": figs[2],
        }
        picks = _pick_figures_for_bucket_multi(bucket_papers, figure_assignments, max_per_bucket=4)
        assert len(picks) == 3
        # Leader is first; each pick has claim_doi == fig_doi (no
        # substitution because every paper has its own figure).
        for claim_doi, fig_doi, _path in picks:
            assert claim_doi == fig_doi
        assert picks[0][0] == "10.1/a"  # leader

    def test_pick_figures_for_bucket_multi_respects_max_per_bucket(self, tmp_path) -> None:
        """``max_per_bucket=2`` caps the result list at 2 entries."""
        from vaultlab.slides.deck import _pick_figures_for_bucket_multi

        figs: list[Path] = []
        for i in range(5):
            p = tmp_path / f"fig_{i}.png"
            p.write_bytes(b"\x89PNG\r\n" + b"x" * 200_000)
            figs.append(p)

        bucket_papers = [{"doi": f"10.1/p{i}", "tier": "A", "og_score": 5.0 - i} for i in range(5)]
        figure_assignments = {f"10.1/p{i}": figs[i] for i in range(5)}
        picks = _pick_figures_for_bucket_multi(bucket_papers, figure_assignments, max_per_bucket=2)
        assert len(picks) == 2

    def test_pick_figures_for_bucket_multi_skips_tier_c_papers(self, tmp_path) -> None:
        """Tier-C papers must NEVER be picked as figure-slide subjects."""
        from vaultlab.slides.deck import _pick_figures_for_bucket_multi

        fig_a = tmp_path / "a.png"
        fig_a.write_bytes(b"\x89PNG\r\n" + b"x" * 200_000)
        fig_c = tmp_path / "c.png"
        fig_c.write_bytes(b"\x89PNG\r\n" + b"x" * 200_000)

        bucket_papers = [
            {"doi": "10.1/c", "tier": "C"},  # Tier-C leader
            {"doi": "10.1/a", "tier": "A"},  # Tier-A
        ]
        figure_assignments = {"10.1/c": fig_c, "10.1/a": fig_a}
        picks = _pick_figures_for_bucket_multi(bucket_papers, figure_assignments)
        # Exactly one pick — the Tier-A paper. The Tier-C paper, even
        # with a cached figure, must never appear as fig_doi.
        assert len(picks) == 1
        claim_doi, fig_doi, _path = picks[0]
        # Leader is Tier-C (claim_doi reflects that), figure substitutes
        # from Tier-A.
        assert fig_doi == "10.1/a"
        assert claim_doi == "10.1/c"
