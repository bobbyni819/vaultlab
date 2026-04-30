"""Stage C of the L4 e2e spatial-tx test.

Acquires figures from PMC OA tar packages for each paper in the corpus,
then composes a 7-slide deck via build_deck_from_lineage_result.
"""
from __future__ import annotations

import os
import pickle
import sys
import time
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

os.environ.pop("ANTHROPIC_API_KEY", None)

from vaultlab.figures.acquisition import acquire_figures_for_corpus
from vaultlab.slides import build_deck_from_lineage_result

PROJECT_SLUG = "spatial-tx-tme-test"
KB_ROOT = Path(r"G:/My Drive/Knowledge/vaultlab")
RESULT_PATH = KB_ROOT / "Output" / PROJECT_SLUG / "stage_b_result.pkl"


def main() -> int:
    started = time.time()
    print("=" * 72)
    print("STAGE C: spatial-tx-tme — figures + deck")
    print("=" * 72)

    with open(RESULT_PATH, "rb") as f:
        state = pickle.load(f)
    result = state["result"]
    corpus = state["corpus"]

    # ---- Figure acquisition ----
    figure_cache_dir = KB_ROOT / "Sources" / "Papers"
    print(f"\nAcquiring figures for {corpus.n_papers} papers...")
    figure_results = acquire_figures_for_corpus(
        corpus,
        cache_dir=figure_cache_dir,
        parallel=2,
    )
    successes = [
        (doi, r) for doi, r in figure_results.items()
        if r.source != "unavailable" and r.figures
    ]
    print(f"Figure acquisition: {len(successes)} / {len(figure_results)} papers got figures")
    figure_assignments = {
        doi: r.figures[0].file_path
        for doi, r in successes
    }
    for doi, r in successes[:10]:
        n_figs = len(r.figures)
        print(f"  {doi}: {n_figs} figs (source={r.source})")

    # ---- Deck composition ----
    print(f"\nComposing deck...")
    deck_path = build_deck_from_lineage_result(
        lineage_result=result,
        speaker="Bobby Y.X. Ni",
        affiliation="Hickey Lab @ Duke BME",
        project_slug=PROJECT_SLUG,
        figure_assignments=figure_assignments,
        kb_root=KB_ROOT,
    )
    print(f"Deck written: {deck_path}")

    # ---- Counts ----
    try:
        from pptx import Presentation
        pres = Presentation(str(deck_path))
        print(f"\nDeck slide count: {len(pres.slides)}")
        for i, s in enumerate(pres.slides, 1):
            # Try to extract a title
            title = ""
            for shape in s.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if text:
                            title = text
                            break
                    if title:
                        break
            print(f"  Slide {i}: {title[:80]}")
    except Exception as exc:
        print(f"(pptx inspection failed: {exc})")

    # Save figure acquisition stats
    fig_stats_path = KB_ROOT / "Output" / PROJECT_SLUG / "figure_stats.txt"
    with open(fig_stats_path, "w", encoding="utf-8") as f:
        f.write(f"Total papers attempted: {len(figure_results)}\n")
        f.write(f"Papers with figures: {len(successes)}\n")
        f.write(f"Total figures: {sum(len(r.figures) for _, r in successes)}\n")
        f.write("\nPer-paper:\n")
        for doi, r in figure_results.items():
            f.write(f"{doi}: source={r.source} figures={len(r.figures)} error={r.error}\n")
    print(f"Figure stats: {fig_stats_path}")

    print(f"\nElapsed: {time.time()-started:.1f}s")
    print("Stage C complete.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
