"""Journal-club deck for Pentimalli & Rajewsky 2025 (Cell Systems).

Bobby is presenting this paper at journal club next week. The deck walks
through all 7 figures with descriptive sentence titles, 3-tier speaker
notes (mental_map / script / extended_walkthrough), and the same Path-A
build-from-plan pipeline used for the multi-lung review deck.

Paper: Pentimalli TM, ..., Rajewsky N. "Combining spatial transcriptomics
and ECM imaging in 3D for mapping cellular interactions in the tumor
microenvironment." Cell Systems 2025;16:101261.
DOI: 10.1016/j.cels.2025.101261

Why this paper for JC: it is the first 3D-multimodal-spatial atlas of an
NSCLC FFPE block (CosMx 1000-plex + SHG ECM imaging + STIM registration)
on a routine clinical sample. The headline is that single-section 2D
spatial transcriptomics SYSTEMATICALLY LOSES the dendritic-cell niche
and T-cell-niche spatial continuity that are recoverable in 3D. This is
also a direct lung-tissue analog of Bobby's thesis multimodal-3D agenda.

Output: G:/My Drive/Knowledge/vaultlab/Output/Decks/journal-club-pentimalli-2026-05-05/
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from vaultlab.slides.audit import audit_deck
from vaultlab.slides.deck import build_from_plan


# ---------------------------------------------------------------------------
# Inline-emphasis post-processor
# ---------------------------------------------------------------------------
# Walks the rendered .pptx, parses markdown-style markers in paragraph text,
# and splits each paragraph into formatted runs.
#
# Markers:
#   **text**      → bold
#   [c]text[/c]   → accent color (warm red-orange) + bold (take-away emphasis)
#
# Applied to every text frame on every slide. Existing run-level formatting
# (size, font name, base color) is preserved. Italic intentionally not used —
# single asterisks would collide with bullet styling, and bold + color give
# enough visual hierarchy for journal-club-grade decks.

def apply_inline_emphasis(pptx_path):
    """Post-process: split runs at markdown markers and apply emphasis."""
    import re
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    ACCENT = RGBColor(0xC8, 0x4B, 0x31)  # warm accent for take-aways / numbers / drug targets
    PATTERN = re.compile(r"(\*\*[^*]+\*\*|\[c\][^\[]+\[/c\])")

    pres = Presentation(pptx_path)
    n_changed = 0

    for slide in pres.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text
                if not text or not PATTERN.search(text):
                    continue
                if not para.runs:
                    continue
                base = para.runs[0]
                base_size = base.font.size
                base_name = base.font.name
                base_bold = base.font.bold
                base_italic = base.font.italic
                try:
                    base_color = base.font.color.rgb
                except Exception:  # noqa: BLE001
                    base_color = None

                tokens = []
                pos = 0
                for m in PATTERN.finditer(text):
                    if m.start() > pos:
                        tokens.append((text[pos:m.start()], None))
                    matched = m.group(0)
                    if matched.startswith("**"):
                        tokens.append((matched[2:-2], "bold"))
                    elif matched.startswith("[c]"):
                        tokens.append((matched[3:-4], "color"))
                    pos = m.end()
                if pos < len(text):
                    tokens.append((text[pos:], None))

                p_elem = para._p
                for child in list(p_elem):
                    if child.tag in (qn("a:r"), qn("a:br")):
                        p_elem.remove(child)

                for txt, fmt in tokens:
                    if not txt:
                        continue
                    r = para.add_run()
                    r.text = txt
                    if base_name:
                        r.font.name = base_name
                    if base_size:
                        r.font.size = base_size
                    if base_bold:
                        r.font.bold = True
                    if base_italic:
                        r.font.italic = True
                    if base_color and fmt != "color":
                        try:
                            r.font.color.rgb = base_color
                        except Exception:  # noqa: BLE001
                            pass
                    if fmt == "bold":
                        r.font.bold = True
                    elif fmt == "color":
                        r.font.color.rgb = ACCENT
                        r.font.bold = True

                n_changed += 1

    # ------- analogy-slide animations (renderer doesn't handle these) -------
    # For every analogy slide: reveal the LEFT-side group (familiar label +
    # body) on click 1, RIGHT-side group (scientific label + body) on click 2.
    # Title, arrow, and citation stay visible from slide load.
    from vaultlab.slides.animations import appear_together_on_click

    n_analogy = 0
    for slide in pres.slides:
        # Detect "analogy" by structure: 5+ text shapes, with a small middle
        # textbox roughly centered and shorter than the side bodies.
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if getattr(s, "has_text_frame", False)]
        if len(text_shapes) < 5:
            continue
        # Heuristic: find shapes by horizontal third of slide.
        sw = pres.slide_width
        sh = pres.slide_height
        title_cutoff = int(sh * 0.13)  # ~1.0" — excludes title (T=0.3) but keeps labels (T≥1.2)
        bottom_cutoff = int(sh * 0.92)  # ~6.9" — excludes citation row
        left_third = sw * 0.40
        right_third = sw * 0.60

        left_shapes = []
        right_shapes = []
        for s in text_shapes:
            if not isinstance(s.top, int) or not isinstance(s.left, int):
                continue
            if s.top < title_cutoff:
                continue  # title row
            if s.top > bottom_cutoff:
                continue  # citation row
            cx = s.left + s.width // 2
            if cx < left_third:
                left_shapes.append(s)
            elif cx > right_third:
                right_shapes.append(s)
            # middle = arrow → leave as always-visible

        if not (len(left_shapes) >= 2 and len(right_shapes) >= 2):
            continue

        try:
            appear_together_on_click(slide, left_shapes, click_index=0)
            appear_together_on_click(slide, right_shapes, click_index=1)
            n_analogy += 1
        except Exception:  # noqa: BLE001
            pass

    pres.save(pptx_path)
    return n_changed, n_analogy

KB = Path("G:/My Drive/Knowledge/vaultlab")
FIG_CACHE = Path("C:/Users/bobby/.cache/vaultlab/_deck_figures_2026_05_03")
SLUG = "10.1016_j.cels.2025.101261"
OUT_DIR = KB / "Output" / "Decks" / "journal-club-pentimalli-2026-05-05"
OUT_DIR.mkdir(parents=True, exist_ok=True)


def _fig(n: int) -> str:
    return str(FIG_CACHE / f"{SLUG}_fig{n}.png")


CITATION = "Pentimalli & Rajewsky, Cell Systems 2025;16:101261"


REFS = [
    "Pentimalli TM et al. Cell Syst 2025;16:101261.",
    "Schurch CM et al. Cell 2020;182:1341.",
    "Goltsev Y et al. Cell 2018;174:968.",
    "Jin S et al. Nat Commun 2021;12:1088.   (CellChat)",
    "Travaglini KJ et al. Nature 2020;587:619.   (Human Lung Cell Atlas)",
    "Kim N et al. Nat Commun 2020;11:2285.   (NSCLC scRNA-seq atlas)",
    "Hickey JW et al. Front Immunol 2021;12:727626.",
    "Sorin M et al. Nature 2023;614:548.",
    "Mi H et al. Brief Bioinform 2024;25:bbae421.",
]


def plan() -> dict:
    return {
        "title": "Combining spatial transcriptomics and ECM imaging in 3D for mapping cellular interactions in the tumor microenvironment",
        "subtitle": "Pentimalli & Rajewsky 2025 — Cell Systems — journal club",
        "topic": "journal-club-pentimalli-2026-05-05",
        "author": "Bobby Y.X. Ni",
        "kb": "vaultlab",
        "theme": "dark",
        "template": "plain",
        "slides": [
            # 1 — Title
            {
                "type": "title",
                "title": "Pentimalli & Rajewsky 2025 — 3D multimodal atlas of one NSCLC",
                "subtitle": "CosMx × SHG × STIM on one FFPE block",
                "author": "Bobby Y.X. Ni  ·  Hickey Lab  ·  Duke BME  ·  Journal club, 2026-05-08",
                "speaker_notes": {
                    "hook": "What if the third dimension is hiding the most clinically actionable structures in your tumour?",
                    "key_claim": "Pentimalli & Rajewsky reconstruct a 3D molecular + ECM atlas of one early-stage NSCLC patient and show that single-section 2D spatial transcriptomics systematically misses the dendritic-cell niche and T-cell niche continuity that are recoverable in 3D.",
                    "transition": "Quick mental map first.",
                },
            },

            # 2 — Why this paper / outline
            {
                "type": "text",
                "title": "Why this paper — first proof that routine 3D multimodal profiling is feasible on clinical FFPE",
                "bullets": [
                    "**SAMPLE** — one patient, early-stage non-small-cell lung cancer (NSCLC), formalin-fixed paraffin-embedded (FFPE) block",
                    "**DESIGN** — [c]34 consecutive 5-µm sections[/c], registered into one 3D coordinate frame via STIM software",
                    "**MODALITIES** — CosMx 1000-plex spatial transcriptomics + label-free second-harmonic generation (SHG) extracellular matrix imaging + H&E",
                    "**SCALE** — [c]340 644 cells[/c] → 18 cell types → 10 multicellular niches",
                    "**HEADLINE** — 3D recovers [c]dendritic-cell niches + T-cell continuity[/c] that 2D analysis erases",
                    "**DRUGGABILITY** — niche-resolved checkpoint signalling: [c]MIF, CCR7, PD-L1, CTLA-4, Tim-3[/c]",
                ],
                "speaker_notes": {
                    "hook": "What does this paper add over Schurch 2020 / Goltsev 2018?",
                    "key_claim": "Earlier neighbourhood-imaging papers were 2D and protein-only. Pentimalli is the first that combines (a) genuine 3D registration over 34 serial sections, (b) 1000-plex transcriptomics, and (c) label-free ECM imaging — on a routine clinical FFPE block, not a research-grade fresh sample.",
                    "evidence": "Cell Systems 2025; co-authored with NanoString (Liang, Gregory) for the CosMx 960-gene panel; Klauschen + Rajewsky + Piccolo lab consortium.",
                    "transition": "Walk through the experimental design first.",
                    "script": (
                        "Three ideas frame the paper. First, the patient is real and the tissue is real — "
                        "an early-stage but aggressive NSCLC, archival FFPE, surgical resection. Nothing "
                        "fresh-frozen, nothing precious. That matters because every clinical pathology "
                        "lab in the world has FFPE blocks; if this workflow only worked on fresh tissue, "
                        "it would never translate. Second, the modalities are complementary rather than "
                        "competing. CosMx gives you which gene each cell expresses; SHG gives you the "
                        "extracellular matrix scaffold those cells live inside; STIM aligns serial "
                        "sections into a coherent 3D coordinate frame. None of those alone is enough. "
                        "Third — and this is the part that matters for our lineage — the 3D analysis "
                        "isn't just decorative. The neighbourhood volumes are 2.28× larger than 2D, "
                        "and that extra volume is exactly where the dendritic-cell niches and the "
                        "T-cell-niche bridges live. Lose 3D, lose the immunology."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is spatial transcriptomics and why is 3D hard?\n\n"
                        "Spatial transcriptomics is a family of methods that measure RNA expression "
                        "while preserving the (x,y) position of each measurement in a tissue section. "
                        "The dominant platforms today are Visium (10x Genomics, barcoded array, ~2-55 µm "
                        "spots), CosMx (NanoString/Bruker, 960-6000 gene single-molecule imaging), MERFISH "
                        "(Vizgen, similar), and Xenium (10x, 300-5000 gene panel). All of them have been "
                        "used almost exclusively in 2D — one section per sample, occasionally a small "
                        "stack of 2-4 sections. Going to true 3D requires three things: (a) consistent "
                        "section quality across many serial cuts (FFPE blocks degrade with each section), "
                        "(b) image registration that handles non-rigid deformations between sections, "
                        "(c) the experimental budget to image many sections (CosMx is ~$3-5k per section "
                        "at this depth, so 34 sections is ~$100-170k of reagent + instrument time).\n\n"
                        "STIM is the registration tool — Software for Spatial Transcriptomics on the "
                        "Image Mosaic, developed in the Rajewsky lab. It uses computer-vision-based "
                        "alignment of consecutive sections, anchored on H&E or DAPI features. The paper "
                        "reports a median post-alignment displacement of ~42 µm — small enough that "
                        "neighbourhoods at 50 µm radius are robust to alignment error.\n\n"
                        "CONTEXT FOR THE LINEAGE — Schurch 2020 in Cell defined the 'cellular neighbourhood' "
                        "concept on 2D CODEX of colorectal cancer. Goltsev 2018 introduced CODEX. Hickey "
                        "2021 systematised the CODEX→single-cell pipeline. Sorin 2023 in Nature scaled "
                        "the neighbourhood approach to 416 LUAD patients on IMC, all 2D. Pentimalli 2025 "
                        "is the first paper that asks: how much of the structure these methods discover "
                        "is artefact of the 2D slicing plane? The answer is: enough to matter clinically."
                    ),
                },
            },

            # 3 — Who built this (text — institution-grouped bullets)
            {
                "type": "text",
                "title": "Who built this — Berlin × Munich × Padua × NanoString consortium",
                "bullets": [
                    "**RAJEWSKY LAB** (Max Delbrück Center, MDC Berlin) — first author Pentimalli; systems biology + spatial omics; built STIM registration tool",
                    "**COSCIA LAB** (MDC Berlin) — rising spatial-proteomics methods lead (ex-Mann lab, Munich)",
                    "**KLAUSCHEN LAB** (Charité + Ludwig Maximilian Univ, LMU Munich) — computational pathology; [c]independent H&E annotations[/c]",
                    "**PICCOLO LAB** (Univ of Padua) — Hippo / YAP-TAZ pathway; mechano-transduction + epithelial-to-mesenchymal transition (EMT) framing",
                    "**NANOSTRING → BRUKER** (Liang, Gregory) — co-designed the CosMx 1000-plex cancer panel",
                ],
                "speaker_notes": {
                    "hook": "Why should you trust this paper before you've read a single figure?",
                    "key_claim": "This isn't a one-lab effort. It's a Berlin-Munich-Padua-NanoString consortium covering systems biology, spatial-proteomics methods, computational pathology, mechano-signalling, and the platform vendor. Cross-validation is built into the author list.",
                    "evidence": "Cell Systems 2025; first author Pentimalli (Rajewsky lab), corresponding Klauschen + Rajewsky; co-authors include Coscia (spatial proteomics) and Piccolo (Hippo/YAP-TAZ).",
                    "transition": "And the field around it has been priming for this paper for ~5 years.",
                    "script": (
                        "Quick credentialling pass before figures. Pentimalli is a Rajewsky-lab "
                        "postdoc; Rajewsky founded systems biology at MDC Berlin and has been pushing "
                        "spatial-omics for the last decade. Klauschen at Munich runs one of the top "
                        "computational-pathology labs and his group does the pathologist H&E "
                        "annotations independently of the molecular analysis — that's the "
                        "orthogonal-validation you'll see on slide 5. Coscia is the rising star in "
                        "spatial-proteomics methods, ex-Mann lab. Piccolo is the Hippo/YAP-TAZ "
                        "guy from Padua and brings the mechano/EMT framing. NanoString — now "
                        "Bruker after the 2024 acquisition — co-developed the 1000-plex CosMx "
                        "panel; the assay is fresh-from-vendor, not arms-length. So when you "
                        "evaluate the methods later, remember: this is a consortium with the "
                        "platform vendor, the pathologist, the spatial-proteomics expert, and "
                        "two systems-biology labs. The cross-checks are in the author list."
                    ),
                    "extended_walkthrough": (
                        "WHY EACH AUTHOR CONTRIBUTION MATTERS — a quick read of the consortium.\n\n"
                        "RAJEWSKY (MDC Berlin) — corresponding author, systems-biology lab. Famous "
                        "for the lncRNA/circRNA work in the 2010s, then pivoted hard into spatial "
                        "transcriptomics + organoids in the 2020s. The lab built STIM, the "
                        "registration tool used here. Reputation: methodologically rigorous, "
                        "publishes in Cell / Nature / Nat Methods.\n\n"
                        "KLAUSCHEN (Munich) — corresponding author, computational pathology. His "
                        "group develops AI-based pathology models and does pathologist-grade tissue "
                        "annotations. The pathologist H&E annotation in Fig 2 is from his lab — "
                        "that's the independent ground-truth check on the molecular niche analysis.\n\n"
                        "COSCIA (MDC Berlin) — spatial-proteomics methods. Ex-Mann lab in Munich. "
                        "Brings the LCMS-MS deep proteomics complement to the CosMx transcriptomics "
                        "(though most of that proteomics work is in supplements / future papers).\n\n"
                        "PICCOLO (Padua) — Hippo/YAP-TAZ pathway. The mechano-transduction framing "
                        "of EMT and ECM stiffness is his contribution. Why YAP-TAZ matters: stiff "
                        "(desmoplastic) ECM activates YAP/TAZ → fibroblast activation → more matrix "
                        "deposition — positive-feedback loop. This is referenced in Fig 5/6/7.\n\n"
                        "PENTIMALLI (Rajewsky lab postdoc) — first author, did most of the wet-lab "
                        "+ analysis. Italian by training, comes through Padua before MDC. Background "
                        "is computational, not just experimental.\n\n"
                        "NANOSTRING / BRUKER (Liang, Gregory) — vendor co-authors. This means the "
                        "1000-plex CosMx panel was co-designed for cancer-relevance, and the "
                        "assay parameters were optimised together rather than the academic group "
                        "buying off-the-shelf reagents. The 0.27% off-target rate from the "
                        "negative-control probes is the result of that co-design.\n\n"
                        "WHAT THIS DOES TO TRUST — collaborative consortium papers in spatial "
                        "biology have a mixed track record. Sometimes the cross-discipline collab "
                        "is meaningful (this paper); sometimes it's a co-author list of "
                        "convenience. Tells: are there multiple corresponding authors from "
                        "different fields (yes), is the methods section detailed enough to "
                        "reproduce (yes — supplementary methods is ~50 pages), and does the "
                        "narrative structure of the paper require all the contributions (yes — "
                        "you cannot tell this story without the pathology, the platform, AND "
                        "the systems biology). All three checks pass, so I take this paper at "
                        "face value."
                    ),
                },
            },

            # 4 — The field — where 3D spatial omics is going (text)
            {
                "type": "text",
                "title": "The field — 3D spatial omics is having a moment",
                "bullets": [
                    "**2018** — CO-Detection by indEXing (CODEX) + imaging mass cytometry (IMC) reach ~50-plex protein imaging",
                    "**2020** — Schürch (colorectal cancer): cellular neighbourhoods predict outcome in 2D CODEX",
                    "**2021-2023** — Hickey CODEX pipeline + Sorin 416-patient lung adenocarcinoma IMC (all 2D)",
                    "**2022-2024** — Visium / CosMx / Xenium scale spatial transcriptomics to 1 000+ genes per section",
                    "**2024** — Human BioMolecular Atlas (HuBMAP) + Human Cell Atlas (HCA) push 3D; Hickey/Agmon couple Vivarium agent-based model with CODEX",
                    "**2025** — [c]Pentimalli: first 3D × multimodal × clinical FFPE[/c]",
                ],
                "speaker_notes": {
                    "hook": "What's the 30-second history of how we got here?",
                    "key_claim": "Spatial omics has been a 7-year arc: 50-plex protein in 2018, neighbourhood concept in 2020, 1000-plex transcriptomics by 2022, 416-patient clinical scale by 2023, 3D atlases by 2024-25. Pentimalli is the first paper that combines all three innovations — 3D + 1000-plex transcriptomics + ECM imaging — on routine clinical FFPE.",
                    "evidence": "Goltsev 2018 Cell (CODEX); Schurch 2020 Cell (CN concept); Hickey 2021 Front Immunol (pipeline); Sorin 2023 Nature (LUAD); HuBMAP/HCA mission 2020+; Pentimalli 2025 Cell Systems (this paper).",
                    "transition": "Now to the figures.",
                    "script": (
                        "Two-minute history. 2018 — Goltsev introduces CODEX, Bodenmiller extends "
                        "IMC; both push multiplexed protein imaging to about 50 markers per section. "
                        "2020 — Schurch publishes the cellular-neighbourhood concept on CRC CODEX, "
                        "which becomes the dominant analytical framework. 2021-23 — the pipeline "
                        "matures: Hickey builds the standard CODEX-to-single-cell-table workflow "
                        "(2021); Sorin scales to 416 LUAD patients on IMC and shows neighbourhood "
                        "structure beats cell frequencies for clinical prediction (2023). All of "
                        "that is 2D. Meanwhile, 2022-24 — spatial transcriptomics scales: Visium, "
                        "CosMx, Xenium each hit 1000+ genes per section. 2024 — HuBMAP and the "
                        "Human Cell Atlas mission start funding 3D efforts at scale. The Hickey/"
                        "Agmon Cell Systems paper integrates CODEX with Vivarium ABM. 2025 — "
                        "Pentimalli is the first to put all three innovations together: 3D, "
                        "1000-plex transcriptomics, multimodal ECM imaging, on a routine FFPE "
                        "block. Each piece existed; the integration didn't until now."
                    ),
                    "extended_walkthrough": (
                        "WHY THIS HISTORY MATTERS FOR EVALUATING THE PAPER — knowing the timeline "
                        "tells you what's genuinely novel vs. what's a refinement.\n\n"
                        "(1) 'Cellular neighbourhoods' — NOT NEW. Concept from Schurch 2020. "
                        "Pentimalli's contribution: extending it to 3D.\n"
                        "(2) 'Multiplexed RNA imaging' — NOT NEW. CosMx existed before this paper. "
                        "Pentimalli's contribution: doing it on 6 sections of one block.\n"
                        "(3) 'Niche-resolved druggable signalling' — PARTIALLY NEW. The CellChat "
                        "spatial-activity score was published in 2D before; doing it in 3D and "
                        "showing the DC-niche network is genuinely new.\n"
                        "(4) 'SHG ECM imaging' — NOT NEW. SHG is decades old in optics. The "
                        "contribution is doing it co-registered with CosMx in 3D.\n"
                        "(5) '3D vs 2D systematic comparison' — NEW. Nobody had run this comparison "
                        "rigorously before. The 2.28× number and the niche-reassignment matrix in "
                        "Fig 3 are the novel quantitative findings.\n"
                        "(6) 'On routine clinical FFPE' — NEW IN COMBINATION. Each previous 3D "
                        "paper used research-grade fresh tissue or model systems. FFPE is the "
                        "gating constraint for clinical translation.\n\n"
                        "WHERE THE FIELD IS HEADED NEXT — three trajectories:\n\n"
                        "(a) Cohort-scale 3D atlases. Pentimalli is n=1; the next step is n=20-100 "
                        "to test reproducibility of the niche grammar. HuBMAP and HCA are funding "
                        "this. Expect a 2026-27 Pentimalli-style paper with multiple patients.\n\n"
                        "(b) Multimodal extension. CosMx + SHG covers transcripts + matrix. Adding "
                        "spatial proteomics (CODEX or MIBI on adjacent sections) gets you proteins. "
                        "Adding metabolomics (MALDI-MSI) gets you metabolites. The 5-modality 3D "
                        "atlas is the obvious next 5-year goal.\n\n"
                        "(c) Functional perturbation. Right now this is descriptive. The next big "
                        "advance is perturbation — CRISPR knockouts of MIF or CCL19 in matched "
                        "model systems and re-running the 3D atlas. Pentimalli explicitly "
                        "acknowledges this gap.\n\n"
                        "WHY A LAB DOING SPATIAL BIOLOGY SHOULD CARE — three concrete things:\n\n"
                        "(1) The 2D-vs-3D systematic-bias result is general. If you're publishing "
                        "anything about niche structure or rare-cell-type localisation from "
                        "single-section 2D, this paper will be cited as evidence your finding "
                        "could be a 2D artefact. Get ahead of it.\n\n"
                        "(2) The CosMx-on-FFPE part is reproducible. If you have a CosMx-equipped "
                        "shared facility (Duke does), the 6-section workflow is doable.\n\n"
                        "(3) The CellChat 3D-spatial activity score is a reusable analytical "
                        "method. You don't need 3D data to use it — works on 2D too — and gives "
                        "more biologically meaningful signalling readouts than gene-expression "
                        "co-occurrence."
                    ),
                },
            },

            # 5 — Section divider
            {
                "type": "section_divider",
                "title": "Now the figures — 7 of them, 1 headline",
            },

            # 6 — fig1 — Experimental design
            {
                "type": "figure",
                "layout": "auto",
                "title": "34 serial sections × CosMx 960-plex × SHG-ECM imaging assembled into one 3D atlas",
                "image_path": _fig(1),
                "caption": "Fig 1 — section layout, H&E, UMAP of 340 k cells, atlas label-transfer.",
                "citation_source": CITATION,
                "bullets": [
                    "**PATIENT** — 1 early-stage NSCLC, archival FFPE block, 16 mm² region of interest",
                    "**DESIGN** — 34 sections × 5 µm; 6 sections imaged with CosMx 1000-plex",
                    "**DATA** — [c]114 M transcripts → 340 644 cells → 18 cell types[/c]",
                ],
                "speaker_notes": {
                    "hook": "How do you actually build a 3D atlas from a single FFPE block?",
                    "key_claim": "The experimental design is a stack of 34 serial 5-µm sections from one FFPE block, with 6 of those sections receiving the full multimodal treatment (CosMx + immunofluorescence + SHG) and the remaining sections used for H&E continuity.",
                    "evidence": "16 mm² ROI selected to capture tumour + cancer-associated stroma + small-caliber airways (anchoring features for 3D registration). 0.27% off-target rate on negative-control probes.",
                    "key_terms": ["FFPE", "CosMx", "STIM", "label transfer", "negative-control probe"],
                    "transition": "From cells to neighbourhoods next.",
                    "script": (
                        "Panel A shows the section layout — 34 sections, 5 µm thick, taken at varying "
                        "z-spacing ranging from 5 to 30 µm so the team can resolve both fine structure "
                        "and longer-range 3D architecture. Sections 3, 4, 10, 12, 16, 22, 28, 34 are "
                        "the multimodal-imaged ones; the rest are H&E for continuity. Panel B is the "
                        "H&E with pathologist annotation of carcinoma, stroma, and normal lung. Panel "
                        "C is the global UMAP of all 340 thousand segmented cells, coloured by 18 "
                        "annotated types — tumour, respiratory epithelium, basal epithelial, alveolar, "
                        "fibroblasts, vascular endothelial, lymphatic endothelial, smooth muscle, "
                        "pericytes, macrophages, monocytes, dendritic cells, mast cells, cytotoxic T, "
                        "regulatory T, B cells, plasma cells, cycling immune. Panel D shows two zooms: "
                        "H&E adjacent to the spatial-transcriptomics colour map, demonstrating they "
                        "carry the same morphological signal. Panel E is the validation — top row is "
                        "label-transfer similarity to the Human Lung Cell Atlas; bottom row is "
                        "similarity to a separate NSCLC reference. Each cell's identity gets confirmed "
                        "by independent reference matching."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is CosMx and what does '960-plex' actually mean?\n\n"
                        "CosMx Spatial Molecular Imager is NanoString's single-molecule spatial-RNA "
                        "platform, now part of Bruker after the 2024 acquisition. Mechanically: the "
                        "tissue is hybridised with up to 6,000 oligonucleotide probes, each carrying a "
                        "unique 64-bit fluorescence barcode. The instrument cycles through 16 rounds of "
                        "fluorescence imaging, decoding the barcodes to identify which probe (and "
                        "therefore which RNA) is at each xy location. The output is a point cloud of "
                        "individual mRNA detections plus segmented cell boundaries. 960-plex means the "
                        "probe panel targets 960 genes — a curated cancer-relevant set including the "
                        "major immune lineage markers, fibroblast markers, EMT markers, and signalling "
                        "axes (cytokines, chemokines, receptors). The 0.27% off-target rate is what "
                        "you read off the negative-control probes — sequences that should match nothing "
                        "in the human transcriptome.\n\n"
                        "WHY 6 SECTIONS, NOT ALL 34? Cost and instrument time. The 34 H&E sections give "
                        "you the morphological backbone for 3D registration; CosMx is the expensive "
                        "molecular layer, run on a sparse subset that still spans the full z-range.\n\n"
                        "WHY 16 mm²? CosMx field-of-view economics. The 1000-plex panel images "
                        "approximately 6 mm² per FOV; 16 mm² is two-three FOVs stitched together. "
                        "Bigger ROIs run into runtime and storage limits on current instrumentation.\n\n"
                        "LABEL TRANSFER — what panel E shows. Rather than calling cell identities "
                        "from de-novo clustering of this dataset alone, the team computes per-cell "
                        "similarity (typically cosine similarity in PCA space) to two reference "
                        "datasets — Travaglini 2020 Nature for healthy lung, Kim 2020 Nat Commun for "
                        "NSCLC — and inherits their labels. The grayscale heat in panel E is the "
                        "similarity score; high grey = high confidence in label transfer. This is "
                        "the standard sanity check: if your alveolar macrophage cluster doesn't look "
                        "like a known alveolar macrophage transcriptome, you have a problem.\n\n"
                        "WHY THE PATIENT MATTERS — early-stage but aggressive NSCLC means surgical "
                        "resection was curative-intent, but the histology showed features predictive "
                        "of recurrence. This is the patient population where understanding the 3D "
                        "tumour microenvironment matters most: the decision is whether to add "
                        "adjuvant chemo or immunotherapy, and that decision is currently made on "
                        "limited information."
                    ),
                },
            },

            # 4 — fig2 — Cellular neighborhoods + 10 niches
            {
                "type": "figure",
                "layout": "auto",
                "title": "10 multicellular niches from 3D neighbourhoods match independent pathologist H&E annotation",
                "image_path": _fig(2),
                "caption": "Fig 2 — z-stack, 50-µm neighbourhood definition, 10 niches, H&E concordance.",
                "citation_source": CITATION,
                "bullets": [
                    "**DEFINITION** — 3D neighbourhood = cells within a 50-µm sphere (z-neighbours included)",
                    "**RESULT** — [c]10 niches[/c] with distinct cell-type compositions (heatmap, panel D)",
                    "**VALIDATION** — niches match [c]independent pathologist H&E[/c] (panel E)",
                ],
                "speaker_notes": {
                    "hook": "Are these 'niches' just clustering artefacts, or do they mean anything?",
                    "key_claim": "10 niches emerge from 3D neighbourhood-vector clustering: tumour core, tumour surface, airways, alveoli, desmoplastic stroma, vascular stroma, smooth muscle, macrophage niches, dendritic-cell niches, T-cell niches. Critically, the niche map agrees with an independent pathologist H&E annotation made WITHOUT seeing the molecular data.",
                    "evidence": "Heatmap (D) shows niche-level cell-type compositions; tumour core is dominated by tumour cells, T-cell niches by cytotoxic + regulatory T cells, etc.",
                    "key_terms": ["cellular neighbourhood", "neighbourhood vector", "niche", "Schurch 2020", "concordance"],
                    "transition": "Niches in hand — now the headline claim about 3D vs 2D.",
                    "script": (
                        "Panel A is the STIM-aligned z-stack rendering — 6 ST sections in their proper "
                        "3D position, with airway lumens visible as the elongated holes spanning "
                        "multiple sections. Panel B is the central definition: for each cell (red), "
                        "find every cell within a 50-micrometre sphere (this includes neighbours from "
                        "sections immediately above and below); count those neighbours by cell type; "
                        "you now have a neighbourhood vector — 18 entries per cell. Cluster all 340k "
                        "neighbourhood vectors with unsupervised methods, and 10 stable niches drop "
                        "out. Panel C is a UMAP of cells coloured by niche assignment — the 10 niches "
                        "occupy distinct regions of UMAP space, telling you the niches are not just "
                        "k-means noise. Panel D is the niche × cell-type heatmap — read across rows "
                        "to see what each niche is made of. Tumour core is mostly tumour cells, T-cell "
                        "niches are enriched for cytotoxic + regulatory T cells. Panel E — and this is "
                        "the orthogonal check — is the spatial map of niches next to a pathologist's "
                        "H&E annotation made without looking at the molecular data. The two agree. "
                        "When two completely independent methods, looking at completely different "
                        "features, draw the same boundaries, the boundaries are real."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is a 'cellular neighbourhood' and where did the concept "
                        "come from?\n\n"
                        "The cellular-neighbourhood (CN) concept was formalised by Christian Schürch "
                        "and Garry Nolan in their 2020 Cell paper on colorectal cancer (Schürch et al. "
                        "Cell 2020;182:1341). The motivation: in oncology, cell FREQUENCIES (e.g., "
                        "'this tumour is 12% Tregs') are weak predictors of outcome, but the SPATIAL "
                        "ARRANGEMENT of those cells matters enormously. A tumour with infiltrating B "
                        "cells in well-organised tertiary lymphoid structures responds completely "
                        "differently to immunotherapy than the same tumour with the same B cells "
                        "scattered uniformly. CNs make 'spatial arrangement' formal: for each cell, "
                        "describe what it's surrounded by (a count vector by cell type), then cluster "
                        "all cells in the cohort by their neighbourhood vectors. The clusters ARE the "
                        "neighbourhood archetypes.\n\n"
                        "WHY 50 µm? It's the typical CODEX/IMC paper choice and it scales to roughly "
                        "the radius at which signalling gradients (chemokines, soluble cytokines) "
                        "remain meaningful before getting washed out by interstitial flow. Smaller "
                        "radii (10-20 µm) capture juxtacrine contact-dependent signalling; 50-100 µm "
                        "captures paracrine signalling; >200 µm starts mixing across distinct tissue "
                        "compartments.\n\n"
                        "WHAT IS A NICHE vs WHAT IS A CN? Terminology is fluid in the field. CNs in "
                        "the strict Schürch sense are the clusters in neighbourhood-vector space. "
                        "Niches in this paper are the spatially-coherent regions of tissue annotated "
                        "by their dominant CN — slightly higher-level. Read 'niche' here as 'the "
                        "tissue-scale region where this CN dominates'.\n\n"
                        "WHY THE PATHOLOGIST AGREEMENT IS THE CRITICAL CHECK — one of the lingering "
                        "doubts about CN analyses is that the clusters are produced by an algorithm "
                        "from features the algorithm chose; how do we know they correspond to real "
                        "tissue structures rather than statistical mirages? Pathologists annotating "
                        "H&E images use morphology — cell shape, tissue density, vessel architecture, "
                        "glandular structure — features the CN algorithm never saw. When the "
                        "pathologist's hand-drawn 'tumour core' boundary aligns with the algorithmic "
                        "'tumour core' niche, the niche is morphologically as well as molecularly "
                        "coherent.\n\n"
                        "READING PANEL D HEATMAP — rows are cell types, columns are niches, colour "
                        "is average count per neighbourhood. The diagonal-block structure tells you "
                        "each niche has a dominant cell-type identity; off-diagonal entries tell you "
                        "the partner cell types. Macrophage niches have ~30 macrophages and partner "
                        "with both tumour cells (immune-suppressive) and dendritic cells (antigen "
                        "presentation). T-cell niches partner cytotoxic + regulatory T cells with "
                        "B cells and lymphatic endothelium — i.e., they look like nascent TLS."
                    ),
                },
            },

            # 5 — fig3 — THE HEADLINE: 3D vs 2D
            {
                "type": "figure",
                "layout": "auto",
                "title": "HEADLINE — 3D neighbourhoods are 2.28× larger; dendritic-cell niches and T-cell continuity exist ONLY in 3D",
                "image_path": _fig(3),
                "caption": "Fig 3 — 3D vs 2D: bigger neighbourhoods, DC niche only in 3D, T-cell bridges.",
                "citation_source": CITATION,
                "bullets": [
                    "**PANEL B** — 3D vs 2D: [c]32 vs 22 cells, 9 vs 7 cell types[/c] (p<0.005)",
                    "**PANEL D** — DC-niche cells reassigned to tumour surface ([c]51%[/c]) or T-cell ([c]24%[/c]) in 2D",
                    "**PANEL E** — T-cell-niche spatial continuity [c]restored only in 3D[/c]",
                ],
                "speaker_notes": {
                    "hook": "Does the third dimension actually buy you anything clinically?",
                    "key_claim": "In 3D, neighbourhoods are 2.28× larger; the dendritic-cell niche EXISTS ONLY when you analyse in 3D; T-cell-niche spatial bridges are recovered only in 3D. 2D analysis is systematically losing the structures that drive immunotherapy response.",
                    "evidence": "Median neighbourhood: 32 vs 22 cells, 9 vs 7 cell types, Chao alpha-diversity 10.5 vs 8 (all p<0.005). 51.2% of DC-niche cells get re-classified as tumour-surface in 2D; 23.6% as T-cell-niche.",
                    "key_terms": ["3D neighbourhood", "Chao alpha-diversity", "niche reassignment", "DC niche", "T-cell continuity"],
                    "transition": "Why does this matter clinically? Niches carry signalling.",
                    "script": (
                        "Panel A schematises the 3D neighbourhood: a sphere of radius 40 µm extends "
                        "across z=±30 µm, picking up neighbours from the section above and below the "
                        "anchor cell. Panel B is the headline statistic, presented as three violin "
                        "plots — number of neighbours, number of cell types, and Chao alpha-diversity. "
                        "All three: 3D wins decisively, p<0.005. The biological consequence appears "
                        "in panels C and D. Panel C is the 2D-restricted clustering with niches "
                        "labelled — note the DENDRITIC-CELL niche is shown but in strikethrough text, "
                        "indicating it does NOT statistically resolve in 2D. Panel D is the cell-by-"
                        "cell reassignment matrix between 3D-niche and 2D-niche labels: 80.4% of "
                        "tumour-core cells stay tumour-core, fine; but only 62.5% of macrophage-niche "
                        "cells stay macrophage-niche, and the dendritic-cell niche disappears entirely "
                        "(blanked out / strikethrough). Where do those DC-niche cells go? Tumour "
                        "surface (51%) or T-cell niches (24%) — they get absorbed into the nearest "
                        "neighbourhood archetype that 2D can statistically support. Panel E shows "
                        "T-cell-niche spatial structure: in 2D it's blue patches scattered across the "
                        "tumour bed; in 3D-rendered view it's a connected bridge structure. Panel F "
                        "shows three 3D-rendered niches — tumour core, tumour surface, dendritic-cell + "
                        "macrophage + T-cell niches — confirming that what looks like point clouds in "
                        "2D becomes coherent 3D anatomy."
                    ),
                    "extended_walkthrough": (
                        "WHY DOES 3D RECOVER MORE STRUCTURE? The geometric intuition: a 2D disc of "
                        "radius r has area πr²; a 3D ball of radius r has volume (4/3)πr³. At matched "
                        "cell density, the 3D ball contains 4r/3 times more cells than the 2D disc — "
                        "for a 50 µm neighbourhood with 200 µm tissue thickness available, that's "
                        "roughly the 2.28× ratio observed. More cells per neighbourhood means the "
                        "cell-type-vector signal-to-noise improves — small populations like dendritic "
                        "cells (typically 1-3% of total) need ≥10 cells in a neighbourhood to be "
                        "statistically distinguishable, and 3D crosses that threshold where 2D "
                        "doesn't.\n\n"
                        "WHY DENDRITIC CELLS, SPECIFICALLY, ARE THE MOST 3D-SENSITIVE — DCs are rare "
                        "(~1-2% of cells in this tumour) AND they cluster sparsely. A 5 µm-thick "
                        "section that happens to slice between two adjacent DCs misses them both. In "
                        "3D, those two DCs are clearly within the same neighbourhood. The DC-niche "
                        "loss in 2D is therefore a sampling artefact compounding cell rarity with "
                        "section thickness.\n\n"
                        "WHY TC-CELL NICHE CONTINUITY IS DIFFERENT — T cells are not rare; they're "
                        "abundant. The 2D-vs-3D issue here isn't 'do we see the cells' but 'do we see "
                        "the spatial continuity'. T-cell niches in this NSCLC are organised into "
                        "long thin bridges that wind through the tumour bed. In 2D each section "
                        "intersects the bridge as a separate ellipse, and the niche-detection "
                        "algorithm sees those ellipses as independent foci. In 3D the bridges are "
                        "connected. This matters because a connected T-cell network is interpreted "
                        "very differently from scattered T-cell foci — connected networks are TLS-like, "
                        "and TLS-like architectures predict immunotherapy response.\n\n"
                        "CLINICAL TAKEAWAY — if your patient gets a single-section 2D spatial "
                        "transcriptomics workup as part of treatment selection (and Visium and CosMx "
                        "are creeping into clinical workflows), and the 2D analysis shows 'no "
                        "tertiary lymphoid structures, no DC niches', the patient may STILL HAVE "
                        "those structures in 3D. The 2D negative is misleading. This implies that "
                        "single-section spatial-omics readouts of immunotherapy-response biomarkers "
                        "are systematically biased toward false negatives.\n\n"
                        "THESIS-WORLD CONNECTION — for our multiscale-simulation programme this is "
                        "the empirical justification for going to 3D in silico. Any model of an "
                        "infected lung that simulates a 3D tissue but compares only to 2D CODEX or "
                        "2D IMC will be 'right' on cell counts but 'wrong' on neighbourhoods. The "
                        "comparison must virtually slice the simulation, not flatten the data."
                    ),
                },
            },

            # 6 — fig4 — Receptor-ligand activity + DC niche druggable signaling
            {
                "type": "figure",
                "layout": "auto",
                "title": "3D receptor-ligand activity scoring exposes a druggable dendritic-cell-niche signalling network",
                "image_path": _fig(4),
                "caption": "Fig 4 — 480-pair activity score; niche ligands; DC-niche druggable network.",
                "citation_source": CITATION,
                "bullets": [
                    "**METHOD** — 480 CellChat receptor-ligand pairs scored at 50-µm radius in 3D",
                    "**PATTERN** — niche-specific ligands (PDGFB → vascular; CCL19 + CXCL9 → DC + T-cell)",
                    "**PAYOFF** — DC-niche druggable network: [c]MIF, CCR7, PD-L1, CTLA-4, Tim-3[/c]",
                ],
                "speaker_notes": {
                    "hook": "If 3D exposes the DC niche, what is the DC niche actually doing?",
                    "key_claim": "Pentimalli scores 480 receptor-ligand pairs in 3D and finds 96 ligands enriched in ≥1 niche. The DC niche carries a clean immune-checkpoint signalling network with multiple druggable nodes — anti-PD-1, anti-CTLA-4, anti-Tim-3, anti-MIF.",
                    "evidence": "CellChat 480-pair database; activity score = sqrt(L_sender × R_receiver) summed within 50 µm; log2FC > 0.5 threshold for enrichment.",
                    "key_terms": ["CellChat", "receptor-ligand activity", "MIF", "CCR7", "PD-L1", "Tim-3"],
                    "transition": "Cells AND signals — but what about the matrix they live in?",
                    "script": (
                        "Panel A is the activity-score formula: for each receptor-ligand pair, multiply "
                        "ligand expression at the sender cell by receptor expression at the receiver "
                        "cell, normalise, sum within 50 µm. This converts a static gene-expression "
                        "snapshot into a directional signalling score. Panel B is the top-5 niche-"
                        "specific ligands across 8 niches — read as a heatmap, blue = high activity. "
                        "PDGFB peaks in vascular stroma (drives pericyte recruitment). AREG / EFNA1 / "
                        "CDH1 peak in tumour surface (autocrine growth + adhesion). CCL19 + CXCL9 peak "
                        "in DC and T-cell niches (immune cell trafficking). Panels C-E show the "
                        "spatial maps for three exemplar ligands: PDGFB tracks the vasculature, AREG "
                        "fills the tumour surface, CCL19 lights up the DC niche regions. Panel F is "
                        "the DC-niche signalling network as a dotplot: rows are sender cell types, "
                        "columns are receptor-ligand pairs. Tumour cells produce MIF that engages "
                        "CXCR4 on DCs and macrophages — recruits them but suppresses APC function. "
                        "Fibroblasts produce CCL19/CCL21 that engages CCR7 on DCs and T cells — "
                        "retention in the niche. DCs produce CD80 engaging CTLA-4, PD-L1 engaging "
                        "PD-1, Galectin-9 (LGALS9) engaging Tim-3 (HAVCR2) — the immune-checkpoint "
                        "suite. Panel G is the network cartoon: tumour → fibroblast → DC → CTL/Treg "
                        "as a chain of cellular control. Every red bar is a node we already drug or "
                        "have a drug in trials for."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is CellChat?\n\n"
                        "CellChat (Jin et al. Nat Commun 2021;12:1088) is a database + inference tool "
                        "for cell-cell communication. The database is a curated set of ~2,000 human "
                        "ligand-receptor (L-R) interactions, with each interaction tagged by signalling "
                        "pathway (e.g., 'TGFβ signalling', 'TNF signalling', 'WNT signalling'). The "
                        "inference part takes a single-cell expression matrix and asks: which L-R pairs "
                        "are 'communicating' in the data, given which cells express the ligand and "
                        "which express the receptor? The Pentimalli paper extends this to spatial: "
                        "instead of asking 'is L expressed somewhere AND R expressed somewhere', they "
                        "ask 'is L expressed in cells WITHIN 50 µm OF cells expressing R'. That spatial "
                        "constraint is what makes the score a meaningful signalling-activity proxy "
                        "rather than a population-level co-expression mean.\n\n"
                        "WHY 480 PAIRS, NOT 2000? The 480 number is the subset of CellChat pairs where "
                        "both ligand AND receptor are on the CosMx 960-gene panel. The full CellChat "
                        "database has receptors and ligands the panel doesn't measure.\n\n"
                        "WHY THE DC NICHE IS THE MONEY SHOT — every modern immuno-oncology drug "
                        "targets a node in this network. anti-PD-1 (pembrolizumab, nivolumab, "
                        "cemiplimab), anti-PD-L1 (atezolizumab, durvalumab, avelumab), anti-CTLA-4 "
                        "(ipilimumab), anti-Tim-3 (sabatolimab, in trials), anti-MIF (CPSI-1306 and "
                        "others, early development), anti-CCR7 antagonists (research-stage). The "
                        "Pentimalli result tells you these drugs are acting on a SPATIALLY-RESOLVED "
                        "niche, not a 'tumour-wide' phenomenon. The clinical implication: response to "
                        "anti-PD-1 should depend on whether the patient HAS a DC niche to disrupt.\n\n"
                        "MIF→CXCR4 IN DETAIL — MIF (Macrophage Migration Inhibitory Factor) is a "
                        "secreted cytokine produced by tumour cells; CXCR4 is a chemokine receptor on "
                        "myeloid cells. The MIF-CXCR4 axis recruits dendritic cells INTO the tumour "
                        "bed but holds them in an immature, antigen-presentation-suppressed state — "
                        "they are present but inactive. Blocking this axis releases the brake on "
                        "DC maturation. Trials are ongoing.\n\n"
                        "CCL19/CCR7 vs PD-L1/PD-1 — the network has both PRO-immune and "
                        "ANTI-immune signalling. CCL19 (from fibroblasts) → CCR7 (on DCs and T cells) "
                        "is recruitment-and-retention — pro-tumour-fighting. PD-L1 (on DCs) → PD-1 "
                        "(on cytotoxic T cells) is the brake — anti-tumour-fighting. The DC niche "
                        "is therefore a battleground: it's where the immune response is being "
                        "organised AND simultaneously where it's being suppressed. Disrupting the "
                        "suppression while preserving the recruitment is the goal of combination "
                        "checkpoint therapy."
                    ),
                },
            },

            # 7 — fig5 — SHG-defined ECM × fibroblast states
            {
                "type": "figure",
                "layout": "auto",
                "title": "SHG-defined ECM compartments couple to 6 fibroblast states with distinct matrix-regulator activities",
                "image_path": _fig(5),
                "caption": "Fig 5 — SHG ECM, 3 compartments, 6 fibroblast states, ECM regulators.",
                "citation_source": CITATION,
                "bullets": [
                    "**MODALITY** — SHG = label-free collagen + elastin imaging ([c]no antibody, no stain[/c])",
                    "**ECM** — 3 compartments: homeostatic / degraded / desmoplastic",
                    "**FIBROBLASTS** — [c]6 transcriptomic states[/c] couple to specific ECM compartments",
                ],
                "speaker_notes": {
                    "hook": "Cells live in matrix — does the matrix follow the cells or do the cells follow the matrix?",
                    "key_claim": "SHG imaging gives label-free resolution of three ECM compartments — homeostatic, degraded, desmoplastic — and these couple to six distinct fibroblast transcriptomic states. ECM regulators (TIMP1, IGF1, INHBA) are spatially organised across compartments.",
                    "evidence": "62,604 fibroblasts; unsupervised clustering → 6 phenotypes; SHG intensities per 50×50 µm neighbourhood → k-means k=3 → 3 ECM compartments.",
                    "key_terms": ["SHG", "second-harmonic generation", "myofibroblast", "desmoplastic", "TIMP1", "IGF1", "INHBA"],
                    "transition": "Where in this matrix does the tumour actually invade?",
                    "script": (
                        "Panel A is the SHG image of section 3 — collagen fibres in green, elastin in "
                        "magenta. Note: this required no antibody, no fluorophore, no stain — SHG is a "
                        "non-linear optical effect that fires only when the laser hits ordered triple-"
                        "helical collagen. Panel B is the per-neighbourhood SHG-intensity scatter — "
                        "x-axis elastin, y-axis collagen. K-means with k=3 picks out three ECM "
                        "compartments. Cluster 1 (yellow, 'homeostatic') is moderate-elastin moderate-"
                        "collagen — normal lung. Cluster 2 (orange, 'degraded') is low-collagen — the "
                        "tumour bed where matrix is being chewed up. Cluster 3 (purple, 'desmoplastic') "
                        "is high-collagen — fibrotic stromal regions. Panel C maps the compartments "
                        "back onto the section. Panel D is the fibroblast UMAP — 6 distinct "
                        "transcriptomic states. Myofibroblasts express FN1, COL11A1, ACTA2 — the "
                        "contractile fibrosis-driving state. 'Activated' fibroblasts express JUN, FOS, "
                        "IGF1 — the early-response state. 'Matrix' fibroblasts express LUM, MGP, "
                        "TIMP1 — the housekeeping ECM-maintenance state. 'Antigen-presenting' "
                        "fibroblasts express CD74 + HLA-DRB1 — these are unusual, they look like "
                        "fibroblasts but expose MHC-II. CCL19+ and CXCL10+ reticular fibroblasts are "
                        "the immune-attracting subsets that probably nucleate the DC and T-cell niches "
                        "we saw in Fig 4. Panel E shows the spatial layout of these states. Panel F "
                        "is the cross-tabulation: which fibroblast states sit in which ECM compartment? "
                        "Myofibroblasts in degraded ECM (where they're tearing it up); activated in "
                        "desmoplastic; matrix in homeostatic. Panels G-I show three exemplar ECM "
                        "regulators — TIMP1 (matrix-state marker), IGF1 (activated-state marker), "
                        "INHBA (myofibroblast marker) — as RNA-density maps that respect compartment "
                        "boundaries."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is second-harmonic generation imaging?\n\n"
                        "SHG is a non-linear optical effect. Two photons of frequency ω hit an ordered "
                        "non-centrosymmetric structure (like a triple-helical collagen fibre) and "
                        "combine into one photon of frequency 2ω. The emission is at exactly half the "
                        "incident wavelength. The catch: SHG only happens for materials that lack "
                        "centrosymmetry — collagen fibrils are the canonical biological example, "
                        "elastin gives weaker but detectable signal at a different wavelength. Because "
                        "no fluorophore is involved, SHG imaging can be acquired on the SAME tissue "
                        "section that's about to be processed for spatial transcriptomics, with no "
                        "interference. That co-acquisition is the methodological trick that makes "
                        "this paper possible.\n\n"
                        "WHY THREE ECM COMPARTMENTS, NOT MORE?\n\n"
                        "K-means with k=3 was chosen empirically — they tried k=2, 3, 4, 5 and k=3 "
                        "produced the most stable clusters (cluster reassignment under perturbation "
                        "was lowest at k=3). The biological labels — 'homeostatic', 'degraded', "
                        "'desmoplastic' — are post-hoc interpretations matching pathological intuition "
                        "about NSCLC stroma:\n"
                        "* Homeostatic: normal lung parenchyma, adjacent to the tumour but not within "
                        "  it. Normal alveolar walls.\n"
                        "* Degraded: matrix that has been actively remodelled by tumour-derived MMPs. "
                        "  Low collagen because it's been digested. Low elastin because it's been "
                        "  fragmented. This is the matrix INSIDE the tumour bed.\n"
                        "* Desmoplastic: high collagen, low elastin. Stiff scarring response of the "
                        "  host to the tumour, deposited at the tumour-stroma interface. This is what "
                        "  makes a tumour 'firm' on palpation.\n\n"
                        "FIBROBLAST HETEROGENEITY — the field has moved away from 'CAF' (cancer-"
                        "associated fibroblast) as a monolithic category. By 2025 we recognise at "
                        "least 4-7 functionally distinct fibroblast states in tumours, with names "
                        "like myCAF (myofibroblast-CAF), iCAF (inflammatory-CAF), apCAF (antigen-"
                        "presenting), and various reticular populations. Pentimalli's 6 states "
                        "match this typology with finer reticular subdivision (CCL19+ and CXCL10+).\n\n"
                        "WHY THE ECM-FIBROBLAST COUPLING MATTERS — TIMP1 (Tissue Inhibitor of "
                        "Metalloproteinases) opposes matrix-degrading MMPs; matrix fibroblasts "
                        "expressing it are protective of the homeostatic compartment. IGF1 (Insulin-"
                        "like Growth Factor 1) drives proliferation of activated fibroblasts and "
                        "their conversion to myofibroblasts; expressed in desmoplastic regions. INHBA "
                        "(Inhibin Beta A, a TGFβ-family ligand) is the canonical myofibroblast-"
                        "activating cytokine — autocrine in this context (myofibroblasts make INHBA "
                        "to maintain themselves). The spatial organisation of these three regulators "
                        "across compartments tells you the matrix biology is being driven from "
                        "within the matrix itself, not from the tumour cells alone.\n\n"
                        "THESIS CONNECTION — for the lung-fibrosis simulation Bobby is building, "
                        "this paper supplies the empirical anchor for HOW MANY fibroblast states a "
                        "model needs (6) and WHAT cytokine outputs distinguish them. INHBA → "
                        "myofibroblast self-maintenance, IGF1 → activated proliferation, TIMP1 → "
                        "matrix protection. Plus the SHG → ECM compartment recipe gives a measurable "
                        "comparison target for any in-silico ECM dynamics."
                    ),
                },
            },

            # 8 — fig6 — Tumor pseudotime + EMT niche
            {
                "type": "figure",
                "layout": "auto",
                "title": "Tumour 3D pseudotime captures EMT — the EMT niche invades collagen-poor ECM",
                "image_path": _fig(6),
                "caption": "Fig 6 — tumour pseudotime captures EMT; EMT niche is collagen-poor.",
                "citation_source": CITATION,
                "bullets": [
                    "**TRAJECTORY** — pseudotime: epithelial → mesenchymal",
                    "**GRADIENT** — highest at invasive front (p<0.0001)",
                    "**SURPRISE** — EMT niche is [c]collagen-POOR[/c]",
                ],
                "speaker_notes": {
                    "hook": "Where exactly in 3D space is the tumour invading from?",
                    "key_claim": "Tumour-cell 3D pseudotime tracks the epithelial-to-mesenchymal transition: cells at the stromal interface are most mesenchymal, and the spatial 'EMT niche' is COLLAGEN-POOR, not desmoplastic — the invading front degrades matrix ahead of itself.",
                    "evidence": "Pseudotime ordering on UMAP using CDH1 / EPCAM (epithelial) vs ITGB6 / COL3A1 (mesenchymal). Pseudotime by niche category: tumour core < tumour surface < stroma-infiltrating, all p<0.0001.",
                    "key_terms": ["pseudotime", "EMT", "CDH1", "ITGB6", "COL3A1", "EMT niche", "invasive front"],
                    "transition": "Now what does the EMT niche express that makes it invasive?",
                    "script": (
                        "Panel A is a 3D rendering of just the tumour-cell population, coloured by "
                        "subtype: yellow tumour-core cells (deep in the tumour), red tumour-surface "
                        "cells (at the tumour-stroma interface), green stroma-infiltrating tumour "
                        "cells (further out, infiltrating the stroma). The colour gradient maps onto a "
                        "biological progression. Panel B is the tumour-cell UMAP, coloured by "
                        "pseudotime rank — early pseudotime in blue/grey, late in red. Pseudotime "
                        "here uses a method like Monocle or Slingshot, ordering cells along a "
                        "trajectory inferred from gene-expression similarity. The four small UMAPs "
                        "below show CDH1 (E-cadherin, epithelial marker), ITGB6 (integrin β6, "
                        "mesenchymal), EPCAM (epithelial), COL3A1 (mesenchymal). They confirm: "
                        "early-pseudotime cells are epithelial, late-pseudotime cells are mesenchymal. "
                        "The trajectory IS EMT. Panel C is the box plot of pseudotime stratified by "
                        "tumour-3D-pseudospace category — core lowest, surface intermediate, "
                        "stroma-infiltrating highest, with all comparisons p<0.0001. So the more "
                        "spatially-invading the cell, the more mesenchymal. Panel D maps pseudotime "
                        "back onto the 2D-projected spatial layout — the rectangle marked 'EMT niche' "
                        "is the highest-pseudotime region. Panel E shows the SHG/collagen signal "
                        "across the same field of view — and in the EMT-niche box, collagen is LOW. "
                        "Panel F is the box-plot of collagen scaled across niche categories: EMT "
                        "niche is decisively collagen-poor compared to homeostatic and desmoplastic. "
                        "This is counter-intuitive — desmoplasia is usually associated with invasion "
                        "— but the resolution here is that desmoplasia is at the tumour BORDER while "
                        "active invasion happens in collagen-poor 'tunnels' the tumour creates ahead "
                        "of itself."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — what is EMT and what is pseudotime?\n\n"
                        "EMT (Epithelial-to-Mesenchymal Transition) is the gene-expression program "
                        "where an epithelial cell loses cell-cell adhesion (E-cadherin / CDH1 down), "
                        "loses apical-basal polarity, and gains migratory / mesenchymal features "
                        "(N-cadherin / CDH2 up, vimentin up, integrin β6 / ITGB6 up, MMP secretion "
                        "up). EMT is a major route by which carcinomas (epithelial-origin tumours) "
                        "become invasive and metastatic. EMT is a continuum, not a binary state — "
                        "tumour cells exist at every position along the epithelial-mesenchymal "
                        "spectrum, and that spectrum is what pseudotime captures.\n\n"
                        "Pseudotime methods (Monocle, Slingshot, PAGA) order cells along a continuous "
                        "trajectory inferred from gene-expression similarity. They start with a UMAP "
                        "or principal-graph embedding, identify a starting state (typically the most "
                        "abundant or most epithelial cluster), and compute a one-dimensional "
                        "coordinate ('pseudotime') for every cell that increases monotonically along "
                        "the inferred trajectory. The ordering is statistical — if the manifold is "
                        "smooth and the trajectory is real, pseudotime captures it.\n\n"
                        "KEY EMT MARKERS — CDH1 (E-cadherin) is the canonical epithelial marker; loss "
                        "is the entry point to EMT. EPCAM is epithelial cell-cell adhesion molecule; "
                        "stable in epithelial cells, lost in mesenchymal. ITGB6 (integrin β6) "
                        "heterodimerises with αv to form αvβ6, which binds RGD motifs on fibronectin "
                        "and TGFβ — gained during EMT, drives invasion. COL3A1 (collagen-3) is "
                        "secreted by mesenchymal cells, contributes to tumour-bed matrix.\n\n"
                        "WHY COLLAGEN-POOR EMT NICHES IS A SURPRISING FINDING — the textbook "
                        "expectation is that EMT happens at the desmoplastic stromal front, where "
                        "tumour cells encounter dense collagen and respond by transitioning. "
                        "Pentimalli's 3D analysis flips this: the EMT-most cells are in collagen-POOR "
                        "regions. The biological interpretation: invading tumour cells DIGEST collagen "
                        "ahead of them as they advance, creating low-collagen tunnels through dense "
                        "stroma. The desmoplastic compartment is the static border; the EMT niche is "
                        "the active invasive front. This distinction is invisible in 2D where the "
                        "two compartments overlap heavily on slicing.\n\n"
                        "CLINICAL RELEVANCE — anti-collagen / anti-MMP therapies have repeatedly "
                        "failed in NSCLC trials. One reason might be that they targeted the WRONG "
                        "compartment — desmoplastic stroma is already deposited and not the active "
                        "invasion site. The EMT niche is where the action is, and it requires a "
                        "different druggable target set."
                    ),
                },
            },

            # 9 — fig7 — Druggable EMT niche
            {
                "type": "figure",
                "layout": "auto",
                "title": "EMT niche has a distinct druggable signalling signature — LGALS1 / IGFBP5 / VEGFA / SPP1 / integrins",
                "image_path": _fig(7),
                "caption": "Fig 7 — EMT-niche markers; tumour-fibroblast-macrophage druggable model.",
                "citation_source": CITATION,
                "bullets": [
                    "**TUMOUR CELLS** — [c]NDRG1 + LGALS1[/c] (galectin-1; immunosuppressive)",
                    "**MYOFIBROBLASTS** — [c]VEGFA + IGFBP5[/c] (angiogenic; modulates IGF axis)",
                    "**MACROPHAGES** — [c]SPP1[/c] (osteopontin; pro-tumour-permissive marker)",
                ],
                "speaker_notes": {
                    "hook": "What is each cell type in the EMT niche actually contributing?",
                    "key_claim": "The EMT niche is a tumour-fibroblast-macrophage cooperative — tumour cells express NDRG1 + LGALS1; myofibroblasts secrete VEGFA + IGFBP5; macrophages secrete SPP1. Each axis is a known druggable target. The integrin 3D-signalling density is highest in the EMT niche, confirming active αvβ6-fibronectin engagement.",
                    "evidence": "Tumour-cell volcano (A) shows EMT-niche-enriched markers; spatial RNA density maps (B-G) localise the marker expression to the rectangular EMT-niche region; integrin signalling activity (H) peaks in the same region.",
                    "key_terms": ["NDRG1", "LGALS1", "galectin-1", "VEGFA", "IGFBP5", "SPP1", "osteopontin", "αvβ6 integrin"],
                    "transition": "What's the take-home for the field — and for the thesis?",
                    "script": (
                        "Panel A is a volcano plot of tumour-cell gene expression, contrasting "
                        "EMT-niche cells against other tumour-cell populations. The orange + yellow "
                        "annotations on the right are the EMT-niche markers — NDRG1, LGALS1, GADD45B, "
                        "MIF, JUNB. NDRG1 is a stress-response gene activated by hypoxia. LGALS1 is "
                        "Galectin-1, a glycan-binding protein known to suppress T-cell responses — so "
                        "the EMT niche is actively immunosuppressive at the tumour-cell level. Panels "
                        "B and C are spatial RNA-density maps for NDRG1 and LGALS1 in tumour cells; "
                        "the dashed rectangle marks the EMT niche, and density peaks there. Panel D "
                        "shows myofibroblast density across the same field — the EMT niche has high "
                        "myofibroblast presence. Panels E and F: VEGFA and IGFBP5 in fibroblasts. "
                        "VEGFA drives angiogenesis (so even though we're in collagen-poor regions, "
                        "they're being vascularised); IGFBP5 modulates IGF availability — pro-tumour. "
                        "Panel G: SPP1 in macrophages. SPP1 is osteopontin, a secreted phosphoprotein "
                        "that promotes tumour invasion and is a marker of macrophage co-option by "
                        "tumour. Panel H is the integrin 3D-signalling density map — highest in the "
                        "EMT niche. αvβ6 (the ITGB6+ pairing we discussed in Fig 6) is the active "
                        "integrin. Panel I is the integrative cartoon: tumour cells (LGALS1 + IGF2 + "
                        "IGFBP5) signal to myofibroblasts (αSMA+, contractile, secrete VEGFA), which "
                        "cooperate with macrophages (VC-, pro-tumour, SPP1+) to support a pre-invasive "
                        "tumour state. Every node is a candidate drug target — anti-galectin (GB1211 "
                        "in trials), anti-VEGFA (bevacizumab), anti-IGF1R (trial-stage), anti-SPP1 "
                        "(research-stage)."
                    ),
                    "extended_walkthrough": (
                        "BACKGROUND — why does this matter beyond NSCLC?\n\n"
                        "The tumour-fibroblast-macrophage cooperative described in panel I is a "
                        "GENERIC pattern that recurs across solid tumours — colorectal, pancreatic, "
                        "breast, lung. Each tumour type has its own version of the EMT niche, with "
                        "slightly different molecular details but the same cell-type architecture. "
                        "What's new here is the spatial resolution: we can actually point at where in "
                        "the tumour this cooperative is operating, in 3D, on a clinical sample.\n\n"
                        "GALECTIN-1 (LGALS1) IN DETAIL — galectin-1 binds β-galactoside-containing "
                        "glycans on the surface of T cells, induces T-cell apoptosis or anergy, "
                        "blocks T-cell receptor clustering. Tumour-cell-derived galectin-1 is one of "
                        "the major mechanisms by which tumours evade T-cell-mediated killing. GB1211 "
                        "(GalecTo) is an oral galectin-3 inhibitor in trials; specific galectin-1 "
                        "inhibitors are in earlier development.\n\n"
                        "SPP1 / OSTEOPONTIN IN DETAIL — SPP1 is a secreted glycoprotein that binds "
                        "integrins (CD44, αvβ3) on tumour and stromal cells, promotes migration, "
                        "matrix degradation, and angiogenesis. SPP1+ macrophages have been identified "
                        "in multiple cancers as a 'tumour-permissive' macrophage state — they look "
                        "anti-inflammatory but actively support tumour invasion. SPP1 is the most "
                        "commonly cited marker of 'pro-tumoural' macrophages in 2024-2025 literature.\n\n"
                        "αvβ6 INTEGRIN PATHWAY — αvβ6 binds RGD motifs on fibronectin and on the "
                        "latent form of TGFβ. When αvβ6 binds latent TGFβ-LAP, it activates TGFβ — "
                        "which then drives further EMT, fibroblast activation, and immunosuppression. "
                        "It's a positive-feedback loop. αvβ6 inhibitors have been tested in fibrotic "
                        "lung disease (idiopathic pulmonary fibrosis); the clinical experience there "
                        "informs cancer development.\n\n"
                        "WHY THE INTEGRATIVE CARTOON IS THE PUNCHLINE — the field has lots of papers "
                        "showing 'this gene is up in this tumour cluster'. What this paper provides "
                        "is a SPATIALLY-CONSISTENT, MULTI-CELL-TYPE story: tumour cells, fibroblasts, "
                        "and macrophages each express a complementary set of markers that ONLY "
                        "make sense as a coordinated module. That module is the EMT niche. That's "
                        "the niche-resolved druggable target.\n\n"
                        "THESIS-WORLD CONNECTION — for an infected lung simulation, replace 'tumour' "
                        "with 'pathogen-infected epithelium' and 'EMT' with 'epithelial activation'. "
                        "The fibroblast-macrophage cooperative around an infection focus is "
                        "structurally analogous; the simulation can use the same niche grammar."
                    ),
                },
            },

            # 10 — Critique / limitations (analogy layout: strengths | limitations)
            {
                "type": "analogy",
                "variant": "side_by_side",
                "title": "Strengths vs limitations — what would you trust?",
                "familiar_label": "What this paper PROVES",
                "familiar_body": (
                    "**3D multimodal on FFPE is feasible + templated.**\n\n"
                    "Single-section 2D loses [c]DC niches + T-cell continuity[/c].\n\n"
                    "Niche grammar (18 types → 10 niches) survives independent H&E.\n\n"
                    "DC-niche druggable network is spatially-resolved."
                ),
                "scientific_label": "What still needs validation",
                "scientific_body": (
                    "**n=1** — 2.28× and 51% from one patient.\n\n"
                    "EMT-collagen-poor — Fig 6F overlap is substantial.\n\n"
                    "CellChat conflates expression w/ signalling — needs [c]CRISPR validation[/c].\n\n"
                    "Generalisation to infection, fibrosis untested."
                ),
                "arrow_text": "vs",
                "citation_source": CITATION,
                "speaker_notes": {
                    "hook": "Where do we land after 7 figures?",
                    "key_claim": "Pentimalli proves the FEASIBILITY of routine 3D multimodal profiling on clinical FFPE and proves the CONSEQUENCES of staying in 2D. The n=1 cohort caveat is real, but the methodological contribution generalises immediately.",
                    "transition": "Open for questions.",
                    "script": (
                        "Two kinds of contributions to call out cleanly. First, the methodological "
                        "contribution: 34-section CosMx + SHG + STIM is now a templated workflow that "
                        "any group with FFPE access and a CosMx instrument can replicate. The cost is "
                        "real (~$150k of reagent + ~6-8 weeks of imaging) but feasible. Second, the "
                        "scientific contribution: 3D analysis is not optional for niche-level questions. "
                        "DC niches and T-cell continuity get systematically erased by single-section "
                        "2D analysis. If you're working on immunotherapy biomarkers, you cannot trust "
                        "single-section spatial readouts. The major limitation is n=1: every quantitative "
                        "claim — the 2.28× ratio, the 51% niche reassignment — comes from one patient. "
                        "Population-level claims await the larger 3D atlas the same group is presumably "
                        "building. For the thesis, the piece I take is the niche grammar: 18 cell types "
                        "→ 10 niches → druggable signalling network. The same grammar should describe "
                        "an infected alveolus — different cell types, different niches, but the same "
                        "logic of 'cells assemble into spatial communities that carry the actionable "
                        "signal'."
                    ),
                    "extended_walkthrough": (
                        "WHAT I'D PUSH BACK ON IN DISCUSSION — three points.\n\n"
                        "(1) The n=1 caveat is more serious than the paper acknowledges. Every panel "
                        "in the figures is from this one patient. If a different patient happened to "
                        "have a tumour with no DC niches, the 'DC niche is critical' framing might "
                        "not generalise. The team are presumably building a 3D atlas with more "
                        "patients; that paper will tell us what's reproducible.\n\n"
                        "(2) The 'EMT niche' is reproducibly observed in 2D papers. Pentimalli's "
                        "claim that it's 'collagen-poor not desmoplastic' is novel and specific. But "
                        "the box plot in Fig 6F shows substantial overlap between EMT and degraded "
                        "compartments; the difference is significant but not dramatic. I'd want to "
                        "see this on more patients before treating it as a strong claim.\n\n"
                        "(3) The CellChat 50-µm activity score conflates expression with signalling. "
                        "Just because L is expressed near R doesn't mean signalling is happening — "
                        "you also need ligand secretion, receptor internalisation, downstream signal. "
                        "The paper acknowledges this but the figures present the activity score as if "
                        "it were a direct readout. Functional validation (CRISPR knockout of MIF or "
                        "CCL19 in a model system) would close this gap.\n\n"
                        "WHAT THIS DOES TO THE THESIS PIPELINE — three concrete updates.\n\n"
                        "(1) Any in-silico tissue model we compare against spatial omics data must "
                        "now virtually-slice the simulation, not flatten the data. Default-2D "
                        "comparisons are biased.\n\n"
                        "(2) The niche grammar transfers: identify cell types, define neighbourhood "
                        "vectors, cluster, validate against pathologist-equivalent independent labels. "
                        "Same recipe for an infected alveolar duct as for an NSCLC.\n\n"
                        "(3) The receptor-ligand activity score with the spatial constraint is a "
                        "reusable method for comparing simulation outputs to real signalling. "
                        "Vivarium agents emitting virtual chemokine fields → simulated activity score "
                        "→ matched against the in-vivo activity score.\n\n"
                        "DISCUSSION QUESTIONS to seed — see next slide."
                    ),
                },
            },

            # 14 — Take-home (quote-style slide — single biggest claim)
            {
                "type": "quote",
                "quote": (
                    "Single-section 2D spatial omics is biased toward [c]false negatives[/c]. "
                    "The **dendritic-cell niches** and **T-cell continuity** that drive "
                    "immunotherapy response are recoverable [c]only in 3D[/c]."
                ),
                "attribution": "the take-home from Pentimalli & Rajewsky 2025",
                "speaker_notes": {
                    "hook": "If you remember one thing, it should be this.",
                    "key_claim": "Single-section 2D spatial omics is systematically biased toward false negatives at the niche level. The structures that predict immunotherapy response are 3D structures.",
                    "transition": "Five questions to seed discussion.",
                    "script": (
                        "If you take only one sentence away from this paper, take this one: "
                        "single-section 2D spatial omics is systematically biased toward "
                        "false negatives at the niche level. The dendritic-cell niches and "
                        "T-cell continuity that drive immunotherapy response are 3D "
                        "structures, and a 5-µm slice through them looks like nothing. "
                        "If your group is publishing niche-level claims from 2D data, you "
                        "now have to defend why your finding isn't an artefact of that "
                        "slicing plane. If your group is reading clinical spatial-omics "
                        "reports, you now have to read negatives as 'absence of signal in "
                        "this slice', not 'absence of structure'. Everything else in the "
                        "paper is downstream of this."
                    ),
                    "extended_walkthrough": (
                        "WHY THIS IS THE CORRECT TAKE-HOME — three reasons.\n\n"
                        "(1) It generalises beyond NSCLC. The geometric argument (2D disc "
                        "vs 3D ball) is independent of disease. Any niche-level claim from "
                        "any tumour, any infection, any fibrosis study built on 2D data is "
                        "vulnerable to the same critique.\n\n"
                        "(2) It is actionable. You can decide today: are you publishing 2D "
                        "niche claims (justify why 2D is enough) or 3D (you have a method "
                        "advantage). The implications for grant writing, paper framing, "
                        "and clinical interpretation are immediate.\n\n"
                        "(3) It is testable. The 51% niche-reassignment number is a "
                        "concrete quantitative prediction that other groups can replicate. "
                        "If they reproduce 30-70% reassignment in their own 3D atlases, "
                        "the claim holds. If they find <10%, the claim was n=1 noise. "
                        "Either way, the field moves forward."
                    ),
                },
            },

            # 15 — Discussion questions
            {
                "type": "text",
                "title": "Discussion seeds — where would you push this paper?",
                "bullets": [
                    "1.  **BUDGET** — would you spend on CosMx 3D / 1 sample, or 2D / 50 samples? Why?",
                    "2.  **ARTEFACT vs BIOLOGY** — is 'DC niche only in 3D' a methods artefact or real? How would you test?",
                    "3.  **GENERALISABILITY** — does the EMT-niche cooperative survive in metastatic disease? Treatment-naïve vs post-immunotherapy?",
                    "4.  **MINIMUM USEFUL 3D** — is 6 sections the floor? 3? 2 + interpolation?",
                    "5.  **CLINICAL PATHOLOGY** — is single-section H&E still enough for diagnostic decisions?",
                ],
                "speaker_notes": {
                    "hook": "Open the floor.",
                    "key_claim": "These five questions span methods, biology, and clinical translation — pick whichever feels most live for the room.",
                    "transition": "References + happy to dig into any figure.",
                },
            },

            # 12 — References
            {
                "type": "references",
                "title": "References",
                "references": REFS,
            },
        ],
    }


def main() -> int:
    out_path = OUT_DIR / "journal-club-pentimalli-2026-05-05.pptx"
    print(f"Building {out_path.name} ...")
    plan_dict = plan()

    missing: list[str] = []
    for s in plan_dict["slides"]:
        ip = s.get("image_path")
        if ip and not Path(ip).exists():
            missing.append(ip)
    if missing:
        print("MISSING figures:")
        for m in missing:
            print(f"  - {m}")
        return 1

    result = build_from_plan(
        plan_dict, out_path, write_marp=False, with_animations=True,
    )
    pptx_path = result["pptx"]
    print(f"  built: {pptx_path}  ({pptx_path.stat().st_size:,} bytes)")

    n_emph, n_analogy = apply_inline_emphasis(pptx_path)
    print(f"  emphasis: applied to {n_emph} paragraphs")
    print(f"  analogy animations: added to {n_analogy} slide(s)")

    audit = audit_deck(out_path)
    print(f"  severity:               {audit.severity}")
    print(f"  overlap_pairs:          {audit.total_overlapping_pairs}")
    print(f"  overflow_shapes:        {audit.total_overflowing_shapes}")
    print(f"  offslide_shapes:        {audit.total_offslide_shapes}")
    print(f"  over_bulleted_slides:   {audit.n_over_bulleted_slides}")
    print(f"  long_titles:            {audit.n_long_titles}")
    print(f"  figure_gap_slides:      {audit.figure_gap_slides}")
    print(f"  thin_slides:            {audit.thin_slides}")
    print(f"  n_total_images:         {audit.n_total_images} / {audit.n_slides} slides")
    if audit.severity == "fail":
        print("AUDIT FAILED")
        for s in audit.per_slide:
            flags = []
            if s.text_overflow_shapes:
                flags.append(f"overflow={s.text_overflow_shapes}")
            if s.offslide_shapes:
                flags.append(f"offslide={s.offslide_shapes}")
            if s.overlapping_shapes:
                flags.append(f"overlap={s.overlapping_shapes}")
            if s.title_too_long:
                flags.append("long_title")
            if s.over_bulleted:
                flags.append("over_bulleted")
            if flags:
                print(f"    slide {s.index} [{s.title[:55]!r}]: {flags}")
        return 2
    print("AUDIT PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
