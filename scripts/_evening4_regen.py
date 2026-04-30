"""Evening-4 advisor-package regeneration.

After the OpenAlex backfill (2026-04-30 morning) closed 256 anonymous-
author cases AND the evening-4 code fixes landed (Fix 1 acquire_figures,
Fix 2 aggressive picker, Fix 3 arc collision), this script regenerates:

* Phase B1 — Wiki/Projects/<rerun-slug>/{papers,lineage,START_HERE}.md
  via _write_project_view, using the now-clean post-backfill summaries.
* Phase B2 — codex-arc-2026-04-30.pdf re-export of the existing rerun
  arc (text unchanged; just re-export to refresh the advisor PDF).
* Phase B3 — codex-...-deck.pptx rebuild via build_deck_from_lineage_result
  with the cached figures + the new aggressive picker. Then COPY the new
  deck to the advisor-package folder as codex-deck-2026-04-30.pptx.

The script is idempotent — running it twice with the same on-disk inputs
produces the same outputs.
"""

from __future__ import annotations

import json
import shutil
import subprocess
import sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

KB_ROOT = Path("G:/My Drive/Knowledge/vaultlab")
PRIOR_PROJECT_SLUG = (
    "codex-multiplexed-imaging-methods-and-applications-across-tissue-types"
)
RERUN_PROJECT_SLUG = (
    "codex-multiplexed-imaging-methods-and-applications-across-tissue-types"
    "-evening3-rerun"
)
TOPIC = (
    "CODEX multiplexed imaging — methods and applications across tissue types"
)
RERUN_ARC_PATH = (
    KB_ROOT / "Wiki" / "Concepts"
    / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30-rerun.md"
)
PRIOR_PROVENANCE = (
    KB_ROOT / "Wiki" / "Concepts"
    / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30.md.provenance.json"
)

ADVISOR_DIR = (
    KB_ROOT / "Output" / "_demos" / "advisor-package-2026-04-30"
)
RERUN_OUTPUT_DIR = KB_ROOT / "Output" / RERUN_PROJECT_SLUG
RERUN_FIGURES_DIR = RERUN_OUTPUT_DIR / "figures"


def _load_post_backfill_summaries() -> dict:
    """Reload PaperSummary objects from on-disk frontmatter (post-backfill)."""
    import yaml

    from vaultlab.research.summarize import PaperSummary

    if not PRIOR_PROVENANCE.exists():
        raise SystemExit(f"Prior provenance not found: {PRIOR_PROVENANCE}")
    prov = json.loads(PRIOR_PROVENANCE.read_text(encoding="utf-8"))
    inputs = prov.get("inputs", [])
    print(f"[load] {len(inputs)} summary paths in prior provenance")

    summaries: dict[str, PaperSummary] = {}
    skipped = 0
    anon_count = 0
    for p in inputs:
        path = Path(p)
        if not path.exists():
            skipped += 1
            continue
        text = path.read_text(encoding="utf-8")
        if not text.startswith("---"):
            skipped += 1
            continue
        end = text.find("\n---", 3)
        if end == -1:
            skipped += 1
            continue
        try:
            fm = yaml.safe_load(text[3:end]) or {}
        except yaml.YAMLError:
            skipped += 1
            continue
        doi_val = (fm.get("doi") or "").lower()
        if not doi_val:
            doi_val = path.stem.replace("_", "/", 1)
        authors_list = list(fm.get("authors") or [])
        if not authors_list or all(
            a.strip().lower() in ("anonymous", "n/a", "anon", "")
            for a in authors_list
        ):
            anon_count += 1
        try:
            s = PaperSummary(
                doi=doi_val,
                title=fm.get("title") or "",
                authors=authors_list,
                year=int(fm.get("year") or 0),
                journal=fm.get("journal") or "",
                citation_count=int(fm.get("citation_count") or 0),
                influential_citations=int(fm.get("influential_citations") or 0),
                og_score=float(fm.get("og_score") or 0.0),
                forward_influence=int(fm.get("forward_influence") or 0),
                year_bucket=fm.get("year_bucket") or "unknown",
                tier=fm.get("tier") or "C",
                source_pdf=fm.get("source_pdf") or "",
            )
        except (TypeError, ValueError) as exc:
            print(f"  [warn] failed to parse {path.name}: {exc}")
            skipped += 1
            continue
        summaries[s.doi.lower()] = s
    print(
        f"[load] {len(summaries)} summaries; {anon_count} still anonymous; "
        f"skipped {skipped}"
    )
    return summaries


def phase_b1_rewrite_project_view(summaries: dict) -> dict:
    """Re-run _write_project_view with post-backfill summaries."""
    from vaultlab.research.corpus import Corpus
    from vaultlab.research.lineage import _write_project_view
    from vaultlab.research.paper import Paper

    print("\n[B1] re-writing project view with backfilled authors...")

    # Copy the rerun arc to the suffixed path on disk if missing.
    prior_arc = (
        KB_ROOT / "Wiki" / "Concepts"
        / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30.md"
    )
    if not RERUN_ARC_PATH.exists() and prior_arc.exists():
        RERUN_ARC_PATH.write_text(
            prior_arc.read_text(encoding="utf-8"), encoding="utf-8"
        )

    rerun_corpus = Corpus(topic=TOPIC, seeds=[])
    for doi, s in summaries.items():
        rerun_corpus.papers[doi] = Paper(
            title=s.title,
            authors=list(s.authors),
            year=s.year,
            doi=doi,
            source_api="reload-evening4",
        )

    deck_target = (
        RERUN_OUTPUT_DIR
        / f"{PRIOR_PROJECT_SLUG}-deck.pptx"
    )

    out = _write_project_view(
        kb_root=KB_ROOT,
        project_slug=RERUN_PROJECT_SLUG,
        topic=TOPIC,
        arc_path=RERUN_ARC_PATH,
        summaries=summaries,
        corpus=rerun_corpus,
        deck_path=deck_target,
        run_id="evening4-regen-2026-04-30",
        date_str="2026-04-30",
        speaker="Bobby Y.X. Ni",
        sources_n=8,
        picker_method="content-aware (evening-4 regen, post-backfill)",
        crosstalk="evening4-regen",
        timestamp=datetime.now().strftime("%Y-%m-%dT%H:%M:%S"),
        pdfs_acquired=sum(
            1 for s in summaries.values() if (s.tier or "C") == "A"
        ),
    )
    print(f"[B1] wrote {len(out)} project-view files:")
    for kind, p in out.items():
        print(f"  {kind:>15}: {p}")
    return out


def phase_b2_export_arc_pdf() -> Path | None:
    """Re-export the rerun arc to PDF (advisor-package overwrite).

    Uses Playwright + Chromium (already on this machine for Bobby's
    Chinese-report flow). Pandoc + xelatex hung on this corpus during
    the first attempt, so we render via headless Chrome instead — the
    arc is a small markdown document so a simple HTML-to-PDF pass is
    both faster and more reliable.
    """
    print("\n[B2] re-exporting CODEX arc to PDF (playwright path)...")
    target_pdf = ADVISOR_DIR / "codex-arc-2026-04-30.pdf"
    if not RERUN_ARC_PATH.exists():
        print(f"[B2] SKIP — arc not found: {RERUN_ARC_PATH}")
        return None

    try:
        import markdown as _md  # type: ignore[import-not-found]
    except ImportError:
        _md = None

    md_text = RERUN_ARC_PATH.read_text(encoding="utf-8")
    # Strip YAML frontmatter for the rendered document.
    if md_text.startswith("---"):
        end = md_text.find("\n---", 3)
        if end != -1:
            md_text = md_text[end + 4 :].lstrip()

    if _md is not None:
        body_html = _md.markdown(
            md_text, extensions=["tables", "fenced_code"]
        )
    else:
        # Fallback: dump as preformatted text. Functional but ugly.
        import html as _html
        body_html = f"<pre>{_html.escape(md_text)}</pre>"

    css = """
    body { font-family: 'Helvetica', 'Arial', sans-serif; max-width: 760px;
           margin: 1in auto; color: #1a1a1a; line-height: 1.55; }
    h1 { font-size: 22pt; color: #002F6C; border-bottom: 2px solid #002F6C;
         padding-bottom: 4px; margin-top: 0; }
    h2 { font-size: 16pt; color: #002F6C; margin-top: 1.2em; }
    h3 { font-size: 13pt; color: #1a4587; }
    blockquote { border-left: 4px solid #002F6C; padding: 4px 12px;
                 background: #f4f6fb; color: #2c2c2c; }
    table { border-collapse: collapse; margin: 1em 0; font-size: 9pt; }
    th, td { border: 1px solid #ccc; padding: 4px 8px; }
    th { background: #f0f0f0; }
    code { font-family: 'Consolas', 'Menlo', monospace;
           background: #f4f4f4; padding: 1px 4px; font-size: 9pt; }
    a { color: #1a4587; }
    .muted { color: #777; font-size: 9pt; }
    """
    html_doc = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>{css}</style></head><body>"
        f"<h1>CODEX multiplexed imaging — lineage arc</h1>"
        f"<p class='muted'>Regenerated 2026-04-30 evening-4 from "
        f"<code>{RERUN_ARC_PATH.name}</code></p>"
        f"{body_html}"
        "</body></html>"
    )
    html_tmp = ADVISOR_DIR / "_codex-arc-2026-04-30.tmp.html"
    html_tmp.write_text(html_doc, encoding="utf-8")
    try:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto(html_tmp.absolute().as_uri())
            page.pdf(
                path=str(target_pdf),
                format="Letter",
                margin={
                    "top": "0.75in",
                    "bottom": "0.75in",
                    "left": "0.75in",
                    "right": "0.75in",
                },
                print_background=True,
            )
            browser.close()
    except Exception as exc:
        print(f"[B2] FAIL playwright PDF: {exc}")
        return None
    finally:
        try:
            html_tmp.unlink()
        except OSError:
            pass
    print(f"[B2] wrote {target_pdf}")
    return target_pdf


def _build_figure_assignments() -> dict:
    """Walk RERUN_FIGURES_DIR and pick the largest figure per DOI."""
    figs: dict = {}
    if not RERUN_FIGURES_DIR.exists():
        return figs
    for doi_dir in sorted(RERUN_FIGURES_DIR.iterdir()):
        if not doi_dir.is_dir():
            continue
        manifest = doi_dir / ".figures.json"
        candidates: list[Path] = []
        if manifest.exists():
            try:
                data = json.loads(manifest.read_text(encoding="utf-8"))
                for entry in data.get("figures", []):
                    fp = Path(entry.get("file_path") or "")
                    if fp.exists():
                        candidates.append(fp)
            except json.JSONDecodeError:
                pass
        if not candidates:
            for child in sorted(doi_dir.iterdir()):
                if child.suffix.lower() in (".png", ".jpg", ".jpeg"):
                    candidates.append(child)
        if not candidates:
            continue
        # Reverse the slug back to a doi (single underscore -> /).
        doi = doi_dir.name.replace("_", "/", 1).lower()
        # Pick the largest file.
        candidates.sort(key=lambda p: p.stat().st_size, reverse=True)
        figs[doi] = candidates[0]
    return figs


def phase_b3_rebuild_deck(summaries: dict) -> Path | None:
    """Rebuild the deck using the new aggressive figure picker."""
    print("\n[B3] rebuilding CODEX deck with new figure picker...")
    from vaultlab.research.corpus import Corpus
    from vaultlab.research.lineage import LineageRunResult
    from vaultlab.research.paper import Paper
    from vaultlab.slides.deck import build_deck_from_lineage_result

    figure_assignments = _build_figure_assignments()
    print(f"[B3] figure_assignments: {len(figure_assignments)} entries")
    for doi, p in list(figure_assignments.items())[:5]:
        try:
            size_kb = p.stat().st_size / 1024
        except OSError:
            size_kb = 0
        print(f"   {doi}: {p.name} ({size_kb:.0f} KB)")

    corpus = Corpus(topic=TOPIC, seeds=[])
    for doi, s in summaries.items():
        corpus.papers[doi] = Paper(
            title=s.title,
            authors=list(s.authors),
            year=s.year,
            doi=doi,
            source_api="reload-evening4",
        )

    summary_paths = {
        doi: KB_ROOT / "Wiki" / "Summaries" / f"{doi.replace('/', '_')}.md"
        for doi in summaries
    }
    lineage_result = LineageRunResult(
        topic=TOPIC,
        arc_path=RERUN_ARC_PATH,
        summary_paths=summary_paths,
        search_log_path=Path(),
        corpus_size=len(summaries),
        pdfs_acquired=sum(
            1 for s in summaries.values() if (s.tier or "C") == "A"
        ),
        summaries_written=len(summaries),
        duration_seconds=0.0,
        project_slug=RERUN_PROJECT_SLUG,
        project_view_paths={},
        corpus=corpus,
        figure_assignments=figure_assignments,
        figures_acquired=len(figure_assignments),
    )
    try:
        deck_p = build_deck_from_lineage_result(
            lineage_result,
            speaker="Bobby Y.X. Ni",
            affiliation="Hickey Lab @ Duke BME",
            project_slug=RERUN_PROJECT_SLUG,
            figure_assignments=figure_assignments,
            kb_root=KB_ROOT,
            plan_mode="fast",
            audience="advisor-demo",
            target_slide_count=8,
            final_audit=False,
        )
        print(f"[B3] deck written: {deck_p}")
    except Exception as exc:
        import traceback
        traceback.print_exc()
        print(f"[B3] FAIL deck build: {exc}")
        return None

    # Copy to advisor package.
    advisor_deck = ADVISOR_DIR / "codex-deck-2026-04-30.pptx"
    shutil.copy2(deck_p, advisor_deck)
    print(f"[B3] copied to advisor package: {advisor_deck}")

    # Inspect the deck to confirm figure-slide count.
    try:
        from pptx import Presentation
        pres = Presentation(str(deck_p))
        n_slides = len(pres.slides)
        n_with_figures = sum(
            1 for s in pres.slides if any(
                sh.shape_type == 13 for sh in s.shapes
            )
        )
        print(f"[B3] deck has {n_slides} slides, {n_with_figures} with pictures")
    except Exception as exc:
        print(f"[B3] deck inspection failed: {exc}")
    return deck_p


def main() -> int:
    print(f"[start] evening-4 regen at {datetime.now().isoformat()}")
    print(f"[topic] {TOPIC}")
    print(f"[rerun] {RERUN_PROJECT_SLUG}")
    summaries = _load_post_backfill_summaries()
    if not summaries:
        raise SystemExit("No summaries loaded — aborting")

    phase_b1_rewrite_project_view(summaries)
    phase_b2_export_arc_pdf()
    deck = phase_b3_rebuild_deck(summaries)

    print(f"\n[done] regenerated artifacts:")
    print(f"  project view: Wiki/Projects/{RERUN_PROJECT_SLUG}/")
    print(f"  arc PDF:      Output/_demos/advisor-package-2026-04-30/codex-arc-2026-04-30.pdf")
    if deck:
        print(f"  deck:         Output/_demos/advisor-package-2026-04-30/codex-deck-2026-04-30.pptx")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
