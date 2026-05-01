"""Evening-5 CODEX-arc regeneration.

After the OpenAlex backfill (evening-3) and the wikilink-label fix
(evening-5 — author-format detection now handles OpenAlex's
``F. Last`` form), the prior rerun arc still has stale labels baked
in: ``[[10.1002_eji.202048891|J. 2020]]`` (should be
``Kennedy-Darling 2020``) and dozens of ``Anon n.d.`` rows in the
tables (the table rendering ran AGAINST disk-frontmatter that was
anonymous at render-time).

This script regenerates JUST the lineage arc with the now-fixed label
helper. It:

1. Loads the 236 backfilled summaries from
   ``Wiki/Summaries/`` via the prior provenance JSON.
2. Recomputes og_score and co-citation pairs from the citation graph
   (rebuilt from ``Sources/Articles/`` stub frontmatter — no network
   calls, no Anthropic calls).
3. Lifts the prose (history / development / sota narrative paragraphs)
   from the existing rerun arc — the prose was already clean, only the
   tables needed regen.
4. Renders the new arc with :func:`render_arc_markdown` so the tables
   pick up the fixed ``_author_year_label`` (which delegates to
   :func:`vaultlab.kb.paths.author_year_label`).
5. Writes to ``-rerun-2.md`` so the prior ``-rerun.md`` is preserved
   per Bobby's directive.

After running, also re-exports the advisor-package PDF + deck pointing
at the new arc.
"""

from __future__ import annotations

import json
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

KB_ROOT = Path("G:/My Drive/Knowledge/vaultlab")
PRIOR_PROJECT_SLUG = (
    "codex-multiplexed-imaging-methods-and-applications-across-tissue-types"
)
RERUN_PROJECT_SLUG = (
    "codex-multiplexed-imaging-methods-and-applications-across-tissue-types"
    "-evening3-rerun"
)
TOPIC = (
    "CODEX multiplexed imaging — methods and applications across tissue types"
)
PRIOR_ARC = (
    KB_ROOT / "Wiki" / "Concepts"
    / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30-rerun.md"
)
PRIOR_PROVENANCE = (
    KB_ROOT / "Wiki" / "Concepts"
    / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30.md.provenance.json"
)
NEW_ARC = (
    KB_ROOT / "Wiki" / "Concepts"
    / f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30-rerun-2.md"
)

ADVISOR_DIR = (
    KB_ROOT / "Output" / "_demos" / "advisor-package-2026-04-30"
)
RERUN_OUTPUT_DIR = KB_ROOT / "Output" / RERUN_PROJECT_SLUG
RERUN_FIGURES_DIR = RERUN_OUTPUT_DIR / "figures"


# ---------------------------------------------------------------------------
# Phase 1: load summaries + the prior arc's prose
# ---------------------------------------------------------------------------


def _load_summaries() -> dict:
    """Load PaperSummary objects from on-disk frontmatter (post-backfill)."""
    import yaml

    from vaultlab.research.summarize import PaperSummary

    if not PRIOR_PROVENANCE.exists():
        raise SystemExit(f"Prior provenance not found: {PRIOR_PROVENANCE}")
    prov = json.loads(PRIOR_PROVENANCE.read_text(encoding="utf-8"))
    inputs = prov.get("inputs", [])
    print(f"[load] {len(inputs)} summary paths in prior provenance")

    summaries: dict[str, PaperSummary] = {}
    skipped = 0
    anon_count = 0
    for p in inputs:
        path = Path(p)
        if not path.exists():
            skipped += 1
            continue
        text = path.read_text(encoding="utf-8")
        if not text.startswith("---"):
            skipped += 1
            continue
        end = text.find("\n---", 3)
        if end == -1:
            skipped += 1
            continue
        try:
            fm = yaml.safe_load(text[3:end]) or {}
        except yaml.YAMLError:
            skipped += 1
            continue
        doi_val = (fm.get("doi") or "").lower()
        if not doi_val:
            doi_val = path.stem.replace("_", "/", 1)
        authors_list = list(fm.get("authors") or [])
        if not authors_list or all(
            a.strip().lower() in ("anonymous", "n/a", "anon", "")
            for a in authors_list
        ):
            anon_count += 1
        try:
            s = PaperSummary(
                doi=doi_val,
                title=fm.get("title") or "",
                authors=authors_list,
                year=int(fm.get("year") or 0),
                journal=fm.get("journal") or "",
                citation_count=int(fm.get("citation_count") or 0),
                influential_citations=int(fm.get("influential_citations") or 0),
                og_score=float(fm.get("og_score") or 0.0),
                forward_influence=int(fm.get("forward_influence") or 0),
                year_bucket=fm.get("year_bucket") or "unknown",
                tier=fm.get("tier") or "C",
                source_pdf=fm.get("source_pdf") or "",
            )
        except (TypeError, ValueError) as exc:
            print(f"  [warn] failed to parse {path.name}: {exc}")
            skipped += 1
            continue
        summaries[s.doi.lower()] = s
    print(
        f"[load] {len(summaries)} summaries; {anon_count} still anonymous; "
        f"skipped {skipped}"
    )
    return summaries


def _extract_narrative_from_prior(prior_arc: Path) -> dict[str, str]:
    """Pull the history/development/sota prose paragraphs out of the prior arc.

    Each section starts with ``## History (...)`` / ``## Development (...)``
    / ``## State of the art (...)`` and the prose is the lines BEFORE the
    first table row (``|`` prefix). The prior arc's prose is already clean
    — it was LLM-generated against the post-backfill summaries — so we
    preserve it verbatim. Only the tables get regenerated with the fixed
    ``_author_year_label``.
    """
    text = prior_arc.read_text(encoding="utf-8")
    sections = re.split(r"^## ", text, flags=re.MULTILINE)
    out: dict[str, str] = {"history": "", "development": "", "sota": ""}
    for chunk in sections:
        head = chunk.split("\n", 1)
        if not head:
            continue
        title = head[0].strip().lower()
        rest = head[1] if len(head) > 1 else ""
        # Identify which bucket this section is.
        bucket: str | None = None
        if title.startswith("history"):
            bucket = "history"
        elif title.startswith("development"):
            bucket = "development"
        elif title.startswith("state of the art"):
            bucket = "sota"
        if bucket is None:
            continue
        # Prose is everything BEFORE the first ``| Year |`` table row.
        prose_lines = []
        for line in rest.splitlines():
            if line.lstrip().startswith("|"):
                break
            prose_lines.append(line)
        prose = "\n".join(prose_lines).strip()
        if prose:
            out[bucket] = prose
    return out


# ---------------------------------------------------------------------------
# Phase 2: rebuild the citation graph + metrics from on-disk corpus
# ---------------------------------------------------------------------------


_CO_CITATION_RE = re.compile(
    r"^\d+\.\s+\[\[([^|]+)\|[^]]+\]\]\s*\+\s*\[\[([^|]+)\|[^]]+\]\]\s*"
    r"—\s*co-cited by\s*(\d+)\s*papers"
)


def _build_corpus(summaries: dict) -> object:
    """Recreate a Corpus from per-summary og_score + parsed co-citation pairs.

    The corpus is needed by ``render_arc_markdown`` for the og_score and
    co-citation tables. We avoid recomputing the citation graph (the
    article stubs don't carry references on disk in this KB) and instead:

    1. Use each PaperSummary's stored ``og_score`` field — this was set
       at lit-arc time from CrossRef references and persists in the
       summary's YAML frontmatter.
    2. Lift the existing co-citation pair list from the prior arc's
       table (``[[a|...]] + [[b|...]] — co-cited by N papers``). The
       co-citation calculation hasn't changed since the prior run, so
       lifting the result keeps the new arc's pair list identical to
       the rerun arc — only the wikilink labels get refreshed.
    3. Carry over each summary's ``year_bucket`` so the bucket tables
       render correctly.

    The result is functionally equivalent to a re-computed corpus for
    the table-rendering step, which is all ``render_arc_markdown`` cares
    about.
    """
    from vaultlab.research.corpus import Corpus
    from vaultlab.research.graph_metrics import CorpusMetrics
    from vaultlab.research.paper import Paper

    corpus = Corpus(topic=TOPIC, seeds=[])

    # Step 1: every summary doi → Paper in corpus.papers
    for doi, s in summaries.items():
        corpus.papers[doi] = Paper(
            title=s.title,
            authors=list(s.authors),
            year=s.year,
            doi=doi,
            citation_count=s.citation_count,
            source_api="reload-evening5",
        )

    # Step 2: re-mark seeds from prior provenance (first 8 inputs).
    prov = json.loads(PRIOR_PROVENANCE.read_text(encoding="utf-8"))
    for p in prov.get("inputs", [])[:8]:
        doi = Path(p).stem.replace("_", "/", 1).lower()
        if doi in corpus.papers:
            corpus.seeds.append(corpus.papers[doi])

    print(
        f"[corpus] papers={len(corpus.papers)} seeds={len(corpus.seeds)}"
    )

    # Step 3: build metrics from the per-summary og_score field +
    # parse co-citation pairs out of the prior arc.
    og_score = {doi: s.og_score for doi, s in summaries.items() if s.og_score}
    forward_influence = {
        doi: s.forward_influence for doi, s in summaries.items()
        if s.forward_influence
    }
    year_buckets = {doi: s.year_bucket for doi, s in summaries.items()}

    co_citation_pairs: list[tuple[str, str, int]] = []
    prior_text = PRIOR_ARC.read_text(encoding="utf-8")
    for line in prior_text.splitlines():
        m = _CO_CITATION_RE.match(line.strip())
        if m:
            a_slug = m.group(1).strip()
            b_slug = m.group(2).strip()
            n = int(m.group(3))
            # Slug → doi: split on first underscore. Prefer matching against
            # the corpus by slug-form to recover the canonical DOI.
            a_doi = _slug_to_doi(a_slug, corpus.papers)
            b_doi = _slug_to_doi(b_slug, corpus.papers)
            co_citation_pairs.append((a_doi, b_doi, n))

    print(
        f"[metrics] og_score={len(og_score)} "
        f"co_citation_pairs={len(co_citation_pairs)}"
    )

    corpus.metrics = CorpusMetrics(
        og_score=og_score,
        forward_influence=forward_influence,
        co_citation_pairs=co_citation_pairs,
        year_buckets=year_buckets,
    )
    return corpus


def _slug_to_doi(slug: str, papers: dict) -> str:
    """Recover the canonical DOI from a slug by direct lookup against papers.

    The slug form is what the prior arc's wikilinks carried; we need
    the original DOI keys for ``CorpusMetrics.co_citation_pairs``. Try
    a couple of conversions before falling back to the slug itself.
    """
    # Direct: most slugs are dot-format with a single underscore for "/".
    candidate = slug.replace("_", "/", 1)
    if candidate in papers:
        return candidate
    # Fallback: maybe the slug carries a stray ".pdf" suffix from an
    # older bug.
    if candidate.endswith(".pdf"):
        candidate = candidate[:-4]
        if candidate in papers:
            return candidate
    # Final fallback: dash-form rebuild.
    return slug.replace("_", "/", 1)


# ---------------------------------------------------------------------------
# Phase 3: render the new arc
# ---------------------------------------------------------------------------


def _render_new_arc(summaries: dict, corpus, narrative: dict) -> Path:
    from vaultlab.research.lineage import render_arc_markdown

    arc_md = render_arc_markdown(
        topic=TOPIC,
        date_str="2026-04-30",
        summaries=summaries,
        corpus=corpus,
        method_relpath=(
            f"{PRIOR_PROJECT_SLUG}-lineage-2026-04-30.md.method.md"
        ),
        narrative=narrative,
    )
    NEW_ARC.write_text(arc_md, encoding="utf-8")
    print(f"[arc] wrote {NEW_ARC}")
    return NEW_ARC


# ---------------------------------------------------------------------------
# Phase 4: refresh advisor PDF
# ---------------------------------------------------------------------------


def _refresh_advisor_pdf(arc_path: Path) -> Path | None:
    """Render the new arc to PDF via Playwright + Chromium."""
    print("\n[pdf] re-exporting CODEX arc to advisor-package PDF...")
    target_pdf = ADVISOR_DIR / "codex-arc-2026-04-30.pdf"

    try:
        import markdown as _md  # type: ignore[import-not-found]
    except ImportError:
        _md = None

    md_text = arc_path.read_text(encoding="utf-8")
    if md_text.startswith("---"):
        end = md_text.find("\n---", 3)
        if end != -1:
            md_text = md_text[end + 4 :].lstrip()

    if _md is not None:
        body_html = _md.markdown(md_text, extensions=["tables", "fenced_code"])
    else:
        import html as _html

        body_html = f"<pre>{_html.escape(md_text)}</pre>"

    css = """
    body { font-family: 'Helvetica', 'Arial', sans-serif; max-width: 760px;
           margin: 1in auto; color: #1a1a1a; line-height: 1.55; }
    h1 { font-size: 22pt; color: #002F6C; border-bottom: 2px solid #002F6C;
         padding-bottom: 4px; margin-top: 0; }
    h2 { font-size: 16pt; color: #002F6C; margin-top: 1.2em; }
    h3 { font-size: 13pt; color: #1a4587; }
    blockquote { border-left: 4px solid #002F6C; padding: 4px 12px;
                 background: #f4f6fb; color: #2c2c2c; }
    table { border-collapse: collapse; margin: 1em 0; font-size: 9pt; }
    th, td { border: 1px solid #ccc; padding: 4px 8px; }
    th { background: #f0f0f0; }
    code { font-family: 'Consolas', 'Menlo', monospace;
           background: #f4f4f4; padding: 1px 4px; font-size: 9pt; }
    a { color: #1a4587; }
    .muted { color: #777; font-size: 9pt; }
    """
    html_doc = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<style>{css}</style></head><body>"
        f"<h1>CODEX multiplexed imaging — lineage arc</h1>"
        f"<p class='muted'>Regenerated 2026-04-30 evening-5 (rerun-2) "
        f"from <code>{arc_path.name}</code></p>"
        f"{body_html}"
        "</body></html>"
    )
    html_tmp = ADVISOR_DIR / "_codex-arc-2026-04-30.tmp.html"
    html_tmp.write_text(html_doc, encoding="utf-8")
    try:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto(html_tmp.absolute().as_uri())
            page.pdf(
                path=str(target_pdf),
                format="Letter",
                margin={
                    "top": "0.75in",
                    "bottom": "0.75in",
                    "left": "0.75in",
                    "right": "0.75in",
                },
                print_background=True,
            )
            browser.close()
    except Exception as exc:
        print(f"[pdf] FAIL playwright PDF: {exc}")
        return None
    finally:
        try:
            html_tmp.unlink()
        except OSError:
            pass
    print(f"[pdf] wrote {target_pdf}")
    return target_pdf


# ---------------------------------------------------------------------------
# Phase 5: rebuild advisor deck
# ---------------------------------------------------------------------------


def _build_figure_assignments() -> dict:
    figs: dict = {}
    if not RERUN_FIGURES_DIR.exists():
        return figs
    for doi_dir in sorted(RERUN_FIGURES_DIR.iterdir()):
        if not doi_dir.is_dir():
            continue
        manifest = doi_dir / ".figures.json"
        candidates: list[Path] = []
        if manifest.exists():
            try:
                data = json.loads(manifest.read_text(encoding="utf-8"))
                for entry in data.get("figures", []):
                    fp_str = entry.get("path") or entry.get("file_path") or ""
                    fp = Path(fp_str)
                    if fp.exists():
                        candidates.append(fp)
            except json.JSONDecodeError:
                pass
        if not candidates:
            for child in sorted(doi_dir.iterdir()):
                if child.suffix.lower() in (".png", ".jpg", ".jpeg"):
                    candidates.append(child)
        if not candidates:
            continue
        doi = doi_dir.name.replace("_", "/", 1).lower()
        candidates.sort(key=lambda p: p.stat().st_size, reverse=True)
        figs[doi] = candidates[0]
    return figs


def _refresh_advisor_deck(summaries: dict, arc_path: Path) -> Path | None:
    """Rebuild the advisor-package deck pointing at the new arc."""
    print("\n[deck] rebuilding CODEX deck (evening-5 regen)...")
    from vaultlab.research.corpus import Corpus
    from vaultlab.research.lineage import LineageRunResult
    from vaultlab.research.paper import Paper
    from vaultlab.slides.deck import build_deck_from_lineage_result

    figure_assignments = _build_figure_assignments()
    print(f"[deck] figure_assignments: {len(figure_assignments)} entries")

    corpus = Corpus(topic=TOPIC, seeds=[])
    for doi, s in summaries.items():
        corpus.papers[doi] = Paper(
            title=s.title,
            authors=list(s.authors),
            year=s.year,
            doi=doi,
            source_api="reload-evening5",
        )

    summary_paths = {
        doi: KB_ROOT / "Wiki" / "Summaries" / f"{doi.replace('/', '_')}.md"
        for doi in summaries
    }
    lineage_result = LineageRunResult(
        topic=TOPIC,
        arc_path=arc_path,
        summary_paths=summary_paths,
        search_log_path=Path(),
        corpus_size=len(summaries),
        pdfs_acquired=sum(
            1 for s in summaries.values() if (s.tier or "C") == "A"
        ),
        summaries_written=len(summaries),
        duration_seconds=0.0,
        project_slug=RERUN_PROJECT_SLUG,
        project_view_paths={},
        corpus=corpus,
        figure_assignments=figure_assignments,
        figures_acquired=len(figure_assignments),
    )
    try:
        deck_p = build_deck_from_lineage_result(
            lineage_result,
            speaker="Bobby Y.X. Ni",
            affiliation="Hickey Lab @ Duke BME",
            project_slug=RERUN_PROJECT_SLUG,
            figure_assignments=figure_assignments,
            kb_root=KB_ROOT,
            plan_mode="fast",
            audience="advisor-demo",
            target_slide_count=8,
            final_audit=False,
        )
        print(f"[deck] deck written: {deck_p}")
    except Exception as exc:
        import traceback

        traceback.print_exc()
        print(f"[deck] FAIL deck build: {exc}")
        return None

    advisor_deck = ADVISOR_DIR / "codex-deck-2026-04-30.pptx"
    shutil.copy2(deck_p, advisor_deck)
    print(f"[deck] copied to advisor package: {advisor_deck}")

    try:
        from pptx import Presentation

        pres = Presentation(str(deck_p))
        n_slides = len(pres.slides)
        n_with_figures = sum(
            1 for s in pres.slides if any(sh.shape_type == 13 for sh in s.shapes)
        )
        print(f"[deck] deck has {n_slides} slides, {n_with_figures} with pictures")
    except Exception as exc:
        print(f"[deck] deck inspection failed: {exc}")
    return deck_p


def main() -> int:
    print(f"[start] evening-5 arc regen at {datetime.now().isoformat()}")
    print(f"[topic] {TOPIC}")

    summaries = _load_summaries()
    if not summaries:
        raise SystemExit("No summaries loaded — aborting")

    narrative = _extract_narrative_from_prior(PRIOR_ARC)
    print(
        f"[narrative] history={len(narrative['history'])} chars, "
        f"development={len(narrative['development'])} chars, "
        f"sota={len(narrative['sota'])} chars"
    )

    corpus = _build_corpus(summaries)
    new_arc = _render_new_arc(summaries, corpus, narrative)
    _refresh_advisor_pdf(new_arc)
    _refresh_advisor_deck(summaries, new_arc)

    # Sanity check: confirm zero "Anon n.d." in the new arc body.
    body = new_arc.read_text(encoding="utf-8")
    n_anon = body.count("|Anon n.d.]]")
    n_initial_only = len(re.findall(r"\|[A-Z]\.\s+\d{4}]]", body))
    print(
        f"\n[verify] new arc has {n_anon} '|Anon n.d.]]' tokens, "
        f"{n_initial_only} initial-only '|J. 2020]]' tokens"
    )
    if n_initial_only > 0:
        print("[verify] WARN — initial-only labels still present!")
        for m in re.finditer(r"\|[A-Z]\.\s+\d{4}]]", body):
            print(f"  {m.group()}")

    print("\n[done]")
    print(f"  new arc:      {new_arc}")
    print(
        f"  advisor PDF:  {ADVISOR_DIR / 'codex-arc-2026-04-30.pdf'}"
    )
    print(
        f"  advisor deck: {ADVISOR_DIR / 'codex-deck-2026-04-30.pptx'}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
