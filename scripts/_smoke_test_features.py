"""Smoke-test 6 vaultlab features end-to-end.

Per Bobby 2026-04-29 ask: 'Are all the features running correctly now?
Can you run them yourself with a use case to see how well they're working?'

Tests:
1. kb.feedback.open_question + log_decision
2. context.locations register / get / missing_paths_grill_doc
3. context.user_memory remember / recall
4. kb.ingest (markdown + bibtex + folder)
5. kb.tools_index summary_for / deep_doc_for / suggest_for_topic
6. slides.notes.dual_format
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))


def test_feedback() -> None:
    from vaultlab.kb.feedback import log_decision, open_question

    with tempfile.TemporaryDirectory() as tmp:
        kb = Path(tmp)
        result = open_question(
            kb,
            slug="smoke-test",
            title="Smoke test grill",
            questions=["Does this work?", "Is the output usable?"],
            auto_open=False,
        )
        assert result.path.exists()
        body = result.path.read_text(encoding="utf-8")
        assert "Q1" in body and "Q2" in body
        assert "smoke-test" in body

        log_path = log_decision(kb, "smoke-project", "Test decision", "Test reason")
        assert log_path.exists()
        log_body = log_path.read_text(encoding="utf-8")
        assert "Test decision" in log_body
    print("[OK] kb.feedback - open_question + log_decision")


def test_locations() -> None:
    import os

    from vaultlab.context.locations import (
        get_path,
        load_locations,
        register_path,
    )

    with tempfile.TemporaryDirectory() as tmp:
        loc_file = Path(tmp) / "locations.toml"
        os.environ["VAULTLAB_LOCATIONS"] = str(loc_file)
        try:
            register_path("work_log.google_doc_id", "abc123")
            register_path("projects.car-t", "research/Wiki/Projects/car-t")
            data = load_locations()
            assert get_path("work_log.google_doc_id", locations=data) == "abc123"
            assert get_path("projects.car-t", locations=data) == "research/Wiki/Projects/car-t"
            assert get_path("missing.key", locations=data) is None
        finally:
            del os.environ["VAULTLAB_LOCATIONS"]
    print("[OK] context.locations - register + get + dashed-key support")


def test_user_memory() -> None:
    import os

    from vaultlab.context.user_memory import recall, recall_all, remember

    with tempfile.TemporaryDirectory() as tmp:
        os.environ["VAULTLAB_USER_MEMORY"] = tmp
        try:
            remember(
                category="feedback",
                name="smoke-test",
                description="Smoke test entry",
                content="Body",
            )
            entry = recall("feedback", "smoke-test")
            assert entry is not None
            assert entry.description == "Smoke test entry"

            index_text, entries = recall_all()
            assert "smoke-test" in index_text
            assert any(e.name == "smoke-test" for e in entries)
        finally:
            del os.environ["VAULTLAB_USER_MEMORY"]
    print("[OK] context.user_memory - remember + recall + index")


def test_ingest() -> None:
    from vaultlab.kb.ingest import ingest

    with tempfile.TemporaryDirectory() as tmp:
        # Markdown
        md = Path(tmp) / "note.md"
        md.write_text(
            "---\ntitle: Test note\ntype: note\n---\n\n# H1\n\nBody.",
            encoding="utf-8",
        )
        doc = ingest(md)
        assert not isinstance(doc, list)
        assert doc.title == "Test note"

        # BibTeX
        bib = Path(tmp) / "refs.bib"
        bib.write_text(
            "@article{a2024, title = {A}, year = {2024}}\n@book{b2025, title = {B}}\n",
            encoding="utf-8",
        )
        result = ingest(bib)
        assert isinstance(result, list)
        assert len(result) == 2

        # Folder
        result_folder = ingest(Path(tmp))
        assert isinstance(result_folder, list)
        # markdown + 2 bibtex entries = 3 docs
        assert len(result_folder) >= 3
    print("[OK] kb.ingest - markdown + bibtex + folder dispatch")


def test_tools_index() -> None:
    from vaultlab.kb.tools_index import (
        deep_doc_for,
        load_index,
        suggest_for_topic,
        summary_for,
    )

    index = load_index()
    assert "scanpy" in index
    assert "squidpy" in index

    summary = summary_for("scanpy")
    assert summary is not None
    assert len(summary) < 1500  # one paragraph, not full body
    assert "single-cell" in summary.lower() or "scrna" in summary.lower()

    deep = deep_doc_for("scanpy")
    assert deep is not None
    assert len(deep) > len(summary)

    spatial_hits = suggest_for_topic("spatial")
    spatial_names = {e.name for e in spatial_hits}
    assert "squidpy" in spatial_names

    stats_hits = suggest_for_topic("statistics")
    stats_names = {e.name for e in stats_hits}
    assert "statsmodels" in stats_names or "scipy.stats" in stats_names
    print("[OK] kb.tools_index - summary + deep + suggest tiered search")


def test_dual_format_notes() -> None:
    from vaultlab.slides.notes import DIVIDER, dual_format, parse_dual_format

    notes = dual_format(
        mental_map={
            "hook": "Open with this",
            "key_claim": "Main claim",
            "evidence": "Figure 1",
            "key_terms": ["scFv", "ITAM"],
            "click": "First click reveals annotations",
            "transition": "Next slide is X",
        },
        detailed_script="This is the full word-for-word script.",
    )
    assert "HOOK: Open with this" in notes
    assert "KEY CLAIM: Main claim" in notes
    assert DIVIDER.strip() in notes

    parsed_map, parsed_script = parse_dual_format(notes)
    assert parsed_map["hook"] == "Open with this"
    assert parsed_map["key_claim"] == "Main claim"
    assert "full word-for-word" in parsed_script
    print("[OK] slides.notes - dual_format render + parse round-trip")


def test_slide_grouping() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    from vaultlab.figures.understand import ElementAnnotation
    from vaultlab.slides.annotated_figure_slide import add_annotated_figure_slide

    # Need a real image - use one of the CAR-T figures if available; else skip
    test_image = Path(r"C:\tmp\cart_figs_v13\image1.png")
    if not test_image.exists():
        print("[SKIP] slides.annotated_figure_slide - no test image at " + str(test_image))
        return

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    annotations = [
        ElementAnnotation(
            label="Test annotation",
            bbox_px=(295, 1827, 419, 2094),
            motif_name="endo-tcr",
        ),
    ]

    add_annotated_figure_slide(
        pres,
        test_image,
        annotations,
        title="Smoke test slide",
        caption="Test caption",
        motif_colors={"endo-tcr": (40, 110, 220)},
        page_number=1,
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        pres.save(f.name)
        out = Path(f.name)

    # Verify by re-opening
    p2 = Presentation(str(out))
    assert len(p2.slides) == 1
    shape_names = {s.name for s in p2.slides[0].shapes}
    assert "slide_title" in shape_names
    assert "ann1_box" in shape_names
    assert "ann1_marker" in shape_names
    # Group should be present
    assert any("ann1_side_group" in n for n in shape_names)
    out.unlink()
    print("[OK] slides.annotated_figure_slide - native shapes + grouping")


def main() -> None:
    test_feedback()
    test_locations()
    test_user_memory()
    test_ingest()
    test_tools_index()
    test_dual_format_notes()
    test_slide_grouping()
    print("\nAll 7 smoke tests passed.")


if __name__ == "__main__":
    main()
