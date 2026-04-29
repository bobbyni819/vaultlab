"""Demo: native-shape annotated PPTX.

Per Bobby 2026-04-29: each annotation must be a separate PowerPoint object
so it can be animated, moved, edited, or removed without re-rendering.

This script builds a deck with:
- Title slide (explanation)
- image1 (TCR/TAA/CAR) - 9 annotations, each as native shapes
- image10 (TME) - 8 annotations from the rigor-protocol redo

In PowerPoint, every annotation appears in the Selection pane as named
shapes (ann1_box, ann1_marker, ann1_label, ann1_side_marker etc.) so user
can apply entrance animations from the Animations tab.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from vaultlab.figures.understand import ElementAnnotation
from vaultlab.figures.understand.whitespace import find_marker_offset
from vaultlab.slides.annotated_figure_slide import (
    HICKEY_LAB_LAYOUT,
    add_annotated_figure_slide,
)
from vaultlab.slides.notes import dual_format
from vaultlab.slides.themes.hickey_lab import (
    hickey_lab_template_path,
    load_hickey_lab_presentation,
)


def auto_offset_annotations(
    image_path: Path,
    annotations: list[ElementAnnotation],
) -> list[ElementAnnotation]:
    """Run the whitespace finder over annotations and update marker_offset_px.

    Per Bobby 2026-04-29 figure-annotation decision tree: don't guess marker
    placement; programmatically pick a nearby whitespace zone.
    """
    placed_marker_bboxes: list[tuple[int, int, int, int]] = []
    out: list[ElementAnnotation] = []
    for ann in annotations:
        offset = find_marker_offset(
            image_path,
            ann.bbox_px,
            avoid_other_bboxes=tuple(placed_marker_bboxes),
        )
        new_ann = ElementAnnotation(
            label=ann.label,
            bbox_px=ann.bbox_px,
            explanation=ann.explanation,
            motif_name=ann.motif_name,
            confidence=ann.confidence,
            use_box=ann.use_box,
            marker_offset_px=offset,
        )
        out.append(new_ann)
        if offset is not None:
            x0, y0 = ann.bbox_px[0] + offset[0], ann.bbox_px[1] + offset[1]
            placed_marker_bboxes.append((x0, y0, x0 + 120, y0 + 120))
    return out

# ---------------------------------------------------------------------------
# Image1 annotations - 9 elements (lifted from v3 script with corrections)
# ---------------------------------------------------------------------------

IMAGE1 = Path(r"C:\tmp\cart_figs_v13\image1.png")

# Per-annotation marker offsets chosen to avoid overlaps + place markers in
# whitespace, NOT on top of motifs. Bobby's 2026-04-29 feedback: be flexible.
# Some annotations also drop the box (use_box=False) when a tight marker
# pointer is cleaner than wrapping the element in a rectangle.

IMAGE1_ANNOTATIONS = [
    # #1 - Endogenous TCR panel a (left edge of panel; whitespace to the left)
    ElementAnnotation(
        label="Endogenous TCR (panel a, competing)",
        bbox_px=(295, 1827, 419, 2094),
        motif_name="endo-tcr",
        marker_offset_px=(-180, -50),  # to the left in whitespace
    ),
    # #2 - Introduced TCR panel a. Bobby 2026-04-29 v5: skip the box; the green
    # color is distinctive enough that just a marker is cleaner.
    ElementAnnotation(
        label="Introduced TCR (panel a)",
        bbox_px=(998, 1853, 1192, 2188),
        motif_name="intro-tcr",
        use_box=False,
    ),
    # #3 - MHC Class I (top of receptor stack; place marker UP and LEFT to clear
    # the introduced-TCR region below)
    ElementAnnotation(
        label="MHC Class I (panel a)",
        bbox_px=(1072, 1666, 1185, 1848),
        motif_name="mhc",
        marker_offset_px=(-220, -150),
    ),
    # #4 - CD3 chains (right of introduced TCR; place marker BELOW the box, in
    # the membrane area)
    ElementAnnotation(
        label="CD3 chains (panel a)",
        bbox_px=(1206, 2100, 1261, 2392),
        motif_name="cd3",
        marker_offset_px=(150, 320),  # below + slightly right
    ),
    # #5 - TAA panel b (narrow horizontal band; box is awkward, skip box and
    # use a marker pointing at it from the LEFT in panel-b whitespace)
    ElementAnnotation(
        label="TAA / aberrantly expressed (panel b)",
        bbox_px=(2627, 1820, 2770, 1900),
        motif_name="taa",
        use_box=False,
        marker_offset_px=(-280, -200),  # well left and above
    ),
    # #6 - Endogenous TCR panel b (move marker LEFT into whitespace, below
    # where #5's marker is)
    ElementAnnotation(
        label="Endogenous TCR (panel b, primary)",
        bbox_px=(2728, 1857, 3309, 2124),
        motif_name="endo-tcr",
        marker_offset_px=(-300, 200),  # left of box, mid-height
    ),
    # #7 - scFv panel c (tall narrow column; default top-left works since
    # the area above scFv is whitespace between tumor cell and the construct)
    ElementAnnotation(
        label="scFv / antigen recognition (panel c)",
        bbox_px=(4558, 1837, 4715, 2430),
        motif_name="endo-tcr",
        marker_offset_px=(-280, 150),  # left of box
    ),
    # #8 - Co-stim panel c (move marker LEFT to clear the construct stack)
    ElementAnnotation(
        label="Co-stim domain (panel c, intracellular)",
        bbox_px=(4558, 2436, 4614, 2589),
        motif_name="intro-tcr",
        marker_offset_px=(-280, 0),
    ),
    # #9 - CD3-zeta panel c (move marker LEFT, separated vertically from #8)
    ElementAnnotation(
        label="CD3-zeta activation (panel c)",
        bbox_px=(4558, 2624, 4620, 2746),
        motif_name="cd3",
        marker_offset_px=(-280, 50),
    ),
]

IMAGE1_COLORS = {
    "endo-tcr": (40, 110, 220),  # electric blue
    "intro-tcr": (60, 200, 60),  # neon green
    "mhc": (130, 70, 180),  # purple
    "cd3": (220, 110, 30),  # orange
    "taa": (210, 50, 70),  # red-magenta
}


# ---------------------------------------------------------------------------
# Image10 annotations - rigor-protocol manual bboxes (image is 2035x1676)
# ---------------------------------------------------------------------------

IMAGE10 = Path(r"C:\tmp\cart_figs_v13\image10.png")

IMAGE10_ANNOTATIONS = [
    ElementAnnotation(
        label="Suppressive cells (MDSC, Treg, TAM)",
        bbox_px=(420, 0, 1380, 510),
        motif_name="suppressive-cells",
        explanation="",
    ),
    ElementAnnotation(
        label="PD-1/PDL-1 + CTLA-4/CD86 detail",
        bbox_px=(20, 360, 720, 800),
        motif_name="checkpoint",
        explanation="",
    ),
    ElementAnnotation(
        label="Soluble inhibitors: IL-10, TGF-beta",
        bbox_px=(1300, 60, 2030, 480),
        motif_name="soluble-inhibitors",
        explanation="",
    ),
    ElementAnnotation(
        label="Tumor antigen heterogeneity",
        bbox_px=(1330, 660, 2030, 1100),
        motif_name="antigen-heterogeneity",
        explanation="",
    ),
    ElementAnnotation(
        label="Metabolic suppression (low O2/pH/nutrients)",
        bbox_px=(820, 670, 1270, 1010),
        motif_name="metabolic",
        explanation="",
    ),
    ElementAnnotation(
        label="Dysregulated vasculature: VCAM/ICAM down",
        bbox_px=(0, 1090, 820, 1530),
        motif_name="vasculature",
        explanation="",
    ),
    ElementAnnotation(
        label="Physical barriers: ECM, CAF, IFP",
        bbox_px=(1260, 1040, 2030, 1620),
        motif_name="physical-barriers",
        explanation="",
    ),
    ElementAnnotation(
        label="Tumor mass (central pink blob, NOT vessel)",
        bbox_px=(420, 180, 1740, 1430),
        motif_name="tumor-mass",
        explanation="",
    ),
]

IMAGE10_COLORS = {
    "suppressive-cells": (180, 60, 200),
    "checkpoint": (220, 80, 60),
    "soluble-inhibitors": (40, 160, 200),
    "antigen-heterogeneity": (220, 100, 30),
    "metabolic": (50, 150, 150),
    "vasculature": (200, 30, 30),
    "physical-barriers": (130, 90, 50),
    "tumor-mass": (220, 130, 150),
}


# ---------------------------------------------------------------------------
# Notes content (concise; expanded later with the dual-format speaker notes
# once Bobby has answered Q10-Q13 in the slide-construction grill)
# ---------------------------------------------------------------------------

IMAGE1_NOTES = dual_format(
    mental_map={
        "hook": "T cells engineered to find tumors come in three architectural flavors.",
        "key_claim": "TCR, TAA, and CAR therapies differ fundamentally in their MHC-dependence.",
        "evidence": "Three-panel BioRender comparison (VanNoy 2025) on this slide.",
        "key_terms": ["scFv", "MHC class I", "ITAM", "CD3-zeta"],
        "click": "9 annotation pairs reveal in order: panel-a items first (#1-4), then panel b (#5-6), then panel c (#7-9).",
        "transition": "Next slide: how the CAR construct is engineered across 5 generations.",
    },
    detailed_script=(
        "Engineered T cells recognize tumor antigens through three distinct architectural "
        "strategies, shown in this three-panel BioRender comparison from VanNoy 2025.\n\n"
        "In panel a, TCR therapy preserves the endogenous T-cell signaling architecture. "
        "An exogenous T-cell receptor is introduced into the cell (the green dimer marked #2), "
        "competing with the endogenous TCR (#1, blue) for the same MHC class I peptide complex "
        "(#3, purple). Both receptors couple through the CD3 chains (#4, orange) for signal "
        "transduction. The therapeutic gain comes from supplying a high-affinity TCR specific "
        "for a known tumor peptide.\n\n"
        "Panel b shows TAA therapy. Here the antigen is an aberrantly expressed self-protein "
        "(#5, the small red bar at the top of the receptor stack), presented via MHC class I to "
        "the endogenous TCR (#6, blue). No exogenous receptor is introduced - the existing T "
        "cell is the effector. The challenge is finding tumor antigens that are both "
        "tumor-restricted and MHC-presentable.\n\n"
        "Panel c shows CAR therapy. The chimeric antigen receptor uses an scFv (#7) to bind "
        "tumor surface antigen directly, MHC-independently. Inside the cell, the construct "
        "carries a co-stimulatory domain (#8, green) and a CD3-zeta-style activation domain "
        "(#9, orange). This MHC-independence is consistent with broader antigen reach but also "
        "with off-target on-tumor toxicity when the antigen is shared with healthy tissue."
    ),
)

IMAGE10_NOTES = dual_format(
    mental_map={
        "hook": "The tumor microenvironment fights back against CAR-T in seven distinct ways.",
        "key_claim": "Solid-tumor CAR-T efficacy is limited by converging suppressive mechanisms, not just antigen escape.",
        "evidence": "8 labeled callout regions surrounding the central tumor mass on this slide.",
        "key_terms": [
            "MDSC",
            "Treg",
            "TAM",
            "PD-1/PDL-1",
            "CTLA-4/CD86",
            "VCAM-1/ICAM-1",
            "ECM",
            "CAF",
            "IFP",
        ],
        "click": "8 annotations enter in clockwise order from suppressive cells (#1) to tumor mass (#8).",
        "transition": "Next slide: how engineering strategies address each of these mechanisms.",
    },
    detailed_script=(
        "Solid-tumor CAR-T efficacy is bounded not by the engineering of the cells but by the "
        "tumor microenvironment they enter. This figure summarizes seven converging suppressive "
        "forces around a central tumor mass.\n\n"
        "The top-center callout (#1) shows suppressive cells - MDSCs, Tregs, and TAMs that engage "
        "the CAR T cell directly through inhibitory ligands. The left circle (#2) zooms into the "
        "checkpoint interaction: PD-1/PDL-1 and CTLA-4/CD86 axes between cancer cells and "
        "dendritic cells, both targets of approved checkpoint inhibitors.\n\n"
        "The top-right callout (#3) shows soluble inhibitors. IL-10 and TGF-beta diffuse from "
        "the suppressive cell cluster, broadly damping immune activation. To the right (#4), "
        "tumor antigen heterogeneity is consistent with antigen-negative escape variants under "
        "single-target CAR pressure.\n\n"
        "The central blue ellipses (#5) capture metabolic suppression: hypoxia, low pH, nutrient "
        "depletion, and accumulated lactate or adenosine. These metabolic conditions limit T-cell "
        "function independently of any inhibitory ligand.\n\n"
        "Bottom-left (#6), dysregulated tumor vasculature with reduced VCAM-1/ICAM-1 limits T-cell "
        "extravasation into the parenchyma. Bottom-right (#7), physical barriers - dense ECM, "
        "cancer-associated fibroblasts, and high interstitial fluid pressure - mechanically "
        "exclude T cells from reaching the tumor cells they target.\n\n"
        "Note on annotation #8: this captures the actual tumor mass, the pink cellular blob, NOT "
        "the red blood vessel running through it - a common visual misread of this kind of figure."
    ),
)


def add_title_slide(pres: Presentation) -> None:
    blank = pres.slide_layouts[6]
    s = pres.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "VaultLab - Native-Shape Annotations"
    run.font.name = "Arial"
    run.font.size = Pt(40)
    run.font.bold = True

    sub = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5))
    sf = sub.text_frame
    sf.word_wrap = True
    text = (
        "Each numbered marker, bounding box, and side label is a SEPARATE "
        "PowerPoint shape - right-click any of them in the Selection pane "
        "(View > Selection Pane) to:\n\n"
        "  - Add entrance animations from Animations tab\n"
        "  - Move them around the slide\n"
        "  - Edit the text directly\n"
        "  - Delete individual ones\n\n"
        "Naming: ann1_box / ann1_marker / ann1_label / ann1_side_marker so all "
        "shapes for one annotation can be selected together."
    )
    for i, line in enumerate(text.split("\n\n")):
        p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(15)
        run.font.color.rgb = RGBColor(60, 60, 60)


SECTIONS = [
    "Background",
    "Antigen recognition",
    "TME challenges",
    "Engineering strategies",
    "Outlook",
]


def main() -> None:
    if hickey_lab_template_path() is not None:
        pres = load_hickey_lab_presentation(theme="dark")
        print(f"Using Hickey Lab template: {hickey_lab_template_path()}")
        layout = HICKEY_LAB_LAYOUT
    else:
        pres = Presentation()
        print("Hickey Lab template not bundled; using plain Presentation.")
        layout = None  # use default

    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    add_title_slide(pres)

    # Auto-pick marker offsets via whitespace finder (no more guessing!)
    image1_anns = auto_offset_annotations(IMAGE1, IMAGE1_ANNOTATIONS)
    image10_anns = auto_offset_annotations(IMAGE10, IMAGE10_ANNOTATIONS)

    kwargs = {} if layout is None else {"layout": layout}

    add_annotated_figure_slide(
        pres,
        IMAGE1,
        image1_anns,
        title="Three architectures of antigen recognition",
        caption="Three-panel comparison from VanNoy 2025 (BioRender). "
        "Each numbered annotation is a native PowerPoint shape.",
        motif_colors=IMAGE1_COLORS,
        notes=IMAGE1_NOTES,
        page_number=2,
        sections=SECTIONS,
        current_section_idx=1,
        **kwargs,
    )

    add_annotated_figure_slide(
        pres,
        IMAGE10,
        image10_anns,
        title="Seven mechanisms of TME suppression against CAR-T",
        caption="8 labeled callouts surrounding the central tumor. "
        "Try View > Selection Pane to see all shapes.",
        motif_colors=IMAGE10_COLORS,
        notes=IMAGE10_NOTES,
        page_number=3,
        sections=SECTIONS,
        current_section_idx=2,
        **kwargs,
    )

    out = Path(r"C:\Users\bobby\Downloads\car_t_decks\figure_understanding_native_v6.pptx")
    pres.save(out)
    print(f"Wrote -> {out}")


if __name__ == "__main__":
    main()
