"""v3 figure-understanding pipeline run on ALL 8 figures.

Per Bobby 2026-04-29 review:
- No leader lines, color-matched markers + labels, white halos, big fonts
- Tumor-cell motif added (light pink/red BioRender tumor blobs)
- Run on all 8 demo figures, build single PPTX

Architecture: shared BioRender palette dictionary; each figure picks the motifs
relevant to it + supplies concept-matching rules. The renderer (v3) gets the
motif->color mapping so on-figure markers + side labels match the actual
visual element colors.
"""

from __future__ import annotations

import sys
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from vaultlab.figures.understand import (
    ColorMotif,
    ElementAnnotation,
    Region,
    extract_regions,
    merge_regions,
)
from vaultlab.figures.understand.render import render_annotated_figure_v3

INPUT_DIR = Path(r"C:\tmp\cart_figs_v13")
OUT_DIR = Path(r"C:\tmp\cart_figs_v13_annotated_v3")
OUT_DIR.mkdir(exist_ok=True)


# ---------------------------------------------------------------------------
# Shared BioRender-style palette - motifs that recur across many figures.
# Each motif also has a display color that the renderer uses for markers + labels.
# ---------------------------------------------------------------------------

# (motif, color)
PALETTE: dict[str, tuple[ColorMotif, tuple[int, int, int]]] = {
    "neon-green": (
        ColorMotif("neon-green", (80, 140), 0.40, 0.40, 0.00003),
        (60, 200, 60),
    ),
    "electric-blue": (
        ColorMotif("electric-blue", (195, 240), 0.30, 0.35, 0.00003),
        (40, 110, 220),
    ),
    "orange": (
        ColorMotif("orange", (15, 40), 0.45, 0.45, 0.00003),
        (220, 110, 30),
    ),
    "red-magenta": (
        ColorMotif("red-magenta", (340, 360), 0.35, 0.40, 0.00003),
        (210, 50, 70),
    ),
    "purple-violet": (
        ColorMotif("purple-violet", (260, 295), 0.18, 0.28, 0.00003),
        (130, 70, 180),
    ),
    "yellow-gold": (
        ColorMotif("yellow-gold", (45, 65), 0.40, 0.55, 0.00003),
        (200, 165, 30),
    ),
    "tumor-pink": (
        # Light pink BioRender tumor cells - low saturation, high value
        ColorMotif("tumor-pink", (340, 360), 0.10, 0.65, 0.0008),
        (220, 130, 150),
    ),
    "cyan-light": (
        ColorMotif("cyan-light", (160, 200), 0.20, 0.55, 0.0001),
        (40, 170, 200),
    ),
}


def motifs_for(names: list[str]) -> list[ColorMotif]:
    return [PALETTE[n][0] for n in names]


def colors_for(names: list[str]) -> dict[str, tuple[int, int, int]]:
    return {n: PALETTE[n][1] for n in names}


# ---------------------------------------------------------------------------
# Per-figure plans
# ---------------------------------------------------------------------------


@dataclass
class FigurePlan:
    image_filename: str
    title: str
    overall_caption: str
    long_explanation: str
    motif_names: list[str]
    # Function takes merged regions, image width/height; returns annotations.
    concept_matcher: Callable[[list[Region], int, int], list[ElementAnnotation]]
    dilation_px: int = 25


def _largest_in(regions: list[Region], motif: str, panel_filter=None) -> Region | None:
    candidates = [r for r in regions if r.motif_name == motif]
    if panel_filter is not None:
        candidates = [r for r in candidates if panel_filter(r)]
    if not candidates:
        return None
    return max(candidates, key=lambda r: r.area_px)


def _all_in(regions: list[Region], motif: str, panel_filter=None) -> list[Region]:
    candidates = [r for r in regions if r.motif_name == motif]
    if panel_filter is not None:
        candidates = [r for r in candidates if panel_filter(r)]
    return candidates


def _box_union(regs: list[Region]) -> tuple[int, int, int, int]:
    return (
        min(r.bbox_px[0] for r in regs),
        min(r.bbox_px[1] for r in regs),
        max(r.bbox_px[2] for r in regs),
        max(r.bbox_px[3] for r in regs),
    )


# Helper builders for panel filters
def panel_third(W: int, idx: int):
    """idx=0,1,2 for left/middle/right thirds."""
    edge = W // 3
    if idx == 0:
        return lambda r: r.bbox_px[0] < edge
    if idx == 1:
        return lambda r: edge <= r.bbox_px[0] < 2 * edge
    return lambda r: r.bbox_px[0] >= 2 * edge


def panel_half_h(H: int, top: bool):
    return (lambda r: r.bbox_px[3] < H // 2) if top else (lambda r: r.bbox_px[1] >= H // 2)


# ---------------------------------------------------------------------------
# CONCEPT MATCHERS - one per figure
# ---------------------------------------------------------------------------


def matcher_image1(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """3-panel: TCR / TAA / CAR T."""
    out: list[ElementAnnotation] = []

    a, b, c = panel_third(W, 0), panel_third(W, 1), panel_third(W, 2)

    # Panel a - TCR therapy
    if r := min(
        (r for r in regions if r.motif_name == "electric-blue" and a(r)),
        default=None,
        key=lambda r: r.bbox_px[0],
    ):
        out.append(
            ElementAnnotation(
                label="Endogenous TCR (panel a, competing)",
                bbox_px=r.bbox_px,
                motif_name="electric-blue",
                explanation="Native αβ TCR competing with the introduced TCR for MHC-peptide.",
                confidence=0.85,
            )
        )
    if r := _largest_in(regions, "neon-green", a):
        out.append(
            ElementAnnotation(
                label="Introduced TCR (panel a)",
                bbox_px=r.bbox_px,
                motif_name="neon-green",
                explanation="Engineered TCR - the therapeutic receptor in TCR therapy.",
                confidence=0.90,
            )
        )
    if r := _largest_in(regions, "purple-violet", a):
        out.append(
            ElementAnnotation(
                label="MHC Class I (panel a)",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="MHC class I + peptide; sits on top of the introduced TCR.",
                confidence=0.85,
            )
        )
    if r := _largest_in(regions, "orange", a):
        out.append(
            ElementAnnotation(
                label="CD3 chains (panel a)",
                bbox_px=r.bbox_px,
                motif_name="orange",
                explanation="CD3 ε/γ/δ subunits - signal transduction.",
                confidence=0.80,
            )
        )
    # Panel b - TAA therapy
    if r := _largest_in(regions, "red-magenta", b):
        out.append(
            ElementAnnotation(
                label="Aberrantly expressed protein (TAA, panel b)",
                bbox_px=r.bbox_px,
                motif_name="red-magenta",
                explanation="Tumor-associated antigen presented via MHC class I.",
                confidence=0.75,
            )
        )
    if r := _largest_in(regions, "electric-blue", b):
        out.append(
            ElementAnnotation(
                label="Endogenous TCR (panel b, primary receptor)",
                bbox_px=r.bbox_px,
                motif_name="electric-blue",
                explanation="In TAA therapy the endogenous TCR is the active receptor.",
                confidence=0.85,
            )
        )
    # Panel c - CAR T
    if r := _largest_in(regions, "electric-blue", c):
        out.append(
            ElementAnnotation(
                label="scFv (panel c, antigen recognition)",
                bbox_px=r.bbox_px,
                motif_name="electric-blue",
                explanation="Single-chain variable fragment - binds antigen MHC-independently.",
                confidence=0.65,
            )
        )
    if r := _largest_in(regions, "neon-green", c):
        out.append(
            ElementAnnotation(
                label="Co-stim domain (panel c)",
                bbox_px=r.bbox_px,
                motif_name="neon-green",
                explanation="Co-stimulatory domain (CD28 / 4-1BB) - intracellular.",
                confidence=0.70,
            )
        )
    if r := _largest_in(regions, "orange", c):
        out.append(
            ElementAnnotation(
                label="CD3-ζ activation (panel c)",
                bbox_px=r.bbox_px,
                motif_name="orange",
                explanation="CD3-ζ-style activation domain - bottom of CAR signaling stack.",
                confidence=0.70,
            )
        )
    # NEW: Tumor cells (Bobby's request)
    tumor_regions = [r for r in regions if r.motif_name == "tumor-pink"]
    if tumor_regions:
        union = _box_union(tumor_regions)
        out.append(
            ElementAnnotation(
                label="Tumor cells (top of all 3 panels)",
                bbox_px=union,
                motif_name="tumor-pink",
                explanation="Pink BioRender tumor-cell blobs at the top of each panel.",
                confidence=0.80,
            )
        )
    return out


def matcher_image2(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """5 generations of CAR. Each generation is a vertical stack of receptors."""
    out: list[ElementAnnotation] = []
    blue = sorted(_all_in(regions, "electric-blue"), key=lambda r: r.bbox_px[0])
    green = sorted(_all_in(regions, "neon-green"), key=lambda r: r.bbox_px[0])
    orange = sorted(_all_in(regions, "orange"), key=lambda r: r.bbox_px[0])
    red = sorted(_all_in(regions, "red-magenta"), key=lambda r: r.bbox_px[0])

    # Approximate columns: leftmost = TCR; next 5 = gen 1-5 CARs
    # This is heuristic - works for typical layouts
    if blue:
        out.append(
            ElementAnnotation(
                label="TCR αβ (reference, leftmost)",
                bbox_px=blue[0].bbox_px,
                motif_name="electric-blue",
                explanation="Native TCR for reference - αβ heterodimer + CD3 + ζ chains.",
                confidence=0.70,
            )
        )
    if red and len(red) >= 1:
        # The TAA at top of CAR scFvs
        top_red = sorted(red, key=lambda r: r.bbox_px[1])[0]
        out.append(
            ElementAnnotation(
                label="TAA (target antigen, top of CAR)",
                bbox_px=top_red.bbox_px,
                motif_name="red-magenta",
                explanation="Tumor-associated antigen at top of CAR construct.",
                confidence=0.65,
            )
        )
    if green:
        out.append(
            ElementAnnotation(
                label="scFv (antigen recognition domain)",
                bbox_px=green[len(green) // 2].bbox_px,
                motif_name="neon-green",
                explanation="scFv variable fragment in each CAR generation.",
                confidence=0.60,
            )
        )
    if orange:
        # CD3-zeta activation domain - lowest orange
        lowest_orange = max(orange, key=lambda r: r.bbox_px[1])
        out.append(
            ElementAnnotation(
                label="CD3-ζ activation domain (bottom of CARs)",
                bbox_px=lowest_orange.bbox_px,
                motif_name="orange",
                explanation="Activation cassette below the membrane in all CAR generations.",
                confidence=0.70,
            )
        )
    if green and len(green) > 1:
        bottom_green = max(green, key=lambda r: r.bbox_px[1])
        out.append(
            ElementAnnotation(
                label="Co-stimulatory domain (gen 2+ addition)",
                bbox_px=bottom_green.bbox_px,
                motif_name="neon-green",
                explanation="Co-stim addition that distinguishes 2nd-gen+ CARs from 1st-gen.",
                confidence=0.55,
            )
        )
    return out


def matcher_image3(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """Signals 1/2/3."""
    out: list[ElementAnnotation] = []
    a, b, c = panel_third(W, 0), panel_third(W, 1), panel_third(W, 2)

    # Panel a - Signal 1 (TAA + CAR)
    if r := _largest_in(regions, "red-magenta", a):
        out.append(
            ElementAnnotation(
                label="Signal 1: TAA antigen",
                bbox_px=r.bbox_px,
                motif_name="red-magenta",
                explanation="TAA bound by CAR scFv -> activation.",
                confidence=0.75,
            )
        )
    if r := _largest_in(regions, "electric-blue", a):
        out.append(
            ElementAnnotation(
                label="Signal 1: scFv (CAR recognition)",
                bbox_px=r.bbox_px,
                motif_name="electric-blue",
                explanation="Antigen recognition domain.",
                confidence=0.65,
            )
        )
    # Panel b - Signal 2 (co-stim)
    if r := _largest_in(regions, "neon-green", b):
        out.append(
            ElementAnnotation(
                label="Signal 2: co-stim domain",
                bbox_px=r.bbox_px,
                motif_name="neon-green",
                explanation="Co-stim signal -> PI3K/AKT, NF-kB, MAPK.",
                confidence=0.70,
            )
        )
    if r := _largest_in(regions, "orange", b):
        out.append(
            ElementAnnotation(
                label="Signal 2: activation domain",
                bbox_px=r.bbox_px,
                motif_name="orange",
                explanation="Activation cascade below membrane.",
                confidence=0.70,
            )
        )
    # Panel c - Signal 3 (cytokines)
    if r := _largest_in(regions, "cyan-light", c):
        out.append(
            ElementAnnotation(
                label="Signal 3: cytokine release",
                bbox_px=r.bbox_px,
                motif_name="cyan-light",
                explanation="IL-2 family cytokines drive proliferation + persistence.",
                confidence=0.60,
            )
        )
    elif r := _largest_in(regions, "electric-blue", c):
        out.append(
            ElementAnnotation(
                label="Signal 3: cytokine field",
                bbox_px=r.bbox_px,
                motif_name="electric-blue",
                explanation="Cytokine release region.",
                confidence=0.55,
            )
        )
    return out


def matcher_image7(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """CAR-M 6-panel."""
    out: list[ElementAnnotation] = []
    top = panel_half_h(H, top=True)
    bottom = panel_half_h(H, top=False)

    if r := _largest_in(regions, "purple-violet", top):
        out.append(
            ElementAnnotation(
                label="M1 macrophage (top)",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="M1-polarized macrophage in cytotoxicity panel A.",
                confidence=0.65,
            )
        )
    if r := _largest_in(regions, "tumor-pink", top):
        out.append(
            ElementAnnotation(
                label="Tumor cell (panel A)",
                bbox_px=r.bbox_px,
                motif_name="tumor-pink",
                explanation="Tumor cell being attacked by CAR-M.",
                confidence=0.70,
            )
        )
    if r := _largest_in(regions, "tumor-pink", bottom):
        out.append(
            ElementAnnotation(
                label="Tumor mass (panel E - infiltration)",
                bbox_px=r.bbox_px,
                motif_name="tumor-pink",
                explanation="Solid tumor with CAR-M infiltrating.",
                confidence=0.70,
            )
        )
    cyan = _all_in(regions, "cyan-light")
    if cyan:
        # Cytokine dots
        union = _box_union(cyan)
        out.append(
            ElementAnnotation(
                label="Cytokine field",
                bbox_px=union,
                motif_name="cyan-light",
                explanation="Released cytokines (IL-6, IL-8, IL-12, TNF-α/β).",
                confidence=0.50,
            )
        )
    return out


def matcher_image8(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """7-step cancer-immunity cycle."""
    out: list[ElementAnnotation] = []
    if r := _largest_in(regions, "tumor-pink"):
        out.append(
            ElementAnnotation(
                label="Tumor mass (center)",
                bbox_px=r.bbox_px,
                motif_name="tumor-pink",
                explanation="Tumor cells releasing antigens at cycle start.",
                confidence=0.70,
            )
        )
    purple = _all_in(regions, "purple-violet")
    for i, r in enumerate(sorted(purple, key=lambda r: r.area_px, reverse=True)[:3]):
        out.append(
            ElementAnnotation(
                label=f"DC / immune cell {i + 1}",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="Immune cell (DC, T cell) in the cycle.",
                confidence=0.50,
            )
        )
    return out


def matcher_image10(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """TME suppressive mechanisms."""
    out: list[ElementAnnotation] = []
    if r := _largest_in(regions, "tumor-pink"):
        out.append(
            ElementAnnotation(
                label="Tumor mass (center)",
                bbox_px=r.bbox_px,
                motif_name="tumor-pink",
                explanation="Central tumor; suppressive mechanisms surround it.",
                confidence=0.75,
            )
        )
    if r := _largest_in(regions, "purple-violet"):
        out.append(
            ElementAnnotation(
                label="Suppressive cells region",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="MDSCs / Tregs / TAMs cluster.",
                confidence=0.55,
            )
        )
    if r := _largest_in(regions, "orange"):
        out.append(
            ElementAnnotation(
                label="CAR T cells",
                bbox_px=r.bbox_px,
                motif_name="orange",
                explanation="CAR T cells trying to engage; orange BioRender icons.",
                confidence=0.55,
            )
        )
    return out


def matcher_image13(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """CAR T/NK/M comparison."""
    out: list[ElementAnnotation] = []
    left = lambda r: r.bbox_px[0] < W // 2
    right = lambda r: r.bbox_px[0] >= W // 2

    if r := _largest_in(regions, "tumor-pink", left):
        out.append(
            ElementAnnotation(
                label="Tumor microenvironment (panel A)",
                bbox_px=r.bbox_px,
                motif_name="tumor-pink",
                explanation="Heterogeneous TME with all cell types.",
                confidence=0.70,
            )
        )
    purple_right = sorted(_all_in(regions, "purple-violet", right), key=lambda r: r.bbox_px[1])
    for i, r in enumerate(purple_right[:3]):
        labels = ["CAR-T cell", "CAR-NK cell", "CAR-M cell"]
        out.append(
            ElementAnnotation(
                label=labels[i] if i < len(labels) else f"Engineered cell {i + 1}",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="Engineered effector cell - panel B row.",
                confidence=0.55,
            )
        )
    if r := _largest_in(regions, "neon-green"):
        out.append(
            ElementAnnotation(
                label="Granzymes / cytotoxic granules",
                bbox_px=r.bbox_px,
                motif_name="neon-green",
                explanation="Cytotoxic granule contents released onto target cells.",
                confidence=0.50,
            )
        )
    return out


def matcher_image14(regions: list[Region], W: int, H: int) -> list[ElementAnnotation]:
    """Notch + syn-Notch."""
    out: list[ElementAnnotation] = []
    top = panel_half_h(H, top=True)
    bottom = panel_half_h(H, top=False)

    if r := _largest_in(regions, "yellow-gold", top):
        out.append(
            ElementAnnotation(
                label="EGF Repeats (NECD top)",
                bbox_px=r.bbox_px,
                motif_name="yellow-gold",
                explanation="Yellow-striped EGF-like repeats - extracellular ligand-binding region.",
                confidence=0.85,
            )
        )
    if r := _largest_in(regions, "red-magenta", top):
        out.append(
            ElementAnnotation(
                label="ANK / NICD region",
                bbox_px=r.bbox_px,
                motif_name="red-magenta",
                explanation="Ankyrin repeats in NICD intracellular cassette.",
                confidence=0.65,
            )
        )
    if r := _largest_in(regions, "purple-violet", top):
        out.append(
            ElementAnnotation(
                label="NICD core",
                bbox_px=r.bbox_px,
                motif_name="purple-violet",
                explanation="Notch intracellular domain - signaling cassette.",
                confidence=0.65,
            )
        )
    if r := _largest_in(regions, "neon-green", bottom):
        out.append(
            ElementAnnotation(
                label="syn-Notch CAR (bottom panel)",
                bbox_px=r.bbox_px,
                motif_name="neon-green",
                explanation="Engineered second-antigen receptor in syn-Notch system.",
                confidence=0.55,
            )
        )
    return out


# ---------------------------------------------------------------------------
# Plans
# ---------------------------------------------------------------------------


PLANS: list[FigurePlan] = [
    FigurePlan(
        image_filename="image1.png",
        title="Mechanisms of Antigen Recognition (TCR / TAA / CAR T)",
        overall_caption="Three-panel comparison of how engineered T cells recognize tumor antigens.",
        long_explanation=(
            "Panel a (TCR): introduced TCR alongside endogenous; both bind MHC-presented peptide. "
            "Panel b (TAA): aberrant self-protein presented via MHC; endogenous TCR is the receptor. "
            "Panel c (CAR T): scFv binds antigen directly, MHC-independent."
        ),
        motif_names=[
            "neon-green",
            "electric-blue",
            "orange",
            "red-magenta",
            "purple-violet",
            "tumor-pink",
        ],
        concept_matcher=matcher_image1,
        dilation_px=25,
    ),
    FigurePlan(
        image_filename="image2.png",
        title="Five Generations of CAR Design",
        overall_caption="Architectural evolution from native TCR through 5th-generation CARs.",
        long_explanation=(
            "TCR (left) for reference. Gen 1: scFv -> CD3-ζ alone. Gen 2: + 1 co-stim. "
            "Gen 3: + 2 co-stim. Gen 4 (TRUCK): cytokine-inducer module. "
            "Gen 5: JAK + STAT3/5 activation domain."
        ),
        motif_names=["neon-green", "electric-blue", "orange", "red-magenta", "purple-violet"],
        concept_matcher=matcher_image2,
        dilation_px=20,
    ),
    FigurePlan(
        image_filename="image3.png",
        title="Three Signals for T-cell Activation",
        overall_caption="T-cell full activation requires three sequential signals.",
        long_explanation=(
            "Signal 1: antigen recognition via TCR / CAR. "
            "Signal 2: co-stimulation drives PI3K/AKT, NF-kB, MAPK. "
            "Signal 3: cytokines (IL-2 family) drive proliferation + persistence."
        ),
        motif_names=["neon-green", "electric-blue", "orange", "red-magenta", "cyan-light"],
        concept_matcher=matcher_image3,
        dilation_px=25,
    ),
    FigurePlan(
        image_filename="image7.png",
        title="CAR-Macrophage (CAR-M) Mechanisms",
        overall_caption="6-panel overview of how CAR-engineered macrophages attack solid tumors.",
        long_explanation=(
            "A: antigen-dependent cytotoxicity. B: TME remodeling + T-cell recruitment. "
            "C: tumor phagocytosis. D: transcription-factor activation -> cytokines. "
            "E: infiltration in solid tumor. F: cell-type legend."
        ),
        motif_names=[
            "neon-green",
            "electric-blue",
            "orange",
            "purple-violet",
            "tumor-pink",
            "cyan-light",
        ],
        concept_matcher=matcher_image7,
        dilation_px=30,
    ),
    FigurePlan(
        image_filename="image8.png",
        title="The Cancer-Immunity Cycle (7 Steps)",
        overall_caption="Iterative loop: antigen release -> presentation -> priming -> trafficking -> "
        "infiltration -> recognition -> tumor death.",
        long_explanation=(
            "Each numbered step (1-7) has its own cytokine + checkpoint context. "
            "Inhibitors at step 7 (PD-1/L1, LAG-3, TIM-3, TGF-β) are therapeutic targets."
        ),
        motif_names=["electric-blue", "orange", "purple-violet", "red-magenta", "tumor-pink"],
        concept_matcher=matcher_image8,
        dilation_px=30,
    ),
    FigurePlan(
        image_filename="image10.png",
        title="TME Suppressive Mechanisms vs CAR-T",
        overall_caption="Multiple converging suppressive forces in the TME degrade CAR-T function.",
        long_explanation=(
            "Suppressive cells (MDSC/Treg/TAM), soluble inhibitors (TGF-β, IL-10), "
            "antigen heterogeneity, dysregulated vasculature, physical barriers (ECM, CAF, IFP), "
            "metabolic suppression (low O2/pH/nutrients, high lactate)."
        ),
        motif_names=["electric-blue", "orange", "purple-violet", "red-magenta", "tumor-pink"],
        concept_matcher=matcher_image10,
        dilation_px=30,
    ),
    FigurePlan(
        image_filename="image13.png",
        title="CAR-T / CAR-NK / CAR-M Comparison",
        overall_caption="How three engineered cell types kill cancer cells in the TME.",
        long_explanation=(
            "CAR-T: granzymes + perforin. CAR-NK: perforin (less GVHD). "
            "CAR-M: cytokines + phagocytosis."
        ),
        motif_names=[
            "neon-green",
            "electric-blue",
            "orange",
            "purple-violet",
            "red-magenta",
            "tumor-pink",
        ],
        concept_matcher=matcher_image13,
        dilation_px=25,
    ),
    FigurePlan(
        image_filename="image14.png",
        title="Notch + syn-Notch CAR-T",
        overall_caption="Native Notch architecture (A) and the syn-Notch CAR-T system (B) using "
        "Notch-style heterodimerization for AND-gate antigen recognition.",
        long_explanation=(
            "A: NECD (EGF Repeats + NRR) + TM + NICD (RAM, ANK, NLS, PEST). "
            "B: 5-step syn-Notch pipeline - first-antigen recognition -> proteolytic TF release -> "
            "CAR transcription -> second-antigen binding -> cytotoxic granule release."
        ),
        motif_names=[
            "neon-green",
            "electric-blue",
            "orange",
            "red-magenta",
            "purple-violet",
            "yellow-gold",
        ],
        concept_matcher=matcher_image14,
        dilation_px=20,
    ),
]


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def process_one(plan: FigurePlan) -> tuple[Path, list[ElementAnnotation]]:
    src = INPUT_DIR / plan.image_filename
    motifs = motifs_for(plan.motif_names)
    color_map = colors_for(plan.motif_names)

    # Read figure dimensions
    from PIL import Image

    with Image.open(src) as im:
        W, H = im.size

    raw = extract_regions(src, motifs, opening_radius=2)
    merged = merge_regions(raw, dilation_px=plan.dilation_px)

    annotations = plan.concept_matcher(merged, W, H)

    out_path = OUT_DIR / f"v3_{plan.image_filename}"
    render_annotated_figure_v3(src, annotations, out_path, motif_colors=color_map)

    return out_path, annotations


def main() -> None:
    results: list[tuple[FigurePlan, Path, list[ElementAnnotation]]] = []
    for plan in PLANS:
        out, anns = process_one(plan)
        print(f"\n=== {plan.image_filename} ({plan.title}) ===")
        print(f"  motifs: {plan.motif_names}")
        print(f"  -> {len(anns)} annotations rendered to {out}")
        for a in anns:
            print(f"    [{a.confidence:.2f}] {a.label}")
        results.append((plan, out, anns))

    # Build PPTX
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    blank = pres.slide_layouts[6]

    s = pres.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VaultLab Figure-Understanding v3 - All 8 Figures"
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True

    sub = s.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0))
    sf = sub.text_frame
    sf.word_wrap = True
    text = (
        "v3 improvements per Bobby's 2026-04-29 review:\n"
        "- No leader lines; markers + boxes only on the figure\n"
        "- Side labels match motif colors (semantic association)\n"
        "- White halos around boxes/markers so any color is visible\n"
        "- Larger label fonts (~3x v2); two-line wrap allowed\n"
        "- NEW: tumor-cell motif (light pink BioRender blobs)\n\n"
        "Each slide's speaker notes show the motif palette + concept matching for that figure."
    )
    for i, line in enumerate(text.split("\n\n")):
        p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(15)

    for plan, out_path, anns in results:
        s = pres.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = plan.title
        run.font.name = "Arial"
        run.font.size = Pt(22)
        run.font.bold = True

        avail_w_in = 12.5
        avail_h_in = 5.6
        with Image.open(out_path) as im:
            aspect = im.width / im.height
        avail_aspect = avail_w_in / avail_h_in
        if aspect > avail_aspect:
            disp_w = Inches(avail_w_in)
            disp_h = Inches(avail_w_in / aspect)
        else:
            disp_h = Inches(avail_h_in)
            disp_w = Inches(avail_h_in * aspect)
        x_in = (13.333 - disp_w / Emu(914400)) / 2
        s.shapes.add_picture(str(out_path), Inches(x_in), Inches(1.0), width=disp_w, height=disp_h)

        cap = s.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7))
        cf = cap.text_frame
        cf.word_wrap = True
        p = cf.paragraphs[0]
        run = p.add_run()
        run.text = plan.overall_caption
        run.font.name = "Arial"
        run.font.size = Pt(13)
        run.font.italic = True

        notes = s.notes_slide.notes_text_frame
        ann_list = "\n".join(
            f"  {i + 1}. [{a.confidence:.2f}] {a.label} - {a.explanation}"
            for i, a in enumerate(anns)
        )
        notes.text = (
            f"{plan.long_explanation}\n\n"
            f"Motif palette used: {', '.join(plan.motif_names)}\n\n"
            f"Annotations rendered ({len(anns)}):\n{ann_list}"
        )

    out_pptx = Path(r"C:\Users\bobby\Downloads\car_t_decks\figure_understanding_demo_v3.pptx")
    pres.save(out_pptx)
    print(f"\nPPTX -> {out_pptx}")


if __name__ == "__main__":
    main()
