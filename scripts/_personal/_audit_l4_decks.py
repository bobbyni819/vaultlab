"""Read both L4 .pptx decks and report exactly what's in each slide."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

DECKS = [
    Path(r"G:\My Drive\Knowledge\vaultlab\Output\codex-cn-test\codex-cellular-neighborhoods-deck.pptx"),
    Path(r"G:\My Drive\Knowledge\vaultlab\Output\spatial-tx-tme-test\spatial-transcriptomics-tumor-microenvironment-deck.pptx"),
]


def audit(deck_path: Path) -> None:
    print("=" * 80)
    print(deck_path.name)
    print("=" * 80)
    pres = Presentation(deck_path)
    print(f"Slide count: {len(pres.slides)}")
    print(f"Slide dims: {pres.slide_width / Emu(914400):.1f} x {pres.slide_height / Emu(914400):.1f} in")

    for i, slide in enumerate(pres.slides, 1):
        print(f"\n--- Slide {i} ---")
        n_pic = sum(1 for s in slide.shapes if s.shape_type == 13)  # PICTURE = 13
        n_text = sum(1 for s in slide.shapes if s.has_text_frame)
        n_total = len(slide.shapes)
        print(f"  shapes: {n_total} total, {n_pic} pictures, {n_text} text")
        for s in slide.shapes:
            label = f"  [{s.shape_type}] name={s.name!r}"
            if s.has_text_frame:
                txt = s.text_frame.text.strip().replace("\n", " | ")
                if txt:
                    label += f"  text={txt[:160]!r}"
            print(label)


for d in DECKS:
    if d.exists():
        audit(d)
    else:
        print(f"MISSING: {d}")
