"""L4 E2E Test - Stage E: verification audit + results doc.

Walks the 9 verification checks from the test plan and writes the
results doc to KB_ROOT/Sources/Notes/L4-e2e-codex-2026-04-29.md.
"""

from __future__ import annotations

import json
import pickle
import re
from pathlib import Path

from vaultlab.kb.paths import (
    article_stub_path,
    concept_path,
    deck_path,
    search_log_path,
    summary_path,
    slugify_topic,
)

TOPIC = "CODEX cellular neighborhoods"
KB_ROOT = Path(r"G:/My Drive/Knowledge/vaultlab")
DATE_STR = "2026-04-29"
STATE_DIR = Path(r"C:/Users/bobby/Downloads/vaultlab/scripts/_l4_state")
PROJECT_SLUG = "codex-cn-test"

TIER_A = [
    "10.1016/j.cell.2018.07.010",
    "10.1126/sciadv.add1166",
    "10.1371/journal.pcbi.1012344",
    "10.1016/j.cell.2024.04.013",
    "10.1038/nmeth.4391",
    "10.1089/cmb.2019.0340",
    "10.1007/s00281-022-00974-0",
    "10.1038/s42003-022-04032-1",
]


def _read(p: Path) -> str:
    return p.read_text(encoding="utf-8") if p.exists() else ""


def main():
    out_lines: list[str] = []

    def log(s=""):
        out_lines.append(s)
        print(s, flush=True)

    log("# L4 End-to-End Test: CODEX Cellular Neighborhoods")
    log("")
    log("Date: 2026-04-29")
    log(f"Topic: {TOPIC}")
    log(f"KB root: {KB_ROOT}")
    log(f"Test plan: G:/My Drive/Knowledge/vaultlab/Sources/Notes/end-to-end-test-plan-2026-04-29.md")
    log("")

    # ------------------------------------------------------------------
    # Load state
    # ------------------------------------------------------------------
    seeds = json.load(open(STATE_DIR / "seeds.json"))
    acq = json.load(open(STATE_DIR / "acq_results.json"))
    figs_path = STATE_DIR / "figure_results.json"
    fig_results = json.load(open(figs_path)) if figs_path.exists() else {}
    with open(STATE_DIR / "corpus.pkl", "rb") as f:
        corpus = pickle.load(f)

    # ------------------------------------------------------------------
    # Check 1: Search coverage
    # ------------------------------------------------------------------
    log("## Check 1: Search coverage")
    log("")
    log(f"- Seeds returned: **{len(seeds)}** (max_seeds=12, all with DOIs)")
    has_schurch_preprint = any(s["doi"].lower() == "10.1101/743989" for s in seeds)
    has_schurch_cell = any(s["doi"].lower() == "10.1016/j.cell.2020.07.005" for s in seeds)
    log(f"- Schurch 2020 (Cell, DOI 10.1016/j.cell.2020.07.005): **{'present' if has_schurch_cell else 'NOT in seeds'}**")
    log(f"- Schurch 2020 (bioRxiv preprint, DOI 10.1101/743989): **{'present' if has_schurch_preprint else 'absent'}**")
    log("- All 12 seeds:")
    for i, s in enumerate(seeds, 1):
        log(f"  {i}. `{s['doi']}` — {s['title'][:80]} ({s['year']}, {s['journal']})")
    log("")

    # ------------------------------------------------------------------
    # Check 2: Citation graph metrics
    # ------------------------------------------------------------------
    log("## Check 2: Citation graph metrics")
    log("")
    metrics = corpus.metrics
    log(f"- Corpus size after CrossRef expansion: **{corpus.n_papers}** papers, **{corpus.n_edges}** edges")
    if metrics:
        top_og = sorted(metrics.og_score.items(), key=lambda kv: kv[1], reverse=True)[:3]
        log("- Top-3 by og_score:")
        for doi, score in top_og:
            p = corpus.papers.get(doi)
            title = p.title[:80] if p and p.title else "(no metadata)"
            year = p.year if p else "?"
            log(f"  - og={score:.3f} | {year} | {title} `[{doi}]`")
        top_fwd = sorted(metrics.forward_influence.items(), key=lambda kv: kv[1], reverse=True)[:3]
        log("- Top-3 by forward_influence:")
        for doi, score in top_fwd:
            p = corpus.papers.get(doi)
            title = p.title[:80] if p and p.title else "(no metadata)"
            year = p.year if p else "?"
            log(f"  - fwd={score} | {year} | {title} `[{doi}]`")
    log("")

    # ------------------------------------------------------------------
    # Check 3: PDF acquisition rate
    # ------------------------------------------------------------------
    log("## Check 3: PDF acquisition rate")
    log("")
    n_with_pdf = sum(1 for a in acq.values() if a.get("pdf_path"))
    sources = {}
    for a in acq.values():
        if a.get("pdf_path"):
            src = a.get("source", "?")
            sources[src] = sources.get(src, 0) + 1
    log(f"- PDFs acquired: **{n_with_pdf} / {len(acq)}** ({100*n_with_pdf/max(len(acq),1):.1f}%)")
    log("- Sources used:")
    for src, n in sorted(sources.items(), key=lambda kv: -kv[1]):
        log(f"  - {src}: {n}")
    seed_pdf = sum(1 for s in seeds if acq.get(s["doi"], {}).get("pdf_path"))
    log(f"- Seeds with PDFs: **{seed_pdf} / {len(seeds)}**")
    log("")

    # ------------------------------------------------------------------
    # Check 4: Per-paper summaries
    # ------------------------------------------------------------------
    log("## Check 4: Per-paper summaries (Tier A)")
    log("")
    for doi in TIER_A:
        sp = summary_path(KB_ROOT, doi)
        body = _read(sp)
        size = len(body)
        # Extract TL;DR (line after `## TL;DR`).
        tldr_match = re.search(r"## TL;DR\s*\n(.*?)(?=\n##|\Z)", body, re.DOTALL)
        tldr = (tldr_match.group(1).strip() if tldr_match else "(NOT FOUND)").split("\n")
        tldr_preview = " ".join(line for line in tldr if line.strip())[:300]
        # Count key findings with [pN] markers.
        findings_section = re.search(r"## Key findings[^\n]*\n(.*?)(?=\n## |\Z)", body, re.DOTALL)
        n_findings = 0
        n_with_marker = 0
        if findings_section:
            lines = [ln for ln in findings_section.group(1).splitlines() if ln.strip().startswith("- ")]
            n_findings = len(lines)
            n_with_marker = sum(1 for ln in lines if re.search(r"\[p\d+\]|\[unknown\]", ln))
        log(f"### `{doi}`")
        log(f"- Path: `{sp.relative_to(KB_ROOT)}`  (exists={sp.exists()}, {size} bytes)")
        log(f"- Key findings: **{n_findings}** total, **{n_with_marker}** with `[pN]` page marker")
        log(f"- TL;DR (first 300 chars): {tldr_preview!r}")
        log("")

    # ------------------------------------------------------------------
    # Check 5: Lineage arc
    # ------------------------------------------------------------------
    log("## Check 5: Lineage arc")
    log("")
    arc = concept_path(KB_ROOT, TOPIC, "lineage", DATE_STR)
    arc_body = _read(arc)
    log(f"- Path: `{arc.relative_to(KB_ROOT)}`  (exists={arc.exists()}, {len(arc_body)} bytes)")
    if arc_body:
        # Pull out the three narrative paragraphs.
        for section in ("History", "Development", "State of the art"):
            m = re.search(rf"## {re.escape(section)}.*?\n\n(.*?)(?=\n##|\n\|)", arc_body, re.DOTALL)
            if m:
                para = m.group(1).strip().split("\n\n")[0]
                log(f"### {section} paragraph")
                log("")
                log(f"> {para}")
                log("")
        # Verify wikilinks resolve.
        links = set(re.findall(r"\[\[([^\]|]+)\|", arc_body))
        log(f"- Wikilinks found: **{len(links)}** unique targets")
        unresolved = []
        for slug in links:
            target = KB_ROOT / "Wiki" / "Summaries" / f"{slug}.md"
            if not target.exists():
                unresolved.append(slug)
        if unresolved:
            log(f"- **UNRESOLVED wikilinks ({len(unresolved)}):**")
            for s in unresolved[:10]:
                log(f"  - `{s}`")
        else:
            log("- All wikilinks resolve to existing Wiki/Summaries/ files")
    log("")

    # ------------------------------------------------------------------
    # Check 6: Figure acquisition
    # ------------------------------------------------------------------
    log("## Check 6: Figure acquisition")
    log("")
    n_with_figs = sum(1 for r in fig_results.values() if r.get("source") != "unavailable" and r.get("n_figures", 0) > 0)
    log(f"- Tier-A papers with figures: **{n_with_figs} / {len(fig_results)}**")
    for doi, r in fig_results.items():
        cap = (r.get("first_caption") or "")[:120]
        log(f"  - `{doi}`: source={r['source']}, n={r['n_figures']}, first_caption={cap!r}")
    log("")

    # ------------------------------------------------------------------
    # Check 7: Deck
    # ------------------------------------------------------------------
    log("## Check 7: Deck")
    log("")
    topic_slug = slugify_topic(TOPIC)
    deck_p = deck_path(KB_ROOT, PROJECT_SLUG, f"{topic_slug}-deck.pptx")
    log(f"- Path: `{deck_p.relative_to(KB_ROOT)}`")
    log(f"- Exists: **{deck_p.exists()}**")
    if deck_p.exists():
        log(f"- Size: **{deck_p.stat().st_size} bytes**")
        try:
            from pptx import Presentation
            pres = Presentation(str(deck_p))
            log(f"- Slide count: **{len(pres.slides)}**")
            log("- Slide titles:")
            for i, s in enumerate(pres.slides, 1):
                title_txt = ""
                for shape in s.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        title_txt = shape.text_frame.text.split("\n")[0]
                        break
                log(f"  {i}. {title_txt[:80]}")
        except Exception as e:
            log(f"- (could not read pptx: {e})")
    log("")

    # ------------------------------------------------------------------
    # Check 8: Provenance
    # ------------------------------------------------------------------
    log("## Check 8: Provenance receipts")
    log("")
    prov_json = arc.with_suffix(arc.suffix + ".provenance.json")
    method_md = arc.with_suffix(arc.suffix + ".method.md")
    log(f"- `{prov_json.name}`: exists={prov_json.exists()} ({prov_json.stat().st_size if prov_json.exists() else 0} bytes)")
    log(f"- `{method_md.name}`: exists={method_md.exists()} ({method_md.stat().st_size if method_md.exists() else 0} bytes)")
    log("")

    # ------------------------------------------------------------------
    # Check 9: Canonical paths
    # ------------------------------------------------------------------
    log("## Check 9: Canonical KB output paths")
    log("")
    expected = [
        ("search log", search_log_path(KB_ROOT, TOPIC, DATE_STR)),
        ("arc", arc),
        ("provenance", prov_json),
        ("method", method_md),
        ("deck", deck_p),
    ]
    for s in seeds:
        expected.append((f"article stub ({s['doi']})", article_stub_path(KB_ROOT, s["doi"])))
    for d in TIER_A:
        expected.append((f"summary ({d})", summary_path(KB_ROOT, d)))
    n_ok = 0
    n_total = 0
    for name, p in expected:
        n_total += 1
        ok = p.exists()
        if ok:
            n_ok += 1
        marker = "OK" if ok else "MISSING"
        log(f"- [{marker}] {name}: `{p.relative_to(KB_ROOT)}`")
    log("")
    log(f"**Path summary: {n_ok}/{n_total} canonical paths exist.**")
    log("")

    # ------------------------------------------------------------------
    # Final verdict
    # ------------------------------------------------------------------
    log("## Final verdict")
    log("")

    text = "\n".join(out_lines) + "\n"
    out = KB_ROOT / "Sources" / "Notes" / "L4-e2e-codex-2026-04-29.md"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(text, encoding="utf-8")
    print(f"\nResults doc: {out}", flush=True)


if __name__ == "__main__":
    main()
