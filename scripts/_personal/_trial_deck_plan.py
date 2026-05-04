"""Trial run for ``vaultlab.workflows.deck_plan`` — synthetic stub callback.

Builds a synthetic CODEX-flavoured corpus (5 Tier-A papers, 3 figure
assignments), defines a stub plan_callback that picks 5 papers across
7 slides, runs ``generate_deck_plan + build_from_plan``, then verifies
the resulting .pptx opens and has 7 slides.

Usage::

    python scripts/_trial_deck_plan.py [kb_root]

If ``kb_root`` is omitted, a temp dir is used.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

from PIL import Image

from vaultlab.research.corpus import Corpus
from vaultlab.research.graph_metrics import CorpusMetrics
from vaultlab.research.paper import Paper
from vaultlab.research.summarize import PaperSummary
from vaultlab.slides import build_from_plan
from vaultlab.workflows.deck_plan import (
    DeckPlanTask,
    generate_deck_plan,
)


def _make_synthetic_inputs(out_dir: Path):
    """Build a synthetic CODEX-flavoured corpus with 5 papers + 3 figures."""
    summaries = {
        "10.1126/science.aar7042": PaperSummary(
            doi="10.1126/science.aar7042",
            title="CODEX multiplexed imaging in tissue",
            authors=["Goltsev Y", "Samusik N", "Kennedy-Darling J"],
            year=2018,
            journal="Cell",
            og_score=0.65,
            forward_influence=42,
            year_bucket="history",
            tier="A",
            tldr=(
                "CODEX uses indexed antibodies cleaved by polymerase to image "
                "40+ proteins in tissue. The method works on FFPE and frozen "
                "tissue. It established multiplexed imaging at single-cell "
                "resolution."
            ),
            key_findings=[
                "First 40+ marker tissue imaging with antibody indexing",
                "Validated on mouse spleen and human tonsil",
                "Single-cell phenotyping resolved B/T/macrophage subsets",
            ],
        ),
        "10.1038/s41587-019-0207-4": PaperSummary(
            doi="10.1038/s41587-019-0207-4",
            title="Spatial cellular neighborhoods in CODEX images",
            authors=["Schurch CM", "Bhate SS"],
            year=2020,
            journal="Cell",
            og_score=0.40,
            forward_influence=18,
            year_bucket="development",
            tier="A",
            tldr=(
                "Identifies recurring 10-cell neighborhoods in CODEX colorectal "
                "tumor images. Neighborhood structure correlates with patient "
                "survival. The framework uses CNN segmentation + clustering."
            ),
            key_findings=[
                "Cellular neighborhoods are reproducible across patients",
                "Neighborhood enrichment predicts survival in CRC",
                "Compositional analysis exposes immune-tumor interface",
            ],
        ),
        "10.1126/science.abd5588": PaperSummary(
            doi="10.1126/science.abd5588",
            title="Whole-organ CODEX of a healthy thymus",
            authors=["Park JE", "Botting RA"],
            year=2020,
            journal="Science",
            og_score=0.30,
            forward_influence=11,
            year_bucket="development",
            tier="A",
            tldr=(
                "Maps thymic stromal cells with single-cell resolution. "
                "Reveals new cortical/medullary niches. Provides reference "
                "atlas for autoimmune disease research."
            ),
            key_findings=[
                "Identified two new thymic epithelial subtypes",
                "Mapped TCR repertoire to spatial niche",
                "Provides public reference atlas",
            ],
        ),
        "10.1038/s41591-022-02101-w": PaperSummary(
            doi="10.1038/s41591-022-02101-w",
            title="CODEX-driven biomarker discovery in melanoma",
            authors=["Lin JR", "Wang S"],
            year=2023,
            journal="Nature Medicine",
            og_score=0.25,
            forward_influence=6,
            year_bucket="sota",
            tier="A",
            tldr=(
                "Uses CODEX of pre-treatment melanoma to predict immune "
                "checkpoint response. Identifies a tertiary lymphoid "
                "structure signature. Validated in two cohorts."
            ),
            key_findings=[
                "TLS signature predicts ICB response (AUC 0.84)",
                "Spatial distance to PD-1+ T cells matters more than density",
                "Multi-cohort validation on n=181 patients",
            ],
        ),
        "10.1038/s41587-024-01001-5": PaperSummary(
            doi="10.1038/s41587-024-01001-5",
            title="Foundation model for CODEX cell typing",
            authors=["Chen T", "Park R"],
            year=2024,
            journal="Nature Biotechnology",
            og_score=0.18,
            forward_influence=2,
            year_bucket="sota",
            tier="A",
            tldr=(
                "Pretrains a vision transformer on 1M CODEX cells. "
                "Yields zero-shot transfer across disease contexts. "
                "Open-sources weights and benchmark."
            ),
            key_findings=[
                "Zero-shot transfer beats supervised baselines on 4/5 tasks",
                "Reduces annotation burden by ~10x",
                "Public weights + 1M-cell benchmark released",
            ],
        ),
    }

    # Build 3 PNG figures (only 3 of 5 papers have cached figures, mimicking
    # real-world acquisition rate).
    fig_dir = out_dir / "figures"
    fig_dir.mkdir(parents=True, exist_ok=True)
    figure_assignments: dict[str, Path] = {}
    for doi, color in [
        ("10.1126/science.aar7042", "red"),
        ("10.1038/s41587-019-0207-4", "blue"),
        ("10.1038/s41587-024-01001-5", "green"),
    ]:
        slug = doi.replace("/", "_").replace(".", "_")
        path = fig_dir / f"{slug}.png"
        Image.new("RGB", (640, 480), color).save(str(path))
        figure_assignments[doi] = path

    # Build the Corpus with metrics.
    papers: dict[str, Paper] = {}
    og_score: dict[str, float] = {}
    forward_influence: dict[str, int] = {}
    year_buckets: dict[str, str] = {}
    for doi, s in summaries.items():
        key = doi.lower()
        papers[key] = Paper(
            doi=doi, title=s.title, authors=list(s.authors),
            year=s.year, journal=s.journal,
        )
        og_score[key] = s.og_score
        forward_influence[key] = s.forward_influence
        year_buckets[key] = s.year_bucket
    metrics = CorpusMetrics(
        og_score=og_score,
        forward_influence=forward_influence,
        co_citation_pairs=[
            ("10.1126/science.aar7042", "10.1038/s41587-019-0207-4", 8),
            ("10.1126/science.aar7042", "10.1126/science.abd5588", 5),
        ],
        year_buckets=year_buckets,
    )
    corpus = Corpus(
        topic="CODEX multiplexed imaging",
        seeds=[],
        papers=papers,
        references={},
        metrics=metrics,
    )
    return corpus, summaries, figure_assignments


def _stub_callback(task: DeckPlanTask) -> dict:
    """Stub plan_callback that picks 5 papers across 7 slides.

    In production this would be the LLM's response after reading
    task.corpus_summaries. For the trial we hand-author a representative
    response that exercises every slide type.
    """
    fig_paths = task.figure_assignments
    foundations = fig_paths.get("10.1126/science.aar7042")
    neighborhoods = fig_paths.get("10.1038/s41587-019-0207-4")
    foundation_model = fig_paths.get("10.1038/s41587-024-01001-5")

    return {
        "story_arc_summary": (
            "Trace CODEX from instrument breakthrough (2018) through "
            "spatial-neighborhood frameworks (2020) to clinical biomarker "
            "discovery and foundation models (2023-2024)."
        ),
        "slides": [
            {
                "type": "title",
                "title": "CODEX multiplexed imaging",
                "subtitle": f"{task.audience.replace('-', ' ').title()} deck",
                "author": task.speaker,
                "speaker_notes": {
                    "mental_map": {
                        "hook": "Imaging just got a lot more multiplexed.",
                        "key_claim": (
                            "CODEX changed how we look at tissue, then "
                            "spatial methods changed how we read it."
                        ),
                        "transition": "Let's start with the original method.",
                    },
                    "detailed_script": (
                        "Hello, today I'll trace the CODEX lineage from the "
                        "2018 method paper to the foundation models we're "
                        "seeing in 2024. Five papers, four key beats."
                    ),
                },
            },
            {"type": "section_divider", "title": "Foundations (2018)"},
            {
                "type": "figure",
                "title": "CODEX establishes multiplexed tissue imaging",
                "image_path": str(foundations),
                "claim_paper_doi": "10.1126/science.aar7042",
                "figure_paper_doi": "10.1126/science.aar7042",
                "caption": "Goltsev et al. 2018: 40+ marker imaging via DNA-indexed antibodies.",
                "bullets": [
                    "[[10.1126_science_aar7042|Goltsev 2018]]: 40+ markers via antibody indexing",
                    "Validated on mouse spleen + human tonsil",
                    "Single-cell phenotyping at tissue scale",
                ],
                "speaker_notes": {
                    "mental_map": {
                        "hook": "How do you image 40 proteins at once?",
                        "key_claim": "Antibody indexing + polymerase cleavage.",
                        "evidence": "Spleen B/T/macrophage panel in this figure.",
                    },
                    "detailed_script": (
                        "The original CODEX paper introduced indexed antibodies "
                        "where each antibody carries a short DNA barcode. A "
                        "polymerase cycles through complementary barcodes, "
                        "imaging 40+ markers without bleaching."
                    ),
                },
            },
            {"type": "section_divider", "title": "Spatial frameworks (2020-2022)"},
            {
                "type": "figure",
                "title": "Neighborhood structure as a clinical readout",
                "image_path": str(neighborhoods),
                "claim_paper_doi": "10.1038/s41587-019-0207-4",
                "figure_paper_doi": "10.1038/s41587-019-0207-4",
                "caption": "Schurch et al. 2020: 10-cell neighborhoods predict CRC survival.",
                "bullets": [
                    "[[10.1038_s41587-019-0207-4|Schurch 2020]]: 10-cell neighborhoods are reproducible",
                    "Neighborhood enrichment predicts CRC survival",
                    "Sets the template for spatial-readout work",
                ],
            },
            {
                "type": "text",
                "title": "Clinical biomarker discovery (2023)",
                "bullets": [
                    "[[10.1038_s41591-022-02101-w|Lin 2023]]: TLS signature predicts ICB response (AUC 0.84)",
                    "Spatial distance to PD-1+ T cells > absolute density",
                    "Validated in two melanoma cohorts (n=181)",
                ],
                "citations": ["10.1038/s41591-022-02101-w"],
                "speaker_notes": {
                    "mental_map": {
                        "key_claim": "Spatial structure beats density for ICB prediction.",
                    },
                },
            },
            {
                "type": "figure",
                "title": "Foundation models for CODEX (2024)",
                "image_path": str(foundation_model),
                "claim_paper_doi": "10.1038/s41587-024-01001-5",
                "figure_paper_doi": "10.1038/s41587-024-01001-5",
                "caption": "Chen et al. 2024: ViT pretrained on 1M CODEX cells, zero-shot transfer.",
                "bullets": [
                    "[[10.1038_s41587-024-01001-5|Chen 2024]]: ViT on 1M cells",
                    "Zero-shot beats supervised on 4/5 tasks",
                    "Reduces annotation burden ~10x",
                ],
            },
        ],
    }


def main() -> int:
    explicit_kb = sys.argv[1] if len(sys.argv) > 1 else None
    cleanup = explicit_kb is None
    if explicit_kb:
        kb_root = Path(explicit_kb)
        kb_root.mkdir(parents=True, exist_ok=True)
    else:
        kb_root = Path(tempfile.mkdtemp(prefix="vaultlab-trial-deck-plan-"))

    print(f"[trial] using kb_root: {kb_root}")
    try:
        corpus, summaries, figure_assignments = _make_synthetic_inputs(kb_root)

        # Run the deck-plan generator with the stub callback
        dict_plan = generate_deck_plan(
            topic="CODEX multiplexed imaging",
            corpus=corpus,
            summaries=summaries,
            figure_assignments=figure_assignments,
            speaker="Bobby Ni",
            affiliation="Hickey Lab @ Duke BME",
            audience="journal-club",
            target_slide_count=7,
            kb_root=kb_root,
            plan_callback=_stub_callback,
        )

        # Print slide structure
        print(f"[trial] story arc: {dict_plan['story_arc_summary']}")
        print(f"[trial] slide count: {len(dict_plan['slides'])}")
        for i, slide in enumerate(dict_plan["slides"], 1):
            stype = slide.get("type", "?")
            title = slide.get("title", "")
            print(f"  {i}. [{stype}] {title}")

        # Render to .pptx
        out = kb_root / "trial-codex-deck.pptx"
        result = build_from_plan(dict_plan, out, write_marp=False)
        out_pptx = result["pptx"]

        if not out_pptx.exists():
            print(f"[trial] FAIL: .pptx not written at {out_pptx}")
            return 1

        size_kb = out_pptx.stat().st_size / 1024
        print(f"[trial] WROTE: {out_pptx} ({size_kb:.1f} KB)")

        # Verify slide count by reopening
        from pptx import Presentation
        prs = Presentation(str(out_pptx))
        n_slides = len(prs.slides)
        print(f"[trial] OPENED: {n_slides} slides in deck")
        if n_slides < 7:
            print(f"[trial] WARN: expected >=7 slides, got {n_slides}")
            return 1

        print("[trial] OK")
        return 0
    finally:
        if cleanup:
            shutil.rmtree(kb_root, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
