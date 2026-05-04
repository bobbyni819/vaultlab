"""Trial run for ``vaultlab.slides.deck.build_deck`` — synthetic 5-paper deck.

Builds a deck plan inline (no /lit-arc invocation needed), composes a .pptx,
verifies it lands at the canonical Output/<project>/<deck>.pptx, prints the
path + open command.

Usage::

    python scripts/_trial_build_deck.py [kb_root]

If ``kb_root`` is omitted, a temp dir is used.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

from vaultlab.kb.paths import deck_path, slugify_doi
from vaultlab.research.lineage import LineageRunResult
from vaultlab.slides import (
    DeckPlan,
    DeckSlide,
    build_deck,
    build_deck_from_lineage_result,
)


def _make_synthetic_kb(kb_root: Path) -> LineageRunResult:
    """Lay down 5 synthetic Wiki/Summaries entries + an arc file."""
    papers = [
        ("10.1234/foundations-1990", 1990, "history", "Foundational discovery"),
        ("10.1234/scaffolding-2002", 2002, "history", "Scaffolding the field"),
        ("10.1234/breakthrough-2014", 2014, "development", "Breakthrough method"),
        ("10.1234/refinement-2019", 2019, "development", "Methodological refinement"),
        ("10.1234/sota-2024", 2024, "sota", "State-of-the-art system"),
    ]
    summary_paths: dict[str, Path] = {}
    for doi, year, bucket, title in papers:
        slug = slugify_doi(doi)
        p = kb_root / "Wiki" / "Summaries" / f"{slug}.md"
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(
            f"---\n"
            f"doi: {doi}\n"
            f"title: {title}\n"
            f"authors: [\"Smith J\", \"Doe A\"]\n"
            f"year: {year}\n"
            f"journal: Journal of Synthetic Science\n"
            f"year_bucket: {bucket}\n"
            f"tier: A\n"
            f"---\n"
            f"\n## TL;DR\n"
            f"In {year}, this paper showed that {title.lower()} could be "
            f"replicated with high fidelity across multiple systems.\n"
            f"\n## Key findings (with [page] provenance)\n"
            f"- Established a new technique for the {bucket} era [p2]\n"
            f"- Demonstrated quantitative improvement over baselines [p4]\n"
            f"- Opened follow-up work on related problems [p7]\n",
            encoding="utf-8",
        )
        summary_paths[doi] = p

    arc_path = kb_root / "Wiki" / "Concepts" / "trial-topic-lineage-2026-04-29.md"
    arc_path.parent.mkdir(parents=True, exist_ok=True)
    arc_path.write_text(
        "---\ntopic: trial topic\n---\n\n"
        "# Lineage: trial topic\n\n"
        "The story begins with foundational work in the early 1990s. "
        "Successive papers refined the basic technique, culminating in "
        "the breakthroughs of the 2010s. Today, state-of-the-art systems "
        "build on every layer of that lineage.\n",
        encoding="utf-8",
    )

    return LineageRunResult(
        topic="trial topic",
        arc_path=arc_path,
        summary_paths=summary_paths,
        search_log_path=Path(),
        corpus_size=len(papers),
        pdfs_acquired=len(papers),
        summaries_written=len(papers),
        duration_seconds=0.0,
    )


def _trial_synthesized_plan(out: Path) -> Path:
    """Quick smoke: handcrafted DeckPlan -> build_deck -> .pptx."""
    plan = DeckPlan(
        title="Trial deck",
        subtitle="Smoke test for build_deck",
        speaker="Bobby Ni",
        affiliation="Hickey Lab @ Duke BME",
        sections=["Background", "Findings", "References"],
        slides=[
            DeckSlide(
                kind="title",
                title="Trial deck",
                content={
                    "subtitle": "Smoke test for build_deck",
                    "speaker": "Bobby Ni",
                    "affiliation": "Hickey Lab @ Duke BME",
                    "date": "2026-04-29",
                },
            ),
            DeckSlide(
                kind="section_intro",
                title="Background",
                content={
                    "section_name": "Background",
                    "key_question": "What did we know coming in?",
                    "bullets": [
                        "Foundational papers from the 1990s",
                        "Methodological refinements through the 2010s",
                    ],
                },
            ),
            DeckSlide(
                kind="bullets",
                title="Key findings",
                content={
                    "bullets": [
                        "Effect size grew across replications",
                        "New mechanism proposed in 2024",
                        "Opens new translational pathway",
                    ],
                    "citations": ["10.1234/foundations-1990", "10.1234/sota-2024"],
                },
            ),
            DeckSlide(
                kind="references",
                title="References",
                content={
                    "refs": [
                        {"n": 1, "citation": "Smith J, Doe A. Foundational. JSS 1990.", "doi": "10.1234/foundations-1990"},
                        {"n": 2, "citation": "Smith J. Breakthrough. JSS 2014.", "doi": "10.1234/breakthrough-2014"},
                        {"n": 3, "citation": "Doe A. SOTA. JSS 2024.", "doi": "10.1234/sota-2024"},
                    ]
                },
            ),
        ],
        theme="hickey_lab",
    )
    return build_deck(plan, out)


def main() -> int:
    keep_kb = len(sys.argv) > 1
    kb_root = Path(sys.argv[1]) if keep_kb else Path(tempfile.mkdtemp(prefix="vaultlab_trial_deck_"))
    print(f"[trial] kb_root: {kb_root}")

    # 1) handcrafted deck via build_deck()
    smoke_out = kb_root / "Output" / "trial" / "smoke-deck.pptx"
    smoke_out.parent.mkdir(parents=True, exist_ok=True)
    smoke_out = _trial_synthesized_plan(smoke_out)
    print(f"[trial] handcrafted deck: {smoke_out} ({smoke_out.stat().st_size:,} bytes)")

    # 2) lineage-driven deck via build_deck_from_lineage_result()
    result = _make_synthetic_kb(kb_root)
    out = build_deck_from_lineage_result(
        result,
        speaker="Bobby Ni",
        affiliation="Hickey Lab @ Duke BME",
        project_slug="trial-deck",
        figure_assignments={},  # no figures -> graceful fallback to bullets
        kb_root=kb_root,
    )
    print(f"[trial] lineage deck:    {out} ({out.stat().st_size:,} bytes)")

    # Inspect
    from pptx import Presentation

    pres = Presentation(str(out))
    print(f"[trial] lineage deck slide count: {len(pres.slides)}")
    for i, s in enumerate(pres.slides, 1):
        names = sorted(sh.name for sh in s.shapes)
        print(f"  slide {i}: {len(s.shapes)} shapes — {names[:4]}...")

    # Canonical path check
    expected = deck_path(kb_root, "trial-deck", "trial-topic-deck.pptx")
    print(f"[trial] canonical match: {out == expected}")

    print()
    print(f"To open: bobby-kb open {out}")

    if not keep_kb:
        # Leave the temp dir for the user to inspect
        print(f"[trial] (temp kb retained at {kb_root}; remove manually if desired)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
