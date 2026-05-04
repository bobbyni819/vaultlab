"""Re-populate all 4 deck originals with figures using the post-fix populator.

Generates -auto + -aspirational variants for each topic-scope combo, audits
each populated deck, and writes an audit report to Output/Reports/.

Run:
    cd vaultlab
    set PYTHONIOENCODING=utf-8
    python scripts/_repopulate_decks_2026_05_03.py
"""

from __future__ import annotations

import datetime as dt
import sys
from dataclasses import dataclass
from pathlib import Path

from vaultlab.slides.audit import _extract_dois_from_arc, audit_deck
from vaultlab.slides.figure_populate import populate_deck_with_figures


def _strip_legacy_corner_stamps(deck_path: Path) -> int:
    """Remove orphan figure stamps left in base decks by prior buggy runs.

    Detects pictures matching the old 4.0x2.5" corner-stamp signature
    (3.5–4.5 in wide, 2.0–3.0 in tall, located in the right or bottom
    half of the slide), removes them along with their accompanying
    8pt/12pt micro-caption directly below. Saves the deck in place.

    Returns the number of pictures removed.
    """
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    removed = 0
    for slide in prs.slides:
        sw = prs.slide_width
        sh = prs.slide_height
        to_remove = []
        for shape in list(slide.shapes):
            if shape.shape_type != 13:
                continue
            try:
                w_in = (shape.width or 0) / 914400
                h_in = (shape.height or 0) / 914400
                l_in = (shape.left or 0) / 914400
                t_in = (shape.top or 0) / 914400
            except (AttributeError, TypeError):
                continue
            sw_in = sw / 914400
            sh_in = sh / 914400
            is_corner_stamp = (
                3.5 <= w_in <= 4.5
                and 2.0 <= h_in <= 3.0
                and (l_in > sw_in / 2 - 1 or t_in > sh_in / 2 + 0.5)
            )
            if is_corner_stamp:
                to_remove.append(shape)
        for sh_obj in to_remove:
            sh_obj._element.getparent().remove(sh_obj._element)
            removed += 1
    if removed:
        prs.save(str(deck_path))
    return removed


KB_ROOT = Path("G:/My Drive/Knowledge/vaultlab")
OUTPUT_ROOT = KB_ROOT / "Output"
ARC_ROOT = KB_ROOT / "Wiki" / "Concepts"
PDF_CACHE = Path("G:/My Drive/Knowledge/vaultlab/Sources/Papers")
FIG_STAGING = Path("C:/Users/bobby/.cache/vaultlab/_deck_figures_2026_05_03")


@dataclass
class Job:
    topic_dir: str
    scope: str
    arc_filename: str
    deck_filename: str


JOBS = [
    Job("multiscale-tissue-simulation-lung-infection", "short",
        "multiscale-tissue-simulation-lung-infection-lineage-short-2026-05-02.md",
        "short-2026-05-02.pptx"),
    Job("multiscale-tissue-simulation-lung-infection", "review",
        "multiscale-tissue-simulation-lung-infection-lineage-review-2026-05-02.md",
        "review-2026-05-02.pptx"),
    Job("spatial-tx-tme", "short",
        "spatial-tx-tme-lineage-short-2026-05-02.md",
        "short-2026-05-02.pptx"),
    Job("spatial-tx-tme", "review",
        "spatial-tx-tme-lineage-review-2026-05-02.md",
        "review-2026-05-02.pptx"),
]


def _format_severity(sev: str) -> str:
    return {"ok": "[OK]   ", "warn": "[WARN] ", "fail": "[FAIL] "}.get(sev, sev)


def main() -> int:
    FIG_STAGING.mkdir(parents=True, exist_ok=True)
    PDF_CACHE.mkdir(parents=True, exist_ok=True)
    reports_dir = OUTPUT_ROOT / "Reports"
    reports_dir.mkdir(parents=True, exist_ok=True)

    rows: list[str] = []
    audit_md_blocks: list[str] = []
    rows.append(
        f"{'topic':45s} {'scope':8s} {'mode':14s} {'figs':>5s} "
        f"{'placeholders':>13s} {'severity':10s} {'deck':40s}"
    )
    rows.append("-" * 140)

    for job in JOBS:
        arc_path = ARC_ROOT / job.arc_filename
        deck_path = OUTPUT_ROOT / job.topic_dir / job.deck_filename

        if not arc_path.exists():
            print(f"SKIP missing arc: {arc_path}")
            continue
        if not deck_path.exists():
            print(f"SKIP missing deck: {deck_path}")
            continue

        ids = _extract_dois_from_arc(arc_path.read_text(encoding="utf-8"))
        print(f"\n=== {job.topic_dir} / {job.scope} ===")
        print(f"  arc:  {arc_path.name}")
        print(f"  deck: {deck_path.name}")
        print(f"  arc citations -> {len(ids)} IDs")

        # One-shot cleanup of any orphan corner-stamp figures left in the
        # base deck by previous buggy runs of figure_populate.
        n_stripped = _strip_legacy_corner_stamps(deck_path)
        if n_stripped:
            print(f"  cleaned {n_stripped} orphan corner-stamp figure(s) from base")

        for mode in ("auto", "aspirational"):
            out = deck_path.with_name(deck_path.stem + f"-{mode}.pptx")
            print(f"  populating {mode}-mode -> {out.name} ...", flush=True)
            try:
                result = populate_deck_with_figures(
                    deck_path,
                    candidate_dois=ids,
                    pdf_cache_dir=PDF_CACHE,
                    figure_staging_dir=FIG_STAGING,
                    mode=mode,
                    out_path=out,
                )
            except Exception as exc:  # noqa: BLE001
                print(f"    ERROR: {exc}")
                continue

            audit = audit_deck(out, arc_path=arc_path)
            sev = _format_severity(audit.severity)
            rows.append(
                f"{job.topic_dir:45s} {job.scope:8s} {mode:14s} "
                f"{result.n_inserted:>5d} {result.n_placeholders:>13d} "
                f"{sev:10s} {out.name:40s}"
            )
            print(
                f"    inserted={result.n_inserted}  "
                f"placeholders={result.n_placeholders}  severity={audit.severity}"
            )
            if result.inserted_dois[:3]:
                print(f"    sample inserts: {result.inserted_dois[:3]}")
            audit_md_blocks.append(audit.to_markdown_report())

    print("\n" + "=" * 140)
    print("Summary table")
    print("=" * 140)
    for r in rows:
        print(r)

    # Persist audit report to Output/Reports/
    today = dt.date.today().isoformat()
    report_path = reports_dir / f"deck-rerun-{today}.md"
    body = [
        f"---", f"date: {today}", f"type: deck-rerun-audit",
        f"trigger: filename-mismatch fix + bobby_slides layout adoption",
        f"---", "",
        f"# Deck rerun audit — {today}", "",
        "## Summary table", "", "```",
    ] + rows + ["```", "", "## Per-deck audit"] + ["", *audit_md_blocks]
    report_path.write_text("\n".join(body), encoding="utf-8")
    print(f"\nAudit report written: {report_path}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
