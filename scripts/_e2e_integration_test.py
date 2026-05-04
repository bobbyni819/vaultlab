"""Pre-L4-run synthetic end-to-end integration test for the vaultlab pipeline.

Exercises the FULL /lit-arc + /build-deck + /lit-report + /onboard-project
pipelines using:

* Real Python code paths (corpus build, picker, summarize, arc render,
  deck render, report assembly, project view writer, onboarding init).
* Stubbed callbacks (picker / reader / narrator / runner / section_writer
  / plan_generator) so no LLM is hit and no tokens spent.
* Test-injected fakes for ``_client`` (search) and ``_acquire`` (PDF
  acquisition) so no network is touched.

Run from the vaultlab repo root::

    python scripts/_e2e_integration_test.py

Exits 0 if every assertion passes; non-zero on the first failure (with a
loud traceback). Bobby reads the report stanza below the test runs to
decide whether to ship the L4 run.
"""

from __future__ import annotations

import json
import logging
import os
import sys
import tempfile
import traceback
from pathlib import Path
from typing import Any

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[attr-defined]
except Exception:
    pass

# Quiet logging so the report section is readable.
logging.basicConfig(
    level=logging.WARNING, format="%(levelname)s %(name)s: %(message)s"
)
# Some pipeline modules log at INFO; keep them out of stdout.
logging.getLogger("vaultlab").setLevel(logging.WARNING)
# Belt-and-suspenders: prove we are not relying on a real API key.
os.environ.pop("ANTHROPIC_API_KEY", None)


# ---------------------------------------------------------------------------
# Fake Paper / search client / acquisition (shared across paths 1 and 3)
# ---------------------------------------------------------------------------


def _seeds():
    """Return a tiny synthetic corpus (3 papers) covering 3 year buckets."""
    from vaultlab.research.paper import Paper

    return [
        Paper(
            title="Foundational programmable cleavage",
            authors=["Jinek M"],
            year=2012,
            journal="Science",
            doi="10.1126/science.1225829",
            citation_count=900,
            source_api="pubmed",
            abstract=(
                "Foundational programmable cleavage of dsDNA via "
                "dual-RNA-guided endonuclease."
            ),
        ),
        Paper(
            title="Cytidine base editor",
            authors=["Komor AC"],
            year=2016,
            journal="Nature",
            doi="10.1038/nature17946",
            citation_count=600,
            source_api="pubmed",
            abstract="Direct C-to-T conversion in genomic DNA without DSBs.",
        ),
        Paper(
            title="Prime editor",
            authors=["Anzalone AV"],
            year=2019,
            journal="Nature",
            doi="10.1038/s41586-019-1711-4",
            citation_count=300,
            source_api="pubmed",
            abstract="Search-and-replace genome editing via reverse transcriptase.",
        ),
    ]


class _FakeClient:
    """Stand-in for ``ResearchClient`` — returns canned seeds."""

    def __init__(self, seeds):
        self._seeds = seeds

    def search(self, query, max_results=20, sources=None):
        return list(self._seeds)


def _fake_fetch_refs(doi: str):
    """Tiny one-layer reference graph: each paper cites the previous one."""
    from vaultlab.research.citation_lookup import Reference

    chain = {
        "10.1126/science.1225829": [],
        "10.1038/nature17946": [Reference(doi="10.1126/science.1225829")],
        "10.1038/s41586-019-1711-4": [
            Reference(doi="10.1126/science.1225829"),
            Reference(doi="10.1038/nature17946"),
        ],
    }
    return chain.get(doi)


def _fake_acquire(corpus, cache_dir, **kwargs):
    """Fake acquisition: write a minimal valid PDF blob for every DOI."""
    from vaultlab.research.acquisition import AcquisitionResult, cache_path_for

    cache_dir = Path(cache_dir)
    cache_dir.mkdir(parents=True, exist_ok=True)
    out: dict[str, AcquisitionResult] = {}
    for doi in corpus.papers:
        target = cache_path_for(doi, cache_dir)
        target.write_bytes(b"%PDF-1.4\n" + b"x" * 4000)
        out[doi] = AcquisitionResult(
            doi=doi, pdf_path=target, source="unpaywall", license="cc-by",
        )
    return out


# ---------------------------------------------------------------------------
# Stub callbacks (deterministic — never hit any LLM)
# ---------------------------------------------------------------------------


def stub_picker(task) -> dict[str, Any]:
    """Top-N candidates by og_score (deterministic)."""
    sorted_cands = sorted(task.candidates, key=lambda c: -c.og_score)[: task.target_n]
    return {
        "picks": [
            {"doi": c.doi, "rank": i, "rationale": "stub-rationale"}
            for i, c in enumerate(sorted_cands, 1)
        ]
    }


def stub_reader(task) -> dict[str, Any]:
    """Canned per-paper summary."""
    return {
        "tldr": (
            f"Synthetic TL;DR for {task.doi}. Sentence two. Sentence three."
        ),
        "why_it_matters": [
            "Synthetic novelty bullet 1",
            "Synthetic novelty bullet 2",
        ],
        "methods_summary": "Synthetic methods paragraph for the integration test.",
        "key_findings": [
            f"Synthetic finding {i} [p{i + 1}]" for i in range(3)
        ],
        "extracted_references": [],
    }


def stub_narrator(task) -> dict[str, Any]:
    """Canned arc paragraphs with wikilinks (one per bucket)."""
    return {
        "history": (
            "History paragraph with [[10.1126_science.1225829|Jinek 2012]] wikilink."
        ),
        "development": (
            "Development paragraph with [[10.1038_nature17946|Komor 2016]] wikilink."
        ),
        "sota": (
            "SOTA paragraph with "
            "[[10.1038_s41586-019-1711-4|Anzalone 2019]] wikilink."
        ),
    }


def stub_runner(meeting, members):
    """Canned per-role outputs for adversarial meetings.

    Handles every meeting purpose:
      * picker  -> synthesizer emits {"picks": [...]} JSON
      * arc     -> synthesizer emits {"history","development","sota"} JSON
      * deck    -> synthesizer emits {"slides":[...]} JSON
      * report  -> synthesizer emits {"section_text",...} JSON
      * audit   -> rigor_auditor emits {"passed":true,"issues":[]} JSON
    """
    topic_str = (meeting.topic or "").lower()
    agenda_text = ""
    if meeting.agenda is not None:
        agenda_text = (meeting.agenda.statement or "").lower()
    # Priority-ordered detection — earlier checks win when statements overlap.
    is_audit = topic_str.startswith("rigor audit")
    # Section meetings: "Write the <label> section (N words) of a deep-research..."
    is_report_section = (
        "of a deep-research review" in agenda_text
        and "section" in agenda_text
    )
    # Picker: "Pick the N BEST papers from the M candidates..." (note: also
    # mentions "literature lineage arc" — must check picker BEFORE arc).
    is_picker = "best papers" in agenda_text
    # Deck plan: "Plan a N-slide deck for '<topic>'..." (also mentions
    # "narrative arc" but is_picker is already false at this point).
    is_deck = "-slide deck" in agenda_text or "slide deck for" in agenda_text
    # Arc: "Write a 3-paragraph lineage arc (history / development / sota)..."
    is_arc = "3-paragraph lineage arc" in agenda_text

    # Resolve which report section is active (label uses spaces, not _).
    section_id = "background"
    section_targets = {
        "background": 650,
        "methods_landscape": 1000,
        "findings": 1250,
        "contradictions": 400,
        "future_directions": 300,
    }
    section_label_to_id = {
        "background": "background",
        "methods landscape": "methods_landscape",
        "findings": "findings",
        "contradictions": "contradictions",
        "future directions": "future_directions",
    }
    if is_report_section:
        for label, sid in section_label_to_id.items():
            # Match "write the <label> section"
            if f"write the {label} section" in agenda_text:
                section_id = sid
                break

    outputs: list[dict[str, Any]] = []
    for role in members:
        rid = (role.id or "").lower()
        if rid == "rigor_auditor" or is_audit:
            payload = {"passed": True, "issues": []}
            outputs.append({"output": json.dumps(payload)})
            continue
        if rid != "synthesizer":
            outputs.append({"output": f"[{role.id}] stub critique"})
            continue

        # Synthesizer: emit purpose-specific JSON.
        # Priority: section > picker > deck > arc (matches the priority of
        # the agenda detection — section statements never collide with the
        # other meeting types, but picker mentions "lineage arc", deck
        # mentions "narrative arc", so order matters).
        if is_report_section:
            target_words = section_targets[section_id]
            base = (
                f"Synthetic {section_id.replace('_', ' ')} body. "
                "[[10.1126_science.1225829|Jinek 2012]] founded the field. "
                "[[10.1038_nature17946|Komor 2016]] extended it. "
                "[[10.1038_s41586-019-1711-4|Anzalone 2019]] reached SOTA. "
            )
            words = base.split()
            full_words: list[str] = []
            while len(full_words) < target_words:
                full_words.extend(words)
            full_text = " ".join(full_words[:target_words])
            payload = {
                "section_text": full_text,
                "claims_with_evidence": [
                    {
                        "claim": f"{section_id} canonical claim",
                        "doi_slugs": ["10.1126_science.1225829"],
                    },
                ],
            }
        elif is_picker:
            payload = {
                "picks": [
                    {
                        "doi": "10.1126/science.1225829",
                        "rank": 1,
                        "rationale": "synth-stub",
                    },
                    {
                        "doi": "10.1038/nature17946",
                        "rank": 2,
                        "rationale": "synth-stub",
                    },
                ]
            }
        elif is_deck:
            payload = {
                "story_arc_summary": "Synthetic adversarial deck arc.",
                "slides": [
                    {
                        "type": "title",
                        "title": "Adversarial Deck",
                        "subtitle": "Synthetic",
                    },
                    {
                        "type": "section_divider",
                        "title": "Background",
                    },
                    {
                        "type": "text",
                        "title": "Foundational",
                        "bullets": [
                            "Programmable cleavage [[10.1126_science.1225829|Jinek 2012]]",
                            "Base editing [[10.1038_nature17946|Komor 2016]]",
                        ],
                        "speaker_notes": {},
                    },
                    {
                        "type": "section_divider",
                        "title": "SOTA",
                    },
                    {
                        "type": "text",
                        "title": "Prime editing",
                        "bullets": [
                            "Search-and-replace [[10.1038_s41586-019-1711-4|Anzalone 2019]]",
                        ],
                        "speaker_notes": {},
                    },
                ],
            }
        elif is_arc:
            payload = {
                "history": (
                    "History paragraph with "
                    "[[10.1126_science.1225829|Jinek 2012]] wikilink."
                ),
                "development": (
                    "Development paragraph with "
                    "[[10.1038_nature17946|Komor 2016]] wikilink."
                ),
                "sota": (
                    "SOTA paragraph with "
                    "[[10.1038_s41586-019-1711-4|Anzalone 2019]] wikilink."
                ),
            }
        else:
            # Unknown synthesizer purpose — emit empty dict.
            payload = {}
        outputs.append({"output": json.dumps(payload)})

    return outputs


def stub_section_writer(task) -> dict[str, Any]:
    """Single-shot section writer fallback (used only when no crosstalk_runner)."""
    word_target = task.target_word_count
    canned = (
        " ".join(["word"] * max(1, word_target))
        + " [[10.1126_science.1225829|Stub 2024]]."
    )
    return {"section_text": canned, "claims_with_evidence": []}


def stub_plan_generator(task) -> dict[str, Any]:
    """Single-shot deck plan callback for the non-adversarial path."""
    return {
        "story_arc_summary": "Stub arc",
        "slides": [
            {
                "type": "title",
                "title": "Test Deck",
                "subtitle": "Synthetic",
            },
            {
                "type": "section_divider",
                "title": "Section 1",
            },
            {
                "type": "text",
                "title": "Findings",
                "bullets": [
                    "Finding A [[10.1126_science.1225829|Jinek 2012]]",
                    "Finding B [[10.1038_nature17946|Komor 2016]]",
                ],
                "speaker_notes": {},
            },
            {
                "type": "section_divider",
                "title": "Section 2",
            },
            {
                "type": "text",
                "title": "More",
                "bullets": [
                    "Detail [[10.1038_s41586-019-1711-4|Anzalone 2019]]",
                ],
                "speaker_notes": {},
            },
        ],
    }


# ---------------------------------------------------------------------------
# Result + bug accumulator
# ---------------------------------------------------------------------------


class TestState:
    def __init__(self):
        self.path_results: dict[str, str] = {}
        self.bugs: list[str] = []
        self.notes: list[str] = []
        self.kb_dirs: dict[str, Path] = {}

    def pass_path(self, name: str):
        self.path_results[name] = "PASS"

    def fail_path(self, name: str, reason: str):
        self.path_results[name] = f"FAIL ({reason[:80]})"
        self.bugs.append(f"[{name}] {reason}")


# ---------------------------------------------------------------------------
# Path 1: /lit-arc with all features on (adversarial picker + arc, content-aware)
# ---------------------------------------------------------------------------


def run_path1(state: TestState, kb_root: Path) -> "Any":
    from vaultlab.research.lineage import run_lit_arc

    project_slug = "e2e-integration-test"

    result = run_lit_arc(
        "test topic",
        kb_root=kb_root,
        project_slug=project_slug,
        depth="balanced",
        max_seeds=3,
        picker_callback=stub_picker,
        reader=stub_reader,
        narrator=stub_narrator,
        picker_mode="adversarial",
        arc_mode="adversarial",
        crosstalk_runner=stub_runner,
        crosstalk_n_rounds=2,
        speaker="Test User",
        max_papers_to_summarize=2,  # force the picker meeting to actually run
        # Test injection: no network, deterministic.
        _client=_FakeClient(_seeds()),
        _fetch_refs=_fake_fetch_refs,
        _acquire=_fake_acquire,
        _today="2026-04-30",
    )

    # Verify outputs.
    assert result.arc_path.exists(), f"arc not written: {result.arc_path}"
    assert result.search_log_path.exists(), (
        f"search log not written: {result.search_log_path}"
    )
    missing_summaries = [str(p) for p in result.summary_paths.values() if not p.exists()]
    assert not missing_summaries, f"summaries missing: {missing_summaries}"

    # Provenance receipts (sidecars next to arc).
    prov_json = result.arc_path.with_name(result.arc_path.name + ".provenance.json")
    prov_md = result.arc_path.with_name(result.arc_path.name + ".method.md")
    assert prov_json.exists(), f"provenance.json missing: {prov_json}"
    assert prov_md.exists(), f"method.md missing: {prov_md}"

    # Project view files (Wiki/Projects/<slug>/).
    project_dir = kb_root / "Wiki" / "Projects" / project_slug
    assert (project_dir / "papers.md").exists(), "papers.md not written"
    assert (project_dir / "lineage.md").exists(), "lineage.md not written"
    assert (project_dir / "decisions-log.md").exists(), (
        "decisions-log.md not written"
    )
    assert (project_dir / "START_HERE.md").exists(), "START_HERE.md not written"

    # Sources/Articles/<slug>.md per seed.
    articles = list((kb_root / "Sources" / "Articles").glob("*.md"))
    assert articles, "no Sources/Articles/<doi>.md stubs written"

    # Sources/Papers/<slug>.pdf cached.
    pdfs = list((kb_root / "Sources" / "Papers").glob("*.pdf"))
    assert pdfs, "no Sources/Papers/<doi>.pdf written"

    state.kb_dirs["path1"] = kb_root
    return result


# ---------------------------------------------------------------------------
# Path 2: /build-deck adversarial plan + final audit
# ---------------------------------------------------------------------------


def run_path2(state: TestState, kb_root: Path, lineage_result):
    from vaultlab.slides import build_deck_from_lineage_result

    deck_path = build_deck_from_lineage_result(
        lineage_result,
        speaker="Test",
        affiliation="Test Lab",
        project_slug="e2e-integration-test",
        figure_assignments={},  # no real figures
        kb_root=kb_root,
        plan_callback=stub_plan_generator,  # exercise non-adversarial path too
        plan_mode="adversarial",  # crosstalk on
        crosstalk_runner=stub_runner,
        final_audit=True,
        audit_strict=False,
        audience="journal-club",
        target_slide_count=5,
    )

    assert deck_path.exists(), f"deck not written: {deck_path}"
    from pptx import Presentation

    pres = Presentation(str(deck_path))
    n_slides = len(pres.slides)
    # The auto-appended references slide should bring this to >=5.
    # Stub plan emits: title + section_divider + text + section_divider + text
    # → 5 LLM-side; then renderer auto-appends references → 6 total.
    assert n_slides >= 5, f"expected >=5 slides, got {n_slides}"
    state.notes.append(f"Path 2 deck slides: {n_slides}")
    return deck_path


# ---------------------------------------------------------------------------
# Path 3: /lit-report 5-section pipeline
# ---------------------------------------------------------------------------


def run_path3(state: TestState, kb_root: Path):
    from vaultlab.research.report import SECTION_ORDER, run_lit_report

    result = run_lit_report(
        "test topic",
        kb_root=kb_root,
        project_slug="e2e-integration-test-report",
        depth="balanced",
        max_seeds=3,
        max_papers_to_summarize=2,
        picker_callback=stub_picker,
        reader=stub_reader,
        section_writer=stub_section_writer,  # fallback path
        crosstalk_runner=stub_runner,         # adversarial path (preferred)
        crosstalk_n_rounds=2,
        speaker="Test",
        audience="graduate-student",
        target_total_words=2000,
        _client=_FakeClient(_seeds()),
        _fetch_refs=_fake_fetch_refs,
        _acquire=_fake_acquire,
        _today="2026-04-30",
    )

    assert result.report_path.exists(), f"report not written: {result.report_path}"
    text = result.report_path.read_text(encoding="utf-8")
    section_h2 = {
        "background": "Background",
        "methods_landscape": "Methods landscape",
        "findings": "Key findings",
        "contradictions": "Contradictions & open questions",
        "future_directions": "Future directions",
    }
    for sec in SECTION_ORDER:
        label = section_h2[sec]
        assert f"## {label}" in text, f"section '{label}' missing"

    # Per-section drafts on disk.
    for sec in SECTION_ORDER:
        p = result.section_paths.get(sec)
        assert p is not None and p.exists(), f"missing per-section draft for {sec}"

    state.notes.append(
        f"Path 3 word count: {result.word_count} (sections "
        f"{result.section_word_counts})"
    )
    return result


# ---------------------------------------------------------------------------
# Path 4: /onboard-project end-to-end
# ---------------------------------------------------------------------------


def run_path4(state: TestState, kb_root: Path, project_path: Path):
    from vaultlab.onboarding import (
        init_project_from_intake,
        render_intake_template,
    )

    project_path.mkdir(parents=True, exist_ok=True)
    intake_md = (
        render_intake_template()
        .replace("YOUR ANSWER:", "YOUR ANSWER: Test research topic", 1)
        .replace(
            "- [ ] Understand a literature field",
            "- [x] Understand a literature field",
        )
        .replace(
            "- [ ] PI / weekly meeting",
            "- [x] PI / weekly meeting",
        )
    )
    intake_path = project_path / "project_intake.md"
    intake_path.write_text(intake_md, encoding="utf-8")

    result = init_project_from_intake(
        intake_path=intake_path,
        kb_root=kb_root,
        project_path=project_path,
        slug="onboarding-test",
    )

    assert (
        kb_root / "Wiki" / "Projects" / "onboarding-test" / "START_HERE.md"
    ).exists(), "onboarding START_HERE.md not written"
    assert (project_path / ".vaultlab-project.json").exists(), (
        ".vaultlab-project.json not written"
    )
    state.notes.append(
        f"Path 4 wrote {len(result.files_written())} files; "
        f"{len(result.follow_up_questions)} follow-up Qs"
    )
    return result


# ---------------------------------------------------------------------------
# KB structure audit: does the synthetic pipeline produce the same shape
# the real L4 run will?
# ---------------------------------------------------------------------------


def audit_kb_structure(state: TestState, kb_root: Path, project_slug: str):
    """Spot-check the canonical L4 KB structure is present."""
    expected_dirs = [
        kb_root / "Sources" / "Articles",
        kb_root / "Sources" / "Papers",
        kb_root / "Sources" / "Notes",
        kb_root / "Wiki" / "Concepts",
        kb_root / "Wiki" / "Summaries",
        kb_root / "Wiki" / "Projects" / project_slug,
        kb_root / "Output" / project_slug,
    ]
    missing: list[str] = []
    for d in expected_dirs:
        if not d.exists():
            missing.append(str(d.relative_to(kb_root)))
    if missing:
        state.bugs.append(
            "[KB structure] Missing canonical directories: " + ", ".join(missing)
        )
        return False
    return True


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    print("=" * 78)
    print("Vaultlab pre-L4 synthetic E2E integration test")
    print("=" * 78)

    state = TestState()

    with tempfile.TemporaryDirectory(prefix="vaultlab_e2e_") as td:
        td_root = Path(td)
        kb_a = td_root / "kb_a"
        kb_b = td_root / "kb_b"
        kb_a.mkdir()
        kb_b.mkdir()
        project_path = td_root / "synthetic_project"

        # ------------------------------------------------------------
        # Path 1: /lit-arc with all features on
        # ------------------------------------------------------------
        print("\n[Path 1] /lit-arc — adversarial picker + arc, content-aware...")
        try:
            lineage_result = run_path1(state, kb_a)
            state.pass_path("Path 1 (/lit-arc)")
            print("  PASS")
        except Exception:
            tb = traceback.format_exc()
            state.fail_path("Path 1 (/lit-arc)", tb.splitlines()[-1])
            state.bugs.append("Path 1 traceback:\n" + tb)
            print("  FAIL")
            print(tb)
            lineage_result = None

        # ------------------------------------------------------------
        # Path 2: /build-deck adversarial plan + final audit
        # ------------------------------------------------------------
        print("\n[Path 2] /build-deck — adversarial plan + final audit...")
        if lineage_result is None:
            state.fail_path(
                "Path 2 (/build-deck)",
                "skipped — Path 1 must succeed first",
            )
            print("  SKIPPED (Path 1 failed)")
        else:
            try:
                run_path2(state, kb_a, lineage_result)
                state.pass_path("Path 2 (/build-deck)")
                print("  PASS")
            except Exception:
                tb = traceback.format_exc()
                state.fail_path("Path 2 (/build-deck)", tb.splitlines()[-1])
                state.bugs.append("Path 2 traceback:\n" + tb)
                print("  FAIL")
                print(tb)

        # ------------------------------------------------------------
        # Path 3: /lit-report 5-section
        # ------------------------------------------------------------
        print("\n[Path 3] /lit-report — 5-section adversarial pipeline...")
        try:
            run_path3(state, kb_b)
            state.pass_path("Path 3 (/lit-report)")
            print("  PASS")
        except Exception:
            tb = traceback.format_exc()
            state.fail_path("Path 3 (/lit-report)", tb.splitlines()[-1])
            state.bugs.append("Path 3 traceback:\n" + tb)
            print("  FAIL")
            print(tb)

        # ------------------------------------------------------------
        # Path 4: /onboard-project
        # ------------------------------------------------------------
        print("\n[Path 4] /onboard-project — intake → init...")
        try:
            run_path4(state, kb_a, project_path)
            state.pass_path("Path 4 (/onboard-project)")
            print("  PASS")
        except Exception:
            tb = traceback.format_exc()
            state.fail_path("Path 4 (/onboard-project)", tb.splitlines()[-1])
            state.bugs.append("Path 4 traceback:\n" + tb)
            print("  FAIL")
            print(tb)

        # ------------------------------------------------------------
        # KB structure audit
        # ------------------------------------------------------------
        print("\n[KB structure] Auditing canonical L4 KB layout under kb_a...")
        ok = audit_kb_structure(state, kb_a, "e2e-integration-test")
        if ok:
            state.notes.append(
                "KB layout matches L4 expectation: Sources/{Articles,Papers,Notes}, "
                "Wiki/{Concepts,Summaries,Projects}, Output/<slug>."
            )
            print("  PASS")
        else:
            print("  FAIL — see bug list")

        # ------------------------------------------------------------
        # Report
        # ------------------------------------------------------------
        print("\n" + "=" * 78)
        print("REPORT")
        print("=" * 78)
        for name, status in state.path_results.items():
            print(f"  {name}: {status}")
        print()
        if state.notes:
            print("Notes:")
            for n in state.notes:
                print(f"  - {n}")
            print()
        if state.bugs:
            print("Bugs / failures (precise):")
            for b in state.bugs:
                # Print first 2 lines of multi-line bug strings to keep output crisp.
                for ln in str(b).splitlines()[:8]:
                    print(f"  {ln}")
                print()
        else:
            print("No bugs surfaced. Pipeline is integration-clean.")

        # Exit with nonzero if any path failed.
        any_fail = any(s.startswith("FAIL") for s in state.path_results.values())
        return 1 if any_fail else 0


if __name__ == "__main__":
    raise SystemExit(main())
