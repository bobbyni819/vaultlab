"""Build the v2 demo PPTX — figure 1 uses extracted regions; rest use v1 plans.

Demonstrates the difference between eye-estimated coords (v1) and color-motif
extracted coords (v2). Slide 2 is the v1 figure 1 (eye-estimated) and slide 3
is the v2 figure 1 (extracted) so Bobby can see the side-by-side improvement.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

# Reuse v1 plans + the new v2-rendered figure 1
sys.path.insert(0, str(Path(__file__).parent))
from figure_understanding_prototype import (  # type: ignore[import-not-found]
    PLANS as V1_PLANS,
)
from figure_understanding_prototype import (
    annotate_figure as v1_annotate,
)

PPTX_OUT = Path(r"C:\Users\bobby\Downloads\car_t_decks\figure_understanding_demo_v2.pptx")
V2_FIG1 = Path(r"C:\tmp\cart_figs_v13_annotated_v2\figure1_v2.png")


def _add_title_slide(pres: Presentation) -> None:
    blank = pres.slide_layouts[6]
    s = pres.slides.add_slide(blank)
    title_box = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VaultLab Figure-Understanding — v2"
    run.font.name = "Arial"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)

    sub = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5))
    sf = sub.text_frame
    sf.word_wrap = True
    sub_text = (
        "v1 (slide 2) used my eye-estimated fractional coordinates -- 3 of 7 boxes "
        "were completely misplaced.\n\n"
        "v2 (slide 3) uses programmatic color-motif extraction "
        "(vaultlab.figures.understand) -- regions found by HSV thresholding + "
        "morphological merging, then mapped to concepts. Boxes sit precisely on the "
        "actual visual elements.\n\n"
        "Slides 4-10 are v1 (no v2 yet) for reference."
    )
    for i, line in enumerate(sub_text.split("\n\n")):
        p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(60, 60, 60)


def _add_figure_slide(
    pres: Presentation,
    title: str,
    image_path: Path,
    caption: str,
    speaker_notes: str,
) -> None:
    blank = pres.slide_layouts[6]
    s = pres.slides.add_slide(blank)

    title_box = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 20, 20)

    avail_w_in = 12.5
    avail_h_in = 5.6
    img = Image.open(image_path)
    aspect = img.width / img.height
    avail_aspect = avail_w_in / avail_h_in
    if aspect > avail_aspect:
        disp_w = Inches(avail_w_in)
        disp_h = Inches(avail_w_in / aspect)
    else:
        disp_h = Inches(avail_h_in)
        disp_w = Inches(avail_h_in * aspect)

    x_in = (13.333 - disp_w / Emu(914400)) / 2
    s.shapes.add_picture(
        str(image_path),
        Inches(x_in),
        Inches(1.0),
        width=disp_w,
        height=disp_h,
    )

    caption_box = s.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7))
    cf = caption_box.text_frame
    cf.word_wrap = True
    p = cf.paragraphs[0]
    run = p.add_run()
    run.text = caption
    run.font.name = "Arial"
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = RGBColor(60, 60, 60)

    s.notes_slide.notes_text_frame.text = speaker_notes


def main() -> None:
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    _add_title_slide(pres)

    # Slide 2: v1 figure 1 (eye-estimated; the BAD one)
    v1_fig1 = v1_annotate(V1_PLANS[0])
    _add_figure_slide(
        pres,
        "Figure 1 (v1) - eye-estimated coords",
        v1_fig1,
        "v1: 3 of 7 boxes are misplaced. Coordinates were guessed from the visual; "
        "no programmatic verification.",
        "v1 protocol: I read the figure visually, estimated fractional coordinates "
        "by eye, drew the boxes. No verification step. Bobby's review (2026-04-29) "
        "found:\n"
        "- #1 (Endogenous TCR): box too large; could shrink top + right.\n"
        "- #2 (Introduced TCR): mostly right; could tighten.\n"
        "- #3 (MHC Class I): COMPLETELY OFF -- floating on a tumor cell instead of "
        "on top of the introduced TCR.\n"
        "- #4 (CD3): mostly right; bottom should extend to capture gamma chains.\n"
        "- #5 (TAA presented): floating; should capture the aberrant protein + MHC + "
        "endogenous TCR in panel b.\n"
        "- #6 (scFv): floating on tumor cell; should be on the panel-c CAR construct.\n"
        "- #7 (CAR signaling domains): too high; should extend downward into the "
        "intracellular region.",
    )

    # Slide 3: v2 figure 1 (color-extracted)
    _add_figure_slide(
        pres,
        "Figure 1 (v2) - color-motif extracted",
        V2_FIG1,
        "v2: pixel-precise boxes from vaultlab.figures.understand. HSV color "
        "thresholding -> connected components -> region merging -> concept matching.",
        "v2 protocol uses vaultlab.figures.understand:\n\n"
        "1. extract_regions: HSV-threshold the image for 5 color motifs (neon green, "
        "electric blue, orange, red-magenta, purple-violet). Returns 60 raw "
        "connected components.\n\n"
        "2. merge_regions: collapse fragments split by dark outlines. With "
        "dilation=25px, 60 raw components merge into 20 element-level regions.\n\n"
        "3. concept matching: pick regions per panel via x-position filter "
        "(panel a: x < W/3; panel b: W/3 <= x < 2W/3; panel c: x >= 2W/3). Apply "
        "domain rules (panel-a leftmost blue = endogenous TCR; panel-a largest "
        "purple = MHC class I; etc.).\n\n"
        "Result: 9 concept annotations with pixel-precise boxes that sit on the "
        "actual visual elements. Compare slide 2 (eye-estimated) to slide 3 "
        "(extracted) -- the same 7 elements Bobby flagged are now correctly placed.\n\n"
        "Code: scripts/figure_understanding_v2.py + src/vaultlab/figures/understand/.",
    )

    # Slides 4-10: v1 figures 2-8 unchanged
    for plan in V1_PLANS[1:]:
        annotated = v1_annotate(plan)
        _add_figure_slide(
            pres,
            f"{plan.title} (v1)",
            annotated,
            plan.overall_caption,
            plan.long_explanation
            + "\n\nNumbered annotations:\n"
            + "\n".join(
                f"{i + 1}. [{a.label}] -- {a.explanation}" for i, a in enumerate(plan.annotations)
            ),
        )

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    pres.save(PPTX_OUT)
    print(f"PPTX written -> {PPTX_OUT}")


if __name__ == "__main__":
    main()
