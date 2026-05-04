"""Infer which figure-layout was applied to a rendered slide.

Bobby's 2026-05-04 ask: read the rendered pptx and report which layout
each slide actually got — so we know layout dispatch produced what we
expected. This is the post-build audit that closes the loop on
"is the deck rendering what I asked for."

Heuristic geometry-based classification by reading shape positions:

- ``figure_top_caption_br`` — figure shape positioned at top spanning
  full slide width (>80%); caption / citation in bottom-right.
- ``figure_with_side_caption`` — figure on left (60-65% width) full
  slide height (>70%); caption stacked in right gutter at top.
- ``figure_above_bullets`` — figure top half full-width; bullets bottom.
- ``figure_only`` — figure horizontally centered, large, no bullets.
- ``default`` (figure_slide) — figure left ~60% width, ~70% height; bullets
  to the right of figure.
- ``unknown`` — slide has a picture but doesn't match above patterns.
- None — slide has no picture (title / text / section_divider).

Returns the inferred layout name as a string.
"""

from __future__ import annotations

from typing import Any


_PIC_TYPE = 13


def infer_slide_layout(slide: Any, slide_w_emu: int, slide_h_emu: int) -> str | None:
    """Infer which figure-layout was used to render this slide.

    Returns the layout name, or ``None`` if the slide has no figure.

    Decision tree based on figure CENTER position + text-shape positions
    around the figure:

    - figure center in upper half AND text in bottom strip → figure_top_caption_br
    - figure center in upper half AND text below in left half → figure_above_bullets
    - figure center left of slide center AND text to the right of figure →
        figure_with_side_caption (figure tall) OR default (figure short)
    - figure center near horizontal middle AND no side/below text →
        figure_only
    """
    pics = [sh for sh in slide.shapes if sh.shape_type == _PIC_TYPE]
    if not pics:
        return None

    pic = pics[0]
    sw_in = slide_w_emu / 914400
    sh_in = slide_h_emu / 914400
    pic_l = (pic.left or 0) / 914400
    pic_t = (pic.top or 0) / 914400
    pic_w = (pic.width or 0) / 914400
    pic_h = (pic.height or 0) / 914400

    fig_cx = pic_l + pic_w / 2
    fig_cy = pic_t + pic_h / 2
    fig_cx_frac = fig_cx / sw_in if sw_in else 0
    fig_cy_frac = fig_cy / sh_in if sh_in else 0
    pic_h_frac = pic_h / sh_in if sh_in else 0
    pic_right_in = pic_l + pic_w
    pic_bottom_in = pic_t + pic_h

    text_shapes = [
        sh for sh in slide.shapes
        if sh.shape_type != _PIC_TYPE and getattr(sh, "has_text_frame", False)
        and (sh.text_frame.text or "").strip()
    ]

    # Text in the right-side gutter (right of figure right edge)
    has_text_right_of_figure = any(
        ((sh.left or 0) / 914400) >= pic_right_in - 0.3
        and ((sh.top or 0) / 914400) > 1.0  # below the title row
        for sh in text_shapes
    )
    # Text below the figure bottom edge in the bottom strip
    has_text_below_figure = any(
        ((sh.top or 0) / 914400) >= pic_bottom_in - 0.2
        for sh in text_shapes
    )
    # Text specifically in the bottom-RIGHT half of the slide
    has_text_in_bottom_right = any(
        ((sh.top or 0) / 914400) >= pic_bottom_in - 0.2
        and (((sh.left or 0) / 914400) + ((sh.width or 0) / 914400) / 2) > sw_in * 0.5
        for sh in text_shapes
    )
    # Text in bottom-LEFT half (bullets in the left gutter for top-BR layout)
    has_text_in_bottom_left = any(
        ((sh.top or 0) / 914400) >= pic_bottom_in - 0.2
        and (((sh.left or 0) / 914400) + ((sh.width or 0) / 914400) / 2) < sw_in * 0.5
        for sh in text_shapes
    )

    # figure_top_caption_br: figure in upper portion, full slide width,
    # caption + citation in bottom-right corner. Identifying signal:
    # figure center horizontally near slide middle AND figure top is
    # near the title-bottom AND there's text content in the bottom strip
    # (specifically bottom-right; bottom-left is OK too if there are
    # bullets there).
    if (
        fig_cy_frac < 0.55
        and 0.35 < fig_cx_frac < 0.65
        and has_text_in_bottom_right
        and pic_h_frac < 0.75
    ):
        return "figure_top_caption_br"

    # figure_above_bullets: figure top half, full-width-ish, bullets below
    # full-width (text below figure but NOT in right gutter — only below).
    if (
        fig_cy_frac < 0.5
        and not has_text_right_of_figure
        and has_text_below_figure
        and pic_h_frac < 0.65
    ):
        return "figure_above_bullets"

    # figure_with_side_caption: figure on left ~60% width, takes >65% slide
    # height (most of the vertical space under the title), text shapes
    # stack in the right gutter.
    if (
        fig_cx_frac < 0.5
        and pic_h_frac >= 0.65
        and has_text_right_of_figure
    ):
        return "figure_with_side_caption"

    # default (figure_slide): figure left ~60% width, ~70% height, bullets
    # right. Less vertical span than side_caption but text still in right.
    if (
        fig_cx_frac < 0.5
        and has_text_right_of_figure
    ):
        return "default"

    # figure_only: figure horizontally centered, no side text, no
    # bottom-strip text (or only caption directly underneath that's
    # narrow centered).
    if 0.35 < fig_cx_frac < 0.65 and not has_text_right_of_figure:
        return "figure_only"

    return "unknown"


def report_deck_layouts(deck_path: Any) -> list[dict[str, Any]]:
    """Return per-slide layout inference for a rendered deck.

    Each entry: ``{"index": N, "title": "...", "layout": "...", "fig_size_in": (w, h)}``.
    Useful for verifying that the auto-layout dispatcher produced what
    we wanted post-build.
    """
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    out: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, 1):
        layout = infer_slide_layout(slide, prs.slide_width, prs.slide_height)
        # Title = first short text shape
        title = ""
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                t = (sh.text_frame.text or "").strip().split("\n")[0]
                if t and len(t) < 120:
                    title = t
                    break
        # Figure size
        fig_size: tuple[float, float] | None = None
        for sh in slide.shapes:
            if sh.shape_type == _PIC_TYPE:
                fig_size = (
                    (sh.width or 0) / 914400,
                    (sh.height or 0) / 914400,
                )
                break
        out.append({
            "index": i,
            "title": title,
            "layout": layout,
            "fig_size_in": fig_size,
        })
    return out


__all__ = ["infer_slide_layout", "report_deck_layouts"]
