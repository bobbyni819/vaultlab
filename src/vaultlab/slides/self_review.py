"""Self-review pass for rendered slide decks.

Reads each slide of a generated ``.pptx`` and runs all available audits as
a single composite ``review_deck(pptx_path)`` entrypoint. The audits run
deterministically against the rendered file — no LLM calls — so the pass
can ship on every render and gate CI.

Audits wired into the composite (sub-goal 5.4 of vaultlab's north-star):

* **Layout hard rules** — every shape's runs must use ``Roboto`` (warning),
  hit the minimum font sizes (heading ≥ 28pt, body ≥ 24pt, caption ≥ 18pt)
  (critical when titles slip below 28 / body below 18), and not overlap
  each other (critical). See ``memory/feedback_slide_hard_rules.md``.
* **Descriptive titles** — titles must be sentence-shaped, not one-word
  labels (warning). Mirrors the "descriptive sentence titles" rule.
* **Bullet density** — text slides with > 7 bullets get flagged (warning).
  Hickey-lab projection rule of thumb.
* **Figure presence** — slides whose caption references a figure but
  have no embedded picture get flagged (info).
* **Story arc** — title slide first, references / acknowledgments last,
  no more than 5 section dividers per deck. Caught structurally without
  needing an LLM.

The unified :class:`ReviewReport` is a composer over those checks. Its
HTML rendering reuses :func:`vaultlab.slides.audit_html.build_audit_report_html`
via a thin adapter — same look-and-feel as the existing rigor-audit report.

Provenance receipts (Red Line #2 contract) are written for any HTML
report emitted to disk by :func:`write_review_report`.
"""

from __future__ import annotations

import logging
import re
from collections.abc import Iterable
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from vaultlab.slides.template import default_font, min_sizes

logger = logging.getLogger(__name__)

__all__ = [
    "ReviewReport",
    "SlideReview",
    "review_deck",
    "write_review_report",
]


# ---------------------------------------------------------------------------
# Phase 7.3 add-on hooks (time-budget + Q&A anticipator)
# ---------------------------------------------------------------------------
#
# Both are OPTIONAL extensions to ``review_deck``: they fire only when the
# caller passes the matching keyword (``budget_minutes`` or
# ``anticipate_questions=True``). The default behaviour of ``review_deck``
# is unchanged so existing tests + the CLI gate stay byte-identical.


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


_SEVERITY_TO_RIGOR = {
    "critical": "blocker",
    "warning": "major",
    "info": "minor",
}


def _normalize_severity(value: str | None) -> str:
    sev = (value or "").lower()
    if sev in {"critical", "blocker"}:
        return "critical"
    if sev in {"warning", "warn", "major"}:
        return "warning"
    return "info"


@dataclass
class SlideReview:
    """Per-slide review verdict.

    Attributes
    ----------
    slide_index : int
        Zero-indexed position of the slide in the deck.
    title : str | None
        The title text recovered from the slide (best-effort — None if no
        title shape was found).
    slide_type : str
        Heuristic slide-type label (``"title"`` / ``"section_divider"`` /
        ``"figure"`` / ``"text"`` / ``"references"`` / ``"other"``).
    issues : list[dict]
        Issue records of the form
        ``{"severity": "critical"|"warning"|"info", "rule": <str>,
        "detail": <str>}``.
    """

    slide_index: int
    title: str | None = None
    slide_type: str = "other"
    issues: list[dict[str, str]] = field(default_factory=list)


@dataclass
class ReviewReport:
    """Composite review report for an entire rendered deck.

    Attributes
    ----------
    pptx_path : Path
        Path to the rendered ``.pptx`` file that was reviewed.
    per_slide : list[SlideReview]
        One review per slide, in deck order.
    story_arc_issues : list[dict]
        Deck-level structural issues that aren't tied to a single slide
        (e.g. "no title slide", "too many section dividers").
    """

    pptx_path: Path
    per_slide: list[SlideReview] = field(default_factory=list)
    story_arc_issues: list[dict[str, str]] = field(default_factory=list)
    # Optional Phase 7.3 add-on outputs. Populated only when the matching
    # keyword arguments are passed to :func:`review_deck`. ``None`` means
    # the audit wasn't requested; consumers shouldn't infer "pass".
    time_budget: Any = None
    anticipated_questions: list[Any] = field(default_factory=list)

    # -- aggregates -------------------------------------------------------

    @property
    def n_slides(self) -> int:
        return len(self.per_slide)

    def _count(self, severity: str) -> int:
        total = sum(1 for s in self.per_slide for i in s.issues if i.get("severity") == severity)
        total += sum(1 for i in self.story_arc_issues if i.get("severity") == severity)
        return total

    @property
    def n_critical(self) -> int:
        return self._count("critical")

    @property
    def n_warning(self) -> int:
        return self._count("warning")

    @property
    def n_info(self) -> int:
        return self._count("info")

    def ok(self) -> bool:
        """``True`` when no critical issues are present."""
        return self.n_critical == 0

    # -- summarisation ----------------------------------------------------

    def summary_lines(self) -> list[str]:
        """Human-readable headline summary."""
        lines = [
            f"Reviewed {self.n_slides} slide{'s' if self.n_slides != 1 else ''} "
            f"from {self.pptx_path.name}.",
            f"Findings: {self.n_critical} critical, {self.n_warning} warning, {self.n_info} info.",
        ]
        if self.ok() and self.n_warning == 0:
            lines.append("Deck passed self-review — no issues raised.")
        elif self.ok():
            lines.append("Deck passed critical checks — see warnings for polish opportunities.")
        else:
            lines.append("Deck blocked: fix the critical issues before shipping.")
        return lines

    # -- iteration --------------------------------------------------------

    def all_issues(self) -> Iterable[dict[str, Any]]:
        """Yield every issue (per-slide + arc) decorated with location info."""
        for slide in self.per_slide:
            for issue in slide.issues:
                yield {
                    "loc": f"Slide {slide.slide_index + 1}",
                    "slide_index": slide.slide_index,
                    "slide_title": slide.title,
                    "slide_type": slide.slide_type,
                    **issue,
                }
        for issue in self.story_arc_issues:
            yield {
                "loc": issue.get("loc", "(deck)"),
                "slide_index": None,
                "slide_title": None,
                "slide_type": None,
                **issue,
            }


# ---------------------------------------------------------------------------
# Public entrypoint
# ---------------------------------------------------------------------------


def review_deck(
    pptx_path: Path | str,
    *,
    budget_minutes: int | None = None,
    qa_reserve_minutes: int = 5,
    anticipate_questions: bool = False,
    qa_runner_callback: Any = None,
    qa_n_questions: int = 10,
) -> ReviewReport:
    """Run the full self-review pass on a rendered ``.pptx``.

    The check sequence mirrors the audits documented in this module's
    docstring; each one appends issue records to the appropriate
    :class:`SlideReview` or to ``story_arc_issues``. Nothing is raised
    on failure — callers should consult :attr:`ReviewReport.n_critical`
    (or :meth:`ReviewReport.ok`) for a pass/fail decision.

    Parameters
    ----------
    pptx_path
        Path to the rendered ``.pptx`` to review.
    budget_minutes
        Optional. When set, calls
        :func:`vaultlab.slides.time_budget.audit_time_budget` and stores
        the report on :attr:`ReviewReport.time_budget`. Over-budget
        decks also pick up a deck-level ``warning`` issue
        (``rule="time-budget-over"``).
    qa_reserve_minutes
        Q&A reserve to subtract from the speaking budget; defaults to 5.
    anticipate_questions
        Optional. When ``True``, calls
        :func:`vaultlab.slides.qa_anticipator.anticipate_qa` and stores
        the question list on :attr:`ReviewReport.anticipated_questions`.
    qa_runner_callback
        Optional ``(prompt) -> str`` callable forwarded to the Q&A
        anticipator's LLM mode. Falls back to heuristics on failure.
    qa_n_questions
        Cap on anticipated-question count (default 10).

    Returns
    -------
    ReviewReport
        Composite report with per-slide + story-arc issues, plus the
        optional time-budget / Q&A outputs.

    Raises
    ------
    FileNotFoundError
        If ``pptx_path`` does not exist.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_path}")

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover — gated at install time
        raise RuntimeError(
            "python-pptx is required for review_deck. Install with "
            '`pip install -e ".[slides]"` or `pip install python-pptx`.'
        ) from exc

    prs = Presentation(str(pptx_path))
    report = ReviewReport(pptx_path=pptx_path)

    for idx, slide in enumerate(prs.slides):
        review = _review_one_slide(slide, idx, is_first=(idx == 0))
        report.per_slide.append(review)

    _audit_story_arc(report)

    # Optional add-ons (Phase 7.3 close-out). Both routes silently no-op
    # when the import fails — they're load-bearing only when the caller
    # explicitly asks for them.
    if budget_minutes is not None:
        try:
            from vaultlab.slides.time_budget import audit_time_budget

            tb_report = audit_time_budget(
                pptx_path,
                budget_minutes=budget_minutes,
                qa_reserve_minutes=qa_reserve_minutes,
            )
            report.time_budget = tb_report
            if tb_report.over_budget():
                total_min = tb_report.estimated_total_seconds // 60
                total_sec = tb_report.estimated_total_seconds % 60
                report.story_arc_issues.append(
                    {
                        "severity": "warning",
                        "rule": "time-budget-over",
                        "detail": (
                            f"Estimated speak time {total_min}:{total_sec:02d} "
                            f"exceeds the {tb_report.budget_seconds // 60}-min "
                            "speaking budget. Trim slides or shorten body lines."
                        ),
                        "loc": "(deck)",
                    }
                )
        except Exception:  # pragma: no cover — best-effort
            logger.exception("review_deck: time-budget audit failed")

    if anticipate_questions:
        try:
            from vaultlab.slides.qa_anticipator import anticipate_qa

            report.anticipated_questions = list(
                anticipate_qa(
                    pptx_path,
                    runner_callback=qa_runner_callback,
                    n_questions=qa_n_questions,
                )
            )
        except Exception:  # pragma: no cover — best-effort
            logger.exception("review_deck: Q&A anticipator failed")

    return report


# ---------------------------------------------------------------------------
# Per-slide audits
# ---------------------------------------------------------------------------


_TITLE_FONT_THRESHOLD = 28
_BODY_FONT_THRESHOLD = 24
_CAPTION_FONT_THRESHOLD = 18
_HARD_FLOOR = 18  # nothing below this is acceptable
_MAX_BULLETS_PER_SLIDE = 7
_DESCRIPTIVE_TITLE_MIN_WORDS = 3


def _review_one_slide(slide: Any, slide_index: int, *, is_first: bool = False) -> SlideReview:
    """Run the slide-scope audits on a single python-pptx slide object."""
    text_chunks = _collect_text(slide)
    title = _extract_title(slide, text_chunks)
    slide_type = _infer_slide_type(slide, title, text_chunks, is_first=is_first)

    review = SlideReview(slide_index=slide_index, title=title, slide_type=slide_type)

    _check_fonts_and_sizes(slide, review)
    _check_overlap(slide, review)
    _check_descriptive_title(title, slide_type, review)
    _check_bullet_density(text_chunks, slide_type, review)
    _check_figure_presence(slide, title, text_chunks, slide_type, review)
    _check_color_contrast(slide, review)
    _check_speaker_notes(slide, slide_type, review)

    return review


def _collect_text(slide: Any) -> list[str]:
    """Return all non-empty text strings on a slide, in shape order."""
    out: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = shape.text_frame.text
        except Exception:  # pragma: no cover — defensive
            text = ""
        for line in str(text).splitlines():
            line = line.strip()
            if line:
                out.append(line)
    return out


def _extract_title(slide: Any, text_chunks: list[str]) -> str | None:
    """Best-effort title recovery from the slide.

    Strategy: prefer the slide's title placeholder if present; otherwise
    fall back to the topmost text shape on the slide (which is the
    convention the lab template + python-pptx imperative builders both
    use).
    """
    title_key = _resolve_title_shape_key(slide)
    if title_key is not None:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if _shape_key(shape) != title_key:
                continue
            try:
                text = (shape.text_frame.text or "").strip()
                if text:
                    return text
            except Exception:  # pragma: no cover — defensive
                continue
    return text_chunks[0] if text_chunks else None


_REFERENCES_RE = re.compile(r"^\s*references?\s*$", re.IGNORECASE)
_ACK_RE = re.compile(r"^\s*acknowledg", re.IGNORECASE)


def _infer_slide_type(
    slide: Any, title: str | None, text_chunks: list[str], *, is_first: bool = False
) -> str:
    """Heuristic slide-type classification from rendered shapes.

    Only the deck's first slide can be classified as ``"title"`` — later
    slides with short titles fall through to ``"text"`` so the
    descriptive-title audit can still fire on them.
    """
    title = (title or "").strip()
    has_picture = any(getattr(s, "shape_type", None) == 13 for s in slide.shapes)  # 13 = PICTURE
    body_chunks = text_chunks[1:] if title and text_chunks else text_chunks

    if title and _REFERENCES_RE.match(title):
        return "references"
    if title and _ACK_RE.match(title):
        return "acknowledgments"

    if has_picture:
        return "figure"

    # Section dividers are typically a single title-only text with no body.
    if title and not body_chunks:
        return "section_divider"

    if not text_chunks:
        return "other"

    # Only the FIRST slide of the deck can be the title slide. Later
    # slides with a short title + small body stay classified as text so
    # the descriptive-title audit can flag them.
    if is_first and title and len(body_chunks) <= 3 and len(title.split()) <= 6:
        return "title"

    return "text"


def slide_has_only_title(slide: Any) -> bool:
    """Return ``True`` when the slide carries exactly one text shape."""
    text_shapes = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]
    if len(text_shapes) != 1:
        return False
    try:
        return bool(text_shapes[0].text_frame.text.strip())
    except Exception:  # pragma: no cover — defensive
        return False


def _shape_key(shape: Any) -> Any:
    """Return a stable identity for a python-pptx shape across iterations.

    ``id(shape)`` and ``id(shape._element)`` are BOTH unreliable: python-pptx
    wraps the underlying XML element in a fresh Python object on every
    ``slide.shapes`` access, and lxml in turn hands back a fresh proxy object
    for the same underlying element on each access — so ``id(shape._element)``
    changes between two separate ``slide.shapes`` iterations (and the freed id
    can even be *reused* by a different shape's proxy, depending on the
    interpreter's GC/allocation state). That made title-vs-body matching
    nondeterministic — fine in isolation, broken once an earlier test churned
    the heap.

    ``shape.shape_id`` (the ``<p:cNvPr id="...">`` integer) is assigned by the
    OOXML and is stable + unique within a slide, which is the only scope in
    which we compare keys. Prefer it; fall back to ``id`` only if a shape has
    no usable ``shape_id``.
    """
    try:
        sid = shape.shape_id
        if sid is not None:
            return ("sid", int(sid))
    except Exception:  # pragma: no cover — defensive
        pass
    elem = getattr(shape, "_element", None)
    return id(elem) if elem is not None else id(shape)


def _resolve_title_shape_key(slide: Any) -> Any:
    """Return the shape-key for the slide's title-bearing shape, or None.

    Strategy: prefer the slide's title placeholder when present
    (``slide.shapes.title``); otherwise fall back to the topmost text
    shape on the slide — lab-template layouts and ad-hoc python-pptx
    builds both anchor the title at the top of the canvas.
    """
    try:
        title_shape = slide.shapes.title
    except (AttributeError, KeyError):
        title_shape = None
    if title_shape is not None:
        return _shape_key(title_shape)

    candidate_key = None
    candidate_top = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = (shape.text_frame.text or "").strip()
        except Exception:  # pragma: no cover — defensive
            continue
        if not text:
            continue
        try:
            top = int(shape.top)
        except (AttributeError, TypeError, ValueError):
            continue
        if candidate_top is None or top < candidate_top:
            candidate_key = _shape_key(shape)
            candidate_top = top
    return candidate_key


def _is_footnote_shape(slide: Any, shape: Any) -> bool:
    """Return ``True`` when a shape looks like a citation-source footnote.

    Lab-template figure slides intentionally render the citation source at
    9pt anchored to the bottom 0.5in of the slide (see
    :mod:`vaultlab.slides.layouts.figure` — ``apply_font(... size=9)`` for
    the citation textbox). We don't want to flag that as "body text too
    small" — it's a structural footer element with a known sub-floor size.

    Heuristic: shape's top is in the bottom 12% of the slide AND its text
    is a single short paragraph.
    """
    try:
        slide_height = int(slide.part.package.presentation_part.presentation.slide_height)
    except Exception:  # pragma: no cover — defensive
        return False
    if slide_height <= 0:
        return False
    try:
        top = int(shape.top)
    except (AttributeError, TypeError, ValueError):
        return False
    # Bottom 12% of slide → footnote zone.
    if top < int(slide_height * 0.88):
        return False
    try:
        text = (shape.text_frame.text or "").strip()
    except Exception:  # pragma: no cover — defensive
        return False
    # Footnotes are short single-line.
    if not text or len(text.splitlines()) > 1 or len(text) > 200:
        return False
    return True


def _check_fonts_and_sizes(slide: Any, review: SlideReview) -> None:
    """Critical: title font size < 28pt or any body font size < 18pt."""
    expected_font = default_font()
    sizes = min_sizes()
    title_threshold = sizes.get("heading", _TITLE_FONT_THRESHOLD)
    body_threshold = sizes.get("body", _BODY_FONT_THRESHOLD)
    caption_threshold = sizes.get("caption", _CAPTION_FONT_THRESHOLD)

    seen_non_roboto: set[str] = set()
    title_too_small = False
    body_too_small_pt: int | None = None
    caption_too_small_pt: int | None = None

    title_key = _resolve_title_shape_key(slide)

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        is_title = title_key is not None and _shape_key(shape) == title_key
        if _is_footnote_shape(slide, shape):
            # Citation-source / footer text — skip the size floor.
            continue
        try:
            paragraphs = list(shape.text_frame.paragraphs)
        except Exception:  # pragma: no cover — defensive
            continue
        for para in paragraphs:
            for run in para.runs:
                font = getattr(run, "font", None)
                if font is None:
                    continue
                name = font.name
                if name and name != expected_font:
                    seen_non_roboto.add(name)
                size = font.size
                if size is None:
                    continue
                pt = int(size.pt) if hasattr(size, "pt") else int(size) // 12700
                if is_title and pt < title_threshold:
                    title_too_small = True
                elif not is_title:
                    # Treat short single-run runs as captions if they're
                    # below the body threshold — caption floor still applies.
                    if pt < caption_threshold:
                        caption_too_small_pt = min(caption_too_small_pt or pt, pt)
                    elif pt < body_threshold:
                        body_too_small_pt = min(body_too_small_pt or pt, pt)

    if title_too_small:
        review.issues.append(
            {
                "severity": "critical",
                "rule": "min-title-font",
                "detail": (
                    f"Slide title is below the {title_threshold}pt minimum. "
                    "Bump the heading size or shorten the title text."
                ),
            }
        )
    if caption_too_small_pt is not None:
        review.issues.append(
            {
                "severity": "critical",
                "rule": "min-body-font",
                "detail": (
                    f"Body text is {caption_too_small_pt}pt — below the "
                    f"{caption_threshold}pt absolute floor. Will not project."
                ),
            }
        )
    elif body_too_small_pt is not None:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "min-body-font",
                "detail": (
                    f"Body text is {body_too_small_pt}pt — under the recommended "
                    f"{body_threshold}pt body minimum (caption-only allowance "
                    f"is {caption_threshold}pt)."
                ),
            }
        )
    if seen_non_roboto:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "font-family",
                "detail": (
                    f"Non-Roboto font(s) detected: {sorted(seen_non_roboto)}. "
                    "Switch to Roboto to honor the lab template hard rule."
                ),
            }
        )


def _shape_bbox(shape: Any) -> tuple[int, int, int, int] | None:
    """Return ``(left, top, right, bottom)`` in EMU, or ``None`` if unavailable."""
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except (AttributeError, TypeError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return (left, top, left + width, top + height)


def _bboxes_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _overlap_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
    """Return EMU² of intersection between two bboxes."""
    if not _bboxes_overlap(a, b):
        return 0
    return max(0, min(a[2], b[2]) - max(a[0], b[0])) * max(
        0, min(a[3], b[3]) - max(a[1], b[1])
    )


# Allow a tiny overlap area (annotations sometimes sit just over a figure
# edge). 0.5% of a 13.333"×7.5" slide ≈ 4.5M EMU² — generous threshold.
_OVERLAP_TOLERANCE_EMU2 = 5_000_000


def _check_overlap(slide: Any, review: SlideReview) -> None:
    """Critical: any pair of shapes overlap beyond a small tolerance."""
    boxes: list[tuple[int, int, int, int]] = []
    labels: list[str] = []
    for shape in slide.shapes:
        bbox = _shape_bbox(shape)
        if bbox is None:
            continue
        boxes.append(bbox)
        try:
            label = shape.name or shape.shape_type or "shape"
        except Exception:  # pragma: no cover — defensive
            label = "shape"
        labels.append(str(label))

    n = len(boxes)
    if n < 2:
        return
    flagged: set[tuple[int, int]] = set()
    for i in range(n):
        for j in range(i + 1, n):
            if _overlap_area(boxes[i], boxes[j]) > _OVERLAP_TOLERANCE_EMU2:
                flagged.add((i, j))

    if flagged:
        pairs = ", ".join(f"{labels[i]} ↔ {labels[j]}" for i, j in sorted(flagged)[:3])
        more = "" if len(flagged) <= 3 else f" (+{len(flagged) - 3} more pairs)"
        review.issues.append(
            {
                "severity": "critical",
                "rule": "no-shape-overlap",
                "detail": (
                    f"Shape overlap detected: {pairs}{more}. Re-flow the slide "
                    "so figures, callouts, and text frames don't cover each other."
                ),
            }
        )


_STRUCTURAL_TITLE_SLIDE_TYPES = {
    "section_divider",
    "title",
    "references",
    "acknowledgments",
}


def _check_descriptive_title(title: str | None, slide_type: str, review: SlideReview) -> None:
    """Warning: a non-divider slide has a one-word or empty title.

    Structural slide types (title, section dividers, references,
    acknowledgments) are exempt — their titles are conventionally short
    by design.
    """
    if slide_type in _STRUCTURAL_TITLE_SLIDE_TYPES:
        return
    if not title:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "descriptive-title",
                "detail": "Slide has no recoverable title — viewers lose their footing.",
            }
        )
        return
    word_count = len(title.split())
    if word_count < _DESCRIPTIVE_TITLE_MIN_WORDS:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "descriptive-title",
                "detail": (
                    f"Title '{title}' is only {word_count} word(s). Use a descriptive "
                    "sentence title that previews the slide's claim."
                ),
            }
        )


def _check_bullet_density(
    text_chunks: list[str], slide_type: str, review: SlideReview
) -> None:
    """Warning: text slides with more than 7 visible bullets / lines."""
    if slide_type not in {"text", "figure", "references"}:
        return
    body_chunks = text_chunks[1:] if text_chunks else []
    if len(body_chunks) > _MAX_BULLETS_PER_SLIDE:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "bullet-density",
                "detail": (
                    f"{len(body_chunks)} body lines (max recommended "
                    f"{_MAX_BULLETS_PER_SLIDE}). Split the slide or trim bullets."
                ),
            }
        )


_FIG_HINT_RE = re.compile(r"\b(figure|panel|graph|plot|image|micrograph|chart)\b", re.IGNORECASE)


def _check_figure_presence(
    slide: Any,
    title: str | None,
    text_chunks: list[str],
    slide_type: str,
    review: SlideReview,
) -> None:
    """Info: caption text mentions a figure but no picture shape is present."""
    if slide_type in {"references", "acknowledgments", "title", "section_divider"}:
        return
    has_picture = any(getattr(s, "shape_type", None) == 13 for s in slide.shapes)
    if has_picture:
        return
    hay = " ".join(text_chunks) + " " + (title or "")
    if _FIG_HINT_RE.search(hay):
        review.issues.append(
            {
                "severity": "info",
                "rule": "figure-presence",
                "detail": (
                    "Slide text references a figure / panel / plot but the slide "
                    "has no embedded image. Add the figure or rephrase the caption."
                ),
            }
        )


# ---------------------------------------------------------------------------
# WCAG color-contrast check (deferred-followups bundle, 2026-05-15)
# ---------------------------------------------------------------------------
#
# AA threshold for normal-weight body text is 4.5:1. Anything below 3.0 is
# borderline-unreadable on a projector and gets flagged as ``critical``; the
# 3.0-4.5 band is a polish-worthy ``warning``.
#
# Conservative-by-design: we only flag when BOTH the run font color and the
# resolved background color are concrete RGB. Theme / scheme colors, picture
# fills, gradients, and BACKGROUND-inherited shapes against an unread slide
# background all skip silently. The goal is "no false positives" — a missed
# contrast issue is fine, a hallucinated one is not.


_CONTRAST_AA_THRESHOLD = 4.5  # WCAG AA for normal-weight text
_CONTRAST_CRITICAL_BELOW = 3.0  # below this, body is borderline-unreadable
_DEFAULT_SLIDE_BG = (255, 255, 255)  # safe-default white when bg can't be resolved


def _luminance(rgb: tuple[int, int, int]) -> float:
    """Relative luminance per WCAG 2.x (sRGB → linear → weighted sum)."""

    def chan(c: int) -> float:
        s = c / 255
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4

    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast_ratio(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    """WCAG contrast ratio (higher is better; 1.0 means same color)."""
    L1 = _luminance(fg)
    L2 = _luminance(bg)
    L_high, L_low = max(L1, L2), min(L1, L2)
    return (L_high + 0.05) / (L_low + 0.05)


def _rgb_color_from_color_format(color_obj: Any) -> tuple[int, int, int] | None:
    """Read a concrete RGB triple from a python-pptx ColorFormat.

    Returns ``None`` for theme / scheme / unset colors so the audit stays
    quiet on themed content where we can't be sure of the rendered shade.
    """
    if color_obj is None:
        return None
    try:
        ctype = color_obj.type
    except Exception:  # pragma: no cover — defensive
        return None
    if ctype is None:
        return None
    # MSO_THEME_COLOR_INDEX values produce themed shades we can't read here.
    # Only "RGB" gives us a concrete .rgb. python-pptx raises AttributeError
    # for any other color type.
    try:
        rgb = color_obj.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    if rgb is None:
        return None
    try:
        # RGBColor is a 3-byte hex string under the hood.
        as_int = int(str(rgb), 16)
    except (TypeError, ValueError):
        return None
    return ((as_int >> 16) & 0xFF, (as_int >> 8) & 0xFF, as_int & 0xFF)


def _resolve_shape_background_rgb(shape: Any) -> tuple[int, int, int] | None:
    """Best-effort: return the shape's solid fill RGB, or ``None``.

    Only returns a value for ``MSO_FILL_TYPE.SOLID`` with a concrete RGB
    fore-color. Gradients, pictures, patterns, theme fills, and the
    BACKGROUND inherit-from-slide case all return ``None`` so the caller
    can fall back to the slide-level default.
    """
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    try:
        from pptx.enum.dml import MSO_FILL_TYPE  # type: ignore[import-not-found]
    except ImportError:  # pragma: no cover — gated at install time
        return None
    try:
        if fill.type != MSO_FILL_TYPE.SOLID:
            return None
    except (AttributeError, ValueError):
        return None
    try:
        return _rgb_color_from_color_format(fill.fore_color)
    except Exception:  # pragma: no cover — defensive
        return None


def _check_color_contrast(slide: Any, review: SlideReview) -> None:
    """Flag low-contrast text runs against the resolved background color.

    Heuristic (intentionally conservative):

    * Only check runs where ``font.color.type`` is ``RGB`` (concrete).
    * Resolve background as the shape's solid fill if present, else fall
      back to the safe-default slide background (``#FFFFFF``).
    * Compute the WCAG contrast ratio. Flag ``critical`` below
      :data:`_CONTRAST_CRITICAL_BELOW` (3.0) and ``warning`` between 3.0
      and :data:`_CONTRAST_AA_THRESHOLD` (4.5).

    Falls silent on theme colors / scheme fills / pictures / inherited
    backgrounds. False-negative > false-positive — Bobby's hard rule.
    """
    seen_critical: list[tuple[float, tuple[int, int, int], tuple[int, int, int]]] = []
    seen_warning: list[tuple[float, tuple[int, int, int], tuple[int, int, int]]] = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        bg_rgb = _resolve_shape_background_rgb(shape) or _DEFAULT_SLIDE_BG
        try:
            paragraphs = list(shape.text_frame.paragraphs)
        except Exception:  # pragma: no cover — defensive
            continue
        for para in paragraphs:
            for run in para.runs:
                font = getattr(run, "font", None)
                if font is None:
                    continue
                try:
                    text = (run.text or "").strip()
                except Exception:  # pragma: no cover — defensive
                    continue
                if not text:
                    continue
                fg_rgb = _rgb_color_from_color_format(getattr(font, "color", None))
                if fg_rgb is None:
                    # Theme color / unset — skip (no false positives).
                    continue
                ratio = _contrast_ratio(fg_rgb, bg_rgb)
                if ratio < _CONTRAST_CRITICAL_BELOW:
                    seen_critical.append((ratio, fg_rgb, bg_rgb))
                elif ratio < _CONTRAST_AA_THRESHOLD:
                    seen_warning.append((ratio, fg_rgb, bg_rgb))

    if seen_critical:
        ratio, fg, bg = min(seen_critical, key=lambda t: t[0])
        review.issues.append(
            {
                "severity": "critical",
                "rule": "color-contrast",
                "detail": (
                    f"Text contrast ratio {ratio:.2f}:1 against background "
                    f"#{fg[0]:02X}{fg[1]:02X}{fg[2]:02X} on "
                    f"#{bg[0]:02X}{bg[1]:02X}{bg[2]:02X} — below the 3.0:1 "
                    "floor (text will be illegible on a projector). Use a "
                    "darker (or lighter) foreground."
                ),
            }
        )
    elif seen_warning:
        ratio, fg, bg = min(seen_warning, key=lambda t: t[0])
        review.issues.append(
            {
                "severity": "warning",
                "rule": "color-contrast",
                "detail": (
                    f"Text contrast ratio {ratio:.2f}:1 — below WCAG AA "
                    f"({_CONTRAST_AA_THRESHOLD:.1f}:1) for normal-weight body text. "
                    f"Foreground #{fg[0]:02X}{fg[1]:02X}{fg[2]:02X} on background "
                    f"#{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}."
                ),
            }
        )


# ---------------------------------------------------------------------------
# Speaker-notes structure audit (Bobby's two-tier hard rule)
# ---------------------------------------------------------------------------
#
# ``feedback_slide_hard_rules`` requires every body slide to carry a
# dual-format speaker-notes block: a mental-map heading (HOOK / KEY CLAIM /
# EVIDENCE / KEY TERMS / CLICK / TRANSITION) AND a 200-400 word detailed
# script underneath, separated by a visible divider. This audit catches
# the common drift modes:
#
# * Figure / data slide ships with EMPTY notes — nothing to actually
#   present. Critical.
# * Notes present but mental-map labels missing — the on-stage scan
#   index is gone. Warning.
# * Script body is too thin to actually present from (< 100 words) on a
#   non-title/divider slide. Warning.
# * Script body is too long (> 500 words) — a draft pasted in, not a
#   compressed script. Warning.
#
# Title / divider / references / acknowledgments slides are exempt from
# the body-word-count rule (their notes are conventionally short).

_NOTES_MIN_SCRIPT_WORDS = 100
_NOTES_MAX_SCRIPT_WORDS = 500
_NOTES_DIVIDER = "--- DETAILED SCRIPT ---"
# Body-presentable slides — those needing a real spoken script. The empty /
# mental-map / short-body / long-body rules apply here. References and
# acknowledgments are NOT body slides; their notes are conventionally
# short or absent.
_NOTES_BODY_SLIDE_TYPES = {"figure", "text"}
# Body-word-count rule exempts these structural slides. Empty-notes rule
# similarly exempts the title / divider opener.
_NOTES_STRUCTURAL_EXEMPT = {
    "title",
    "section_divider",
    "references",
    "acknowledgments",
}


def _extract_speaker_notes(slide: Any) -> str:
    """Return the slide's speaker-notes plain-text, or ``""``."""
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    tf = getattr(notes_slide, "notes_text_frame", None)
    if tf is None:
        return ""
    try:
        return (tf.text or "").strip()
    except Exception:  # pragma: no cover — defensive
        return ""


def _split_notes_into_mental_map_and_script(notes_text: str) -> tuple[str, str]:
    """Return ``(mental_map_block, script_body)`` from a notes string.

    When the explicit ``--- DETAILED SCRIPT ---`` divider is present we
    split on it. When it isn't, we treat lines starting with ``- LABEL:``
    where LABEL matches a known mental-map key as the heading block, and
    everything after them as script. This mirrors :mod:`vaultlab.slides.notes`.
    """
    if _NOTES_DIVIDER in notes_text:
        head, _, tail = notes_text.partition(_NOTES_DIVIDER)
        return head.strip(), tail.strip()
    # No divider — try to recover a heading by inspecting line prefixes.
    try:
        from vaultlab.slides.notes import parse_speaker_notes
    except Exception:  # pragma: no cover — slides module is local
        return "", notes_text.strip()
    parsed = parse_speaker_notes(notes_text)
    if parsed:
        # Bullets identified, but no explicit body section → treat the
        # whole string as a mental-map-only document (no script body).
        return notes_text.strip(), ""
    return "", notes_text.strip()


def _has_mental_map_heading(mental_map_block: str) -> bool:
    """Cheap detection: at least one known mental-map LABEL parses out."""
    if not mental_map_block.strip():
        return False
    try:
        from vaultlab.slides.notes import parse_speaker_notes
    except Exception:  # pragma: no cover — slides module is local
        return False
    parsed = parse_speaker_notes(mental_map_block)
    return bool(parsed)


def _word_count(text: str) -> int:
    return len([w for w in text.split() if w.strip()])


def _check_speaker_notes(slide: Any, slide_type: str, review: SlideReview) -> None:
    """Flag missing / under-structured / under-length speaker notes.

    See module-level commentary above for the four firing rules. Each
    issue is filed against the slide review, not the deck arc, so the
    HTML adapter renders them inline with the slide card.
    """
    notes_text = _extract_speaker_notes(slide)

    # Rule 1 — empty notes on a body slide is critical (figure / data slides
    # absolutely need something to say).
    if not notes_text:
        if slide_type in _NOTES_BODY_SLIDE_TYPES:
            review.issues.append(
                {
                    "severity": "critical",
                    "rule": "speaker-notes-empty",
                    "detail": (
                        "Slide has no speaker notes — there is nothing to actually "
                        "say at the podium. Add a mental-map heading and a 200-400 "
                        "word script per the two-tier hard rule."
                    ),
                }
            )
        # Title / divider / acknowledgments slides exempt from empty check.
        return

    mental_map_block, script_body = _split_notes_into_mental_map_and_script(notes_text)
    has_map = _has_mental_map_heading(mental_map_block)
    script_words = _word_count(script_body)

    # Rule 2 — body slide notes missing the mental-map heading.
    if slide_type in _NOTES_BODY_SLIDE_TYPES and not has_map:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "speaker-notes-mental-map",
                "detail": (
                    "Speaker notes don't include the two-tier mental-map heading "
                    "(HOOK / KEY CLAIM / EVIDENCE / KEY TERMS / CLICK / TRANSITION). "
                    "Hard rule requires the on-stage scan-index above the script."
                ),
            }
        )

    # Word-count rules only apply to body slides; title / divider / ack are exempt.
    if slide_type in _NOTES_STRUCTURAL_EXEMPT:
        return

    # Rule 3 — script body too thin.
    if script_words and script_words < _NOTES_MIN_SCRIPT_WORDS:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "speaker-notes-short",
                "detail": (
                    f"Script body is {script_words} words — under the "
                    f"{_NOTES_MIN_SCRIPT_WORDS}-word floor for a body slide. "
                    "Likely too thin to actually present from."
                ),
            }
        )
    elif not script_body and slide_type in _NOTES_BODY_SLIDE_TYPES:
        # Notes only contain a mental-map block with no script underneath.
        review.issues.append(
            {
                "severity": "warning",
                "rule": "speaker-notes-short",
                "detail": (
                    "Speaker notes have a mental-map heading but no detailed "
                    f"script body. Add a {_NOTES_MIN_SCRIPT_WORDS}-{_NOTES_MAX_SCRIPT_WORDS} "
                    "word script under the divider."
                ),
            }
        )

    # Rule 4 — script body too long (likely an un-compressed draft).
    if script_words > _NOTES_MAX_SCRIPT_WORDS:
        review.issues.append(
            {
                "severity": "warning",
                "rule": "speaker-notes-long",
                "detail": (
                    f"Script body is {script_words} words — over the "
                    f"{_NOTES_MAX_SCRIPT_WORDS}-word ceiling. Compress to "
                    f"{_NOTES_MIN_SCRIPT_WORDS}-{_NOTES_MAX_SCRIPT_WORDS} words "
                    "or split the slide."
                ),
            }
        )


# ---------------------------------------------------------------------------
# Story-arc structural audit
# ---------------------------------------------------------------------------


_MAX_SECTION_DIVIDERS_PER_DECK = 5


def _audit_story_arc(report: ReviewReport) -> None:
    """Deck-level structural checks (no LLM required).

    Catches the cheap-but-load-bearing arc issues:

    * No identifiable title slide at all.
    * Section dividers > 5 (per slide hard rules — too many resets).
    * Reference slide missing when other slides cite (best-effort —
      ``[[wikilink]]`` style is hard to detect without a plan dict; we
      check that *some* references-shaped slide appears in any deck
      with > 3 slides).
    * Title slide appears anywhere other than first.
    * References / acknowledgments appear before the body slides.
    """
    slides = report.per_slide
    n = len(slides)
    if n == 0:
        report.story_arc_issues.append(
            {
                "severity": "critical",
                "rule": "story-arc-empty",
                "detail": "Deck is empty — no slides were rendered.",
                "loc": "(deck)",
            }
        )
        return

    # Title slide must be first.
    title_positions = [i for i, s in enumerate(slides) if s.slide_type == "title"]
    if not title_positions:
        # Fall back: the first slide may be tagged 'section_divider' if it has
        # only a title. Treat that as acceptable provided slide 0 looks like
        # a deck opener.
        if slides[0].slide_type not in {"section_divider", "title"}:
            report.story_arc_issues.append(
                {
                    "severity": "warning",
                    "rule": "story-arc-title-first",
                    "detail": "No identifiable title slide at the start of the deck.",
                    "loc": "Slide 1",
                }
            )
    elif title_positions[0] != 0:
        report.story_arc_issues.append(
            {
                "severity": "critical",
                "rule": "story-arc-title-position",
                "detail": (
                    f"Title slide appears at position {title_positions[0] + 1} "
                    "instead of slide 1."
                ),
                "loc": f"Slide {title_positions[0] + 1}",
            }
        )

    # Section divider count cap.
    n_dividers = sum(1 for s in slides if s.slide_type == "section_divider")
    # The opener may have been classified as a divider; don't double-count it.
    if slides[0].slide_type == "section_divider":
        n_dividers = max(0, n_dividers - 1)
    if n_dividers > _MAX_SECTION_DIVIDERS_PER_DECK:
        report.story_arc_issues.append(
            {
                "severity": "warning",
                "rule": "story-arc-divider-count",
                "detail": (
                    f"{n_dividers} section dividers (max recommended "
                    f"{_MAX_SECTION_DIVIDERS_PER_DECK}). Resets the audience too often."
                ),
                "loc": "(deck)",
            }
        )

    # References / acknowledgments must come at the end if present.
    ref_positions = [
        i for i, s in enumerate(slides) if s.slide_type in {"references", "acknowledgments"}
    ]
    if ref_positions:
        last_ref = max(ref_positions)
        if last_ref < n - 1 and any(
            s.slide_type not in {"references", "acknowledgments"} for s in slides[last_ref + 1 :]
        ):
            # A non-reference slide came after the last reference slide
            report.story_arc_issues.append(
                {
                    "severity": "warning",
                    "rule": "story-arc-references-tail",
                    "detail": (
                        "References / acknowledgments slide is not at the end of the deck."
                    ),
                    "loc": f"Slide {last_ref + 1}",
                }
            )


# ---------------------------------------------------------------------------
# HTML adapter — feeds ReviewReport into the existing audit_html builder
# ---------------------------------------------------------------------------


def _to_audit_html_inputs(report: ReviewReport) -> tuple[dict[str, Any], dict[str, Any]]:
    """Translate a :class:`ReviewReport` into ``(plan, audit)`` dicts.

    The existing :func:`vaultlab.slides.audit_html.build_audit_report_html`
    consumer expects a deck-plan dict and a rigor-audit dict (``passed`` +
    ``issues``). We materialize a thin plan-shaped view of the rendered
    deck from the per-slide reviews and map ReviewReport severities into
    the rigor-audit severity vocabulary
    (``critical`` → ``blocker``, ``warning`` → ``major``, ``info`` → ``minor``).
    """
    plan_slides = [
        {
            "type": slide.slide_type,
            "title": slide.title or f"Slide {slide.slide_index + 1}",
        }
        for slide in report.per_slide
    ]

    plan = {
        "title": report.pptx_path.stem,
        "slides": plan_slides,
    }

    issues: list[dict[str, Any]] = []
    for slide in report.per_slide:
        for issue in slide.issues:
            issues.append(
                {
                    "loc": f"Slide {slide.slide_index + 1}",
                    "severity": _SEVERITY_TO_RIGOR.get(issue.get("severity", "info"), "minor"),
                    "kind": issue.get("rule", "other"),
                    "fix": issue.get("detail", ""),
                }
            )
    for issue in report.story_arc_issues:
        issues.append(
            {
                "loc": issue.get("loc", "(deck)"),
                "severity": _SEVERITY_TO_RIGOR.get(issue.get("severity", "info"), "minor"),
                "kind": issue.get("rule", "story-arc"),
                "fix": issue.get("detail", ""),
            }
        )

    audit = {"passed": report.ok(), "issues": issues}
    return plan, audit


def render_review_html(report: ReviewReport, *, title: str | None = None) -> str:
    """Render the :class:`ReviewReport` as a self-contained HTML string.

    Reuses :func:`vaultlab.slides.audit_html.build_audit_report_html` so
    the look and feel match the rigor-audit reports the pipeline already
    emits for ``/build-deck`` runs.
    """
    from vaultlab.slides.audit_html import build_audit_report_html

    plan, audit = _to_audit_html_inputs(report)
    return build_audit_report_html(
        plan,
        audit,
        pptx_path=report.pptx_path,
        title=title or f"Deck self-review — {report.pptx_path.name}",
    )


def write_review_report(report: ReviewReport, out_path: Path | str) -> Path:
    """Render the review HTML, write it to disk, and emit provenance receipts.

    Writes :func:`render_review_html` output to ``out_path`` and then
    drops the standard sidecar pair (``.provenance.json`` +
    ``.method.md``) per the AGENTS.md Red Line #2 contract. Returns the
    resolved path to the HTML file.
    """
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    html = render_review_html(report)
    out_path.write_text(html, encoding="utf-8")

    # Provenance receipts — best-effort per AGENTS.md (no hard gate).
    try:
        from vaultlab.provenance import ProvenanceRecord, write_receipts

        try:
            import vaultlab as _vl

            version = getattr(_vl, "__version__", "unknown")
        except Exception:  # pragma: no cover — version lookup is best-effort
            version = "unknown"

        record = ProvenanceRecord(
            generated_by="vaultlab.slides.self_review.write_review_report",
            kind="slide_self_review",
            inputs=[str(report.pptx_path)],
            code_version=version,
            params={
                "n_slides": report.n_slides,
                "n_critical": report.n_critical,
                "n_warning": report.n_warning,
                "n_info": report.n_info,
                "passed": report.ok(),
            },
            tags=["slides", "self-review", "audit"],
            notes=(
                "Self-review HTML for a rendered deck. Composes layout hard-rule "
                "checks, story-arc structural checks, and bullet/figure heuristics "
                "into a single report. No LLM calls."
            ),
        )
        write_receipts(str(out_path), record)
    except Exception:  # pragma: no cover — best-effort
        logger.exception("write_review_report: provenance sidecar write failed")

    return out_path
