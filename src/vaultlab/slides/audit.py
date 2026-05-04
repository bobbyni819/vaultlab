"""Deck audit — self-evaluation step for slide-deck outputs.

Runs after a deck is composed to surface gaps before declaring "done":

* Slides with **0 images** when the layout intended a figure-or-bullets
  slot (the build_deck_from_lineage_result composer falls back to
  bullets when figure-fetch fails — the audit catches that silent
  degradation).
* Slides that are **text-only and very thin** (e.g. <50 chars across
  all text frames) — usually a layout failure, not a deliberate empty.
* Citation-vs-figure mismatch: papers cited in the arc that *would
  have* yielded a slide-feature figure but weren't fetched (so the
  user can manually fetch them).

The audit returns a structured ``DeckAuditResult`` and a
human-readable markdown report. Bobby's 2026-05-02 ask was: "have
you actually read everything here and evaluate it" — this is the
"have you actually" check.

Usage::

    from vaultlab.slides.audit import audit_deck

    result = audit_deck(
        deck_path,
        arc_path=arc_path,                # optional: cross-check arc citations
        figure_staging_dir=Path("/tmp/_deck_figures"),  # optional: figures available locally
    )
    print(result.to_markdown_report())
    if result.severity >= "warn":
        print(result.manual_fetch_shopping_list())
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# Import at module top so tests can patch it cleanly via
# vaultlab.slides.audit.Presentation. python-pptx is already a vaultlab
# dependency so the import never fails in normal usage.
try:
    from pptx import Presentation
except ImportError:  # pragma: no cover — defensive only
    Presentation = None  # type: ignore[assignment]

logger = logging.getLogger(__name__)


_THIN_SLIDE_TEXT_THRESHOLD = 50  # chars
_PICTURE_SHAPE_TYPE = 13  # python-pptx MSO_SHAPE_TYPE.PICTURE


@dataclass
class SlideAudit:
    """Per-slide audit record."""

    index: int
    title: str
    n_images: int
    n_text_chars: int
    is_thin_text: bool
    is_figure_intended: bool  # heuristic: title contains "figure", "image", section_intro layout, etc.
    figure_gap: bool  # True if figure_intended but n_images == 0
    is_section_divider: bool = False  # title-only chapter transition (no body); never a figure target
    overlapping_shapes: int = 0  # count of shapes whose bbox intersects another shape's bbox by >50%
    text_overflow_shapes: int = 0  # estimated text height > container height by >10%
    offslide_shapes: int = 0  # shapes whose bbox extends past slide canvas edges
    title_too_long: bool = False  # title >100 chars (will wrap awkwardly)
    over_bulleted: bool = False  # >7 bullets in any single text frame


@dataclass
class DeckAuditResult:
    """Audit findings for one deck."""

    deck_path: Path
    n_slides: int
    n_total_images: int
    slides_with_images: int
    text_only_slides: int
    thin_slides: int
    figure_gap_slides: int  # slides where a figure was intended but missing
    per_slide: list[SlideAudit] = field(default_factory=list)
    citations_in_arc: list[str] = field(default_factory=list)  # DOIs cited in the arc

    @property
    def total_overlapping_pairs(self) -> int:
        return sum(s.overlapping_shapes for s in self.per_slide)

    @property
    def total_overflowing_shapes(self) -> int:
        return sum(s.text_overflow_shapes for s in self.per_slide)

    @property
    def total_offslide_shapes(self) -> int:
        return sum(s.offslide_shapes for s in self.per_slide)

    @property
    def n_over_bulleted_slides(self) -> int:
        return sum(1 for s in self.per_slide if s.over_bulleted)

    @property
    def n_long_titles(self) -> int:
        return sum(1 for s in self.per_slide if s.title_too_long)

    @property
    def severity(self) -> str:
        """One of "ok", "warn", "fail"."""
        # Hard fails: shapes off the slide canvas or many overlap pairs
        if self.total_offslide_shapes >= 1:
            return "fail"
        if self.total_overlapping_pairs >= 3:
            return "fail"
        if self.total_overflowing_shapes >= 2:
            return "fail"
        if self.n_total_images == 0 and self.n_slides >= 5:
            return "fail"
        # Warns: any overflow or overlap, structural issues
        if self.figure_gap_slides > 0 or self.thin_slides > 1:
            return "warn"
        if self.total_overlapping_pairs > 0:
            return "warn"
        if self.total_overflowing_shapes > 0:
            return "warn"
        if self.n_over_bulleted_slides > 0:
            return "warn"
        if self.n_long_titles > 0:
            return "warn"
        if self.n_total_images < self.n_slides // 4:
            return "warn"
        return "ok"

    def to_markdown_report(self) -> str:
        """Render the audit as a human-readable markdown block."""
        lines: list[str] = []
        lines.append(f"### Audit — `{self.deck_path.name}`")
        lines.append(f"")
        sev_emoji = {"ok": "✅", "warn": "⚠️", "fail": "❌"}.get(self.severity, "?")
        lines.append(f"**Severity:** {sev_emoji} `{self.severity}`")
        lines.append(f"")
        lines.append(f"- {self.n_slides} slides, {self.n_total_images} total images")
        lines.append(f"- {self.slides_with_images} slides with at least one image")
        lines.append(f"- {self.text_only_slides} text-only slides")
        lines.append(f"- {self.thin_slides} thin slides (< {_THIN_SLIDE_TEXT_THRESHOLD} chars text)")
        lines.append(f"- {self.figure_gap_slides} slides where a figure was intended but missing")
        lines.append(f"")
        if self.figure_gap_slides:
            lines.append(f"**Figure gaps:**")
            for s in self.per_slide:
                if s.figure_gap:
                    lines.append(f"  * Slide {s.index}: {s.title or '(no title)'}")
            lines.append("")
        return "\n".join(lines)

    def manual_fetch_shopping_list(self, max_items: int = 10) -> str:
        """Render a 'please manually fetch figures from these papers' list.

        For now this is a stub — the citation-to-figure-source map needs
        the arc's DOI list to be passed in via ``citations_in_arc``. When
        present, this surfaces the top-N most-cited DOIs as "high-value
        figures the user could manually pull and drop into the deck."
        """
        if not self.citations_in_arc:
            return "(no arc citations supplied → cannot suggest manual fetches)"
        lines = ["**Recommended manual figure fetches**"]
        lines.append("")
        lines.append("These DOIs are cited in the arc and have ≥1 figure-gap slide. ")
        lines.append("Manually grab Figure 1 from each, drop into the deck:")
        lines.append("")
        for doi in self.citations_in_arc[:max_items]:
            lines.append(f"- `{doi}` — https://doi.org/{doi}")
        return "\n".join(lines)


def audit_deck(
    deck_path: Path | str,
    *,
    arc_path: Path | str | None = None,
    figure_intended_titles: tuple[str, ...] = (
        "figure", "image", "panel", "schematic", "diagram", "overview",
    ),
) -> DeckAuditResult:
    """Audit a .pptx deck for figure / content gaps.

    Args:
        deck_path: Path to the .pptx file.
        arc_path: Optional path to the source arc markdown — when given,
            the audit extracts cited DOIs and includes them in the
            manual-fetch shopping list.
        figure_intended_titles: Lowercased substrings; a slide whose
            title contains any of these is considered "figure-intended".
            Slides matching this pattern with 0 images are flagged as
            ``figure_gap``.

    Returns:
        :class:`DeckAuditResult`.
    """
    deck_path = Path(deck_path)
    prs = Presentation(str(deck_path))

    per_slide: list[SlideAudit] = []
    n_total_images = 0
    slides_with_images = 0
    text_only_slides = 0
    thin_slides = 0
    figure_gap_slides = 0

    for i, slide in enumerate(prs.slides, 1):
        n_images = sum(
            1 for sh in slide.shapes if sh.shape_type == _PICTURE_SHAPE_TYPE
        )
        title = ""
        all_text = ""
        text_shape_chars: list[int] = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt = sh.text_frame.text or ""
                all_text += txt + "\n"
                if not title and txt.strip():
                    title = txt.strip().split("\n")[0][:80]
                if txt.strip():
                    text_shape_chars.append(len(txt.strip()))
        n_text_chars = len(all_text.strip())
        # Section dividers are intentionally short — don't count as thin.
        # (We compute is_section_divider AFTER text_shape_chars below; this
        # initial is_thin gets overridden once is_section_divider is known.)
        is_thin = n_text_chars < _THIN_SLIDE_TEXT_THRESHOLD and n_images == 0

        # Detect section_divider: title-only slide with no body content.
        # Specifically: exactly one non-empty text shape, that shape's text
        # is short (<120 chars), AND no images. These are chapter transitions
        # ("BACKGROUND", "1. Origins", "Take-aways") and should NEVER receive
        # post-populated figures — figures belong on content slides.
        is_section_divider = (
            n_images == 0
            and len([c for c in text_shape_chars if c > 5]) == 1
            and text_shape_chars[0] < 120
        )
        # Section dividers are intentional, not "thin"
        if is_section_divider:
            is_thin = False

        title_lc = title.lower()
        is_figure_intended = any(
            kw in title_lc for kw in figure_intended_titles
        )
        if title_lc.startswith(("history", "development", "state of",
                                "background", "methods", "results")):
            is_figure_intended = True
        if re.match(r"^\s*\d+[.\-]?\s+\S", title) or re.match(
            r"^\s*\d+\s*[-]\s*\d+[.\-]?\s+\S", title
        ):
            is_figure_intended = True
        review_paper_kws = (
            "introduction", "foundation", "framework", "early method",
            "seminal", "refinement", "instrumentation",
            "application", "specialised", "specialized",
            "state of the art", "sota", "limitation", "future direction",
            "thesis", "discussion", "synthesis",
        )
        if any(k in title_lc for k in review_paper_kws):
            is_figure_intended = True

        # Section dividers are NEVER figure-gaps — they're intentional
        # transitions. The figure goes on the *content* slide that follows.
        figure_gap = (
            is_figure_intended and n_images == 0 and not is_section_divider
        )

        # References slides have many refs by design — not a quality issue
        is_references = title_lc.startswith("references") or title_lc.startswith("selected references")

        # Shape-overlap detection: count pairs of shapes whose bboxes
        # overlap by >50% of the smaller shape's area.
        n_overlap = _count_overlapping_shape_pairs(slide)

        # Text-overflow + off-slide checks. Section dividers wrap large
        # 48pt titles by design and look fine; references list many entries
        # by design — exempt both from these warnings.
        if is_section_divider or is_references:
            n_overflow = 0
        else:
            n_overflow = _count_text_overflow_shapes(slide)
        n_offslide = _count_offslide_shapes(slide, prs.slide_width, prs.slide_height)

        # Title length + bullet density. References slide intentionally has
        # many bullets — don't flag.
        title_long = len(title) > 100
        if is_references or is_section_divider:
            over_bul = False
        else:
            over_bul = _has_over_bulleted_textbox(slide, threshold=7)

        per_slide.append(SlideAudit(
            index=i,
            title=title,
            n_images=n_images,
            n_text_chars=n_text_chars,
            is_thin_text=is_thin,
            is_figure_intended=is_figure_intended,
            figure_gap=figure_gap,
            is_section_divider=is_section_divider,
            overlapping_shapes=n_overlap,
            text_overflow_shapes=n_overflow,
            offslide_shapes=n_offslide,
            title_too_long=title_long,
            over_bulleted=over_bul,
        ))

        n_total_images += n_images
        if n_images:
            slides_with_images += 1
        else:
            text_only_slides += 1
        if is_thin:
            thin_slides += 1
        if figure_gap:
            figure_gap_slides += 1

    citations: list[str] = []
    if arc_path:
        try:
            arc_text = Path(arc_path).read_text(encoding="utf-8", errors="replace")
            citations = _extract_dois_from_arc(arc_text)
        except OSError:
            pass

    return DeckAuditResult(
        deck_path=deck_path,
        n_slides=len(prs.slides),
        n_total_images=n_total_images,
        slides_with_images=slides_with_images,
        text_only_slides=text_only_slides,
        thin_slides=thin_slides,
        figure_gap_slides=figure_gap_slides,
        per_slide=per_slide,
        citations_in_arc=citations,
    )


def _count_overlapping_shape_pairs(slide: Any) -> int:
    """Count shape pairs whose bounding boxes overlap >50% of the smaller area.

    Uses python-pptx position attributes (left/top/width/height in EMU).
    Catches the "all body shapes stacked at the same coords" failure mode
    of post-populating figures into existing slides.
    """
    boxes = []
    for sh in slide.shapes:
        l = getattr(sh, "left", None)
        t = getattr(sh, "top", None)
        w = getattr(sh, "width", None)
        h = getattr(sh, "height", None)
        # Real EMU values from python-pptx are concrete ints >0. Reject
        # anything else (None, MagicMock from tests, unset placeholders)
        # to keep the audit robust under mocking.
        if not all(isinstance(x, int) and x is not True and x is not False
                   for x in (l, t, w, h)):
            continue
        if w <= 0 or h <= 0:
            continue
        boxes.append((l, t, l + w, t + h, w * h))

    overlap_pairs = 0
    for i in range(len(boxes)):
        l1, t1, r1, b1, a1 = boxes[i]
        for j in range(i + 1, len(boxes)):
            l2, t2, r2, b2, a2 = boxes[j]
            ix = max(0, min(r1, r2) - max(l1, l2))
            iy = max(0, min(b1, b2) - max(t1, t2))
            inter = ix * iy
            if inter <= 0:
                continue
            smaller = min(a1, a2)
            if smaller > 0 and inter / smaller > 0.5:
                overlap_pairs += 1
    return overlap_pairs


def _count_text_overflow_shapes(slide: Any) -> int:
    """Estimate how many text frames have content that exceeds container height.

    Uses a heuristic: avg char width ≈ 0.55 × font_size_pt; line height ≈
    1.4 × font_size_pt. Counts wrapped lines per paragraph from chars-per-line
    and sums × line-height. Flags when estimated > container × 1.1.

    Catches the "bullets clip off the bottom" failure mode users see when
    24pt body content runs longer than its 5-inch box.
    """
    n = 0
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            text = sh.text_frame.text or ""
        except Exception:  # noqa: BLE001
            continue
        if not text.strip():
            continue
        w = getattr(sh, "width", None)
        h = getattr(sh, "height", None)
        if not (isinstance(w, int) and isinstance(h, int)) or w <= 0 or h <= 0:
            continue
        width_in = w / 914400
        height_in = h / 914400

        # Determine dominant font size (Pt) — read first run's size, default 18
        font_size_pt = 18
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_size_pt = run.font.size.pt
                    break
            if font_size_pt != 18:
                break

        avg_char_w_in = (font_size_pt * 0.55) / 72
        chars_per_line = max(1, int(width_in / avg_char_w_in))
        line_height_in = (font_size_pt * 1.4) / 72

        total_lines = 0
        for para in sh.text_frame.paragraphs:
            ptext = para.text or ""
            if not ptext:
                total_lines += 1
                continue
            wrap = max(1, (len(ptext) + chars_per_line - 1) // chars_per_line)
            total_lines += wrap

        estimated_h_in = total_lines * line_height_in
        if estimated_h_in > height_in * 1.1:
            n += 1
    return n


def _count_offslide_shapes(slide: Any, slide_w: int, slide_h: int) -> int:
    """Count shapes whose bbox extends past the slide canvas edges."""
    n = 0
    for sh in slide.shapes:
        l = getattr(sh, "left", None)
        t = getattr(sh, "top", None)
        w = getattr(sh, "width", None)
        h = getattr(sh, "height", None)
        if not all(isinstance(x, int) for x in (l, t, w, h)):
            continue
        # Allow tiny float-error tolerance: 0.05" = ~46k EMU
        tol = 46_000
        if (l + w) > slide_w + tol or (t + h) > slide_h + tol or l < -tol or t < -tol:
            n += 1
    return n


def _has_over_bulleted_textbox(slide: Any, threshold: int = 7) -> bool:
    """True if any text frame contains >threshold non-empty paragraphs."""
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            paragraphs = sh.text_frame.paragraphs
        except Exception:  # noqa: BLE001
            continue
        non_empty = sum(1 for p in paragraphs if (p.text or "").strip())
        if non_empty > threshold:
            return True
    return False


def _extract_dois_from_arc(text: str) -> list[str]:
    """Pull fetchable paper identifiers out of arc markdown.

    Recognises:
    1. Inline DOI: ``10.1038/s41586-...``
    2. Wikilink with slash: ``[[10.1038/s41586-...]]``
    3. Wikilink with underscore-slash: ``[[10.1038_s41586-...]]``
       — sluggified DOIs that use ``_`` instead of ``/`` between the
       registrant and the suffix.
    4. Wikilink with author-year suffix: ``[[10.1038_s41586-x|Author Year]]``
       — strip the ``|...`` part.
    5. Paperclip native IDs: ``[[PMC<digits>]]``, ``[[PMC<digits>|Author Year]]``,
       ``[[arx_...]]``, ``[[bio_...]]``, ``[[med_...]]`` — these are
       directly fetchable via paperclip lookup without DOI resolution.

    Returns a deduplicated list in arc-order. PMC/arx/bio/med IDs are
    returned as-is (fetcher must distinguish via prefix).
    """
    seen: set[str] = set()
    out: list[str] = []

    # Pattern 1+2: standard 10.<digits>/<rest>
    doi_re = re.compile(r"10\.\d{3,9}/[\w.\-/_:;()]+", re.I)
    for m in doi_re.finditer(text):
        doi = m.group(0).rstrip(".,;:)")
        if doi.lower() not in seen:
            seen.add(doi.lower())
            out.append(doi)

    # Pattern 3+4: wikilink-slug with underscore between registrant and
    # suffix, possibly followed by ``|<author-year>``.
    # Format: [[10.1038_s41586-...|Author Year]] → 10.1038/s41586-...
    wiki_doi_re = re.compile(
        r"\[\[\s*(10\.\d{3,9}_[^\]|]+?)\s*(?:\|[^\]]+)?\s*\]\]", re.I,
    )
    for m in wiki_doi_re.finditer(text):
        slug = m.group(1)
        doi = slug.replace("_", "/", 1).rstrip(".,;:)")
        if doi.lower() not in seen:
            seen.add(doi.lower())
            out.append(doi)

    # Pattern 5: paperclip native IDs in wikilinks
    pcl_re = re.compile(
        r"\[\[\s*((?:PMC\d+|arx_[\w.]+|bio_\w+|med_\w+))\s*"
        r"(?:\|[^\]]+)?\s*\]\]", re.I,
    )
    for m in pcl_re.finditer(text):
        pid = m.group(1)
        # Normalize: paperclip IDs are case-insensitive but PMC is uppercase
        if pid.upper().startswith("PMC"):
            pid = pid.upper()
        if pid not in seen:
            seen.add(pid)
            out.append(pid)

    return out


__all__ = [
    "DeckAuditResult",
    "SlideAudit",
    "audit_deck",
]
