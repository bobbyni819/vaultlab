"""Hickey Lab theme for vaultlab.slides.

Ported from ``bobby_slides._template`` (Bobby's existing personal-toolkit
implementation). Bundles the Hickey Lab .pptx template with vaultlab so any
user with access to the lab template gets the right colors / fonts / layouts.

Theme palette (extracted from ``2023_09_19_Template.pptx``):

- accent1: ``#29AF8C`` (teal-green primary)
- accent2: ``#97BE49`` (lime)
- accent3: ``#3D9CCC`` (blue)
- accent4: ``#7C60C6`` (purple)
- accent5: ``#C9492C`` (red-orange)
- accent6: ``#D58C2E`` (gold)
- dk2: ``#44546A`` (dark slate)
- lt2: ``#E7E6E6`` (light gray)

Default font: Roboto. Min sizes: heading 28pt, body 24pt, caption 18pt.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

# Ported palette - extracted from the Hickey Lab .pptx template
_LAB_THEME_COLORS_HEX: dict[str, str] = {
    "accent1": "29AF8C",
    "accent2": "97BE49",
    "accent3": "3D9CCC",
    "accent4": "7C60C6",
    "accent5": "C9492C",
    "accent6": "D58C2E",
    "dk2": "44546A",
    "lt2": "E7E6E6",
}

_BUNDLED_TEMPLATE_PATH = Path(__file__).parent / "_assets" / "hickey_lab_template.pptx"

_DEFAULT_FONT = "Roboto"

# Min sizes (projector-readable per Bobby 2026-04-29 grill rules)
_MIN_SIZES: dict[str, int] = {
    "heading": 28,
    "body": 24,
    "caption": 18,
}

# The lab template has 3 slide masters with different color maps:
#   Master 0 - dark background
#   Master 1 - light background
#   Master 2 - dark section-divider variant
_THEME_MASTER_INDEX: dict[str, int] = {
    "dark": 0,
    "light": 1,
    "default": 0,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def hickey_lab_template_path() -> Path | None:
    """Return the bundled Hickey Lab template path, or ``None`` if missing."""
    return _BUNDLED_TEMPLATE_PATH if _BUNDLED_TEMPLATE_PATH.exists() else None


def load_hickey_lab_presentation(theme: str = "dark", strip_starter_slides: bool = True) -> Any:
    """Create a python-pptx Presentation initialized from the Hickey Lab template.

    Parameters
    ----------
    theme
        ``"dark"`` (default) or ``"light"``. Selects which slide master.
    strip_starter_slides
        If True (default), removes the example slides in the template
        (Instructions, Citation, etc.) but preserves layouts + masters.
    """
    from pptx import Presentation

    path = hickey_lab_template_path()
    if path is None:
        raise FileNotFoundError(
            f"Hickey Lab template not bundled. Expected at {_BUNDLED_TEMPLATE_PATH}"
        )

    pres = Presentation(str(path))

    if strip_starter_slides:
        _strip_existing_slides(pres)

    pres._vaultlab_master_index = _THEME_MASTER_INDEX.get(theme.lower(), 0)
    pres._vaultlab_theme_name = "hickey-lab"
    pres._vaultlab_theme_variant = theme.lower()
    return pres


def hickey_lab_colors() -> dict[str, Any]:
    """Return the Hickey Lab palette as ``RGBColor`` instances."""
    from pptx.dml.color import RGBColor

    return {name: RGBColor.from_string(hex_val) for name, hex_val in _LAB_THEME_COLORS_HEX.items()}


def hickey_lab_colors_hex() -> dict[str, str]:
    """Return the Hickey Lab palette as hex strings."""
    return dict(_LAB_THEME_COLORS_HEX)


def hickey_lab_font() -> str:
    """The Hickey Lab default font (Roboto)."""
    return _DEFAULT_FONT


def hickey_lab_min_sizes() -> dict[str, int]:
    """Min font sizes for headings (28), body (24), captions (18)."""
    return dict(_MIN_SIZES)


# ---------------------------------------------------------------------------
# Internal
# ---------------------------------------------------------------------------


def _strip_existing_slides(pres: Any) -> None:
    """Remove all slides; preserve layouts + masters."""
    rels_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    xml_slides = pres.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for slide_id_elem in slides_to_remove:
        rel_id = slide_id_elem.attrib[f"{{{rels_ns}}}id"]
        pres.part.drop_rel(rel_id)
        xml_slides.remove(slide_id_elem)


__all__ = [
    "hickey_lab_colors",
    "hickey_lab_colors_hex",
    "hickey_lab_font",
    "hickey_lab_min_sizes",
    "hickey_lab_template_path",
    "load_hickey_lab_presentation",
]
