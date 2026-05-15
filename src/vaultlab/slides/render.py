"""PPTX renderer — turns a :class:`Deck` into a ``.pptx`` via python-pptx.

This is the **only** module in vaultlab.slides that depends on python-pptx.
Everything else operates on the dataclasses in ``deck.py`` and ``layouts/``.
"""

from __future__ import annotations

from pathlib import Path

from vaultlab.provenance import ProvenanceRecord, write_receipts
from vaultlab.slides.deck import Deck, Slide
from vaultlab.slides.layouts import LayoutSpec, get_layout
from vaultlab.slides.themes import Theme, get_theme


class RenderError(Exception):
    """Raised when rendering fails (missing dependency, missing figure, etc.)."""


def render_pptx(deck: Deck, output_path: str | Path) -> Path:
    """Render a :class:`Deck` to a ``.pptx`` file.

    Parameters
    ----------
    deck
        The deck to render.
    output_path
        Output ``.pptx`` path. Parent directories created as needed.

    Returns
    -------
    Path
        Resolved path to the written ``.pptx``.

    Raises
    ------
    RenderError
        If python-pptx is not installed, or a referenced figure cannot be found.
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
        from pptx.util import Inches, Pt  # type: ignore[import-not-found]
    except ImportError as e:
        raise RenderError(
            "python-pptx is required for slides rendering. Install with "
            '`pip install -e ".[slides]"` or `pip install python-pptx`.'
        ) from e

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    theme = get_theme(deck.theme)
    pres = Presentation()
    # Standard 16:9 widescreen
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    blank_layout = pres.slide_layouts[6]  # python-pptx index 6 = Blank

    for slide in deck.slides:
        layout = get_layout(slide.layout)
        ppt_slide = pres.slides.add_slide(blank_layout)
        _render_slide(
            ppt_slide,
            slide=slide,
            layout=layout,
            theme=theme,
            slide_width=pres.slide_width,
            slide_height=pres.slide_height,
            working_dir=deck.working_dir,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
        )

    # Stash provenance into the document properties so reviewers can trace
    # decks back to vaultlab runs
    if deck.metadata:
        try:
            pres.core_properties.comments = "\n".join(f"{k}: {v}" for k, v in deck.metadata.items())
        except (AttributeError, ValueError):  # python-pptx version skew — non-fatal
            pass

    pres.save(out)

    # Audit-manifest contract (red line #2: no silent failures).
    # Every artifact-producing entrypoint writes provenance receipts;
    # see vaultlab/.claude/goals/vaultlab-north-star.md.
    record = ProvenanceRecord(
        generated_by="vaultlab.slides.render.render_pptx",
        kind="slide_deck",
        inputs=[],
        params={
            "n_slides": len(deck.slides),
            "theme": deck.theme,
            **{k: v for k, v in (deck.metadata or {}).items() if isinstance(v, (str, int, float, bool))},
        },
    )
    write_receipts(str(out), record)
    return out


# ---------------------------------------------------------------------------
# Internal: render one slide
# ---------------------------------------------------------------------------


def _render_slide(
    ppt_slide,
    *,
    slide: Slide,
    layout: LayoutSpec,
    theme: Theme,
    slide_width,
    slide_height,
    working_dir: Path | None,
    Inches,
    Pt,
    RGBColor,
) -> None:
    # Title
    if slide.title:
        _add_text(
            ppt_slide,
            text=slide.title,
            box=layout.title_box,
            slide_width=slide_width,
            slide_height=slide_height,
            font_name=theme.title_font,
            font_size_pt=theme.title_size_pt,
            color_rgb=theme.title_color_rgb,
            bold=True,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
        )

    # Body — interpretation depends on layout
    if layout.name == "title":
        if slide.subtitle and layout.body_box is not None:
            _add_text(
                ppt_slide,
                text=slide.subtitle,
                box=layout.body_box,
                slide_width=slide_width,
                slide_height=slide_height,
                font_name=theme.body_font,
                font_size_pt=theme.body_size_pt,
                color_rgb=theme.body_color_rgb,
                bold=False,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
            )
    elif layout.name == "content_with_bullets":
        if slide.bullets and layout.body_box is not None:
            _add_bullets(
                ppt_slide,
                bullets=slide.bullets,
                box=layout.body_box,
                slide_width=slide_width,
                slide_height=slide_height,
                font_name=theme.body_font,
                font_size_pt=theme.bullet_size_pt,
                color_rgb=theme.body_color_rgb,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
            )
    elif layout.name == "figure_with_caption":
        if slide.figure_path and layout.figure_box is not None:
            _add_figure(
                ppt_slide,
                figure_path=slide.figure_path,
                box=layout.figure_box,
                slide_width=slide_width,
                slide_height=slide_height,
                working_dir=working_dir,
                Inches=Inches,
            )
        if slide.caption and layout.body_box is not None:
            _add_text(
                ppt_slide,
                text=slide.caption,
                box=layout.body_box,
                slide_width=slide_width,
                slide_height=slide_height,
                font_name=theme.body_font,
                font_size_pt=theme.caption_size_pt,
                color_rgb=theme.body_color_rgb,
                bold=False,
                Inches=Inches,
                Pt=Pt,
                RGBColor=RGBColor,
            )

    # Speaker notes
    if slide.speaker_notes:
        notes = ppt_slide.notes_slide.notes_text_frame
        notes.text = slide.speaker_notes


# ---------------------------------------------------------------------------
# Drawing primitives — kept tiny + parameterized so the renderer is testable
# ---------------------------------------------------------------------------


def _add_text(
    ppt_slide,
    *,
    text,
    box,
    slide_width,
    slide_height,
    font_name,
    font_size_pt,
    color_rgb,
    bold,
    Inches,
    Pt,
    RGBColor,
) -> None:
    left = int(slide_width * box.x)
    top = int(slide_height * box.y)
    width = int(slide_width * box.width)
    height = int(slide_height * box.height)
    txbox = ppt_slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)


def _add_bullets(
    ppt_slide,
    *,
    bullets,
    box,
    slide_width,
    slide_height,
    font_name,
    font_size_pt,
    color_rgb,
    Inches,
    Pt,
    RGBColor,
) -> None:
    left = int(slide_width * box.x)
    top = int(slide_height * box.y)
    width = int(slide_width * box.width)
    height = int(slide_height * box.height)
    txbox = ppt_slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(*color_rgb)


def _add_figure(
    ppt_slide,
    *,
    figure_path,
    box,
    slide_width,
    slide_height,
    working_dir,
    Inches,
) -> None:
    fig = Path(figure_path)
    if not fig.is_absolute() and working_dir is not None:
        fig = working_dir / fig
    if not fig.exists():
        raise RenderError(f"Figure not found: {fig}")
    left = int(slide_width * box.x)
    top = int(slide_height * box.y)
    width = int(slide_width * box.width)
    height = int(slide_height * box.height)
    ppt_slide.shapes.add_picture(str(fig), left, top, width=width, height=height)


__all__ = ["RenderError", "render_pptx"]
