"""Populate an existing deck with figures, producing auto + aspirational variants.

Bobby's 2026-05-02 ask: "create two versions of the decks — 1 deck with
all figures we can actually extract automatically, 2nd deck with what we
hope for — a mix of figures we can already extract + placeholders for
figures the user needs to fetch manually."

Two modes:

* ``auto``         — insert only figures we can actually fetch
                     (paperclip JPGs OR PyMuPDF on cached PDFs). Slides
                     where no figure is obtainable keep their bullet
                     fallback. The deck is "honest about what we have."
* ``aspirational`` — same as auto, PLUS for slides where the cited
                     paper is high-impact but figure-fetch failed, drop a
                     "FIGURE NEEDED" call-out text box on the slide
                     pointing to the publisher URL. The deck is "what
                     we'd want — fetch the missing figures manually to
                     make this slide complete."

Both variants are saved as separate .pptx files (``-auto.pptx`` and
``-aspirational.pptx``) so the user can pick which to share. Audit
runs on both before declaring done.
"""

from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


_PICTURE_SHAPE_TYPE = 13


@dataclass
class FigureSlot:
    """One figure-gap slide that needs filling."""

    slide_index: int
    slide_title: str
    candidate_dois: list[str] = field(default_factory=list)


@dataclass
class FigurePopulateResult:
    """Outcome of populating a deck with figures."""

    deck_path: Path
    mode: str  # "auto" | "aspirational"
    n_inserted: int = 0
    n_placeholders: int = 0
    inserted_dois: list[str] = field(default_factory=list)
    placeholder_dois: list[str] = field(default_factory=list)


def _paperclip_env() -> dict[str, str]:
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["MSYS_NO_PATHCONV"] = "1"
    return env


def _list_paperclip_figures(
    paper_id: str, *, paperclip_binary: str, timeout: int = 30,
) -> list[str]:
    """Return the actual filenames in paperclip's figures dir for a paper.

    Paperclip's `ls /papers/<id>/figures/` prints all files on a single
    line separated by 2+ spaces, then a few metadata lines. Filters out
    metadata lines and returns the file list.

    Returns empty list on any failure.
    """
    try:
        r = subprocess.run(
            [paperclip_binary, "ls", f"/papers/{paper_id}/figures/"],
            capture_output=True, timeout=timeout, env=_paperclip_env(),
        )
    except (subprocess.TimeoutExpired, OSError):
        return []
    if r.returncode != 0:
        return []

    raw = r.stdout
    text = None
    for enc in ("utf-8", "cp1252", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        return []

    files: list[str] = []
    # Strict filename pattern: word-chars / dots / dashes, ending in a 2-5 char extension.
    # Rejects "file_b.jpg · trailing-junk" or any token with embedded whitespace
    # / cp1252 separators that survived decoding.
    name_re = re.compile(r"^[\w.\-]+\.[a-zA-Z]{2,5}$")
    seen: set[str] = set()
    for line in text.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith(("(", "[", "💡", "ERR:", "INFO:")):
            continue
        if "—" in s or "->" in s or "read-only" in s.lower():
            continue
        # Split on any whitespace and validate each candidate strictly
        for c in s.split():
            c = c.strip()
            if name_re.match(c) and c not in seen:
                seen.add(c)
                files.append(c)
    return files


_FIG_PRIORITY_PATTERNS = (
    # Highest: main first-figure indicators across publishers
    re.compile(r"(?i)(?:^|[._\-])fig(?:ure)?[._\-]?0*1[._\-]?(?:html|main)?\."),
    re.compile(r"(?i)(?:^|[._\-])g0*0?1[._\-]?(?:html)?\."),  # PMC publisher slug
    re.compile(r"(?i)_fig0*1\."),                              # bioRxiv style
    # Mid: any figure_N or fig_N (small N preferred via sort below)
    re.compile(r"(?i)(?:^|[._\-])fig(?:ure)?[._\-]?\d+\."),
    re.compile(r"(?i)(?:^|[._\-])g\d+\."),
    re.compile(r"(?i)_fig\d+\."),
)
_PREFERRED_EXTS = (".jpg", ".jpeg", ".png")
_CONVERTIBLE_EXTS = (".tif", ".tiff", ".webp", ".bmp", ".gif")


def _is_main_figure_candidate(fname: str) -> bool:
    """Reject equation glyphs, table images, supplementary, etc."""
    low = fname.lower()
    bad = (
        "equ", "ieq", "scheme", "logo", "icon", "thumb",
        "supp", "_si_", "_sup_", "tbl", "table", "graphabs",
    )
    return not any(b in low for b in bad)


def _pick_main_figure(filenames: list[str]) -> str | None:
    """Pick the best 'main figure 1' candidate from a paperclip ls listing.

    Strategy:
        1. Filter out non-image files and known-bad slugs (equations, icons).
        2. Walk priority regex patterns; for each pattern find all matches.
        3. Within a pattern hit, prefer JPG/PNG over TIFF, prefer larger
           publisher slugs (`_HTML.jpg`) over generic `.gif`.
        4. Fall back to the first usable image if no pattern hits.
    """
    images = [
        f for f in filenames
        if f.lower().endswith(_PREFERRED_EXTS + _CONVERTIBLE_EXTS)
        and _is_main_figure_candidate(f)
    ]
    if not images:
        return None

    def _ext_rank(f: str) -> int:
        low = f.lower()
        for i, ext in enumerate(_PREFERRED_EXTS):
            if low.endswith(ext):
                return i
        for i, ext in enumerate(_CONVERTIBLE_EXTS):
            if low.endswith(ext):
                return len(_PREFERRED_EXTS) + i
        return 99

    for pat in _FIG_PRIORITY_PATTERNS:
        hits = [f for f in images if pat.search(f)]
        if not hits:
            continue
        # Sort: best ext first, then prefer "_HTML" variants over plain
        hits.sort(key=lambda f: (_ext_rank(f), 0 if "_html" in f.lower() else 1, f))
        return hits[0]

    images.sort(key=lambda f: (_ext_rank(f), f))
    return images[0]


def _normalize_image_to_jpg(raw: bytes, src_ext: str, out_path: Path) -> Path | None:
    """Save raw bytes to ``out_path`` as JPG. Convert if not already JPG/PNG.

    python-pptx accepts JPG and PNG; TIFF / WEBP / BMP / GIF must be
    converted first. Returns the path the bytes ended up at, or None on
    failure (silently swallows Pillow errors).
    """
    src_ext = src_ext.lower().lstrip(".")
    if src_ext in ("jpg", "jpeg", "png"):
        out_path.write_bytes(raw)
        return out_path
    try:
        from io import BytesIO
        from PIL import Image
        img = Image.open(BytesIO(raw))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        # Force a .jpg suffix on the output
        if out_path.suffix.lower() not in (".jpg", ".jpeg"):
            out_path = out_path.with_suffix(".jpg")
        img.save(out_path, format="JPEG", quality=88)
        return out_path
    except Exception as exc:  # noqa: BLE001
        logger.warning("Pillow conversion failed for %s: %s", src_ext, exc)
        return None


def fetch_figure_from_paperclip(
    paper_id: str,
    *,
    cache_dir: Path,
    paperclip_binary: str | None = None,
    min_bytes: int = 8_000,
) -> Path | None:
    """Download a 'main figure' from paperclip's virtual filesystem.

    First lists the paper's ``figures/`` directory to discover the
    publisher-specific filename (PMC slugs like ``MOL2-19-3465-g001.jpg``,
    bioRxiv versioned names like ``690313v1_fig1.tif``, arXiv canonical
    ``figure_1.jpg``), picks the best 'figure 1' candidate, fetches it,
    and converts TIFF / WEBP / BMP / GIF to JPG via Pillow so python-pptx
    can ingest it.

    Returns local cache path on success, ``None`` on failure.
    """
    if not paperclip_binary:
        paperclip_binary = shutil.which("paperclip")
    if not paperclip_binary:
        return None
    cache_dir = Path(cache_dir)
    cache_dir.mkdir(parents=True, exist_ok=True)
    out_path = cache_dir / f"{paper_id}_figure_1.jpg"
    if out_path.exists() and out_path.stat().st_size >= min_bytes:
        return out_path

    files = _list_paperclip_figures(paper_id, paperclip_binary=paperclip_binary)
    if not files:
        return None

    pick = _pick_main_figure(files)
    if not pick:
        return None

    try:
        r = subprocess.run(
            [paperclip_binary, "cat", f"/papers/{paper_id}/figures/{pick}"],
            capture_output=True, timeout=120, env=_paperclip_env(),
        )
    except (subprocess.TimeoutExpired, OSError):
        return None
    if r.returncode != 0 or len(r.stdout) < min_bytes:
        return None

    src_ext = "." + pick.rsplit(".", 1)[-1]
    saved = _normalize_image_to_jpg(r.stdout, src_ext, out_path)
    if saved and saved.exists() and saved.stat().st_size >= min_bytes:
        return saved
    return None


def fetch_figure_from_local_pdf(
    pdf_path: Path,
    *,
    cache_dir: Path,
) -> Path | None:
    """Extract figure_1 from a locally-cached PDF via PyMuPDF.

    Returns local path to the largest extracted figure, None on failure.
    """
    if not pdf_path.exists():
        return None
    cache_dir = Path(cache_dir)
    cache_dir.mkdir(parents=True, exist_ok=True)

    try:
        from vaultlab.research.figures import extract_figures
    except ImportError:
        return None

    try:
        records = extract_figures(
            pdf_path=pdf_path,
            output_dir=cache_dir,
            min_dimension=300,
            min_bytes=10_000,
            write_metadata=False,
        )
    except Exception as exc:  # noqa: BLE001
        logger.warning("PyMuPDF extract failed for %s: %s", pdf_path, exc)
        return None

    if not records:
        return None
    # Sort by min_dimension desc — pick the largest figure (most likely to be a real figure, not an icon)
    records.sort(key=lambda r: r.get("min_dimension", 0), reverse=True)
    return Path(records[0]["path"])


def _doi_slug(doi: str) -> str:
    return doi.lower().replace("/", "_")


def _try_resolve_paperclip_id_for_doi(doi: str, *, paperclip_binary: str) -> str | None:
    """Map a DOI to a paperclip paper-id via lookup.

    Returns the paperclip ID (e.g. ``arx_2107.07953``, ``PMC...``) or None.
    """
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["MSYS_NO_PATHCONV"] = "1"
    try:
        r = subprocess.run(
            [paperclip_binary, "lookup", "doi", doi.strip().lower()],
            capture_output=True,
            timeout=30,
            env=env,
        )
    except (subprocess.TimeoutExpired, OSError):
        return None
    if r.returncode != 0:
        return None
    out = r.stdout.decode("utf-8", errors="replace")
    # Output format includes "<id> · <source> · <date>" line
    m = re.search(r"\b(arx_\S+|PMC\d+|bio_\w+|med_\w+)\b", out)
    return m.group(1) if m else None


def populate_deck_with_figures(
    deck_path: Path | str,
    *,
    candidate_dois: list[str],
    pdf_cache_dir: Path | str,
    figure_staging_dir: Path | str,
    mode: str = "auto",
    out_path: Path | str | None = None,
) -> FigurePopulateResult:
    """Insert figures into a deck's figure-gap slides.

    Args:
        deck_path: Existing .pptx file to populate.
        candidate_dois: Ordered list of DOIs cited in the arc, most-
            important first. We try to resolve each to a paperclip
            paper-id and fetch its figure_1.jpg.
        pdf_cache_dir: Where local PDFs live for the PyMuPDF fallback.
            Used when paperclip doesn't have the paper.
        figure_staging_dir: Where downloaded figures are cached so
            multiple deck variants can share them.
        mode: "auto" (insert only fetchable) or "aspirational" (also
            add FIGURE-NEEDED placeholders for unfetchable papers).
        out_path: Where to save the populated deck. Defaults to
            ``<deck_path>-<mode>.pptx`` next to the input.

    Returns:
        :class:`FigurePopulateResult`.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    deck_path = Path(deck_path)
    pdf_cache_dir = Path(pdf_cache_dir)
    figure_staging_dir = Path(figure_staging_dir)
    figure_staging_dir.mkdir(parents=True, exist_ok=True)

    if out_path is None:
        out_path = deck_path.with_name(deck_path.stem + f"-{mode}.pptx")
    else:
        out_path = Path(out_path)

    paperclip_binary = shutil.which("paperclip")

    # Identify figure-gap slides — figure-intended content slides without
    # an existing image. Section dividers are excluded by the audit (they
    # are intentional chapter transitions, not figure targets).
    from vaultlab.slides.audit import audit_deck
    audit = audit_deck(deck_path)
    gap_slides = [
        s for s in audit.per_slide
        if s.figure_gap and not s.is_section_divider
    ]
    if not gap_slides:
        logger.info("No figure-gap slides in %s; nothing to populate", deck_path)
        # Still copy to the output path for consistency
        shutil.copy(str(deck_path), str(out_path))
        return FigurePopulateResult(deck_path=out_path, mode=mode)

    # Resolve DOIs OR paperclip native IDs → figure paths. Cache so
    # multiple deck variants share the staging dir.
    resolved: dict[str, Path | None] = {}
    for ident in candidate_dois:
        if ident in resolved:
            continue
        fig_path: Path | None = None

        # Native paperclip ID (PMC<digits>, arx_, bio_, med_)? Fetch
        # directly without DOI lookup.
        is_pcl_native = bool(re.match(
            r"^(PMC\d+|arx_|bio_|med_)", ident, re.I,
        ))
        if is_pcl_native and paperclip_binary:
            fig_path = fetch_figure_from_paperclip(
                ident, cache_dir=figure_staging_dir,
                paperclip_binary=paperclip_binary,
            )
        # Otherwise treat as DOI — try paperclip lookup → figure_1
        elif paperclip_binary:
            pid = _try_resolve_paperclip_id_for_doi(
                ident, paperclip_binary=paperclip_binary,
            )
            if pid:
                fig_path = fetch_figure_from_paperclip(
                    pid, cache_dir=figure_staging_dir,
                    paperclip_binary=paperclip_binary,
                )

        # PyMuPDF fallback on cached PDF (only for DOI inputs)
        if fig_path is None and not is_pcl_native:
            slug = _doi_slug(ident)
            for candidate in (
                pdf_cache_dir / f"{slug}.pdf",
                pdf_cache_dir / f"{slug.replace('.', '-')}.pdf",
            ):
                if candidate.exists():
                    fig_path = fetch_figure_from_local_pdf(
                        candidate, cache_dir=figure_staging_dir,
                    )
                    if fig_path is not None:
                        break
        resolved[ident] = fig_path
        if len(resolved) >= len(gap_slides) * 3:
            break

    available_dois = [d for d, p in resolved.items() if p is not None]
    unavailable_dois = [d for d, p in resolved.items() if p is None]

    # Open the deck and insert
    prs = Presentation(str(deck_path))
    inserted_dois: list[str] = []
    placeholder_dois: list[str] = []

    available_iter = iter(available_dois)
    unavail_iter = iter(unavailable_dois)

    # Use the proper layout primitives shared with vaultlab.slides.layouts.figure
    from vaultlab.slides.layouts._helpers import (
        add_picture_fit, apply_font, sizes as _sizes,
    )
    SIZES = _sizes()

    def _classify_existing_layout(slide):
        """Categorize a figure-gap slide so we know how to place the figure.

        Returns ``(layout_kind, title_sh, dominant_body, decoration_shapes)``:

        - ``layout_kind``: "title_only" (only short title text, no body) or
          "title_plus" (title + dominant body block).
        - ``title_sh``: the title shape (first short non-empty text).
        - ``dominant_body``: the SINGLE largest non-title text shape — that's
          the bullet block we resize/relocate. ``None`` if title-only.
        - ``decoration_shapes``: page numbers, navigation footers, version
          labels — text shapes we leave UNTOUCHED. Resizing these is what
          caused the multi-shape stack-overlap bug on 2026-05-03.

        The dominant body is picked by largest area among non-title text
        shapes, which reliably picks the bullet block over page-number
        shapes (typically <0.5 sq in vs >20 sq in for bullets).
        """
        title_sh = None
        candidates = []  # (shape, area)
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = (sh.text_frame.text or "").strip()
            if not txt:
                continue
            if title_sh is None and len(txt) < 120:
                title_sh = sh
                continue
            try:
                area = (sh.width or 0) * (sh.height or 0)
            except (AttributeError, TypeError):
                area = 0
            candidates.append((sh, area))

        if not candidates:
            return "title_only", title_sh, None, []
        candidates.sort(key=lambda x: x[1], reverse=True)
        dominant = candidates[0][0]
        decorations = [sh for sh, _ in candidates[1:]]
        return "title_plus", title_sh, dominant, decorations

    for s in gap_slides:
        slide = prs.slides[s.index - 1]
        sw_in = prs.slide_width / 914400
        sh_in = prs.slide_height / 914400

        layout_kind, title_sh, dominant_body, decoration_shapes = _classify_existing_layout(slide)

        # Compute figure box per layout.
        fig_top_in = 1.2
        cap_h_in = 0.4
        cit_h_in = 0.4
        cap_gap = 0.05

        if layout_kind == "title_plus" and dominant_body is not None:
            # Hickey-lab style: figure on LEFT, bullets on RIGHT. Only the
            # SINGLE dominant body shape (the bullet block) gets resized;
            # decoration shapes (page numbers, navigation, footers) stay
            # untouched to avoid the stacked-shape overlap bug.
            fig_left_in = 0.4
            fig_w_in = sw_in * 0.58
            fig_h_in = sh_in - fig_top_in - cap_gap - cap_h_in - cit_h_in - 0.1
            if fig_h_in < 3.0:
                fig_h_in = 3.0
            body_left_in = fig_left_in + fig_w_in + 0.3
            body_top_in = fig_top_in
            body_w_in = sw_in - body_left_in - 0.3
            body_h_in = fig_h_in
            dominant_body.left = Inches(body_left_in)
            dominant_body.top = Inches(body_top_in)
            dominant_body.width = Inches(body_w_in)
            dominant_body.height = Inches(body_h_in)
            # Decoration shapes are NOT touched.
        else:
            # title-only: figure dominates centrally below title.
            fig_left_in = 0.5
            fig_w_in = sw_in - 1.0
            fig_h_in = sh_in - fig_top_in - cap_gap - cap_h_in - cit_h_in - 0.1

        # Try to insert a real figure
        try:
            doi = next(available_iter)
            fig_path = resolved[doi]
            add_picture_fit(
                slide, str(fig_path),
                Inches(fig_left_in), Inches(fig_top_in),
                Inches(fig_w_in), Inches(fig_h_in),
            )

            # Caption — italic 12pt Roboto, directly below figure
            cap_top_in = fig_top_in + fig_h_in + cap_gap
            cx = slide.shapes.add_textbox(
                Inches(fig_left_in), Inches(cap_top_in),
                Inches(fig_w_in), Inches(cap_h_in),
            )
            cx.text_frame.text = f"Figure 1 — {doi}"
            cx.text_frame.word_wrap = True
            apply_font(cx.text_frame, size=12, pres=prs)
            for para in cx.text_frame.paragraphs:
                for run in para.runs:
                    run.font.italic = True

            # Citation footer — 9pt Roboto, bottom of slide
            cit = slide.shapes.add_textbox(
                Inches(0.3), Inches(sh_in - cit_h_in),
                Inches(sw_in - 0.6), Inches(cit_h_in - 0.05),
            )
            cit.text_frame.text = (
                f"Source: https://doi.org/{doi}" if "/" in doi
                else f"Source: paperclip /papers/{doi}/"
            )
            apply_font(cit.text_frame, size=9, pres=prs)

            inserted_dois.append(doi)
            continue
        except StopIteration:
            pass

        # Aspirational: full-size FIGURE NEEDED placeholder where the figure
        # would have gone. Visible call-to-action box, not a corner stamp.
        if mode == "aspirational":
            try:
                doi = next(unavail_iter)
            except StopIteration:
                continue
            ph = slide.shapes.add_textbox(
                Inches(fig_left_in), Inches(fig_top_in),
                Inches(fig_w_in), Inches(fig_h_in),
            )
            tf = ph.text_frame
            tf.word_wrap = True
            tf.text = (
                f"📥 FIGURE NEEDED\n\n"
                f"Manually fetch Figure 1 from:\n"
                f"https://doi.org/{doi}\n\n"
                f"Drop the image here to complete this slide."
            )
            apply_font(tf, size=18, bold=True, pres=prs)
            for para in tf.paragraphs:
                try:
                    from pptx.enum.text import PP_ALIGN
                    para.alignment = PP_ALIGN.CENTER
                except Exception:
                    pass
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xAA, 0x33, 0x33)
            placeholder_dois.append(doi)

    prs.save(str(out_path))
    return FigurePopulateResult(
        deck_path=out_path,
        mode=mode,
        n_inserted=len(inserted_dois),
        n_placeholders=len(placeholder_dois),
        inserted_dois=inserted_dois,
        placeholder_dois=placeholder_dois,
    )


__all__ = [
    "FigurePopulateResult",
    "FigureSlot",
    "fetch_figure_from_paperclip",
    "fetch_figure_from_local_pdf",
    "populate_deck_with_figures",
]
