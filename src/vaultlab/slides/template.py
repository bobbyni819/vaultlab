"""Template loaders for vaultlab.slides — Hickey Lab + plain-theme presentations.

Lifted from ``bobby_slides._template`` (bobby-tools, 2026-04). Provides:

- :func:`load_template` — initialize a python-pptx Presentation from the
  bundled Hickey Lab template, with a chosen theme variant ("dark" / "light").
- :func:`load_plain_presentation` — black or white background, no template,
  for non-lab decks.
- :func:`default_font` / :func:`min_sizes` — projector-readable sizing
  (heading 28pt, body 24pt, caption 18pt).
- :func:`theme_colors` / :func:`theme_colors_hex` — Hickey Lab palette.

The bundled template lives at
``vaultlab/slides/themes/_assets/hickey_lab_template.pptx`` and is shipped
with the ``vaultlab`` package.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

# Lab theme colors extracted from 2023_09_19_Template.pptx
_LAB_THEME_COLORS_HEX: dict[str, str] = {
    "accent1": "29AF8C",  # teal/green primary
    "accent2": "97BE49",  # lime
    "accent3": "3D9CCC",  # blue
    "accent4": "7C60C6",  # purple
    "accent5": "C9492C",  # red/orange
    "accent6": "D58C2E",  # gold
    "dk2": "44546A",      # dark slate
    "lt2": "E7E6E6",      # light gray
}

_BUNDLED_TEMPLATE_PATH = (
    Path(__file__).parent / "themes" / "_assets" / "hickey_lab_template.pptx"
)

_DEFAULT_FONT = "Roboto"

_MIN_SIZES: dict[str, int] = {
    "heading": 28,
    "body": 24,
    "caption": 18,
}

# The Hickey Lab template has 3 slide masters with different color maps:
#   Master 0 - bg1=dk1 (DARK background, light text)
#   Master 1 - bg1=lt1 (LIGHT background, dark text)
#   Master 2 - bg1=dk1 (DARK section-divider variant)
_THEME_MASTER_INDEX: dict[str, int] = {
    "dark": 0,
    "light": 1,
    "default": 0,
}


def lab_template_path() -> Path | None:
    """Return the bundled Hickey Lab template path, or ``None`` if missing."""
    return _BUNDLED_TEMPLATE_PATH if _BUNDLED_TEMPLATE_PATH.exists() else None


def load_plain_presentation(theme: str = "dark") -> Any:
    """Create a plain Presentation with no lab template — black or white bg.

    Use this for decks that should NOT use the Hickey lab template (e.g.
    talks outside Hickey lab work). Returns a python-pptx Presentation
    with the slide background set to black (``theme="dark"``) or white
    (``"light"``) on the master.
    """
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    # 16:9 widescreen
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    bg_hex = "000000" if theme.lower() == "dark" else "FFFFFF"
    _set_master_background(pres, bg_hex)

    pres._vaultlab_master_index = None  # signals "use default layouts"
    pres._vaultlab_plain_theme = theme.lower()
    return pres


def _set_master_background(pres: Any, hex_color: str) -> None:
    """Set the slide master's background to a solid color (hex without #)."""
    from lxml import etree

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    bg_xml = (
        f'<p:bg xmlns:p="{p_ns}" xmlns:a="{a_ns}">'
        f'<p:bgPr><a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f'<a:effectLst/></p:bgPr>'
        f'</p:bg>'
    )
    for master in pres.slide_masters:
        sld_root = master.element
        cSld = sld_root.find(f"{{{p_ns}}}cSld")
        if cSld is None:
            continue
        existing = cSld.find(f"{{{p_ns}}}bg")
        if existing is not None:
            cSld.remove(existing)
        cSld.insert(0, etree.fromstring(bg_xml))


def load_template(
    path: Path | str | None = None,
    strip_starter_slides: bool = True,
    theme: str = "dark",
) -> Any:
    """Load a ``.pptx`` template into a python-pptx Presentation.

    Args:
        path: Explicit template path. If ``None``, loads the bundled Hickey
            Lab template.
        strip_starter_slides: If ``True`` (default), removes any pre-existing
            slides in the template. The Hickey lab template ships with 6
            example slides we don't want in user output. Layouts and masters
            are preserved.
        theme: ``"dark"`` (default) or ``"light"``. Selects which slide
            master new slides will use.

    Returns:
        A python-pptx Presentation initialized from the template. The
        chosen master index is attached as ``pres._vaultlab_master_index``
        so layout primitives know which master to draw layouts from.

    Raises:
        FileNotFoundError: If no template is available.
        ImportError: If ``python-pptx`` is not installed.
    """
    from pptx import Presentation

    if path is None:
        resolved = lab_template_path()
        if resolved is None:
            raise FileNotFoundError(
                "Hickey Lab template not bundled. Expected at:\n"
                f"  {_BUNDLED_TEMPLATE_PATH}"
            )
    else:
        resolved = Path(path)
        if not resolved.exists():
            raise FileNotFoundError(f"Template not found: {resolved}")

    pres = Presentation(str(resolved))

    if strip_starter_slides:
        _strip_existing_slides(pres)

    pres._vaultlab_master_index = _THEME_MASTER_INDEX.get(theme.lower(), 0)
    pres._vaultlab_theme_variant = theme.lower()
    return pres


def _strip_existing_slides(pres: Any) -> None:
    """Remove all slides from a Presentation, preserving layouts + masters."""
    rels_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    xml_slides = pres.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for slide_id_elem in slides_to_remove:
        rel_id = slide_id_elem.attrib[f"{{{rels_ns}}}id"]
        pres.part.drop_rel(rel_id)
        xml_slides.remove(slide_id_elem)


def theme_colors() -> dict[str, Any]:
    """Return the Hickey Lab theme colors as ``RGBColor`` instances."""
    from pptx.dml.color import RGBColor

    return {
        name: RGBColor.from_string(hex_val)
        for name, hex_val in _LAB_THEME_COLORS_HEX.items()
    }


def theme_colors_hex() -> dict[str, str]:
    """Return the Hickey Lab theme colors as hex strings."""
    return dict(_LAB_THEME_COLORS_HEX)


def default_font() -> str:
    """The Hickey Lab default font (Roboto)."""
    return _DEFAULT_FONT


def min_sizes() -> dict[str, int]:
    """Min font sizes (heading 28, body 24, caption 18) in points."""
    return dict(_MIN_SIZES)


__all__ = [
    "default_font",
    "lab_template_path",
    "load_plain_presentation",
    "load_template",
    "min_sizes",
    "theme_colors",
    "theme_colors_hex",
]
