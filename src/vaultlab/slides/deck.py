"""Slide and Deck data classes — declarative deck representation.

Backend-independent. The renderer (``vaultlab.slides.render``) is the only
module that talks to ``python-pptx``; everything else operates on these
dataclasses.

This module also hosts the **higher-level multi-slide composer** —
:class:`DeckSlide`, :class:`DeckPlan`, :func:`build_deck`, and
:func:`build_deck_from_lineage_result`. The composer wraps the low-level
``Slide``/``Deck``/``render_pptx`` pipeline plus the annotated-figure-slide
primitive into a single function that turns a structured "deck plan"
(title / section_intro / figure / bullets / references slides) into a
finished ``.pptx``. It is the prerequisite for the L4 e2e test that turns
a ``run_lit_arc`` corpus into a journal-club deck.
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from vaultlab.research.lineage import LineageRunResult

logger = logging.getLogger(__name__)

# Layout names supported by the starter layouts module
SUPPORTED_LAYOUTS: frozenset[str] = frozenset(
    {"title", "content_with_bullets", "figure_with_caption"}
)


@dataclass
class Slide:
    """One slide in a deck.

    Fields are deliberately layout-agnostic — each layout reads only the fields
    it needs. Unused fields are ignored.

    Attributes
    ----------
    layout
        Layout name from ``SUPPORTED_LAYOUTS``.
    title
        Slide title (rendered as the H1 of the slide).
    subtitle
        Subtitle for ``title`` layout. Ignored elsewhere.
    bullets
        Bullet-list content for ``content_with_bullets``.
    figure_path
        Path to a ``.png`` / ``.jpg`` for ``figure_with_caption``. Resolved
        relative to ``Deck.working_dir`` if not absolute.
    caption
        Caption text for ``figure_with_caption``.
    speaker_notes
        Optional speaker notes (rendered into the .pptx notes panel).
    """

    layout: str
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    figure_path: str | None = None
    caption: str = ""
    speaker_notes: str = ""

    def __post_init__(self) -> None:
        if self.layout not in SUPPORTED_LAYOUTS:
            raise ValueError(
                f"Unsupported layout {self.layout!r}. Supported: {sorted(SUPPORTED_LAYOUTS)}"
            )


@dataclass
class Deck:
    """A complete slide deck.

    Attributes
    ----------
    title
        Deck title (used in the .pptx properties + as a default for the
        opening title slide).
    slides
        Ordered list of slides.
    theme
        Theme name. Defaults to ``"default"``.
    working_dir
        Used to resolve relative ``figure_path`` values.
    metadata
        Free-form provenance attached to the deck — kept in the file
        properties (``vaultlab_provenance`` key when supported by the backend).
    """

    title: str
    slides: list[Slide] = field(default_factory=list)
    theme: str = "default"
    working_dir: Path | None = None
    metadata: dict[str, str] = field(default_factory=dict)

    def add(self, slide: Slide) -> None:
        """Append a slide. Useful for building decks programmatically."""
        self.slides.append(slide)

    def __len__(self) -> int:  # convenience
        return len(self.slides)


# ---------------------------------------------------------------------------
# Multi-slide composer (DeckSlide / DeckPlan / build_deck)
# ---------------------------------------------------------------------------

# Slide kinds the composer understands. These are higher-level than the
# layout names above — a single composer kind may pick one of several
# concrete renderers (e.g. ``figure`` uses the annotated-figure primitive).
SUPPORTED_DECK_SLIDE_KINDS: frozenset[str] = frozenset(
    {"title", "section_intro", "figure", "bullets", "references"}
)


@dataclass(frozen=True)
class DeckSlide:
    """One slide spec for the multi-slide composer.

    The ``content`` dict is interpreted per-kind:

    - ``title`` — ``{"subtitle": str, "speaker": str, "affiliation": str,
      "date": str}``
    - ``section_intro`` — ``{"section_name": str, "key_question": str,
      "bullets": list[str]}`` (1-3 bullets)
    - ``figure`` — ``{"figure_path": Path | str, "annotations":
      list[ElementAnnotation], "motif_colors": dict[str, tuple[int,int,int]],
      "caption": str, "citation_doi": str}``
    - ``bullets`` — ``{"bullets": list[str], "citations": list[str]}``
      (citations are DOI strings rendered as ``[N]`` markers)
    - ``references`` — ``{"refs": list[dict]}`` where each dict is
      ``{"n": int, "citation": str, "doi": str}``

    Unused keys per-kind are ignored — the dict shape is intentionally
    permissive so callers can pass extra context without breaking.
    """

    kind: str
    title: str
    content: dict[str, Any] = field(default_factory=dict)
    notes: str = ""

    def __post_init__(self) -> None:
        if self.kind not in SUPPORTED_DECK_SLIDE_KINDS:
            raise ValueError(
                f"Unsupported DeckSlide kind {self.kind!r}. "
                f"Supported: {sorted(SUPPORTED_DECK_SLIDE_KINDS)}"
            )


@dataclass(frozen=True)
class DeckPlan:
    """Full deck structure consumed by :func:`build_deck`."""

    title: str
    subtitle: str
    speaker: str
    affiliation: str
    sections: list[str] = field(default_factory=list)
    slides: list[DeckSlide] = field(default_factory=list)
    theme: str = "hickey_lab"


def build_deck(
    plan: DeckPlan,
    output_path: Path | str,
    *,
    citations: dict[str, dict[str, Any]] | None = None,
) -> Path:
    """Compose a multi-slide ``.pptx`` from a :class:`DeckPlan`.

    Slide-by-slide dispatch:

    - ``title`` / ``section_intro`` / ``bullets`` / ``references`` are
      rendered with python-pptx text primitives directly (these don't need
      annotation magic — bullets and headings are fine in plain text boxes).
    - ``figure`` slides delegate to
      :func:`vaultlab.slides.annotated_figure_slide.add_annotated_figure_slide`.

    Theming:

    - ``plan.theme == "hickey_lab"`` initialises the presentation from the
      bundled Hickey Lab template (preserves logos, color theme, masters)
      and uses the ``HICKEY_LAB_LAYOUT`` for figure slides.
    - Any other theme falls back to a vanilla 16:9 presentation with the
      ``DEFAULT`` layout.

    Returns the written ``.pptx`` path.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    citations = citations or {}

    # Lazy imports — keep the dataclass module light when callers only need
    # the data shapes (e.g. unit tests that inspect a DeckPlan).
    from pptx import Presentation
    from pptx.util import Inches

    from vaultlab.slides.annotated_figure_slide import (
        DEFAULT,
        HICKEY_LAB_LAYOUT,
        SlideLayout,
        add_annotated_figure_slide,
    )

    if plan.theme == "hickey_lab":
        try:
            from vaultlab.slides.themes.hickey_lab import (
                load_hickey_lab_presentation,
            )

            pres = load_hickey_lab_presentation(theme="dark")
            layout: SlideLayout = HICKEY_LAB_LAYOUT
        except FileNotFoundError:
            logger.warning("Hickey Lab template not bundled; falling back to default theme.")
            pres = Presentation()
            pres.slide_width = Inches(13.333)
            pres.slide_height = Inches(7.5)
            layout = DEFAULT
    else:
        pres = Presentation()
        pres.slide_width = Inches(13.333)
        pres.slide_height = Inches(7.5)
        layout = DEFAULT

    for idx, ds in enumerate(plan.slides):
        page_number = idx + 1
        section_idx = _current_section_idx(ds, plan.sections)
        if ds.kind == "title":
            _add_title_slide(pres, ds, plan, layout)
        elif ds.kind == "section_intro":
            _add_section_intro_slide(pres, ds, layout)
        elif ds.kind == "bullets":
            _add_bullets_slide(
                pres,
                ds,
                layout,
                page_number=page_number,
                sections=plan.sections,
                section_idx=section_idx,
                citations=citations,
            )
        elif ds.kind == "figure":
            _add_figure_slide(
                pres,
                ds,
                layout,
                page_number=page_number,
                sections=plan.sections,
                section_idx=section_idx,
                add_fn=add_annotated_figure_slide,
            )
        elif ds.kind == "references":
            _add_references_slide(pres, ds, layout)
        else:  # pragma: no cover — guarded by DeckSlide.__post_init__
            raise ValueError(f"Unknown DeckSlide kind: {ds.kind}")

    pres.save(str(output_path))
    return output_path


def _format_author_lastname(author: str) -> str:
    """Extract a Vancouver-style last-name surface form from a free-form author string.

    Thin shim over :func:`vaultlab.kb.paths.format_author_lastname` —
    that helper is the single source of truth for surname extraction
    across NCBI ``"Last F"`` / OpenAlex ``"J. Kennedy-Darling"`` /
    Vancouver ``"Last, First"`` / western ``"First Last"`` formats and
    normalizes unicode hyphens to ASCII at the same time.

    Pre-2026-04-30 evening-5 this had its own copy of the logic which
    DIDN'T handle OpenAlex's "F. Last" format — produced ``"J. 2020"``
    instead of ``"Kennedy-Darling 2020"``.
    """
    from vaultlab.kb.paths import format_author_lastname

    return format_author_lastname(author)


def _format_citation_label(
    citation: dict[str, Any] | None,
    n: int,
    fallback: str = "",
) -> str:
    """Render a Vancouver-style ``[N] Last Year`` footnote label.

    Used by the bullets-slide citations footer. ``citation`` is the dict
    populated by :func:`_collect_citations_from_summaries` (or callers
    passing equivalent shape).  Falls back to ``fallback`` (typically the
    DOI) when authors are missing entirely.
    """
    authors = (citation or {}).get("authors") or []
    last = ""
    for a in authors:
        last = _format_author_lastname(a)
        if last:
            break
    if not last:
        if not fallback:
            logger.warning(
                "citation %s has no parseable author; using 'Anon' fallback",
                fallback or "(unknown doi)",
            )
            last = "Anon"
        else:
            return f"[{n}] {fallback}"
    year = (citation or {}).get("year") or ""
    return f"[{n}] {last} {year}".strip()


def _write_deck_provenance(
    deck_pptx_path: Path,
    *,
    lineage_result: LineageRunResult,
    project: str,
    speaker: str,
    affiliation: str,
    audience: str,
    target_slide_count: int,
    plan_mode_used: str,
    audit_status: str,
    figure_assignments: dict[str, Path] | None,
) -> None:
    """Write the deck's ``.provenance.json`` + ``.method.md`` sidecars.

    Implements F-6 from the pipeline-integration-map audit: per AGENTS.md
    Invariant 3 every output writes provenance, but
    :func:`build_deck_from_lineage_result` historically shipped a
    ``.pptx`` with no receipt. This helper closes that gap. Best-effort:
    receipt-write failures are logged and never raised so deck shipping
    isn't gated on optional metadata.
    """
    try:
        from vaultlab.provenance import ProvenanceRecord, write_receipts

        record = ProvenanceRecord(
            generated_by="vaultlab.slides.deck.build_deck_from_lineage_result",
            project=project,
            topic=lineage_result.topic,
            kind="slide_deck",
            inputs=[str(p) for p in lineage_result.summary_paths.values()],
            related_outputs=[str(lineage_result.arc_path)] if lineage_result.arc_path else [],
            params={
                "speaker": speaker,
                "affiliation": affiliation,
                "audience": audience,
                "target_slide_count": target_slide_count,
                "plan_mode": plan_mode_used,
                "audit_status": audit_status or "n/a",
                "n_figure_assignments": len(figure_assignments or {}),
                "n_summaries": len(lineage_result.summary_paths),
            },
            tags=["deck", "lit-arc"],
            notes=(
                "Composed via build_deck_from_lineage_result. "
                f"Arc source: {lineage_result.arc_path}"
            ),
        )
        write_receipts(deck_pptx_path, record)
    except Exception:
        logger.exception("write_receipts failed for deck %s", deck_pptx_path)


def build_deck_from_lineage_result(
    lineage_result: LineageRunResult,
    *,
    speaker: str,
    affiliation: str = "Hickey Lab @ Duke BME",
    project_slug: str | None = None,
    figure_assignments: dict[str, Path] | None = None,
    kb_root: Path | None = None,
    plan_callback: Any = None,
    audience: str = "journal-club",
    target_slide_count: int = 7,
    plan_mode: str = "fast",
    crosstalk_runner: Any = None,
    crosstalk_n_rounds: int = 3,
    final_audit: bool = False,
    audit_strict: bool = False,
    run_dir: Path | None = None,
) -> Path:
    """Take a ``/lit-arc`` result and synthesize a ~7-slide deck.

    Two paths:

    * **v0.1 fast path** — ``plan_callback=None`` (default). Reads each
      per-paper summary's frontmatter to recover ``year_bucket``, then
      composes mechanically:

      1. Title (lineage topic)
      2. Section intro: Background
      3. Figure (history bucket) — IF a figure assignment exists for a
         history-bucket DOI; otherwise dropped and replaced with bullets.
      4. Section intro: Development
      5. Bullets: SOTA findings — TL;DR snippets from sota-bucket papers
      6. Section intro: Synthesis — pulled from the arc narrative when
         available; otherwise a stub instructing the user to fill it in.
      7. References — Vancouver-style 2-column

    * **v0.1+ rigor path** — ``plan_callback`` set to a
      :data:`vaultlab.workflows.deck_plan.PlanGeneratorCallback`. The
      callback receives a :class:`DeckPlanTask` (with full corpus
      summaries, metrics, figure assignments) and returns Claude Code's
      typed JSON plan. The resulting plan is rendered via
      :func:`vaultlab.slides.build_from_plan` (the dict-plan renderer)
      so we get the richer slide-type set
      (``title / section_divider / figure / multi_figure / text /
      references``) instead of the 5-kind composer.

    Output is routed through :func:`vaultlab.kb.paths.deck_path` so the
    .pptx lands at ``<kb_root>/Output/<project>/<topic>-deck.pptx``.

    Per AGENTS.md Invariant 3 (every output writes provenance), this
    function also drops ``<deck>.pptx.provenance.json`` and
    ``<deck>.pptx.method.md`` next to the ``.pptx`` via
    :func:`vaultlab.provenance.write_receipts` (F-6 in the
    pipeline-integration-map audit).
    """
    from vaultlab.kb.paths import deck_path, slugify_topic

    # G-2 fix (option b): if project_slug wasn't threaded explicitly, walk
    # up from cwd looking for ``.vaultlab-project.json`` and adopt its
    # slug. Aligns with the "state-aware, additive, read-before-write"
    # memory rule — explicit kwarg still wins, but a forgetful slash-
    # command body no longer creates a parallel ``Wiki/Projects/<slug>/``.
    if project_slug is None:
        try:
            from vaultlab.onboarding import load_project_config_from_cwd

            _cfg = load_project_config_from_cwd()
        except Exception:  # pragma: no cover — never break a run
            logger.exception("load_project_config_from_cwd failed")
            _cfg = None
        if _cfg is not None and getattr(_cfg, "slug", ""):
            project_slug = _cfg.slug
            logger.info(
                "auto-discovered project_slug=%s from .vaultlab-project.json (cwd=%s)",
                project_slug,
                Path.cwd(),
            )

    # Multi-tenant KB-root resolution (Layer A, 2026-04-30): kb_root is now
    # optional. Falls through env-var → vaultlab config → bobby_kb compat →
    # first-run prompt. See run_lit_arc / run_lit_report for the matching
    # blocks; resolve_kb_root() is the single canonical resolver.
    if kb_root is None:
        from vaultlab.context.locations import resolve_kb_root

        kb_root = resolve_kb_root()

    project = project_slug or "lit-arc"
    deck_name = f"{slugify_topic(lineage_result.topic)}-deck.pptx"
    out = deck_path(Path(kb_root), project, deck_name)
    out.parent.mkdir(parents=True, exist_ok=True)

    # Track which path was taken so the provenance receipt can record it.
    plan_mode_used = "fast"
    audit_status: str = ""

    # Adversarial plan path — multi-agent crosstalk over plan generation.
    if plan_mode == "adversarial" and crosstalk_runner is not None:
        from vaultlab.workflows.crosstalk import (
            adversarial_deck_plan_meeting,
            rigor_audit,
            write_crosstalk_artifacts,
        )

        summaries = _summaries_to_paper_summaries(
            _read_summary_frontmatters(lineage_result.summary_paths)
        )
        ct = adversarial_deck_plan_meeting(
            topic=lineage_result.topic,
            summaries=summaries,
            figure_assignments=figure_assignments or {},
            target_slide_count=target_slide_count,
            n_rounds=crosstalk_n_rounds,
            runner_callback=crosstalk_runner,
        )
        if run_dir is not None:
            try:
                write_crosstalk_artifacts(ct, run_dir=Path(run_dir))
            except Exception:
                logger.exception("write_crosstalk_artifacts (deck-plan) failed")

        # Render the dict-plan from the synthesizer's final_output.
        from vaultlab.workflows.deck_plan import (
            prepare_deck_plan_task,
            render_plan_from_response,
        )

        # F-13: prefer the live corpus carried on the lineage result so we
        # keep ``co_citation_pairs`` / ``seeds`` / ``references`` instead of
        # the empty stand-ins ``_synthetic_corpus_from_summaries`` produces.
        corpus = lineage_result.corpus or _synthetic_corpus_from_summaries(
            topic=lineage_result.topic,
            summaries=summaries,
        )
        task = prepare_deck_plan_task(
            topic=lineage_result.topic,
            corpus=corpus,
            summaries=summaries,
            figure_assignments=figure_assignments or {},
            speaker=speaker,
            affiliation=affiliation,
            audience=audience,
            target_slide_count=target_slide_count,
            kb_root=Path(kb_root),
        )
        dict_plan = render_plan_from_response(task, ct.final_output or {})

        plan_mode_used = "adversarial"

        # Final-gate rigor audit before .pptx ships.
        if final_audit:
            text = _render_plan_for_audit(dict_plan)
            audit = rigor_audit(
                document=text,
                summaries=summaries,
                audit_kind="deck",
                runner_callback=crosstalk_runner,
            )
            blockers = [i for i in audit.get("issues", []) if i.get("severity") == "blocker"]
            if blockers and audit_strict:
                raise RuntimeError(
                    f"rigor_audit found {len(blockers)} blocker issue(s) "
                    "and audit_strict=True; refusing to ship deck."
                )
            audit_status = (
                "failed"
                if blockers
                else ("passed_with_warnings" if not audit.get("passed", True) else "passed")
            )
            if not audit.get("passed", True):
                # Prepend a warning slide noting the audit issues.
                dict_plan.setdefault("slides", []).insert(
                    1,
                    {
                        "type": "text",
                        "title": "[Audit warnings]",
                        "bullets": [
                            f"{i.get('severity', '?')}: {i.get('loc', '?')} — {i.get('fix', '')}"
                            for i in audit.get("issues", [])[:6]
                        ],
                    },
                )

        out_pptx = build_from_plan(dict_plan, out, write_marp=False)["pptx"]
        _write_deck_provenance(
            out_pptx,
            lineage_result=lineage_result,
            project=project,
            speaker=speaker,
            affiliation=affiliation,
            audience=audience,
            target_slide_count=target_slide_count,
            plan_mode_used=plan_mode_used,
            audit_status=audit_status,
            figure_assignments=figure_assignments,
        )
        return out_pptx

    if plan_callback is not None:
        # v0.1+ rigor path — content-aware deck plan via Claude Code.
        # Lazy imports so the slides module stays light when the deck-
        # plan workflow isn't being used.
        from vaultlab.workflows.deck_plan import generate_deck_plan

        summaries = _summaries_to_paper_summaries(
            _read_summary_frontmatters(lineage_result.summary_paths)
        )
        # F-13: prefer the live corpus from the lineage result.
        corpus = lineage_result.corpus or _synthetic_corpus_from_summaries(
            topic=lineage_result.topic,
            summaries=summaries,
        )
        dict_plan = generate_deck_plan(
            topic=lineage_result.topic,
            corpus=corpus,
            summaries=summaries,
            figure_assignments=figure_assignments or {},
            speaker=speaker,
            affiliation=affiliation,
            audience=audience,
            target_slide_count=target_slide_count,
            kb_root=Path(kb_root),
            plan_callback=plan_callback,
        )

        plan_mode_used = "plan_callback"

        # Optional rigor audit on the single-shot path too.
        if final_audit and crosstalk_runner is not None:
            from vaultlab.workflows.crosstalk import rigor_audit

            text = _render_plan_for_audit(dict_plan)
            audit = rigor_audit(
                document=text,
                summaries=summaries,
                audit_kind="deck",
                runner_callback=crosstalk_runner,
            )
            blockers = [i for i in audit.get("issues", []) if i.get("severity") == "blocker"]
            if blockers and audit_strict:
                raise RuntimeError(
                    f"rigor_audit found {len(blockers)} blocker issue(s) "
                    "and audit_strict=True; refusing to ship deck."
                )
            audit_status = (
                "failed"
                if blockers
                else ("passed_with_warnings" if not audit.get("passed", True) else "passed")
            )
            if not audit.get("passed", True):
                dict_plan.setdefault("slides", []).insert(
                    1,
                    {
                        "type": "text",
                        "title": "[Audit warnings]",
                        "bullets": [
                            f"{i.get('severity', '?')}: {i.get('loc', '?')} — {i.get('fix', '')}"
                            for i in audit.get("issues", [])[:6]
                        ],
                    },
                )

        out_pptx = build_from_plan(dict_plan, out, write_marp=False)["pptx"]
        _write_deck_provenance(
            out_pptx,
            lineage_result=lineage_result,
            project=project,
            speaker=speaker,
            affiliation=affiliation,
            audience=audience,
            target_slide_count=target_slide_count,
            plan_mode_used=plan_mode_used,
            audit_status=audit_status,
            figure_assignments=figure_assignments,
        )
        return out_pptx

    plan = _plan_from_lineage(
        lineage_result,
        speaker=speaker,
        affiliation=affiliation,
        figure_assignments=figure_assignments or {},
    )
    out_pptx = build_deck(plan, out, citations=_collect_citations_from_summaries(lineage_result))
    _write_deck_provenance(
        out_pptx,
        lineage_result=lineage_result,
        project=project,
        speaker=speaker,
        affiliation=affiliation,
        audience=audience,
        target_slide_count=target_slide_count,
        plan_mode_used=plan_mode_used,  # "fast" — initial default
        audit_status=audit_status,
        figure_assignments=figure_assignments,
    )
    return out_pptx


def _render_plan_for_audit(dict_plan: dict[str, Any]) -> str:
    """Flatten a dict-plan into plain text suitable for rigor_audit input."""
    lines: list[str] = [f"# {dict_plan.get('title', '')}"]
    for slide in dict_plan.get("slides", []):
        lines.append("")
        lines.append(f"## [{slide.get('type', '?')}] {slide.get('title', '')}")
        if slide.get("subtitle"):
            lines.append(slide["subtitle"])
        for b in slide.get("bullets", []) or []:
            lines.append(f"- {b}")
        if slide.get("caption"):
            lines.append(slide["caption"])
        if slide.get("citation_source"):
            lines.append(f"_(source: {slide['citation_source']})_")
        for ref in slide.get("references", []) or []:
            lines.append(f"- {ref}")
    return "\n".join(lines)


def _summaries_to_paper_summaries(
    summaries_by_doi: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    """Convert frontmatter dicts to :class:`PaperSummary` instances.

    Used by :func:`build_deck_from_lineage_result` when routing through
    the content-aware deck-plan generator.  Imports
    :class:`vaultlab.research.summarize.PaperSummary` lazily so the
    slides module stays decoupled from the research module unless this
    path is exercised.
    """
    from vaultlab.research.summarize import PaperSummary

    out: dict[str, PaperSummary] = {}
    for doi, fm in summaries_by_doi.items():
        s = PaperSummary(
            doi=fm.get("doi", "") or doi,
            title=fm.get("title", "") or "",
            authors=list(fm.get("authors", []) or []),
            year=int(fm.get("year") or 0) if fm.get("year") else 0,
            journal=fm.get("journal", "") or "",
            citation_count=int(fm.get("citation_count") or 0) if fm.get("citation_count") else 0,
            og_score=float(fm.get("og_score") or 0.0),
            forward_influence=int(fm.get("forward_influence") or 0),
            year_bucket=fm.get("year_bucket", "unknown") or "unknown",
            tier=fm.get("tier", "C") or "C",
            tldr=(fm.get("tldr") or "").strip(),
            key_findings=list(fm.get("key_findings", []) or []),
        )
        out[doi] = s
    return out


def _synthetic_corpus_from_summaries(
    *,
    topic: str,
    summaries: dict[str, Any],
) -> Any:
    """Build a minimal :class:`Corpus` for the deck-plan generator.

    The deck-plan generator only reads ``corpus.papers``,
    ``corpus.metrics.og_score``, ``corpus.metrics.forward_influence``,
    ``corpus.metrics.co_citation_pairs``, and
    ``corpus.metrics.year_buckets``.  We synthesize those from the
    on-disk summaries when no live Corpus is available (e.g. when the
    deck is being built from a previous-day's ``Wiki/Summaries/``).
    """
    from vaultlab.research.corpus import Corpus
    from vaultlab.research.graph_metrics import CorpusMetrics
    from vaultlab.research.paper import Paper

    papers: dict[str, Any] = {}
    og_score: dict[str, float] = {}
    forward_influence: dict[str, int] = {}
    year_buckets: dict[str, str] = {}
    seeds: list[Paper] = []
    for doi, s in summaries.items():
        key = (doi or s.doi or "").lower()
        if not key:
            continue
        paper = Paper(
            doi=s.doi or doi,
            title=s.title,
            authors=list(s.authors or []),
            year=int(s.year) if s.year else 0,
            journal=s.journal,
        )
        papers[key] = paper
        og_score[key] = float(getattr(s, "og_score", 0.0) or 0.0)
        forward_influence[key] = int(getattr(s, "forward_influence", 0) or 0)
        bucket = getattr(s, "year_bucket", "unknown") or "unknown"
        year_buckets[key] = bucket
    metrics = CorpusMetrics(
        og_score=og_score,
        forward_influence=forward_influence,
        co_citation_pairs=[],
        year_buckets=year_buckets,
    )
    return Corpus(
        topic=topic,
        seeds=seeds,
        papers=papers,
        references={},
        metrics=metrics,
    )


# ---------------------------------------------------------------------------
# Internal: per-kind slide builders
# ---------------------------------------------------------------------------


def _current_section_idx(ds: DeckSlide, sections: list[str]) -> int | None:
    """Find which section a section_intro slide belongs to (for banner highlight)."""
    if ds.kind != "section_intro":
        return None
    name = ds.content.get("section_name", "")
    if name in sections:
        return sections.index(name)
    return None


def _add_title_slide(pres, ds: DeckSlide, plan: DeckPlan, layout) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    from vaultlab.slides.annotated_figure_slide import (
        muted_text_color_for_theme,
        text_color_for_theme,
    )

    blank = pres.slide_layouts[6]
    s = pres.slides.add_slide(blank)

    # Big centered title
    box = s.shapes.add_textbox(
        Inches(0.5),
        Inches(2.5),
        Inches(layout.slide_w_in - 1.0),
        Inches(1.5),
    )
    box.name = "title_main"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ds.title or plan.title
    r.font.name = "Arial"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    subtitle = ds.content.get("subtitle") or plan.subtitle
    speaker = ds.content.get("speaker") or plan.speaker
    affiliation = ds.content.get("affiliation") or plan.affiliation
    date_str = ds.content.get("date", "")

    sub_lines = [x for x in (subtitle, speaker, affiliation, date_str) if x]
    if sub_lines:
        sub = s.shapes.add_textbox(
            Inches(0.5),
            Inches(4.2),
            Inches(layout.slide_w_in - 1.0),
            Inches(2.5),
        )
        sub.name = "title_subtitle"
        tf2 = sub.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(sub_lines):
            para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = line
            run.font.name = "Arial"
            run.font.size = Pt(20 if i == 0 else 16)
            run.font.color.rgb = RGBColor(*muted_text_color_for_theme(layout.theme_variant))

    if ds.notes:
        s.notes_slide.notes_text_frame.text = ds.notes


def _add_section_intro_slide(pres, ds: DeckSlide, layout) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    from vaultlab.slides.annotated_figure_slide import (
        muted_text_color_for_theme,
        text_color_for_theme,
    )

    blank = pres.shapes if False else pres.slide_layouts[6]
    s = pres.slides.add_slide(pres.slide_layouts[6])

    section_name = ds.content.get("section_name", "")
    key_question = ds.content.get("key_question", "")
    bullets = list(ds.content.get("bullets", []))[:3]

    # Section name (small, top)
    if section_name:
        sn = s.shapes.add_textbox(
            Inches(0.6),
            Inches(0.5),
            Inches(layout.slide_w_in - 1.2),
            Inches(0.6),
        )
        sn.name = "section_name"
        tf = sn.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = section_name.upper()
        r.font.name = "Arial"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0, 102, 204)

    # Title
    title_box = s.shapes.add_textbox(
        Inches(0.6),
        Inches(1.2),
        Inches(layout.slide_w_in - 1.2),
        Inches(1.5),
    )
    title_box.name = "section_title"
    ttf = title_box.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    r = p.add_run()
    r.text = ds.title
    r.font.name = "Arial"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    # Key question
    if key_question:
        kq = s.shapes.add_textbox(
            Inches(0.6),
            Inches(2.9),
            Inches(layout.slide_w_in - 1.2),
            Inches(0.9),
        )
        kq.name = "section_key_question"
        ktf = kq.text_frame
        ktf.word_wrap = True
        p = ktf.paragraphs[0]
        r = p.add_run()
        r.text = key_question
        r.font.name = "Arial"
        r.font.size = Pt(20)
        r.font.italic = True
        r.font.color.rgb = RGBColor(*muted_text_color_for_theme(layout.theme_variant))

    # Bullets
    if bullets:
        bb = s.shapes.add_textbox(
            Inches(0.8),
            Inches(4.0),
            Inches(layout.slide_w_in - 1.6),
            Inches(2.8),
        )
        bb.name = "section_bullets"
        btf = bb.text_frame
        btf.word_wrap = True
        for i, b in enumerate(bullets):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            run = para.add_run()
            run.text = f"•  {b}"
            run.font.name = "Arial"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    if ds.notes:
        s.notes_slide.notes_text_frame.text = ds.notes


def _add_bullets_slide(
    pres,
    ds: DeckSlide,
    layout,
    *,
    page_number: int,
    sections: list[str],
    section_idx: int | None,
    citations: dict[str, dict[str, Any]],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    from vaultlab.slides.annotated_figure_slide import (
        _add_page_number,
        _add_section_banner,
        muted_text_color_for_theme,
        text_color_for_theme,
    )

    s = pres.slides.add_slide(pres.slide_layouts[6])

    # Title
    tbox = s.shapes.add_textbox(
        Inches(0.4),
        Inches(0.15),
        Inches(layout.slide_w_in - 0.8),
        Inches(layout.title_h_in - 0.1),
    )
    tbox.name = "slide_title"
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ds.title
    r.font.name = "Arial"
    r.font.size = Pt(layout.title_font_pt)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    bullets = list(ds.content.get("bullets", []))
    cite_dois = list(ds.content.get("citations", []))

    bb = s.shapes.add_textbox(
        Inches(0.8),
        Inches(layout.title_h_in + 0.3),
        Inches(layout.slide_w_in - 1.6),
        Inches(layout.slide_h_in - layout.title_h_in - layout.footer_h_in - 0.5),
    )
    bb.name = "slide_bullets"
    btf = bb.text_frame
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        run = para.add_run()
        run.text = f"•  {b}"
        run.font.name = "Arial"
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    # Citations footnote line
    if cite_dois:
        labels = []
        for n, doi in enumerate(cite_dois, 1):
            cite = citations.get(doi)
            labels.append(_format_citation_label(cite, n, fallback=doi))
        cb = s.shapes.add_textbox(
            Inches(0.8),
            Inches(layout.slide_h_in - layout.footer_h_in - 0.5),
            Inches(layout.slide_w_in - 1.6),
            Inches(0.4),
        )
        cb.name = "slide_citations_footer"
        ctf = cb.text_frame
        ctf.word_wrap = True
        p = ctf.paragraphs[0]
        run = p.add_run()
        run.text = "  ".join(labels)
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*muted_text_color_for_theme(layout.theme_variant))

    _add_page_number(s, page_number, layout)
    if sections:
        _add_section_banner(s, sections=sections, current_idx=section_idx, layout=layout)

    if ds.notes:
        s.notes_slide.notes_text_frame.text = ds.notes


def _add_figure_slide(
    pres,
    ds: DeckSlide,
    layout,
    *,
    page_number: int,
    sections: list[str],
    section_idx: int | None,
    add_fn,
) -> None:
    figure_path = ds.content.get("figure_path")
    if figure_path is None:
        raise ValueError(
            f"figure-kind DeckSlide '{ds.title}' missing figure_path; "
            "use build_deck_from_lineage_result for graceful fallback."
        )
    annotations = list(ds.content.get("annotations", []))
    motif_colors = ds.content.get("motif_colors") or {}
    caption = ds.content.get("caption", "")

    add_fn(
        pres,
        image_path=figure_path,
        annotations=annotations,
        title=ds.title,
        caption=caption,
        motif_colors=motif_colors,
        layout=layout,
        notes=ds.notes,
        page_number=page_number,
        sections=sections or None,
        current_section_idx=section_idx,
    )


def _add_references_slide(pres, ds: DeckSlide, layout) -> None:
    """Vancouver-style references in two columns."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    from vaultlab.slides.annotated_figure_slide import (
        text_color_for_theme,
    )

    s = pres.slides.add_slide(pres.slide_layouts[6])

    # Title
    tbox = s.shapes.add_textbox(
        Inches(0.4),
        Inches(0.15),
        Inches(layout.slide_w_in - 0.8),
        Inches(layout.title_h_in - 0.1),
    )
    tbox.name = "slide_title"
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ds.title or "References"
    r.font.name = "Arial"
    r.font.size = Pt(layout.title_font_pt)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    refs = list(ds.content.get("refs", []))
    if not refs:
        return

    # Two columns
    col_w = (layout.slide_w_in - 1.2) / 2.0
    col_h = layout.slide_h_in - layout.title_h_in - 0.5
    col_y = layout.title_h_in + 0.2

    half = (len(refs) + 1) // 2
    cols = [refs[:half], refs[half:]]

    for ci, col in enumerate(cols):
        cb = s.shapes.add_textbox(
            Inches(0.6 + ci * (col_w + 0.2)),
            Inches(col_y),
            Inches(col_w),
            Inches(col_h),
        )
        cb.name = f"refs_col_{ci}"
        ctf = cb.text_frame
        ctf.word_wrap = True
        for i, ref in enumerate(col):
            para = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            n = ref.get("n", "?")
            citation = ref.get("citation", "")
            doi = ref.get("doi", "")
            text = f"[{n}] {citation}"
            if doi:
                text += f"  doi:{doi}"
            run = para.add_run()
            run.text = text
            run.font.name = "Arial"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    # Bug 6: og_score methodology footer on the References slide.
    footer_box = s.shapes.add_textbox(
        Inches(0.6),
        Inches(layout.slide_h_in - 0.4),
        Inches(layout.slide_w_in - 1.2),
        Inches(0.3),
    )
    footer_box.name = "refs_methodology_footer"
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = (
        "og_score: Kessler 1963 bibliographic coupling — fraction of seed "
        "papers that cite each candidate. See vaultlab/docs/methodology.md."
    )
    fr.font.name = "Arial"
    fr.font.size = Pt(8)
    fr.font.italic = True
    fr.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))

    if ds.notes:
        s.notes_slide.notes_text_frame.text = ds.notes


# ---------------------------------------------------------------------------
# Internal: lineage-result -> DeckPlan
# ---------------------------------------------------------------------------


def _plan_from_lineage(
    lineage_result: LineageRunResult,
    *,
    speaker: str,
    affiliation: str,
    figure_assignments: dict[str, Path],
) -> DeckPlan:
    """Pure-Python composer (no LLM call) that turns a lineage result into a DeckPlan."""
    from datetime import date as _date

    from vaultlab.slides.notes import dual_format

    # Read summaries from disk if they exist; fall back gracefully when not.
    summaries = _read_summary_frontmatters(lineage_result.summary_paths)
    bucketed = _bucket_by_year(summaries)
    bucketed = _fill_empty_buckets(bucketed, summaries)
    arc_text = _read_arc_narrative(lineage_result.arc_path)

    sections = ["Background", "Development", "Synthesis", "References"]
    slides: list[DeckSlide] = []
    cited_dois: set[str] = set()  # tracks every DOI surfaced anywhere in the deck

    # 1. Title
    slides.append(
        DeckSlide(
            kind="title",
            title=f"Lineage: {lineage_result.topic}",
            content={
                "subtitle": "Journal-club deck",
                "speaker": speaker,
                "affiliation": affiliation,
                "date": _date.today().isoformat(),
            },
            notes=dual_format(
                mental_map={
                    "hook": f"Today we trace the lineage of {lineage_result.topic}.",
                    "key_claim": (
                        f"This corpus has {lineage_result.corpus_size} papers; "
                        f"we'll walk through history, development, and the SOTA."
                    ),
                    "transition": "Let's start with the foundational work.",
                },
                detailed_script=(
                    f"Hello, I'm {speaker}, and today I'm presenting a "
                    f"lineage view of {lineage_result.topic}. The corpus assembled "
                    f"by /lit-arc covers {lineage_result.corpus_size} papers, "
                    f"of which {lineage_result.pdfs_acquired} have full-text "
                    f"PDFs available. We'll move from history through "
                    f"development to the state of the art."
                ),
            ),
        )
    )

    # 2. Section intro: Background
    history_papers = bucketed.get("history", [])
    history_bullets = [_one_line_label(s) for s in history_papers[:3]]
    for s in history_papers[:3]:
        if s.get("doi"):
            cited_dois.add(s["doi"])
    slides.append(
        DeckSlide(
            kind="section_intro",
            title="Background",
            content={
                "section_name": "Background",
                "key_question": (f"What foundational work established {lineage_result.topic}?"),
                # Bobby 2026-04-30: never ship the placeholder text — bucket
                # fallback above guarantees history_papers is non-empty when
                # *any* corpus papers exist, but fall back to a generic
                # framing line if the entire summaries dict is empty.
                "bullets": history_bullets
                or [f"Establishing the foundations of {lineage_result.topic}."],
            },
            notes=dual_format(
                mental_map={
                    "key_claim": "These are the OG papers everyone cites.",
                    "transition": "Now let's see how the field evolved.",
                }
            ),
        )
    )

    # Fix 2 (2026-04-30 evening-4): allocate a global figure budget across
    # buckets BEFORE we start emitting slides. Previously the deck only
    # ever produced ONE figure-slide (history bucket leader). Now every
    # Tier-A paper that has a representative figure can become a figure
    # slide, capped globally at _FIGURE_TOTAL_CAP (default 8) so the deck
    # doesn't blow out.
    development_papers = bucketed.get("development", [])
    sota_papers = bucketed.get("sota", [])
    figure_budget = _allocate_figure_budget(
        history_papers=history_papers,
        development_papers=development_papers,
        sota_papers=sota_papers,
        figure_assignments=figure_assignments,
        summaries=summaries,
    )

    def _emit_figure_slide(
        pick: tuple[str, str, Path],
        *,
        bucket_title: str,
    ) -> None:
        """Append a figure slide for ``pick`` and update ``cited_dois``.

        Centralised so the same caption + substitution + manifest logic
        runs for every bucket (history / development / sota), not just
        history (which used to be the only path).
        """
        claim_doi, fig_doi, fig_path = pick
        cited_dois.add(fig_doi)
        if claim_doi and claim_doi != fig_doi:
            cited_dois.add(claim_doi)
        fig_label = _label_for_doi(summaries, fig_doi)
        is_substituted = bool(claim_doi) and claim_doi != fig_doi
        if is_substituted:
            manifest_entry = _read_figures_manifest_for(fig_doi, fig_path)
            caption = _compose_substitution_caption(
                summaries.get(fig_doi),
                fig_doi,
                figure_label=(manifest_entry or {}).get("label", "") or "",
                figure_caption=(manifest_entry or {}).get("caption", "") or "",
            )
        else:
            caption_summary = _summary_caption(summaries.get(fig_doi))
            caption_parts = [fig_label]
            if caption_summary:
                caption_parts.append(caption_summary)
            caption = " — ".join(caption_parts)
        slides.append(
            DeckSlide(
                kind="figure",
                title=bucket_title,
                content={
                    "figure_path": fig_path,
                    "annotations": [],
                    "motif_colors": {},
                    "caption": caption,
                    "citation_doi": claim_doi or fig_doi,
                    "claim_paper_doi": claim_doi or fig_doi,
                    "figure_paper_doi": fig_doi,
                },
                notes=dual_format(
                    mental_map={
                        "hook": "Look at this canonical figure.",
                        "evidence": (
                            f"Substituted figure from {fig_label}"
                            if is_substituted
                            else f"Figure from {fig_label}"
                        ),
                    }
                ),
            )
        )

    # 3. Figure(s) or bullets — history bucket
    # Bobby 2026-04-30 (Bug #5): figure-slide subjects must be Tier-A
    # papers (LLM-read full text). Tier-C stubs without summaries leak
    # off-topic figures (e.g. Gjerstorff 2006 cancer figure on a spatial-tx
    # deck) when they happen to have a cached PMC figure.
    history_picks = figure_budget.get("history", [])
    if history_picks:
        for pick in history_picks:
            _emit_figure_slide(pick, bucket_title="Foundational findings")
    else:
        # Drop the figure slide; replace with bullets from history TL;DRs.
        # Bug #4: never ship "(no history-bucket summaries available)" —
        # _fill_empty_buckets above guaranteed at least one history paper
        # if the corpus is non-empty.
        history_bullet_dois = [s.get("doi", "") for s in history_papers[:5] if s.get("doi")]
        cited_dois.update(history_bullet_dois)
        slides.append(
            DeckSlide(
                kind="bullets",
                title="Foundational findings",
                content={
                    "bullets": [_bullet_from_summary(s) for s in history_papers[:5]]
                    or [
                        f"No prior work catalogued in this corpus for {lineage_result.topic}.",
                    ],
                    "citations": history_bullet_dois,
                },
                notes=dual_format(
                    mental_map={
                        "key_claim": "These are the foundational findings.",
                    }
                ),
            )
        )

    # 4. Section intro: Development
    development_bullets = [_one_line_label(s) for s in development_papers[:3]]
    for s in development_papers[:3]:
        if s.get("doi"):
            cited_dois.add(s["doi"])
    slides.append(
        DeckSlide(
            kind="section_intro",
            title="Development",
            content={
                "section_name": "Development",
                "key_question": "How did the field evolve?",
                "bullets": development_bullets
                or [
                    f"Tracing how {lineage_result.topic} developed.",
                ],
            },
        )
    )

    # 4b. Figure(s) — development bucket (Fix 2, 2026-04-30 evening-4).
    # Each Tier-A development paper that has a cached figure gets its own
    # slide, subject to the global cap.
    for pick in figure_budget.get("development", []):
        _emit_figure_slide(pick, bucket_title="Development milestone")

    # 5. Bullets: SOTA findings
    sota_bullets = [_bullet_from_summary(s) for s in sota_papers[:5]]
    sota_dois = [s.get("doi", "") for s in sota_papers[:5] if s.get("doi")]
    cited_dois.update(sota_dois)
    slides.append(
        DeckSlide(
            kind="bullets",
            title="State of the art",
            content={
                "bullets": sota_bullets
                or [
                    f"Current state of the art for {lineage_result.topic}.",
                ],
                "citations": sota_dois,
            },
            notes=dual_format(
                mental_map={
                    "key_claim": "Here's where the field stands today.",
                    "transition": "Let's pull it together.",
                }
            ),
        )
    )

    # 5b. Figure(s) — SOTA bucket (Fix 2, 2026-04-30 evening-4).
    for pick in figure_budget.get("sota", []):
        _emit_figure_slide(pick, bucket_title="State-of-the-art result")

    # 6. Section intro: Synthesis (use arc narrative when present)
    # Bobby 2026-04-30 (Bug #3): the synthesis slide was previously
    # rendering YAML frontmatter ('topic: ... | date: ... | seeds: 12')
    # as bullets because _arc_bullets walked the file top-down. Now we
    # prefer a "## Synthesis" heading inside the arc, then fall back to
    # the last narrative paragraph.
    synthesis_bullets = _synthesis_bullets_from_arc(arc_text) if arc_text else []
    if not synthesis_bullets:
        synthesis_bullets = [
            "Synthesis pending — re-run /lit-arc with ANTHROPIC_API_KEY for a full narrative."
        ]
    slides.append(
        DeckSlide(
            kind="section_intro",
            title="Synthesis",
            content={
                "section_name": "Synthesis",
                "key_question": "What's the through-line?",
                "bullets": synthesis_bullets,
            },
        )
    )

    # 7. References — Bobby 2026-04-30 (Bug #2): only the DOIs *cited
    # somewhere in this deck*, not the entire 610-paper corpus.
    refs = _build_references(summaries, cited_dois=cited_dois)
    slides.append(
        DeckSlide(
            kind="references",
            title="References",
            content={"refs": refs},
        )
    )

    return DeckPlan(
        title=f"Lineage: {lineage_result.topic}",
        subtitle="Journal-club deck",
        speaker=speaker,
        affiliation=affiliation,
        sections=sections,
        slides=slides,
        theme="hickey_lab",
    )


# ---------------------------------------------------------------------------
# Internal: read summary frontmatter / arc text
# ---------------------------------------------------------------------------


def _read_summary_frontmatters(
    summary_paths: dict[str, Path],
) -> dict[str, dict[str, Any]]:
    """Parse YAML frontmatter from each Wiki/Summaries/<doi>.md file.

    Returns ``{doi: {fm_dict, plus 'tldr', 'key_findings' parsed from body}}``.
    Missing files are skipped silently — the caller falls back to bullets.
    """
    import yaml

    out: dict[str, dict[str, Any]] = {}
    for doi, path in summary_paths.items():
        if not path or not path.exists():
            continue
        try:
            text = path.read_text(encoding="utf-8")
        except OSError:
            continue
        # Frontmatter block between '---' delimiters
        if not text.startswith("---"):
            out[doi] = {"doi": doi}
            continue
        end = text.find("\n---", 3)
        if end == -1:
            out[doi] = {"doi": doi}
            continue
        try:
            fm = yaml.safe_load(text[3:end]) or {}
        except yaml.YAMLError:
            fm = {}
        body = text[end + 4 :]
        # Extract TL;DR (first non-empty line under "## TL;DR")
        tldr = _extract_section(body, "## TL;DR")
        findings = _extract_bullet_section(body, "## Key findings")
        rec = {"doi": doi, **fm, "tldr": tldr, "key_findings": findings}
        out[doi] = rec
    return out


def _extract_section(body: str, heading: str) -> str:
    """Return text under ``heading`` until the next ``## `` heading, stripped."""
    import re

    pat = re.compile(
        rf"^{re.escape(heading)}\s*\n(.*?)(?=^## |\Z)",
        re.MULTILINE | re.DOTALL,
    )
    m = pat.search(body)
    if not m:
        return ""
    return m.group(1).strip()


def _extract_bullet_section(body: str, heading: str) -> list[str]:
    text = _extract_section(body, heading)
    if not text:
        return []
    return [
        line[2:].strip()
        for line in text.splitlines()
        if line.startswith("- ") and "_(none)_" not in line and "_(empty)_" not in line
    ][:5]


def _bucket_by_year(
    summaries: dict[str, dict[str, Any]],
) -> dict[str, list[dict[str, Any]]]:
    out: dict[str, list[dict[str, Any]]] = {
        "history": [],
        "development": [],
        "sota": [],
        "unknown": [],
    }
    for s in summaries.values():
        bucket = s.get("year_bucket", "unknown")
        out.setdefault(bucket, []).append(s)
    # Sort each bucket by year ascending
    for k in out:
        out[k].sort(key=lambda d: d.get("year", 0) or 0)
    return out


def _read_arc_narrative(arc_path: Path) -> str:
    if not arc_path or not arc_path.exists():
        return ""
    try:
        return arc_path.read_text(encoding="utf-8")
    except OSError:
        return ""


def _arc_bullets(arc_text: str) -> list[str]:
    """Backwards-compat alias kept for any external callers.

    New code should use :func:`_synthesis_bullets_from_arc` instead — it
    skips YAML frontmatter and prefers an explicit ``## Synthesis``
    heading.
    """
    return _synthesis_bullets_from_arc(arc_text)


def _strip_yaml_frontmatter(text: str) -> str:
    """Drop a leading ``---\\n...\\n---`` YAML block from arc-narrative text.

    The original ``_arc_bullets`` walked from line 0 and ended up taking
    the YAML keys (``topic: ...``, ``date: ...``, ``seeds: 12``) as
    bullets. Stripping the frontmatter at parse time prevents that.
    """
    if not text.startswith("---"):
        return text
    end = text.find("\n---", 3)
    if end == -1:
        return text
    return text[end + 4 :].lstrip("\n")


def _synthesis_bullets_from_arc(arc_text: str) -> list[str]:
    """Extract 3-5 synthesis bullets from a lineage-arc markdown.

    Order of preference:

    1. An explicit ``## Synthesis`` heading — split its body into sentences.
    2. The *last* paragraph of the arc (typically the most synthesizing) —
       split into sentences.
    3. A first-paragraph sentence walk as a final fallback.

    YAML frontmatter is stripped before parsing so we never emit
    ``topic: ... | date: ... | seeds: 12`` as bullets.
    """
    if not arc_text:
        return []
    body = _strip_yaml_frontmatter(arc_text)

    # 1. Explicit "## Synthesis" section -------------------------------------
    synthesis = _extract_section(body, "## Synthesis")
    if synthesis:
        sentences = _sentences_from_paragraph(synthesis)
        if sentences:
            return sentences[:5]

    # 2. Last narrative paragraph ----------------------------------------------
    paragraphs = _narrative_paragraphs(body)
    if paragraphs:
        last = paragraphs[-1]
        sentences = _sentences_from_paragraph(last)
        if sentences:
            return sentences[:5]

    # 3. First-paragraph fallback (similar to legacy behaviour but with the
    #    YAML and bullet-line guards in place) -------------------------------
    bullets: list[str] = []
    for line in body.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith(("#", ">", "|", "---", "- ", "* ")):
            continue
        # Skip simple ``key: value`` lines that survived frontmatter stripping
        if _looks_like_yaml_kv_line(stripped):
            continue
        first_period = stripped.find(". ")
        if first_period > 30:
            bullets.append(stripped[: first_period + 1])
        else:
            bullets.append(stripped[:200])
        if len(bullets) >= 3:
            break
    return bullets


def _looks_like_yaml_kv_line(line: str) -> bool:
    """Heuristic: 'topic: foo' or 'seeds: 12' looking lines."""
    if ":" not in line:
        return False
    head, _, _ = line.partition(":")
    head = head.strip()
    if not head:
        return False
    # Short alphanum + underscore key with no whitespace -> looks YAML-ish.
    return len(head) <= 32 and " " not in head and head.replace("_", "").replace("-", "").isalnum()


def _narrative_paragraphs(body: str) -> list[str]:
    """Return list of non-empty narrative paragraphs (paragraph = >=1 prose line).

    Skips: headings, blockquotes, table rows, bullet lists, fence lines.
    """
    paragraphs: list[str] = []
    current: list[str] = []
    for line in body.splitlines():
        stripped = line.strip()
        if not stripped:
            if current:
                paragraphs.append(" ".join(current).strip())
                current = []
            continue
        if stripped.startswith(("#", ">", "|", "---", "```", "- ", "* ")):
            if current:
                paragraphs.append(" ".join(current).strip())
                current = []
            continue
        if _looks_like_yaml_kv_line(stripped):
            continue
        current.append(stripped)
    if current:
        paragraphs.append(" ".join(current).strip())
    return [p for p in paragraphs if p]


def _sentences_from_paragraph(paragraph: str) -> list[str]:
    """Split a paragraph into 1-5 reasonable-length sentences."""
    import re

    parts = re.split(r"(?<=[.!?])\s+", paragraph.strip())
    out: list[str] = []
    for p in parts:
        p = p.strip()
        if not p:
            continue
        # Trim very long ones to keep slide bullets readable
        out.append(p if len(p) <= 220 else p[:217] + "...")
    return out


def _one_line_label(summary: dict[str, Any]) -> str:
    authors = summary.get("authors") or []
    last = ""
    for a in authors:
        last = _format_author_lastname(a)
        if last:
            break
    if not last:
        last = "Anon"
        if authors:
            logger.warning(
                "no parseable author in %s; falling back to 'Anon'",
                summary.get("doi") or "(unknown doi)",
            )
    year = summary.get("year") or "n.d."
    title = summary.get("title", "") or summary.get("doi", "")
    short = title if len(title) <= 70 else title[:67] + "..."
    return f"{last} {year} — {short}"


def _label_for_doi(summaries: dict[str, dict[str, Any]], doi: str) -> str:
    s = summaries.get(doi)
    if not s:
        return doi
    return _one_line_label(s)


def _summary_caption(summary: dict[str, Any] | None) -> str:
    if not summary:
        return ""
    tldr = (summary.get("tldr") or "").strip()
    if tldr:
        return tldr.split("\n")[0][:200]
    return ""


def _bullet_from_summary(summary: dict[str, Any]) -> str:
    findings = summary.get("key_findings") or []
    if findings:
        return findings[0][:200]
    tldr = (summary.get("tldr") or "").strip()
    if tldr:
        return tldr.split("\n")[0][:200]
    return _one_line_label(summary)


def _pick_figure_for_bucket(
    bucket_papers: list[dict[str, Any]],
    figure_assignments: dict[str, Path],
    *,
    summaries: dict[str, dict[str, Any]] | None = None,
) -> tuple[str, str, Path] | None:
    """Pick a Tier-A figure for a bucket, with substitution support.

    Returns ``(claim_doi, figure_doi, figure_path)`` where:

    - ``claim_doi`` is the bucket leader (``bucket_papers[0]``) — the
      paper whose claim the slide is *about*.
    - ``figure_doi`` is the paper whose figure is actually shown; equals
      ``claim_doi`` when the leader has its own figure, or the DOI of a
      later Tier-A bucket paper whose figure substitutes in.
    - ``figure_path`` is the on-disk figure file.

    Bobby 2026-04-30 (Bug #5): filter to Tier-A only. Tier-C papers are
    abstract-only stubs without LLM-read content; using one as a figure-slide
    foundational paper put off-topic figures (e.g. Gjerstorff 2006 cancer
    figure on a spatial-tx deck) into the deck whenever the corpus had a
    cached PMC figure for the wrong DOI.

    Bobby 2026-04-30 (figure substitution): when the leader has no
    figure, walk down the bucket and substitute the first Tier-A paper
    that does. The renderer composes a "Substituted figure from <author>
    <year>" caption so the audience knows the figure source differs from
    the claim source. If no Tier-A paper in the bucket has a figure,
    return ``None`` so the caller falls back to a bullets slide.
    """
    summaries = summaries if summaries is not None else {}
    if not bucket_papers:
        return None
    claim_doi = (bucket_papers[0].get("doi") or "").strip()
    for s in bucket_papers:
        doi = (s.get("doi") or "").strip()
        if not doi:
            continue
        # Honor the tier from the original bucket_papers entry first; fall
        # back to the summaries dict when the caller passed bucketed
        # summaries that already include the tier field.
        tier = (s.get("tier") or summaries.get(doi, {}).get("tier") or "").upper()
        if tier and tier != "A":
            continue
        path = figure_assignments.get(doi)
        if path is not None and Path(path).exists():
            return claim_doi or doi, doi, Path(path)
    return None


# Fix 2 (2026-04-30 evening-4): aggressive figure picker
#
# The original ``_pick_figure_for_bucket`` returned at most ONE figure per
# bucket, so a corpus with 43 cached figures across 3 Tier-A papers
# produced AT MOST 3 figure-slides (one per bucket) — and in practice only
# 1 (since most bucket leaders had no figure assignment). The new
# ``_pick_figures_for_bucket_multi`` walks every Tier-A paper in the
# bucket, picks the largest available figure per paper (skipping decorative
# crops < 100 KB), and returns a list of figure picks. The caller decides
# how many to render based on the global cap.

_FIGURE_MIN_BYTES = 100 * 1024  # 100 KB — skip decorative crops/labels
_FIGURE_TOTAL_CAP = 8  # global ceiling on figure-slides per deck


def _figure_size_bytes(path: Path) -> int:
    """Return file size in bytes, or 0 if the file is missing."""
    try:
        return Path(path).stat().st_size
    except OSError:
        return 0


def _list_cached_figures_for_doi(
    doi: str,
    figure_assignments: dict[str, Path],
) -> list[Path]:
    """Return all cached figures on disk for ``doi``.

    The ``figure_assignments`` map only carries ONE Path per DOI (the
    "default pick" the orchestrator chose). To pick the LARGEST figure
    we also enumerate sibling files in the same directory and consult
    the ``.figures.json`` manifest when present.
    """
    seed = figure_assignments.get(doi)
    if seed is None:
        return []
    seed = Path(seed)
    parent = seed.parent
    if not parent.exists():
        return [seed] if seed.exists() else []
    # Prefer the manifest order (only files that the acquirer actually
    # extracted as figures, not stray PDFs/XML in the same dir).
    manifest = parent / ".figures.json"
    figs: list[Path] = []
    if manifest.exists():
        try:
            data = json.loads(manifest.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            data = {}
        for entry in data.get("figures", []):
            fp = Path(entry.get("file_path") or "")
            if fp.exists():
                figs.append(fp)
    if not figs:
        # No manifest — fall back to scanning common image extensions.
        exts = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".gif"}
        for child in sorted(parent.iterdir()):
            if child.is_file() and child.suffix.lower() in exts:
                figs.append(child)
    if seed not in figs and seed.exists():
        figs.append(seed)
    return figs


def _pick_largest_figure_for_doi(
    doi: str,
    figure_assignments: dict[str, Path],
    *,
    min_size_bytes: int = _FIGURE_MIN_BYTES,
) -> Path | None:
    """Pick the largest figure on disk for ``doi``.

    Logic:

    * When multiple figures are cached for the DOI, prefer the largest
      figure that is at least ``min_size_bytes`` (filters decorative
      panel labels / cropped legends < 100 KB).
    * When ZERO figures meet the size threshold, fall back to the largest
      cached figure overall — better one small figure than no figure
      (Tier-A papers have something to show even when only a small
      thumbnail was extracted).
    * When no figures are cached at all for the DOI, returns ``None``.
    """
    figs = _list_cached_figures_for_doi(doi, figure_assignments)
    if not figs:
        return None
    sized = [(f, _figure_size_bytes(f)) for f in figs]
    above_threshold = [t for t in sized if t[1] >= min_size_bytes]
    if above_threshold:
        above_threshold.sort(key=lambda t: t[1], reverse=True)
        return above_threshold[0][0]
    # Soft fallback — use whatever's cached.
    sized.sort(key=lambda t: t[1], reverse=True)
    return sized[0][0]


def _pick_figures_for_bucket_multi(
    bucket_papers: list[dict[str, Any]],
    figure_assignments: dict[str, Path],
    *,
    summaries: dict[str, dict[str, Any]] | None = None,
    max_per_bucket: int = 4,
    min_size_bytes: int = _FIGURE_MIN_BYTES,
) -> list[tuple[str, str, Path]]:
    """Pick MULTIPLE figures for a bucket — one per Tier-A paper that has one.

    Returns a list of ``(claim_doi, figure_doi, figure_path)`` tuples.
    The first entry pairs the bucket leader (claim paper) with the best
    Tier-A figure (its own when available, otherwise a substitute). Each
    subsequent Tier-A paper that has a representative figure produces an
    additional entry where ``claim_doi == figure_doi``.

    Tier-C papers are NEVER used as figure-slide subjects (Bug #5).
    Figures < ``min_size_bytes`` are skipped (likely decorative crops or
    panel labels rather than the main result figure).

    Cap with ``max_per_bucket`` — when more Tier-A papers than slots
    have a figure, the highest-og_score papers win.
    """
    summaries = summaries if summaries is not None else {}
    if not bucket_papers:
        return []
    leader_row = bucket_papers[0]
    claim_doi = (leader_row.get("doi") or "").strip()
    leader_tier = (leader_row.get("tier") or summaries.get(claim_doi, {}).get("tier") or "").upper()
    leader_is_tier_a = leader_tier in ("", "A")

    # Collect Tier-A papers in the bucket along with the largest figure
    # available for each (or None when the paper has no usable figure).
    tier_a: list[tuple[str, dict[str, Any], Path | None]] = []
    for s in bucket_papers:
        doi = (s.get("doi") or "").strip()
        if not doi:
            continue
        tier = (s.get("tier") or summaries.get(doi, {}).get("tier") or "").upper()
        if tier and tier != "A":
            continue
        fig = _pick_largest_figure_for_doi(
            doi,
            figure_assignments,
            min_size_bytes=min_size_bytes,
        )
        tier_a.append((doi, s, fig))

    if not tier_a:
        return []

    # Order: bucket leader first (when leader IS Tier-A), then remaining
    # Tier-A papers by og_score descending so when we exceed
    # max_per_bucket the most influential papers' figures win.
    leader = next((t for t in tier_a if t[0] == claim_doi), None)
    rest = [t for t in tier_a if t[0] != claim_doi]
    rest.sort(
        key=lambda t: float(t[1].get("og_score") or summaries.get(t[0], {}).get("og_score") or 0.0),
        reverse=True,
    )

    picks: list[tuple[str, str, Path]] = []

    # 1. Leader slot.
    if leader_is_tier_a and leader is not None and leader[2] is not None:
        # Leader is Tier-A and has its own figure.
        picks.append((leader[0], leader[0], leader[2]))
    elif leader_is_tier_a and leader is not None:
        # Leader is Tier-A but has no figure of its own — substitute
        # from the highest-ranked Tier-A paper that does.
        sub = next((t for t in rest if t[2] is not None), None)
        if sub is not None:
            picks.append((leader[0], sub[0], sub[2]))
            rest = [t for t in rest if t[0] != sub[0]]
    else:
        # Leader is Tier-C (or non-A) — keep claim_doi pointing at the
        # leader (the slide is still ABOUT the leader paper) but the
        # figure MUST come from a Tier-A bucket member. This preserves
        # the contract from the original ``_pick_figure_for_bucket`` so
        # _plan_from_lineage still surfaces the leader's claim while
        # showing a topical figure.
        sub = next((t for t in rest if t[2] is not None), None)
        if sub is None:
            # See if a Tier-A paper IS in tier_a but happened to be
            # ordered first (claim_doi == first Tier-A doi).
            sub = next((t for t in tier_a if t[2] is not None), None)
            if sub is not None:
                # Drop it from rest so it doesn't double-count.
                rest = [t for t in rest if t[0] != sub[0]]
        else:
            rest = [t for t in rest if t[0] != sub[0]]
        if sub is not None:
            picks.append((claim_doi, sub[0], sub[2]))

    # 2. Additional slots — one per remaining Tier-A paper with a figure,
    # capped at max_per_bucket.
    for doi, _row, fig in rest:
        if fig is None:
            continue
        if len(picks) >= max_per_bucket:
            break
        picks.append((doi, doi, fig))

    return picks


def _allocate_figure_budget(
    *,
    history_papers: list[dict[str, Any]],
    development_papers: list[dict[str, Any]],
    sota_papers: list[dict[str, Any]],
    figure_assignments: dict[str, Path],
    summaries: dict[str, dict[str, Any]],
    total_cap: int = _FIGURE_TOTAL_CAP,
    min_size_bytes: int = _FIGURE_MIN_BYTES,
) -> dict[str, list[tuple[str, str, Path]]]:
    """Distribute a global figure-slide budget across the three buckets.

    Returns ``{"history": [...], "development": [...], "sota": [...]}``.
    When the per-bucket maximum across all buckets exceeds ``total_cap``,
    we trim from buckets with more papers proportionally so the deck
    doesn't blow out (e.g. 5 history + 5 dev + 5 sota = 15 -> capped at
    8 = roughly 3+3+2).
    """
    # Compute the maximum each bucket COULD produce.
    raw = {
        "history": _pick_figures_for_bucket_multi(
            history_papers,
            figure_assignments,
            summaries=summaries,
            max_per_bucket=total_cap,  # we'll cap at the end
            min_size_bytes=min_size_bytes,
        ),
        "development": _pick_figures_for_bucket_multi(
            development_papers,
            figure_assignments,
            summaries=summaries,
            max_per_bucket=total_cap,
            min_size_bytes=min_size_bytes,
        ),
        "sota": _pick_figures_for_bucket_multi(
            sota_papers,
            figure_assignments,
            summaries=summaries,
            max_per_bucket=total_cap,
            min_size_bytes=min_size_bytes,
        ),
    }
    total_available = sum(len(v) for v in raw.values())
    if total_available <= total_cap:
        return raw

    # Need to trim. Round-robin removal from the bucket with the most
    # remaining picks, working from the LOWEST-og_score figure forward
    # (so the figure-with-the-highest-OG paper survives).
    # Implementation: rebuild each bucket sorted by og_score descending,
    # then pop from the end of whichever bucket is currently largest
    # until the total fits within total_cap.
    def _og(doi: str) -> float:
        return float((summaries.get(doi) or {}).get("og_score") or 0.0)

    out = {b: sorted(picks, key=lambda t: _og(t[1]), reverse=True) for b, picks in raw.items()}
    while sum(len(v) for v in out.values()) > total_cap:
        # Pick the bucket with the most picks; ties broken by alphabetic
        # bucket order for determinism.
        bucket_name = max(out, key=lambda b: (len(out[b]), -ord(b[0])))
        if not out[bucket_name]:
            break
        out[bucket_name].pop()  # drop the lowest-og_score pick in that bucket
    return out


def _read_figures_manifest_for(
    figure_doi: str,
    figure_path: Path,
) -> dict[str, Any] | None:
    """Read the ``.figures.json`` manifest sitting alongside a figure file.

    Returns the matching figure entry (``{figure_id, file_path, caption,
    label, panels}``) if the path appears in the manifest; ``None``
    otherwise.  Used to recover the original Elsevier/PMC caption when
    composing a substitution caption — the figure paper's TL;DR may not
    describe the specific panel being shown.
    """
    import json as _json

    del figure_doi  # the manifest lives in figure_path.parent
    manifest = Path(figure_path).parent / ".figures.json"
    if not manifest.exists():
        return None
    try:
        data = _json.loads(manifest.read_text(encoding="utf-8"))
    except (OSError, _json.JSONDecodeError):
        return None
    target = str(figure_path)
    for f in data.get("figures", []):
        if f.get("file_path") == target:
            return f
    return None


def _compose_substitution_caption(
    figure_summary: dict[str, Any] | None,
    figure_doi: str,
    *,
    figure_label: str = "",
    figure_caption: str = "",
) -> str:
    """Build the caption for a slide whose figure was substituted from another paper.

    The reader needs to know the figure isn't from the paper the slide is
    about, so the caption is prefixed with ``Substituted figure from
    [[<doi-slug>|<author> <year>]]: ...`` where the wikilink resolves to
    the substituted paper's summary.  The body text prefers the figure
    paper's summary TL;DR; when that's missing (Tier-C source) it falls
    back to the figure's own caption from the ``.figures.json`` manifest.
    """
    from vaultlab.kb.paths import slugify_doi

    label = _one_line_label(figure_summary) if figure_summary else figure_doi
    slug = slugify_doi(figure_doi) if figure_doi else ""
    body = ""
    if figure_summary:
        tldr = (figure_summary.get("tldr") or "").strip()
        if tldr:
            body = tldr.split("\n")[0][:200]
    if not body:
        if figure_label and figure_caption:
            body = f"{figure_label}: {figure_caption}"[:240]
        elif figure_caption:
            body = figure_caption[:240]
        else:
            body = "(figure source has limited summary content)"
    if slug:
        prefix = f"Substituted figure from [[{slug}|{label}]]"
    else:
        prefix = f"Substituted figure from {label}"
    return f"{prefix}: {body}"


def _bullets_from_substituted_figure(
    claim_summary: dict[str, Any] | None,
    figure_summary: dict[str, Any] | None,
    *,
    n: int = 5,
) -> list[str]:
    """Mix bullets from claim paper (~60%) and figure paper (~40%).

    When a figure is substituted, the slide's bullets should incorporate
    findings from the figure paper too — the slide is "figure-based" and
    the figure is now from a different paper.  When the figure paper is
    Tier-C (no key_findings), we degrade gracefully to claim-only
    bullets.
    """
    claim_findings = list((claim_summary or {}).get("key_findings") or [])
    fig_findings = list((figure_summary or {}).get("key_findings") or [])
    if not fig_findings:
        if claim_findings:
            return [b[:200] for b in claim_findings[:n]]
        tldr = ((claim_summary or {}).get("tldr") or "").strip()
        return [tldr.split("\n")[0][:200]] if tldr else []
    n_claim = max(1, (n * 3 + 4) // 5)  # ~ceil(n*0.6)
    n_fig = max(0, n - n_claim)
    out: list[str] = []
    for b in claim_findings[:n_claim]:
        out.append(b[:200])
    for b in fig_findings[:n_fig]:
        out.append(b[:200])
    return out[:n]


def _fill_empty_buckets(
    bucketed: dict[str, list[dict[str, Any]]],
    summaries: dict[str, dict[str, Any]],
) -> dict[str, list[dict[str, Any]]]:
    """Fill empty year-buckets so the deck never ships placeholder text.

    Bobby 2026-04-30 (Bug #4): when the year-bucket quartile algorithm
    puts 0 papers in ``history`` (because the corpus tail is older than
    the topic's seminal era) the deck shipped a hard-coded
    ``"(no history-bucket papers in corpus)"`` line. This helper falls
    back to the oldest N papers (history) / newest N papers (sota) /
    middle-year papers (development) from *all* summaries when a bucket
    is empty.
    """
    out = {k: list(v) for k, v in bucketed.items()}
    by_year = sorted(
        (s for s in summaries.values() if s.get("doi")),
        key=lambda d: d.get("year", 0) or 0,
    )
    if not by_year:
        return out

    n_fill = 5

    if not out.get("history"):
        # Oldest N
        out["history"] = by_year[:n_fill]
    if not out.get("sota"):
        out["sota"] = list(reversed(by_year[-n_fill:]))
    if not out.get("development"):
        # Pick from the middle of the year range; if the corpus is too small
        # just reuse history+sota as the dev bucket.
        if len(by_year) >= 3:
            mid_lo = len(by_year) // 4
            mid_hi = max(mid_lo + 1, (3 * len(by_year)) // 4)
            out["development"] = by_year[mid_lo:mid_hi][:n_fill]
        else:
            out["development"] = list(by_year)
    return out


def _build_references(
    summaries: dict[str, dict[str, Any]],
    cited_dois: set[str] | None = None,
) -> list[dict[str, Any]]:
    """Build the deck-references list.

    Pre-2026-04-30 this function returned every paper in the corpus
    (610 papers for spatial-tx) which made the references slide
    illegible — the deck only actually cites ~10 of them. The
    ``cited_dois`` arg now restricts the output to DOIs referenced
    *somewhere* in the deck (bullets-slide citations, figure-slide
    citation_doi, etc.). When ``cited_dois`` is None we keep the legacy
    behaviour for back-compat with callers that don't track citations.
    """
    items = summaries.items()
    if cited_dois is not None:
        cited_set = {d for d in cited_dois if d}
        items = [(d, s) for d, s in items if d in cited_set]
    refs: list[dict[str, Any]] = []
    for n, (doi, s) in enumerate(
        sorted(
            items,
            key=lambda kv: (kv[1].get("year", 0) or 0, kv[0]),
        ),
        start=1,
    ):
        authors = s.get("authors") or []
        if authors:
            if len(authors) > 3:
                authors_part = ", ".join(authors[:3]) + ", et al."
            else:
                authors_part = ", ".join(authors)
        else:
            authors_part = "Anon."
        title = s.get("title", "") or "(no title)"
        year = s.get("year", "") or ""
        journal = s.get("journal", "") or ""
        citation = f"{authors_part} {title}. {journal} {year}.".strip()
        refs.append({"n": n, "citation": citation, "doi": doi})
    return refs


def _collect_citations_from_summaries(
    lineage_result: LineageRunResult,
) -> dict[str, dict[str, Any]]:
    summaries = _read_summary_frontmatters(lineage_result.summary_paths)
    return {
        doi: {
            "title": s.get("title", ""),
            "authors": s.get("authors", []) or [],
            "year": s.get("year", ""),
            "journal": s.get("journal", ""),
        }
        for doi, s in summaries.items()
    }


# ---------------------------------------------------------------------------
# Plan-driven builder (lifted from bobby_slides._builder.build_from_plan)
# ---------------------------------------------------------------------------
#
# This is the SECOND deck builder in this module. It accepts a flexible
# **dict** plan (rather than the typed :class:`DeckPlan` above) and
# dispatches each slide kind to the lifted imperative layout primitives in
# :mod:`vaultlab.slides.layouts`.
#
# When to use which:
#
# - :func:`build_deck` — when you have a typed :class:`DeckPlan` (used by
#   the lineage flow). Renders the 5-kind composer (title /
#   section_intro / figure / bullets / references) with the
#   annotated-figure-slide primitive.
#
# - :func:`build_from_plan` (this function) — when you have a richer dict
#   plan with heterogeneous slide types (figure / multi_figure / quote /
#   two_figure / section_divider / references / text / title) and want
#   lab-template imperative styling.  This is the path the **L4 deck
#   quality** work expects: an LLM-driven plan generator emits this dict
#   shape, ``build_from_plan`` renders it deterministically.

# Slide types the dict-plan builder understands.
SUPPORTED_PLAN_SLIDE_TYPES: frozenset[str] = frozenset(
    {
        "title",
        "section_divider",
        "figure",
        "two_figure",
        "quote",
        "multi_figure",
        "text",
        "references",
    }
)


def build_from_plan(
    plan: dict[str, Any],
    output: Path | str,
    write_marp: bool = True,
    kb_log: Any = None,
    with_animations: bool = False,
    theme: str | None = None,
    template: str = "lab",
) -> dict[str, Path]:
    """Render a structured deck-plan dict to ``.pptx`` (+ optional Marp .md).

    Lifted from ``bobby_slides._builder.build_from_plan`` (bobby-tools,
    2026-04). The dict-plan shape::

        {
            "title": "...",
            "author": "...",
            "subtitle": "...",       # optional
            "topic": "...",          # optional, used for KB logging
            "kb": "...",             # optional, KB name for log entry
            "theme": "dark"|"light", # optional, default "dark"
            "template": "lab"|"plain",  # optional, default "lab"
            "slides": [
                {"type": "title", ...},
                {"type": "section_divider", "title": "..."},
                {"type": "figure", "image_path": "...", ...,
                 "annotations": [...]},   # see vaultlab.slides.annotate
                {"type": "multi_figure", "figures": [...]},
                {"type": "two_figure", ...},
                {"type": "quote", "quote": "...", "attribution": "..."},
                {"type": "text", "title": "...", "bullets": [...]},
                {"type": "references", "references": [...]},
            ]
        }

    Each slide may carry ``"speaker_notes"`` (a mental-map dict — see
    :mod:`vaultlab.slides.notes`).  Bullets in figure slides may be plain
    strings OR dicts with embedded annotations (Option A); the builder
    normalizes them into slide-level annotations automatically.

    Args:
        plan: deck-plan dict (see schema).
        output: path to the ``.pptx`` output.
        write_marp: if ``True`` (default), also write a Marp ``.md`` mirror.
        kb_log: optional :class:`vaultlab.slides.kb_reader.KBReader` — when
            provided, appends an entry to ``<kb>/_Log.md`` and writes a
            deck-plan record to ``<kb>/Output/Reports/``.
        with_animations: if ``True``, auto-applies entrance animations to
            text slides (bullet-by-bullet reveal) and multi_figure slides
            (panel build-up). Title and section dividers stay static.
        theme: explicit ``"dark"`` / ``"light"``, overrides ``plan["theme"]``.
        template: ``"lab"`` (default — Hickey lab template) or ``"plain"``.

    Returns:
        Dict with keys ``"pptx"`` (Path) and optionally ``"marp"`` (Path)
        and ``"report"`` (Path).
    """
    # Lazy imports — keep the module light when callers only use the typed
    # DeckPlan path.
    from vaultlab.slides.animations import bullet_reveal, panel_buildup
    from vaultlab.slides.annotate import add_annotations
    from vaultlab.slides.layouts import (
        add_figure_above_bullets_slide,
        add_figure_only_slide,
        add_figure_slide,
        add_multi_figure_slide,
        add_quote_slide,
        add_references_slide,
        add_section_divider,
        add_text_slide,
        add_title_slide,
        add_two_figure_compare_slide,
    )
    from vaultlab.slides.notes import attach_to_slide
    from vaultlab.slides.template import load_plain_presentation, load_template

    chosen_theme = theme or plan.get("theme", "dark")
    chosen_template = plan.get("template", template)
    if chosen_template == "plain":
        pres = load_plain_presentation(theme=chosen_theme)
    else:
        pres = load_template(theme=chosen_theme)

    slides_plan = plan.get("slides", [])

    for slide_spec in slides_plan:
        slide_spec = _normalize_bullet_annotations(slide_spec)

        stype = slide_spec.get("type", "text")
        notes = slide_spec.get("speaker_notes")
        slide = None

        if stype == "title":
            slide = add_title_slide(
                pres,
                slide_spec.get("title", ""),
                subtitle=slide_spec.get("subtitle", ""),
                author=slide_spec.get("author", ""),
            )
        elif stype == "section_divider":
            slide = add_section_divider(pres, slide_spec.get("title", ""))
        elif stype == "figure":
            layout = slide_spec.get("layout", "default")
            if layout == "figure_only":
                slide = add_figure_only_slide(
                    pres,
                    image_path=slide_spec.get("image_path", ""),
                    title=slide_spec.get("title", ""),
                    caption=slide_spec.get("caption", ""),
                    citation_source=slide_spec.get("citation_source", ""),
                )
            elif layout == "figure_above_bullets":
                slide = add_figure_above_bullets_slide(
                    pres,
                    image_path=slide_spec.get("image_path", ""),
                    title=slide_spec.get("title", ""),
                    caption=slide_spec.get("caption", ""),
                    bullets=slide_spec.get("bullets"),
                    citation_source=slide_spec.get("citation_source", ""),
                )
            else:
                slide = add_figure_slide(
                    pres,
                    image_path=slide_spec.get("image_path", ""),
                    title=slide_spec.get("title", ""),
                    caption=slide_spec.get("caption", ""),
                    bullets=slide_spec.get("bullets"),
                    citation_source=slide_spec.get("citation_source", ""),
                )
        elif stype == "two_figure":
            slide = add_two_figure_compare_slide(
                pres,
                left_image=slide_spec.get("left_image", ""),
                right_image=slide_spec.get("right_image", ""),
                title=slide_spec.get("title", ""),
                left_label=slide_spec.get("left_label", ""),
                right_label=slide_spec.get("right_label", ""),
                left_caption=slide_spec.get("left_caption", ""),
                right_caption=slide_spec.get("right_caption", ""),
                citation_source=slide_spec.get("citation_source", ""),
            )
        elif stype == "quote":
            slide = add_quote_slide(
                pres,
                quote=slide_spec.get("quote", ""),
                attribution=slide_spec.get("attribution", ""),
            )
        elif stype == "multi_figure":
            slide = add_multi_figure_slide(
                pres,
                figures=slide_spec.get("figures", []),
                title=slide_spec.get("title", ""),
            )
        elif stype == "text":
            slide = add_text_slide(
                pres,
                title=slide_spec.get("title", ""),
                bullets=slide_spec.get("bullets", []),
            )
        elif stype == "references":
            slide = add_references_slide(
                pres,
                references=slide_spec.get("references", []),
                title=slide_spec.get("title", "References"),
            )
        else:
            # Unknown type — skip silently; never break a deck render
            continue

        if notes:
            attach_to_slide(slide, notes)

        annotations = slide_spec.get("annotations")
        if annotations and slide is not None and stype == "figure":
            try:
                pictures = [s for s in slide.shapes if s.shape_type == 13]
                if pictures:
                    add_annotations(
                        slide,
                        pictures[0],
                        annotations,
                        with_animations=with_animations,
                    )
            except Exception:
                pass

        if with_animations and slide is not None:
            try:
                _auto_animate_slide(slide, slide_spec, stype, bullet_reveal, panel_buildup)
            except Exception:
                # Animation is best-effort — never break deck rendering
                pass

    out_pptx = Path(output)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(out_pptx))

    result: dict[str, Path] = {"pptx": out_pptx}

    if write_marp:
        from vaultlab.slides.marp import write_marp as _write_marp

        marp_path = out_pptx.with_suffix(".md")
        _write_marp(plan, marp_path)
        result["marp"] = marp_path

    if kb_log is not None:
        report_path = _write_plan_kb_report(plan, out_pptx, kb_log)
        result["report"] = report_path
        topic = plan.get("topic") or plan.get("title", "untitled")
        kb_log.append_log(
            action="compile",
            title=f"Slide deck — {topic}",
            body=f"Generated {len(slides_plan)} slides → {out_pptx.name}",
            pages=[report_path.stem],
        )

    return result


def _normalize_bullet_annotations(slide_spec: dict[str, Any]) -> dict[str, Any]:
    """Support bullet-embedded annotations (Option A).

    A bullet may be a plain string OR a dict like::

        {"text": "Gen 2: ...", "click": 1,
         "annotation": {"type": "rect", "bbox": [...], "color": "FF5252"}}

    This function extracts annotations from such bullets, sets their
    ``click_index`` based on the bullet's ``click`` (or position), and
    merges into the slide-level ``"annotations"`` list. Returns a NEW
    slide_spec with bullets simplified to plain strings.
    """
    bullets = slide_spec.get("bullets")
    if not bullets:
        return slide_spec

    needs_normalize = any(isinstance(b, dict) for b in bullets)
    if not needs_normalize:
        return slide_spec

    new_bullets: list[str] = []
    extracted_anns: list[dict[str, Any]] = []
    for i, b in enumerate(bullets):
        if isinstance(b, dict):
            text = b.get("text", "")
            new_bullets.append(text)
            ann = b.get("annotation")
            if ann:
                ann = dict(ann)
                if "click_index" not in ann:
                    ann["click_index"] = b.get("click", i)
                extracted_anns.append(ann)
        else:
            new_bullets.append(b)

    out = dict(slide_spec)
    out["bullets"] = new_bullets
    if extracted_anns:
        existing = list(slide_spec.get("annotations", []))
        out["annotations"] = existing + extracted_anns
    return out


def _auto_animate_slide(
    slide: Any,
    slide_spec: dict[str, Any],
    stype: str,
    bullet_reveal: Any,
    panel_buildup: Any,
) -> None:
    """Apply default animations based on slide type.

    Rules:
      - text slide with >1 bullets → bullet-by-bullet reveal.
      - multi_figure slide with >1 figure → panel build-up.
      - figure slide with bullets → bullet reveal on the bullets text frame.
      - title, section_divider, references → no animations.
    """
    if stype == "text":
        bullets = slide_spec.get("bullets", [])
        if len(bullets) > 1:
            bullet_shape = slide.shapes[-1]
            if bullet_shape.has_text_frame:
                bullet_reveal(slide, bullet_shape.text_frame)
        return

    if stype == "multi_figure":
        figures = slide_spec.get("figures", [])
        if len(figures) > 1:
            groups = getattr(slide, "_vaultlab_panel_groups", None)
            if groups and len(groups) > 1:
                panel_buildup(slide, groups)
            else:
                picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
                if len(picture_shapes) > 1:
                    panel_buildup(slide, picture_shapes)
        return

    if stype == "figure":
        bullets = slide_spec.get("bullets")
        if bullets and len(bullets) > 1:
            picture_left_edges = [s.left for s in slide.shapes if s.shape_type == 13]
            min_picture_left = min(picture_left_edges) if picture_left_edges else 0
            best_match = None
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                if len(tf.paragraphs) < len(bullets):
                    continue
                if shape.left > min_picture_left:
                    best_match = tf
                    break
            if best_match is not None:
                bullet_reveal(slide, best_match)


def _write_plan_kb_report(plan: dict[str, Any], pptx_path: Path, kb_log: Any) -> Path:
    """Write a deck-plan record to ``<kb>/Output/Reports/``."""
    import json
    from datetime import datetime as _datetime

    date = _datetime.now().strftime("%Y-%m-%d")
    topic = plan.get("topic") or plan.get("title", "untitled")
    slug = "".join(c if c.isalnum() else "-" for c in topic.lower()).strip("-")[:60]
    filename = f"{date}-{slug}-deck.md"

    slides_plan = plan.get("slides", [])
    slide_summary_lines = []
    for i, s in enumerate(slides_plan, 1):
        stype = s.get("type", "?")
        title = s.get("title", "")
        slide_summary_lines.append(f"{i}. **{stype}** — {title}")

    sources_used = sorted(
        {s.get("citation_source", "") for s in slides_plan if s.get("citation_source")}
        | {
            f.get("citation_source", "")
            for s in slides_plan
            if s.get("type") == "multi_figure"
            for f in s.get("figures", [])
            if f.get("citation_source")
        }
    )
    sources_used = [src for src in sources_used if src]

    content = f"""---
title: "Slide deck — {topic}"
created: {date}
type: deck-report
status: COMPLETE
tags: [deck, slides, generated]
---

# Slide deck — {topic}

Generated by `vaultlab.slides.build_from_plan` on {date}.

## Output files

- PPTX: `{pptx_path}`
- Author: {plan.get("author", "")}
- Slides: {len(slides_plan)}

## Slide structure

{chr(10).join(slide_summary_lines)}

## Sources cited

{chr(10).join(f"- {s}" for s in sources_used) if sources_used else "_None_"}

## Full plan (JSON)

```json
{json.dumps(plan, indent=2, default=str)}
```
"""
    return kb_log.write_report(filename, content)
