"""Build a slide where every annotation is a NATIVE python-pptx shape.

Bobby's 2026-04-29 ask: stop baking annotations into the PNG. Each numbered
marker, bounding box, and side label should be its own PowerPoint object so
the user can animate, move, edit, or delete it without re-rendering.

Architecture
------------

For each annotation we add three named native shapes to the slide:

- ``ann{i}_box``   - Rectangle outline around the figure element (transparent fill,
                     colored outline, white shadow/halo).
- ``ann{i}_marker``- Small filled square at the box's top-left corner with the
                     numbered text in white.
- ``ann{i}_label`` - Text box in the right-side gutter column with the label
                     text colored to match the motif.

Naming makes shapes addressable for animations: PowerPoint's animation pane
will list them as ``ann1_box``, ``ann1_marker``, ``ann1_label`` etc., so a user
can build "appear all 3 together on click" or whatever sequence they want.

Coordinate mapping
------------------

Source figures are pixel arrays of varying sizes (e.g. 2035×1676). Slides are
EMU (914400 EMU = 1 inch). The figure is placed in the slide's content area
with aspect-preserved scaling, centered horizontally; annotation bboxes (given
in source-pixel coordinates) are mapped through the same scale factor +
top-left offset so they land on the correct visual element.
"""

from __future__ import annotations

from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from vaultlab.figures.understand.models import ElementAnnotation

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class SlideLayout:
    """Geometry + sizing for one annotated-figure slide.

    Tuned per Bobby 2026-04-29 v3 / v4 review feedback.
    """

    slide_w_in: float = 13.333
    slide_h_in: float = 7.5
    title_h_in: float = 0.85
    caption_h_in: float = 0.55
    footer_h_in: float = 0.40
    figure_area_x_in: float = 0.20
    figure_area_w_in: float = 10.20
    gutter_x_in: float = 10.55
    gutter_w_in: float = 2.65
    marker_size_in: float = 0.24
    marker_font_pt: int = 11
    label_font_pt: int = 11
    title_font_pt: int = 26
    caption_font_pt: int = 12
    footer_font_pt: int = 9
    # Bobby 2026-04-29 v7 review: 'boxes should be slightly larger than the
    # motif so they are not right on the boundaries'. Pads the box by N
    # source pixels on all sides when drawing.
    bbox_padding_px: int = 30
    # Theme variant: "light" (white bg, dark text) or "dark" (dark bg, light text)
    theme_variant: str = "light"
    # Per Bobby's Hickey-template-logo-zones memory: when using Hickey Lab masters,
    # the section banner must shrink to avoid the bottom-flanking lab + Duke logos.
    banner_left_margin_in: float = 0.4
    banner_right_margin_in: float = 1.1  # leaves room for page number


DEFAULT = SlideLayout()


HICKEY_LAB_LAYOUT = SlideLayout(
    theme_variant="dark",
    banner_left_margin_in=1.2,  # avoid Hickey lab logo
    banner_right_margin_in=1.4,  # avoid Duke logo
)


# Theme-aware text-color helpers
def text_color_for_bg(bg_rgb: tuple[int, int, int]) -> tuple[int, int, int]:
    """Return readable text color (white or near-black) for a given background.

    Per Bobby 2026-04-29 v4 review + theme-aware-font-color memory.
    """
    r, g, b = bg_rgb
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return (240, 240, 240) if luminance < 130 else (30, 30, 30)


def text_color_for_theme(variant: str) -> tuple[int, int, int]:
    """Quick text color lookup by theme variant ('dark' or 'light')."""
    if variant == "dark":
        return (240, 240, 240)
    return (30, 30, 30)


def muted_text_color_for_theme(variant: str) -> tuple[int, int, int]:
    """Muted secondary text color (footer, inactive section pills)."""
    return (180, 185, 195) if variant == "dark" else (110, 115, 125)


def card_fill_for_theme(variant: str) -> tuple[int, int, int]:
    """Background fill for inactive section pills."""
    return (50, 55, 65) if variant == "dark" else (245, 246, 248)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def add_annotated_figure_slide(
    pres: Presentation,
    image_path: str | Path,
    annotations: Sequence[ElementAnnotation],
    *,
    title: str,
    caption: str = "",
    motif_colors: dict[str, tuple[int, int, int]] | None = None,
    layout: SlideLayout = DEFAULT,
    notes: str = "",
    page_number: int | None = None,
    sections: Sequence[str] | None = None,
    current_section_idx: int | None = None,
) -> None:
    """Add one slide with figure + native-shape annotations + footer.

    Parameters
    ----------
    pres
        The presentation to append to (must already have slide_width/height set).
    image_path
        Path to the source figure image (no pre-baked annotations).
    annotations
        Concept-to-region pairings.
    title
        Slide title (rendered in title text box at top, centered).
    caption
        Optional italic caption below the figure.
    motif_colors
        Mapping from ``motif_name`` to RGB tuple. Falls back to a default
        palette when ``motif_name`` is missing.
    layout
        Geometric layout constants.
    notes
        Speaker notes text.
    page_number
        If set, draws "<n>" in the right side of the footer (Bobby Q2 2026-04-29).
    sections
        If set, draws a section banner across the full bottom: one rectangle
        per section, current section highlighted via ``current_section_idx``.
    current_section_idx
        Index into ``sections`` of the current section (highlighted).
    """
    from PIL import Image

    blank = pres.slide_layouts[6]
    s = pres.slides.add_slide(blank)

    _add_title(s, title, layout)

    src_w, src_h = Image.open(Path(image_path)).size
    img_x_in, img_y_in, img_w_in, img_h_in = _placed_figure_geometry(src_w, src_h, layout)
    s.shapes.add_picture(
        str(image_path),
        Inches(img_x_in),
        Inches(img_y_in),
        width=Inches(img_w_in),
        height=Inches(img_h_in),
    )

    if caption:
        _add_caption(s, caption, layout)

    if annotations:
        _add_annotations(
            s,
            annotations=annotations,
            src_w=src_w,
            src_h=src_h,
            img_x_in=img_x_in,
            img_y_in=img_y_in,
            img_w_in=img_w_in,
            img_h_in=img_h_in,
            motif_colors=motif_colors or {},
            layout=layout,
        )

    # Footer (page number + optional section banner)
    if page_number is not None:
        _add_page_number(s, page_number, layout)
    if sections:
        _add_section_banner(
            s,
            sections=sections,
            current_idx=current_section_idx,
            layout=layout,
        )

    if notes:
        s.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# Geometry
# ---------------------------------------------------------------------------


# Minimum clearance between the title's bottom edge and the picture's
# top edge. Bobby 2026-04-30 (Bug #6 / title-overlap): the legacy 0.10"
# offset gave only 0.05" clearance after the title text box's
# ``title_h_in - 0.1`` height adjustment, which let figure pictures
# overlap wrapping titles.
#
# The picture is placed at ``title_h_in + TITLE_PICTURE_CLEARANCE_IN``;
# the title box itself ends at ``0.15 + (title_h_in - 0.10) = title_h_in
# + 0.05``. Real clearance is therefore ``TITLE_PICTURE_CLEARANCE_IN -
# 0.05``. Setting the constant to 0.30" yields 0.25" of *actual* gap —
# which is the audit-required minimum.
TITLE_PICTURE_CLEARANCE_IN: float = 0.30


def _placed_figure_geometry(
    src_w: int, src_h: int, layout: SlideLayout
) -> tuple[float, float, float, float]:
    """Compute the figure's position + size in inches on the slide.

    Aspect-preserve fit inside the figure-area rectangle; center horizontally
    within the column.

    Reserves :data:`TITLE_PICTURE_CLEARANCE_IN` of vertical clearance
    between the title's bottom edge and the picture's top edge so a
    wrapping multi-line title never overlaps the figure.
    """
    avail_h = (
        layout.slide_h_in
        - layout.title_h_in
        - TITLE_PICTURE_CLEARANCE_IN
        - layout.caption_h_in
        - layout.footer_h_in
        - 0.05
    )
    avail_w = layout.figure_area_w_in
    aspect_src = src_w / src_h
    aspect_avail = avail_w / avail_h
    if aspect_src > aspect_avail:
        w = avail_w
        h = avail_w / aspect_src
    else:
        h = avail_h
        w = avail_h * aspect_src
    x = layout.figure_area_x_in + (avail_w - w) / 2
    y = layout.title_h_in + TITLE_PICTURE_CLEARANCE_IN
    return x, y, w, h


def _src_bbox_to_inches(
    bbox: tuple[int, int, int, int],
    src_w: int,
    src_h: int,
    img_x_in: float,
    img_y_in: float,
    img_w_in: float,
    img_h_in: float,
) -> tuple[float, float, float, float]:
    """Map a source-pixel bbox to slide inches (x, y, w, h)."""
    x0, y0, x1, y1 = bbox
    sx = img_w_in / src_w
    sy = img_h_in / src_h
    return (
        img_x_in + x0 * sx,
        img_y_in + y0 * sy,
        (x1 - x0) * sx,
        (y1 - y0) * sy,
    )


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------


_FALLBACK_PALETTE = [
    (0, 102, 204),
    (220, 50, 47),
    (38, 139, 210),
    (181, 137, 0),
    (133, 153, 0),
    (108, 113, 196),
    (203, 75, 22),
    (211, 54, 130),
    (42, 161, 152),
]


def _color_for(
    ann: ElementAnnotation, idx: int, motif_colors: dict[str, tuple[int, int, int]]
) -> tuple[int, int, int]:
    if ann.motif_name in motif_colors:
        return motif_colors[ann.motif_name]
    return _FALLBACK_PALETTE[idx % len(_FALLBACK_PALETTE)]


def _darken_for_text(rgb: tuple[int, int, int]) -> tuple[int, int, int]:
    r, g, b = rgb
    if 0.299 * r + 0.587 * g + 0.114 * b > 180:
        return (max(0, r - 80), max(0, g - 80), max(0, b - 80))
    return rgb


def _resolve_padding(
    override: int | tuple[int, int, int, int] | None,
    default_scalar: int,
) -> tuple[int, int, int, int]:
    """Return ``(top, right, bottom, left)`` padding in source pixels.

    Resolution order: per-annotation override (scalar broadcast OR explicit
    4-tuple) > layout default scalar (broadcast).
    """
    if override is None:
        return (default_scalar, default_scalar, default_scalar, default_scalar)
    if isinstance(override, int):
        return (override, override, override, override)
    return override


def _add_title(s, title: str, layout: SlideLayout) -> None:
    from pptx.enum.text import PP_ALIGN

    box = s.shapes.add_textbox(
        Inches(0.4),
        Inches(0.15),
        Inches(layout.slide_w_in - 0.8),
        Inches(layout.title_h_in - 0.1),
    )
    box.name = "slide_title"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER  # Bobby Q1 2026-04-29: centered title
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(layout.title_font_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*text_color_for_theme(layout.theme_variant))


def _add_caption(s, caption: str, layout: SlideLayout) -> None:
    # Caption sits above the footer
    caption_top = layout.slide_h_in - layout.footer_h_in - layout.caption_h_in - 0.02
    box = s.shapes.add_textbox(
        Inches(0.4),
        Inches(caption_top),
        Inches(layout.slide_w_in - 0.8),
        Inches(layout.caption_h_in),
    )
    box.name = "slide_caption"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = caption
    run.font.name = "Arial"
    run.font.size = Pt(layout.caption_font_pt)
    run.font.italic = True
    run.font.color.rgb = RGBColor(*muted_text_color_for_theme(layout.theme_variant))


def _add_page_number(s, page_number: int, layout: SlideLayout) -> None:
    """Right-aligned page number in the footer."""
    from pptx.enum.text import PP_ALIGN

    box = s.shapes.add_textbox(
        Inches(layout.slide_w_in - 1.0),
        Inches(layout.slide_h_in - layout.footer_h_in + 0.05),
        Inches(0.8),
        Inches(layout.footer_h_in - 0.10),
    )
    box.name = "slide_page_number"
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_number)
    run.font.name = "Arial"
    run.font.size = Pt(layout.footer_font_pt)
    run.font.color.rgb = RGBColor(*muted_text_color_for_theme(layout.theme_variant))


def _add_section_banner(s, *, sections, current_idx: int | None, layout: SlideLayout) -> None:
    """Section-banner pills - modern style.

    Per Bobby 2026-04-29 v3 review ('look more modern'):
    - Rounded-rectangle pills (not sharp)
    - No outline (clean)
    - Soft fills: very light card for inactive, cobalt for active
    - Inactive: muted slate text; active: white bold text
    - Comfortable gap between pills
    """
    from pptx.enum.text import PP_ALIGN

    n = len(sections)
    if n == 0:
        return

    left_margin = layout.banner_left_margin_in
    right_margin = layout.banner_right_margin_in
    gap_in = 0.10
    avail = layout.slide_w_in - left_margin - right_margin
    total_gap = gap_in * (n - 1)
    rect_w = (avail - total_gap) / max(n, 1)
    rect_h = 0.32
    rect_y = layout.slide_h_in - layout.footer_h_in + 0.04

    for i, name in enumerate(sections):
        x = left_margin + i * (rect_w + gap_in)
        rect = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(rect_y),
            Inches(rect_w),
            Inches(rect_h),
        )
        rect.name = f"section_banner_{i}"
        # Crank the corner radius so it looks like a pill
        try:
            rect.adjustments[0] = 0.5
        except (IndexError, AttributeError):
            pass

        rect.fill.solid()
        if i == current_idx:
            rect.fill.fore_color.rgb = RGBColor(0, 102, 204)  # cobalt active
            text_color = RGBColor(255, 255, 255)
        else:
            rect.fill.fore_color.rgb = RGBColor(*card_fill_for_theme(layout.theme_variant))
            text_color = RGBColor(*muted_text_color_for_theme(layout.theme_variant))
        rect.line.fill.background()  # no border

        tf = rect.text_frame
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.name = "Arial"
        run.font.size = Pt(layout.footer_font_pt + 1)
        run.font.bold = i == current_idx
        run.font.color.rgb = text_color


def _add_leader_line(
    s,
    *,
    marker_cx_in: float,
    marker_cy_in: float,
    bbox_x_in: float,
    bbox_y_in: float,
    bbox_w_in: float,
    bbox_h_in: float,
    color: tuple[int, int, int],
    name: str,
) -> None:
    """Thin connector from the bbox edge nearest the marker to the marker center.

    Used when the marker is offset far from its element (e.g., placed in a
    distant whitespace region). Picks the bbox edge midpoint closest to the
    marker so the line stays short and doesn't cross figure content.
    """
    bbox_left = bbox_x_in
    bbox_right = bbox_x_in + bbox_w_in
    bbox_top = bbox_y_in
    bbox_bottom = bbox_y_in + bbox_h_in
    bbox_cx = bbox_x_in + bbox_w_in / 2
    bbox_cy = bbox_y_in + bbox_h_in / 2

    # Pick which edge of the bbox to connect from (whichever side the marker is on)
    if abs(marker_cx_in - bbox_cx) > abs(marker_cy_in - bbox_cy):
        anchor_y = bbox_cy
        anchor_x = bbox_right if marker_cx_in > bbox_cx else bbox_left
    else:
        anchor_x = bbox_cx
        anchor_y = bbox_bottom if marker_cy_in > bbox_cy else bbox_top

    line = s.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(anchor_x),
        Inches(anchor_y),
        Inches(marker_cx_in),
        Inches(marker_cy_in),
    )
    line.name = name
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(1.0)


def _group_shapes(slide, shapes_to_group, group_name: str) -> None:
    """Wrap a list of shapes into a single PowerPoint group.

    Per Bobby 2026-04-29 v6 review: side-marker + label should be one click,
    not two. python-pptx has no high-level group API; we build the XML and
    move the shape elements into it.
    """
    if len(shapes_to_group) < 2:
        return

    sp_tree = slide.shapes._spTree

    min_x = min(s.left for s in shapes_to_group)
    min_y = min(s.top for s in shapes_to_group)
    max_x = max(s.left + s.width for s in shapes_to_group)
    max_y = max(s.top + s.height for s in shapes_to_group)

    # Pick a unique-ish id (current count + 100 to avoid collisions)
    next_id = len(sp_tree.findall(qn("p:sp"))) + len(sp_tree.findall(qn("p:grpSp"))) + 100

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"

    grp_xml = (
        f'<p:grpSp xmlns:p="{p_ns}" xmlns:a="{a_ns}">'
        f"<p:nvGrpSpPr>"
        f'<p:cNvPr id="{next_id}" name="{group_name}"/>'
        f"<p:cNvGrpSpPr/>"
        f"<p:nvPr/>"
        f"</p:nvGrpSpPr>"
        f"<p:grpSpPr>"
        f"<a:xfrm>"
        f'<a:off x="{min_x}" y="{min_y}"/>'
        f'<a:ext cx="{max_x - min_x}" cy="{max_y - min_y}"/>'
        f'<a:chOff x="{min_x}" y="{min_y}"/>'
        f'<a:chExt cx="{max_x - min_x}" cy="{max_y - min_y}"/>'
        f"</a:xfrm>"
        f"</p:grpSpPr>"
        f"</p:grpSp>"
    )
    grp_sp = etree.fromstring(grp_xml)

    for shape in shapes_to_group:
        sp_element = shape._element
        sp_tree.remove(sp_element)
        grp_sp.append(sp_element)

    sp_tree.append(grp_sp)


def _add_annotations(
    s,
    *,
    annotations: Sequence[ElementAnnotation],
    src_w: int,
    src_h: int,
    img_x_in: float,
    img_y_in: float,
    img_w_in: float,
    img_h_in: float,
    motif_colors: dict[str, tuple[int, int, int]],
    layout: SlideLayout,
) -> None:
    n = len(annotations)
    # Numerical order in the gutter per Bobby 2026-04-29 v6 review:
    # users expect 1, 2, 3... reading down; box-y ordering caused odd sequences
    # like 3, 5, 1, 7, ... when annotations were spatially scattered.
    sorted_anns = list(enumerate(annotations))

    gutter_top = layout.title_h_in + 0.15
    gutter_bottom = layout.slide_h_in - layout.caption_h_in - 0.05
    spacing = (gutter_bottom - gutter_top) / max(n, 1)

    for slot, (orig_idx, ann) in enumerate(sorted_anns):
        color = _color_for(ann, orig_idx, motif_colors)
        text_color = _darken_for_text(color)
        num = str(orig_idx + 1)

        x_in, y_in, w_in, h_in = _src_bbox_to_inches(
            ann.bbox_px,
            src_w,
            src_h,
            img_x_in,
            img_y_in,
            img_w_in,
            img_h_in,
        )

        # 1. Bounding-box outline (skip if ann.use_box=False).
        # Padding: per-annotation override > layout default. Tuple = asymmetric
        # (top, right, bottom, left) per Bobby 2026-04-29 v8 ask. Shape may be
        # rectangle or circle/ellipse depending on `bbox_shape`.
        if ann.use_box:
            sx = img_w_in / src_w
            sy = img_h_in / src_h
            pad_top, pad_right, pad_bottom, pad_left = _resolve_padding(
                ann.bbox_padding_px, layout.bbox_padding_px
            )
            box_shape_kind = MSO_SHAPE.OVAL if ann.bbox_shape == "circle" else MSO_SHAPE.RECTANGLE
            box = s.shapes.add_shape(
                box_shape_kind,
                Inches(x_in - pad_left * sx),
                Inches(y_in - pad_top * sy),
                Inches(w_in + (pad_left + pad_right) * sx),
                Inches(h_in + (pad_top + pad_bottom) * sy),
            )
            box.name = f"ann{orig_idx + 1}_box"
            box.fill.background()
            box.line.color.rgb = RGBColor(*color)
            box.line.width = Pt(2.5)
            box.text_frame.text = ""

        # 2. Numbered marker - default top-left, but ann.marker_offset_px
        # can shift it to nearby whitespace to avoid collisions with other
        # markers or to clear important figure content underneath.
        marker_size = layout.marker_size_in
        sx = img_w_in / src_w
        sy = img_h_in / src_h
        if ann.marker_offset_px is not None:
            dx_px, dy_px = ann.marker_offset_px
            marker_x_in = x_in + dx_px * sx
            marker_y_in = y_in + dy_px * sy
        else:
            marker_x_in = x_in
            marker_y_in = max(0.05, y_in - marker_size - 0.02)

        # Leader line: when the marker is FAR from the bbox (>15% of figure
        # width on either axis), draw a thin connector so the eye can still
        # trace the marker back to its element. Bobby 2026-04-29 v8: labels
        # may roam, but the viewer needs a visual link.
        far_threshold_in = 0.15 * img_w_in
        marker_cx = marker_x_in + marker_size / 2
        marker_cy = marker_y_in + marker_size / 2
        bbox_cx = x_in + w_in / 2
        bbox_cy = y_in + h_in / 2
        is_far = (
            abs(marker_cx - bbox_cx) > far_threshold_in
            or abs(marker_cy - bbox_cy) > far_threshold_in
        )
        # Draw a leader line whenever the marker is far from the bbox center,
        # regardless of whether use_box is True. When use_box=True the box
        # itself is a visual anchor; when far + boxed, the leader still
        # connects the eye to the marker.
        if is_far:
            _add_leader_line(
                s,
                marker_cx_in=marker_cx,
                marker_cy_in=marker_cy,
                bbox_x_in=x_in,
                bbox_y_in=y_in,
                bbox_w_in=w_in,
                bbox_h_in=h_in,
                color=color,
                name=f"ann{orig_idx + 1}_leader",
            )

        marker = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(marker_x_in),
            Inches(marker_y_in),
            Inches(marker_size),
            Inches(marker_size),
        )
        marker.name = f"ann{orig_idx + 1}_marker"
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(*color)
        marker.line.color.rgb = RGBColor(255, 255, 255)
        marker.line.width = Pt(1.5)
        mtf = marker.text_frame
        mtf.margin_top = Emu(0)
        mtf.margin_bottom = Emu(0)
        mtf.margin_left = Emu(0)
        mtf.margin_right = Emu(0)
        mp = mtf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN

        mp.alignment = PP_ALIGN.CENTER
        mr = mp.add_run()
        mr.text = num
        mr.font.name = "Arial"
        mr.font.bold = True
        mr.font.size = Pt(layout.marker_font_pt)
        mr.font.color.rgb = RGBColor(255, 255, 255)

        # 3. Side label in gutter - marker square + label text
        label_y_in = gutter_top + slot * spacing
        side_marker = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(layout.gutter_x_in),
            Inches(label_y_in),
            Inches(marker_size),
            Inches(marker_size),
        )
        side_marker.name = f"ann{orig_idx + 1}_side_marker"
        side_marker.fill.solid()
        side_marker.fill.fore_color.rgb = RGBColor(*color)
        side_marker.line.color.rgb = RGBColor(255, 255, 255)
        side_marker.line.width = Pt(1.5)
        smt = side_marker.text_frame
        smt.margin_top = Emu(0)
        smt.margin_bottom = Emu(0)
        smt.margin_left = Emu(0)
        smt.margin_right = Emu(0)
        smp = smt.paragraphs[0]
        smp.alignment = PP_ALIGN.CENTER
        smr = smp.add_run()
        smr.text = num
        smr.font.name = "Arial"
        smr.font.bold = True
        smr.font.size = Pt(layout.marker_font_pt)
        smr.font.color.rgb = RGBColor(255, 255, 255)

        # Label text box
        label_box = s.shapes.add_textbox(
            Inches(layout.gutter_x_in + marker_size + 0.10),
            Inches(label_y_in - 0.05),
            Inches(layout.gutter_w_in - marker_size - 0.15),
            Inches(spacing - 0.05),
        )
        label_box.name = f"ann{orig_idx + 1}_label"
        ltf = label_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = ann.label
        lr.font.name = "Arial"
        lr.font.size = Pt(layout.label_font_pt)
        lr.font.bold = True
        lr.font.color.rgb = RGBColor(*text_color)

        # Group the side-marker + label so the user can move/animate them
        # together with one click (Bobby 2026-04-29 v6 ask).
        _group_shapes(s, [side_marker, label_box], f"ann{orig_idx + 1}_side_group")


__all__ = ["DEFAULT", "SlideLayout", "add_annotated_figure_slide"]
