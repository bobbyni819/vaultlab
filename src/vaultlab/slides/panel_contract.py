"""Manuscript-style PowerPoint panel layout contracts.

This module separates two things agents often blur together:

* an intended panel layout contract with slots, gutters, and font floors;
* the geometry extracted from an actual PPTX slide.

Both are deterministic and safe for CI; no PowerPoint COM automation is used.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

PanelSeverity = Literal["pass", "warn", "fail"]

_EMU_PER_INCH = 914400
_SEVERITY_RANK: dict[PanelSeverity, int] = {"pass": 0, "warn": 1, "fail": 2}


@dataclass(frozen=True)
class PanelSlot:
    """One image panel slot in slide inches."""

    letter: str
    image_path: str
    slot_in: list[float]
    letter_box_in: list[float] | None = None
    title_box_in: list[float] | None = None
    allow_text_overlap: bool = False


@dataclass(frozen=True)
class PanelLayoutContract:
    """Geometry contract for a manuscript-style composite slide."""

    figure_id: str
    slide_width_in: float
    slide_height_in: float
    panels: list[PanelSlot]
    min_gutter_in: float = 0.04
    panel_letter_font_pt: float = 14.0
    min_effective_text_pt: float = 5.5


@dataclass(frozen=True)
class PanelLayoutIssue:
    """One panel-contract audit finding."""

    rule: str
    severity: PanelSeverity
    detail: str
    panel: str | None = None


@dataclass(frozen=True)
class PanelLayoutAudit:
    """Aggregate panel-contract audit result."""

    overall_severity: PanelSeverity
    issues: list[PanelLayoutIssue] = field(default_factory=list)

    @property
    def n_fail(self) -> int:
        return sum(1 for issue in self.issues if issue.severity == "fail")

    @property
    def n_warn(self) -> int:
        return sum(1 for issue in self.issues if issue.severity == "warn")

    def ok(self) -> bool:
        return self.n_fail == 0


@dataclass(frozen=True)
class SlideShapeGeometry:
    """Geometry for one native PPTX shape."""

    kind: str
    name: str
    x_in: float
    y_in: float
    w_in: float
    h_in: float
    text: str = ""
    font_sizes_pt: list[float] = field(default_factory=list)
    media_width_px: int | None = None
    media_height_px: int | None = None


@dataclass(frozen=True)
class SlideGeometry:
    """Geometry recovered from one PPTX slide."""

    pptx_path: str
    slide_number: int
    slide_width_in: float
    slide_height_in: float
    shapes: list[SlideShapeGeometry]


def audit_panel_layout_contract(contract: PanelLayoutContract) -> PanelLayoutAudit:
    """Audit a manuscript composite layout contract."""

    issues: list[PanelLayoutIssue] = []
    if contract.slide_width_in <= 0 or contract.slide_height_in <= 0:
        issues.append(
            PanelLayoutIssue(
                rule="slide-size",
                severity="fail",
                detail="slide dimensions must be positive",
            )
        )

    if contract.panel_letter_font_pt < 12.0:
        issues.append(
            PanelLayoutIssue(
                rule="panel-letter-font",
                severity="fail",
                detail=(
                    f"panel letter font {contract.panel_letter_font_pt:.1f} pt is below "
                    "the 12 pt floor"
                ),
            )
        )

    for panel in contract.panels:
        if not _box_is_valid(panel.slot_in):
            issues.append(
                PanelLayoutIssue(
                    rule="panel-slot-shape",
                    severity="fail",
                    detail=f"panel {panel.letter} slot must be [x, y, w, h] with w/h > 0",
                    panel=panel.letter,
                )
            )
            continue
        if not _box_inside(panel.slot_in, contract.slide_width_in, contract.slide_height_in):
            issues.append(
                PanelLayoutIssue(
                    rule="panel-slot-bounds",
                    severity="fail",
                    detail=f"panel {panel.letter} slot extends outside the slide",
                    panel=panel.letter,
                )
            )

    for i, panel_a in enumerate(contract.panels):
        if not _box_is_valid(panel_a.slot_in):
            continue
        for panel_b in contract.panels[i + 1 :]:
            if not _box_is_valid(panel_b.slot_in):
                continue
            overlap = _overlap_area(panel_a.slot_in, panel_b.slot_in)
            if overlap > 0:
                issues.append(
                    PanelLayoutIssue(
                        rule="panel-slot-overlap",
                        severity="fail",
                        detail=(
                            f"panels {panel_a.letter} and {panel_b.letter} overlap by "
                            f"{overlap:.3f} square inches"
                        ),
                        panel=f"{panel_a.letter},{panel_b.letter}",
                    )
                )
            elif _nearest_gap(panel_a.slot_in, panel_b.slot_in) < contract.min_gutter_in:
                issues.append(
                    PanelLayoutIssue(
                        rule="panel-gutter",
                        severity="warn",
                        detail=(
                            f"panels {panel_a.letter} and {panel_b.letter} are closer "
                            f"than the {contract.min_gutter_in:.3f} in gutter"
                        ),
                        panel=f"{panel_a.letter},{panel_b.letter}",
                    )
                )

    return PanelLayoutAudit(overall_severity=_aggregate(issues), issues=issues)


def extract_pptx_slide_geometry(pptx_path: Path | str, *, slide_number: int = 1) -> SlideGeometry:
    """Extract native shape geometry from one PPTX slide.

    ``slide_number`` is 1-indexed to match PowerPoint's UI.
    """

    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency-gated
        raise RuntimeError(
            "python-pptx is required for PPTX geometry extraction. Install vaultlab[slides]."
        ) from exc

    path = Path(pptx_path)
    prs = Presentation(str(path))
    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"slide_number {slide_number} outside deck with {len(prs.slides)} slides")
    slide = prs.slides[slide_number - 1]
    shapes = [_shape_geometry(shape) for shape in slide.shapes]
    return SlideGeometry(
        pptx_path=str(path),
        slide_number=slide_number,
        slide_width_in=_emu_to_in(prs.slide_width),
        slide_height_in=_emu_to_in(prs.slide_height),
        shapes=shapes,
    )


def _shape_geometry(shape: Any) -> SlideShapeGeometry:
    kind = "picture" if getattr(shape, "shape_type", None) == 13 else "shape"
    text = ""
    font_sizes: list[float] = []
    if getattr(shape, "has_text_frame", False):
        try:
            text = str(shape.text_frame.text or "").strip()
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    size = getattr(run.font, "size", None)
                    if size is not None:
                        font_sizes.append(float(size.pt))
        except Exception:
            text = ""

    media_width_px: int | None = None
    media_height_px: int | None = None
    if kind == "picture":
        try:
            media_width_px, media_height_px = shape.image.size
        except Exception:
            media_width_px = None
            media_height_px = None

    return SlideShapeGeometry(
        kind=kind,
        name=str(getattr(shape, "name", "")),
        x_in=_emu_to_in(getattr(shape, "left", 0)),
        y_in=_emu_to_in(getattr(shape, "top", 0)),
        w_in=_emu_to_in(getattr(shape, "width", 0)),
        h_in=_emu_to_in(getattr(shape, "height", 0)),
        text=text,
        font_sizes_pt=font_sizes,
        media_width_px=media_width_px,
        media_height_px=media_height_px,
    )


def _box_is_valid(box: list[float]) -> bool:
    return len(box) == 4 and box[2] > 0 and box[3] > 0


def _box_inside(box: list[float], slide_width: float, slide_height: float) -> bool:
    x, y, w, h = box
    return x >= 0 and y >= 0 and x + w <= slide_width and y + h <= slide_height


def _overlap_area(a: list[float], b: list[float]) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


def _nearest_gap(a: list[float], b: list[float]) -> float:
    """Smallest axis-aligned gap between two boxes that share an x- or y-band.

    NOTE: when two boxes share neither a horizontal nor a vertical projection
    (they are placed strictly diagonally), this returns ``inf`` and the gutter
    check is skipped for that pair -- corner-to-corner proximity is not measured.
    Manuscript panel grids are axis-aligned, so this limitation is acceptable.
    """

    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    horizontal_projection = ay < by + bh and by < ay + ah
    vertical_projection = ax < bx + bw and bx < ax + aw
    gaps: list[float] = []
    if horizontal_projection:
        gaps.append(max(bx - (ax + aw), ax - (bx + bw), 0.0))
    if vertical_projection:
        gaps.append(max(by - (ay + ah), ay - (by + bh), 0.0))
    return min(gaps) if gaps else float("inf")


def _aggregate(issues: list[PanelLayoutIssue]) -> PanelSeverity:
    if not issues:
        return "pass"
    return max((issue.severity for issue in issues), key=lambda severity: _SEVERITY_RANK[severity])


def _emu_to_in(value: Any) -> float:
    return float(value) / _EMU_PER_INCH


__all__ = [
    "PanelLayoutAudit",
    "PanelLayoutContract",
    "PanelLayoutIssue",
    "PanelSlot",
    "SlideGeometry",
    "SlideShapeGeometry",
    "audit_panel_layout_contract",
    "extract_pptx_slide_geometry",
]

