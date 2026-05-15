"""Slide-deck time-budget audit.

Estimate per-slide and total presentation time for a rendered ``.pptx``,
then compare against a target ``time_minutes`` budget. Flags decks that
will not fit inside ``budget_minutes - qa_reserve_minutes`` of speaking
time.

Heuristics (research-presentation context, validated against ~10 min
Hickey-lab journal-club rhythm):

* **Title / divider** — 15s. Audience reads the title and you start.
* **Bullets only** (≤5 bullets, no figure) — 30-45s.
* **Bullets-heavy** (>5 bullets, no figure) — 45-75s.
* **Bullets + figure** — 60-90s. The figure adds a beat per panel.
* **Figure-heavy / data slide** (figure but minimal text) — 90-120s.
* **Methods / equation slide** — 60s. Equations need a pause but the
  audience usually accepts the math at face value.
* **References / acknowledgments** — 20s. People scan these.
* **Discussion / Q&A prompt** — handled out-of-band: the
  ``qa_reserve_minutes`` argument subtracts a fixed Q&A block from the
  effective speaking budget.

The estimator is intentionally conservative — slightly LONGER than the
practiced run-time so a deck that audits "under" almost always lands
under in the room. Bobby would rather have 30s of slack than a missed
take-home slide.

The module exposes a single public entrypoint
:func:`audit_time_budget` returning a :class:`TimeBudgetReport`.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

__all__ = [
    "SlideTimeEstimate",
    "TimeBudgetReport",
    "audit_time_budget",
]


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass
class SlideTimeEstimate:
    """Time-budget estimate for a single slide.

    Attributes
    ----------
    slide_index : int
        Zero-indexed deck position of the slide.
    title : str | None
        Best-effort title text (None when no text shape was recovered).
    kind : str
        Heuristic classification —
        ``"title"`` / ``"section_divider"`` / ``"bullets"`` /
        ``"figure"`` / ``"methods"`` / ``"references"`` /
        ``"acknowledgments"`` / ``"discussion"`` / ``"other"``.
    estimate_seconds : int
        Estimated speaking time, in whole seconds.
    rationale : str
        Short human-readable explanation of how the estimate was reached.
    """

    slide_index: int
    title: str | None
    kind: str
    estimate_seconds: int
    rationale: str


@dataclass
class TimeBudgetReport:
    """Time-budget audit for an entire rendered deck.

    The ``budget_minutes`` argument is the TOTAL slot duration (e.g.
    ``10`` for a 10-minute talk). The ``qa_reserve_minutes`` is the
    amount of that slot reserved for audience Q&A — the effective
    speaking budget is ``budget_minutes - qa_reserve_minutes``.
    """

    deck_path: Path
    budget_minutes: int
    qa_reserve_minutes: int = 5
    per_slide: list[SlideTimeEstimate] = field(default_factory=list)

    @property
    def estimated_total_seconds(self) -> int:
        return sum(s.estimate_seconds for s in self.per_slide)

    @property
    def estimated_total_minutes(self) -> float:
        return self.estimated_total_seconds / 60.0

    @property
    def budget_seconds(self) -> int:
        """Effective speaking budget (total slot minus Q&A reserve)."""
        return (self.budget_minutes - self.qa_reserve_minutes) * 60

    def over_budget(self) -> bool:
        return self.estimated_total_seconds > self.budget_seconds

    def summary_lines(self) -> list[str]:
        """Human-readable headline summary."""
        verdict = "OVER" if self.over_budget() else "under"
        return [
            (
                f"Time budget: {self.budget_minutes} min slot "
                f"(− {self.qa_reserve_minutes} min Q&A reserve "
                f"= {self.budget_seconds // 60} min speaking)."
            ),
            (
                f"Estimated speak time: "
                f"{self.estimated_total_seconds // 60} min "
                f"{self.estimated_total_seconds % 60:02d}s "
                f"({self.estimated_total_seconds}s) — {verdict} budget."
            ),
        ]


# ---------------------------------------------------------------------------
# Audit
# ---------------------------------------------------------------------------


# Per-kind heuristic estimates (seconds). Tuned against ~10 min Hickey-lab
# journal-club rhythm — slightly conservative on purpose.
_KIND_BASE_SECONDS = {
    "title": 15,
    "section_divider": 15,
    "bullets": 40,
    "bullets_heavy": 60,
    "figure": 75,
    "figure_heavy": 100,
    "methods": 60,
    "references": 20,
    "acknowledgments": 20,
    "discussion": 30,
    "other": 30,
}


_REFERENCES_RE = re.compile(r"^\s*references?\s*$", re.IGNORECASE)
_ACK_RE = re.compile(r"^\s*acknowledg", re.IGNORECASE)
_DISCUSSION_RE = re.compile(r"^\s*(discussion|q\s*&\s*a|questions?)\b", re.IGNORECASE)
_METHODS_RE = re.compile(r"\b(methods?|equations?|algorithm|protocol)\b", re.IGNORECASE)


def audit_time_budget(
    pptx_path: Path | str,
    *,
    budget_minutes: int,
    qa_reserve_minutes: int = 5,
) -> TimeBudgetReport:
    """Estimate per-slide + total speaking time for a rendered deck.

    Parameters
    ----------
    pptx_path
        Path to the rendered ``.pptx`` to audit.
    budget_minutes
        Total presentation slot in minutes (e.g. ``10`` for a 10-min
        talk slot).
    qa_reserve_minutes
        Minutes to reserve for Q&A; subtracted from ``budget_minutes``
        to give the effective speaking budget.

    Returns
    -------
    TimeBudgetReport
        Per-slide estimates + aggregate verdict via
        :meth:`TimeBudgetReport.over_budget`.

    Raises
    ------
    FileNotFoundError
        If ``pptx_path`` doesn't exist.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_path}")

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover — gated at install time
        raise RuntimeError(
            "python-pptx is required for audit_time_budget. Install with "
            '`pip install -e ".[slides]"` or `pip install python-pptx`.'
        ) from exc

    prs = Presentation(str(pptx_path))
    report = TimeBudgetReport(
        deck_path=pptx_path,
        budget_minutes=budget_minutes,
        qa_reserve_minutes=qa_reserve_minutes,
    )

    for idx, slide in enumerate(prs.slides):
        estimate = _estimate_one_slide(slide, idx, is_first=(idx == 0))
        report.per_slide.append(estimate)

    return report


# ---------------------------------------------------------------------------
# Internals
# ---------------------------------------------------------------------------


def _collect_text(slide: Any) -> list[str]:
    """Return non-empty, stripped text lines from every text shape on the slide."""
    out: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = shape.text_frame.text
        except Exception:  # pragma: no cover — defensive
            text = ""
        for line in str(text).splitlines():
            line = line.strip()
            if line:
                out.append(line)
    return out


def _has_picture(slide: Any) -> bool:
    """Return ``True`` when any shape on the slide is a picture (shape_type 13)."""
    return any(getattr(s, "shape_type", None) == 13 for s in slide.shapes)


def _topmost_title(slide: Any, chunks: list[str]) -> str | None:
    """Best-effort title recovery — placeholder, then topmost text shape."""
    try:
        title_shape = slide.shapes.title
    except (AttributeError, KeyError):
        title_shape = None
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        try:
            text = (title_shape.text_frame.text or "").strip()
            if text:
                return text
        except Exception:  # pragma: no cover — defensive
            pass

    # Fall back: topmost text shape on the slide.
    best_top: int | None = None
    best_text: str | None = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = (shape.text_frame.text or "").strip()
        except Exception:  # pragma: no cover — defensive
            continue
        if not text:
            continue
        try:
            top = int(shape.top)
        except (AttributeError, TypeError, ValueError):
            continue
        if best_top is None or top < best_top:
            best_top = top
            best_text = text.splitlines()[0].strip() if text.splitlines() else text.strip()
    if best_text:
        return best_text
    return chunks[0] if chunks else None


def _classify_slide(
    slide: Any, title: str | None, body_chunks: list[str], *, is_first: bool
) -> str:
    """Return the slide kind used to look up base seconds.

    Decision order:

    1. Explicit reference / acknowledgments / discussion title matches.
    2. Methods / equation title hints.
    3. First slide with a short title-only body → ``"title"``.
    4. Picture-bearing slide → ``"figure"`` or ``"figure_heavy"``.
    5. Title-only (no body) non-first slide → ``"section_divider"``.
    6. Body-heavy text slide → ``"bullets"`` or ``"bullets_heavy"``.
    7. Fall through → ``"other"``.
    """
    title_str = (title or "").strip()
    n_bullets = len(body_chunks)

    if title_str:
        if _REFERENCES_RE.match(title_str):
            return "references"
        if _ACK_RE.match(title_str):
            return "acknowledgments"
        if _DISCUSSION_RE.match(title_str):
            return "discussion"
        if _METHODS_RE.search(title_str):
            return "methods"

    if _has_picture(slide):
        # Figure with few or no bullets → figure-heavy data slide.
        if n_bullets <= 1:
            return "figure_heavy"
        return "figure"

    if is_first and title_str and n_bullets <= 3 and len(title_str.split()) <= 6:
        return "title"

    if title_str and n_bullets == 0:
        return "section_divider"

    if n_bullets > 5:
        return "bullets_heavy"

    if n_bullets > 0:
        return "bullets"

    return "other"


def _public_kind(internal_kind: str) -> str:
    """Collapse ``bullets_heavy`` / ``figure_heavy`` to the public kind labels."""
    if internal_kind == "bullets_heavy":
        return "bullets"
    if internal_kind == "figure_heavy":
        return "figure"
    return internal_kind


def _estimate_one_slide(slide: Any, slide_index: int, *, is_first: bool) -> SlideTimeEstimate:
    chunks = _collect_text(slide)
    title = _topmost_title(slide, chunks)
    # Body chunks = everything except the title text.
    body_chunks = [c for c in chunks if c != title] if title else list(chunks)
    internal_kind = _classify_slide(slide, title, body_chunks, is_first=is_first)
    base = _KIND_BASE_SECONDS.get(internal_kind, _KIND_BASE_SECONDS["other"])

    # Add a small bullet-density adjustment for bullets-only slides:
    # 5s per bullet beyond 3, capped at +30s.
    if internal_kind == "bullets" and len(body_chunks) > 3:
        base += min(30, 5 * (len(body_chunks) - 3))

    rationale = (
        f"Classified as '{_public_kind(internal_kind)}' "
        f"({len(body_chunks)} body line{'s' if len(body_chunks) != 1 else ''}, "
        f"figure={'yes' if _has_picture(slide) else 'no'}). "
        f"Base {_KIND_BASE_SECONDS.get(internal_kind, 30)}s."
    )

    return SlideTimeEstimate(
        slide_index=slide_index,
        title=title,
        kind=_public_kind(internal_kind),
        estimate_seconds=base,
        rationale=rationale,
    )
