"""Figure annotation overlays — slide-level circles, rectangles, arrows, labels.

Lifted from ``bobby_slides._annotations`` (bobby-tools, 2026-04). Adds
overlay shapes on top of a figure shape. Coordinates are normalized 0-1
relative to the FIGURE BOUNDING BOX (not the whole slide), so annotations
stay aligned even if the figure size changes.

This is the **slide-level** annotation API used by
:func:`vaultlab.slides.deck.build_from_plan` for plan-driven decks. It is
distinct from ``vaultlab.figures.understand`` (figure-level color-motif
detection / element extraction) — that module produces the structured
``ElementAnnotation`` data that the annotated-figure-slide primitive uses.

Example annotation specs::

    [
        {"type": "circle", "x": 0.3, "y": 0.45, "r": 0.05,
         "color": "FF0000", "click_index": 1, "label": "Signal 2"},
        {"type": "arrow", "from": [0.2, 0.5], "to": [0.4, 0.6],
         "color": "FFEB3B", "click_index": 2},
        {"type": "rect", "bbox": [0.1, 0.1, 0.3, 0.3],
         "color": "00FF00", "click_index": 0},
        {"type": "label", "x": 0.5, "y": 0.85,
         "text": "Key area", "color": "FFFFFF"},
    ]

If ``click_index`` is set, the annotation appears on the (click_index)-th
click during slideshow. Otherwise, the annotation is visible from slide
entry.
"""

from __future__ import annotations

from typing import Any

# Default colors (no leading #)
_DEFAULT_COLOR = "FF0000"  # red
_DEFAULT_LINE_WEIGHT_PT = 3.0
_DEFAULT_LABEL_FONT_SIZE = 16


def add_annotations(
    slide: Any,
    figure_shape: Any,
    annotations: list[dict[str, Any]],
    with_animations: bool = True,
) -> list[Any]:
    """Add overlay annotations on a figure.

    Args:
        slide: python-pptx Slide.
        figure_shape: the picture shape to annotate (provides bbox).
        annotations: list of annotation spec dicts (see module docstring).
        with_animations: if True, ties annotations with click_index to
            entrance animations.

    Returns:
        List of created shape objects (in spec order). May be empty if
        ``figure_shape`` is invalid or annotations list is empty.
    """
    if not annotations or figure_shape is None:
        return []

    fig_left = figure_shape.left
    fig_top = figure_shape.top
    fig_width = figure_shape.width
    fig_height = figure_shape.height

    created: list[Any] = []
    for spec in annotations:
        atype = spec.get("type", "").lower()
        try:
            if atype == "circle":
                shape = _add_circle(slide, spec, fig_left, fig_top, fig_width, fig_height)
            elif atype == "oval":
                shape = _add_oval(slide, spec, fig_left, fig_top, fig_width, fig_height)
            elif atype == "rect":
                shape = _add_rect(slide, spec, fig_left, fig_top, fig_width, fig_height)
            elif atype == "arrow":
                shape = _add_arrow(slide, spec, fig_left, fig_top, fig_width, fig_height)
            elif atype == "label":
                shape = _add_label(slide, spec, fig_left, fig_top, fig_width, fig_height)
            else:
                continue
            if shape is not None:
                created.append(shape)
        except Exception:
            # Annotation is best-effort — never crash the deck
            continue

    if with_animations and created:
        _animate_annotations(slide, list(zip(created, annotations)))

    return created


# ---------------------------------------------------------------------------
# Shape constructors
# ---------------------------------------------------------------------------


def _add_circle(slide, spec, fig_left, fig_top, fig_width, fig_height):
    """Round circle at (x, y) with radius r in figure units (smaller dim)."""
    from pptx.enum.shapes import MSO_SHAPE
    x = float(spec.get("x", 0.5))
    y = float(spec.get("y", 0.5))
    r = float(spec.get("r", 0.05))
    cx = fig_left + int(x * fig_width)
    cy = fig_top + int(y * fig_height)
    diameter = int(2 * r * min(fig_width, fig_height))
    left = cx - diameter // 2
    top = cy - diameter // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    _style_outline_shape(shape, spec)
    return shape


def _add_oval(slide, spec, fig_left, fig_top, fig_width, fig_height):
    """Ellipse covering a normalized bbox. Use for non-square round shapes."""
    from pptx.enum.shapes import MSO_SHAPE
    bbox = spec.get("bbox", [0.1, 0.1, 0.3, 0.3])
    x0, y0, x1, y1 = (float(v) for v in bbox)
    left = fig_left + int(x0 * fig_width)
    top = fig_top + int(y0 * fig_height)
    width = int((x1 - x0) * fig_width)
    height = int((y1 - y0) * fig_height)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    _style_outline_shape(shape, spec)
    return shape


def _add_rect(slide, spec, fig_left, fig_top, fig_width, fig_height):
    from pptx.enum.shapes import MSO_SHAPE
    bbox = spec.get("bbox", [0.1, 0.1, 0.3, 0.3])
    x0, y0, x1, y1 = (float(v) for v in bbox)
    left = fig_left + int(x0 * fig_width)
    top = fig_top + int(y0 * fig_height)
    width = int((x1 - x0) * fig_width)
    height = int((y1 - y0) * fig_height)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _style_outline_shape(shape, spec)
    return shape


def _add_arrow(slide, spec, fig_left, fig_top, fig_width, fig_height):
    from pptx.enum.shapes import MSO_CONNECTOR
    src = spec.get("from", [0.2, 0.2])
    dst = spec.get("to", [0.4, 0.4])
    x0 = fig_left + int(float(src[0]) * fig_width)
    y0 = fig_top + int(float(src[1]) * fig_height)
    x1 = fig_left + int(float(dst[0]) * fig_width)
    y1 = fig_top + int(float(dst[1]) * fig_height)
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0, y0, x1, y1)
    _style_arrow(connector, spec)
    return connector


def _add_label(slide, spec, fig_left, fig_top, fig_width, fig_height):
    """Add a text label.

    placement options:
      - ``"on_figure"`` (default) — x/y normalized 0-1 within figure bbox
      - ``"above"`` — sits above the figure in the slide margin
      - ``"below"`` — sits below the figure (above any caption)
      - ``"left"`` — sits in the left margin alongside the figure
      - ``"right"`` — sits in the right margin alongside the figure
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    text = str(spec.get("text", ""))
    if not text:
        return None

    placement = spec.get("placement", "on_figure")
    width_in = float(spec.get("width_in", 1.5))
    height_in = float(spec.get("height_in", 0.5))
    margin = int(914400 * 0.05)  # 0.05" gap between figure edge and label

    if placement == "above":
        left = fig_left
        top = fig_top - int(914400 * height_in) - margin
        width = fig_width
        height = int(914400 * height_in)
    elif placement == "below":
        left = fig_left
        top = fig_top + fig_height + margin
        width = fig_width
        height = int(914400 * height_in)
    elif placement == "left":
        left = max(0, fig_left - int(914400 * width_in) - margin)
        top = fig_top + fig_height // 2 - int(914400 * height_in) // 2
        width = int(914400 * width_in)
        height = int(914400 * height_in)
    elif placement == "right":
        left = fig_left + fig_width + margin
        top = fig_top + fig_height // 2 - int(914400 * height_in) // 2
        width = int(914400 * width_in)
        height = int(914400 * height_in)
    else:
        x = float(spec.get("x", 0.5))
        y = float(spec.get("y", 0.5))
        cx = fig_left + int(x * fig_width)
        cy = fig_top + int(y * fig_height)
        left = cx - int(914400 * width_in / 2)
        top = cy - int(914400 * height_in / 2)
        width = int(914400 * width_in)
        height = int(914400 * height_in)

    tx = slide.shapes.add_textbox(left, top, width, height)
    tx.text_frame.text = text
    tx.text_frame.word_wrap = True
    color_hex = spec.get("color", "FF5252")
    font_size = int(spec.get("font_size", _DEFAULT_LABEL_FONT_SIZE))
    for para in tx.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            try:
                run.font.color.rgb = RGBColor.from_string(color_hex)
            except Exception:
                pass
        if placement in ("above", "below"):
            try:
                from pptx.enum.text import PP_ALIGN
                para.alignment = PP_ALIGN.CENTER
            except Exception:
                pass
    return tx


# ---------------------------------------------------------------------------
# Style helpers
# ---------------------------------------------------------------------------


def _style_outline_shape(shape, spec):
    """Apply outline-only styling (no fill, colored stroke)."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    color_hex = spec.get("color", _DEFAULT_COLOR)
    weight_pt = float(spec.get("weight_pt", _DEFAULT_LINE_WEIGHT_PT))
    try:
        shape.fill.background()
    except Exception:
        pass
    try:
        shape.line.color.rgb = RGBColor.from_string(color_hex)
        shape.line.width = Pt(weight_pt)
    except Exception:
        pass


def _style_arrow(connector, spec):
    """Apply arrow head + colored stroke to a connector."""
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    color_hex = spec.get("color", _DEFAULT_COLOR)
    weight_pt = float(spec.get("weight_pt", _DEFAULT_LINE_WEIGHT_PT))
    try:
        connector.line.color.rgb = RGBColor.from_string(color_hex)
        connector.line.width = Pt(weight_pt)
    except Exception:
        pass
    try:
        ln = connector.line._get_or_add_ln()
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        head_xml = (
            f'<a:headEnd xmlns:a="{a_ns}" type="none"/>'
            f'<a:tailEnd xmlns:a="{a_ns}" type="triangle" w="med" len="med"/>'
        )
        for elem in (
            ln.findall(f"{{{a_ns}}}headEnd")
            + ln.findall(f"{{{a_ns}}}tailEnd")
        ):
            ln.remove(elem)
        for elem in etree.fromstring(f"<root xmlns:a=\"{a_ns}\">{head_xml}</root>"):
            ln.append(elem)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Animation wiring
# ---------------------------------------------------------------------------


def _animate_annotations(slide, shape_spec_pairs):
    """Tie annotations to click animations based on click_index."""
    from vaultlab.slides.animations import appear_on_click

    for shape, spec in shape_spec_pairs:
        click_index = spec.get("click_index")
        if click_index is None:
            continue
        try:
            appear_on_click(slide, shape, click_index=int(click_index))
        except Exception:
            pass


__all__ = ["add_annotations"]
