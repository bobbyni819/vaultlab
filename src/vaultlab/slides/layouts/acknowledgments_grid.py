"""Acknowledgments-grid slide layout — closing-slide contributor grid.

Use case: the last slide of any deck — list contributors / collaborators /
funders in a clean grid. Each cell shows ``name`` (bold), optional
``role``, and optional ``affiliation``.

Auto-picks a column count based on the number of people:
    1-3   → 1 row
    4     → 2x2
    5-6   → 2x3
    7-9   → 3x3
    10-12 → 3x4
    >12   → 4 columns × ceil(N/4) rows

Honors the vaultlab hard rules: Roboto, min sizes 28/24/18, no overlap.
"""

from __future__ import annotations

from collections.abc import Iterable
from math import ceil
from typing import Any

from vaultlab.slides.layouts._helpers import (
    Inches,
    apply_font,
    ensure_blank_layout,
    sizes,
)


def _pick_grid(n: int) -> tuple[int, int]:
    """Return (rows, cols) for *n* people."""
    if n <= 1:
        return 1, 1
    if n <= 3:
        return 1, n
    if n == 4:
        return 2, 2
    if n <= 6:
        return 2, 3
    if n <= 9:
        return 3, 3
    if n <= 12:
        return 3, 4
    cols = 4
    rows = ceil(n / cols)
    return rows, cols


def add_acknowledgments_grid_slide(
    pres: Any,
    people: Iterable[tuple[str, str, str]],
    title: str = "Acknowledgments",
) -> Any:
    """Add a closing-slide grid of contributors.

    Args:
        pres: python-pptx Presentation.
        people: Iterable of ``(name, role, affiliation)`` tuples. ``role``
            and ``affiliation`` may be empty strings.
        title: Slide title (defaults to ``"Acknowledgments"``).

    Returns:
        The slide object.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(0.9)
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    people_list = list(people)
    if not people_list:
        return slide

    rows, cols = _pick_grid(len(people_list))

    grid_top_in = 1.4
    grid_left_in = 0.5
    grid_width_in = sw_in - 1.0
    grid_height_in = sh_in - grid_top_in - 0.5
    cell_w_in = grid_width_in / cols
    cell_h_in = grid_height_in / rows

    # Pick per-cell font sizes that respect the 18pt caption floor.
    # The grid gets denser as the population grows.
    if rows * cols <= 4:
        name_size = sizes_d["body"]  # 24pt
        sub_size = sizes_d["caption"]  # 18pt
    elif rows * cols <= 9:
        name_size = max(sizes_d["caption"], 20)  # 20pt name, still >= 18
        sub_size = sizes_d["caption"]  # 18pt
    else:
        name_size = sizes_d["caption"]  # 18pt
        sub_size = sizes_d["caption"]  # 18pt

    for i, person in enumerate(people_list[: rows * cols]):
        name, role, affiliation = (list(person) + ["", "", ""])[:3]
        r, c = divmod(i, cols)
        cell_left_in = grid_left_in + c * cell_w_in + 0.05
        cell_top_in = grid_top_in + r * cell_h_in + 0.05
        cell_w_inner = cell_w_in - 0.1
        cell_h_inner = cell_h_in - 0.1

        bx = slide.shapes.add_textbox(
            Inches(cell_left_in),
            Inches(cell_top_in),
            Inches(cell_w_inner),
            Inches(cell_h_inner),
        )
        tf = bx.text_frame
        tf.word_wrap = True
        tf.text = name
        # Apply the name (first paragraph) styling.
        if tf.paragraphs and tf.paragraphs[0].runs:
            for run in tf.paragraphs[0].runs:
                run.font.name = "Roboto"
                from pptx.util import Pt as _Pt

                run.font.size = _Pt(name_size)
                run.font.bold = True

        # Role and affiliation as subsequent paragraphs.
        sub_lines: list[str] = []
        if role:
            sub_lines.append(role)
        if affiliation:
            sub_lines.append(affiliation)
        for line in sub_lines:
            p = tf.add_paragraph()
            p.text = line
            for run in p.runs:
                run.font.name = "Roboto"
                from pptx.util import Pt as _Pt

                run.font.size = _Pt(sub_size)

        # Auto-pick text color from the plain-theme hint (if any).
        try:
            from pptx.dml.color import RGBColor

            plain_theme = getattr(pres, "_vaultlab_plain_theme", None)
            if plain_theme == "dark":
                color = RGBColor(0xFF, 0xFF, 0xFF)
            elif plain_theme == "light":
                color = RGBColor(0x00, 0x00, 0x00)
            else:
                color = None
            if color is not None:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color
        except Exception:
            pass

    return slide


__all__ = ["add_acknowledgments_grid_slide"]
