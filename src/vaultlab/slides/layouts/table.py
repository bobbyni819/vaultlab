"""Tabular-data slide layout — render a list-of-rows or DataFrame.

Use case: experimental conditions, parameter sweeps, comparison data
that is too dense for bullets but too small for an appendix.

Layout:
    - Title at top (heading size, bold).
    - Native python-pptx table below.
    - Header row bold + accent-1 fill, alternating row fills for readability.
    - If the body has more than ``max_body_rows`` rows, the layout
      collapses to a single "see appendix" notice instead of trying to
      cram everything onto one slide.

Honors the vaultlab hard rules: Roboto, min sizes 28/24/18, no overlap.
"""

from __future__ import annotations

from typing import Any

from vaultlab.slides.layouts._helpers import (
    Inches,
    apply_font,
    ensure_blank_layout,
    sizes,
)

# Header fill — Hickey Lab accent1 (teal).
_HEADER_FILL_HEX = "29AF8C"
# Subtle alternating-row fill.
_ALT_ROW_FILL_HEX = "F2F2F2"


def _coerce_rows(rows: Any) -> list[list[str]]:
    """Accept a list-of-rows OR a pandas DataFrame and return list-of-rows.

    Strings are preserved as-is; non-strings are coerced via ``str()``.
    For a DataFrame, the header row is taken from ``df.columns``.
    """
    # pandas DataFrame check (cheap, by attribute — avoids importing pandas).
    if hasattr(rows, "columns") and hasattr(rows, "itertuples"):
        try:
            header = [str(c) for c in rows.columns]
            body = [[str(v) for v in row] for row in rows.itertuples(index=False, name=None)]
            return [header] + body
        except Exception:
            pass
    out: list[list[str]] = []
    for r in rows:
        out.append([str(c) for c in r])
    return out


def add_table_slide(
    pres: Any,
    rows: Any,
    title: str = "",
    max_body_rows: int = 10,
    appendix_message: str = "Table truncated — see appendix for full data",
) -> Any:
    """Add a slide containing a styled table.

    Args:
        pres: python-pptx Presentation.
        rows: Either a list-of-rows (first row treated as header) or a
            pandas DataFrame.
        title: Slide title.
        max_body_rows: Maximum number of body rows (excluding header) to
            display. If exceeded, the table is replaced with the
            ``appendix_message``.
        appendix_message: Notice displayed when the table is too long.

    Returns:
        The slide object.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(0.9)
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    coerced = _coerce_rows(rows)
    if not coerced:
        return slide

    header = coerced[0]
    body = coerced[1:]
    n_cols = len(header)

    # Fallback for oversized tables.
    if len(body) > max_body_rows:
        msg = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(sh_in / 2 - 0.5),
            Inches(sw_in - 1.0),
            Inches(1.0),
        )
        msg.text_frame.text = appendix_message
        msg.text_frame.word_wrap = True
        apply_font(msg.text_frame, size=sizes_d["body"], pres=pres)
        try:
            from pptx.enum.text import PP_ALIGN

            for para in msg.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.italic = True
        except Exception:
            pass
        return slide

    n_rows = 1 + len(body)
    tbl_top_in = 1.5
    tbl_left_in = 0.5
    tbl_width_in = sw_in - 1.0
    # Cap table height to leave bottom breathing room.
    tbl_height_in = min(sh_in - tbl_top_in - 0.6, 0.6 * n_rows + 0.4)

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=n_cols,
        left=Inches(tbl_left_in),
        top=Inches(tbl_top_in),
        width=Inches(tbl_width_in),
        height=Inches(tbl_height_in),
    )
    table = shape.table

    from pptx.dml.color import RGBColor

    header_rgb = RGBColor.from_string(_HEADER_FILL_HEX)
    alt_rgb = RGBColor.from_string(_ALT_ROW_FILL_HEX)

    # Header row.
    for col_i, cell_text in enumerate(header[:n_cols]):
        cell = table.cell(0, col_i)
        cell.text = cell_text
        try:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_rgb
        except Exception:
            pass
        apply_font(
            cell.text_frame,
            size=sizes_d["body"],
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
        )

    # Body rows.
    for row_i, row in enumerate(body, start=1):
        for col_i in range(n_cols):
            cell = table.cell(row_i, col_i)
            cell.text = row[col_i] if col_i < len(row) else ""
            if row_i % 2 == 0:
                # alternate-row shading on even body rows (rows 2, 4, ...)
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_rgb
                except Exception:
                    pass
            apply_font(cell.text_frame, size=sizes_d["caption"], pres=pres)

    return slide


__all__ = ["add_table_slide"]
