"""Shared helpers for the lifted slide-layout primitives.

Lifted from ``bobby_slides._layout`` (bobby-tools, 2026-04). These helpers
encapsulate font / size / color application, blank-layout selection, and
aspect-preserved picture insertion.

Conventions:

- ``pres`` carries optional ``_vaultlab_master_index`` (set by
  :func:`vaultlab.slides.template.load_template`) — used to choose layouts
  from the dark or light master.
- ``pres._vaultlab_plain_theme`` (``"dark"`` / ``"light"`` / ``None``) lets
  the font helper auto-pick text color when no explicit color is supplied.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from vaultlab.slides.template import default_font, min_sizes

_FONT = default_font()
_SIZES = min_sizes()


def Pt(n: int) -> Any:
    """Lazy import of ``pptx.util.Pt`` to avoid module-level ImportError."""
    from pptx.util import Pt as _Pt
    return _Pt(n)


def Inches(n: float) -> Any:
    from pptx.util import Inches as _Inches
    return _Inches(n)


def Emu(n: int) -> Any:
    from pptx.util import Emu as _Emu
    return _Emu(n)


def ensure_blank_layout(pres: Any) -> Any:
    """Return a blank slide layout from the chosen slide master.

    If ``load_template`` set ``pres._vaultlab_master_index``, draw layouts
    from that master (so dark vs light theme is honored). Otherwise fall
    back to ``pres.slide_layouts`` (the active master's layouts).
    """
    master_idx = getattr(pres, "_vaultlab_master_index", None)
    if master_idx is not None and 0 <= master_idx < len(pres.slide_masters):
        layouts = pres.slide_masters[master_idx].slide_layouts
    else:
        layouts = pres.slide_layouts

    for layout in layouts:
        if layout.name.lower() in ("blank", "blank slide"):
            return layout
    # Fallback: last layout in the master is conventionally blank
    return list(layouts)[-1]


def apply_font(
    text_frame: Any,
    size: int,
    bold: bool = False,
    color: Any = None,
    pres: Any = None,
) -> None:
    """Apply Roboto + size + optional bold/color to all runs in a text frame.

    If ``pres`` is provided and has ``_vaultlab_plain_theme`` set (plain
    template mode), text color is auto-set to white on dark or black on
    light when no explicit color is given.
    """
    auto_color = None
    if color is None and pres is not None:
        plain_theme = getattr(pres, "_vaultlab_plain_theme", None)
        if plain_theme == "dark":
            from pptx.dml.color import RGBColor
            auto_color = RGBColor(0xFF, 0xFF, 0xFF)
        elif plain_theme == "light":
            from pptx.dml.color import RGBColor
            auto_color = RGBColor(0x00, 0x00, 0x00)
    final_color = color if color is not None else auto_color

    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.name = _FONT
            run.font.size = Pt(size)
            if bold:
                run.font.bold = True
            if final_color is not None:
                run.font.color.rgb = final_color


def add_picture_fit(
    slide: Any,
    image_path: str,
    left: Any,
    top: Any,
    max_width: Any,
    max_height: Any,
) -> Any | None:
    """Add a picture sized to fit within ``max_width`` × ``max_height``.

    Preserves aspect ratio and centers the image within the bounding box.
    Returns the picture shape, or ``None`` if the image file doesn't exist
    or can't be opened.
    """
    p = Path(image_path)
    if not p.exists():
        return None
    try:
        from PIL import Image as _Image
        with _Image.open(p) as im:
            iw, ih = im.size
        if iw <= 0 or ih <= 0:
            return None
        img_aspect = iw / ih
        # max_width/height are EMU; ratio is unitless
        box_aspect = float(max_width) / float(max_height)
        if img_aspect > box_aspect:
            target_w = max_width
            target_h = int(max_width / img_aspect)
            offset_top = top + (max_height - target_h) // 2
            return slide.shapes.add_picture(
                str(p), left, offset_top, width=target_w, height=target_h
            )
        target_h = max_height
        target_w = int(max_height * img_aspect)
        offset_left = left + (max_width - target_w) // 2
        return slide.shapes.add_picture(
            str(p), offset_left, top, width=target_w, height=target_h
        )
    except Exception:
        return None


def sizes() -> dict[str, int]:
    """Return the min-size dict (heading 28, body 24, caption 18)."""
    return dict(_SIZES)


# Char-width factor for Roboto Bold at heading sizes. Matches the audit's
# heading-tier factor in slides.audit._count_text_overflow_shapes so the
# title-box-height estimator and the overflow-detector agree on whether
# a given title wraps.
_HEADING_CHAR_FACTOR = 0.45  # at ≥22pt; matches audit
_HEADING_LINE_HEIGHT_FACTOR = 1.30  # line-height as fraction of font_size


def estimate_title_lines(
    title: str,
    box_width_in: float,
    *,
    font_size_pt: int = 28,
) -> int:
    """Estimate how many lines the title will wrap to in its box.

    PowerPoint wraps at word boundaries, so we count chars-per-line
    conservatively then divide. Falls back to 1 for very short titles.
    """
    if not title:
        return 1
    char_w_in = font_size_pt * _HEADING_CHAR_FACTOR / 72
    chars_per_line = max(1, int(box_width_in / char_w_in))
    n_lines = (len(title) + chars_per_line - 1) // chars_per_line
    return max(1, n_lines)


def estimate_title_box_height(
    title: str,
    box_width_in: float,
    *,
    font_size_pt: int = 28,
    min_height_in: float = 0.7,
) -> float:
    """Estimate the minimum title-box height (in inches) to fit ``title``.

    Pairs with the layout's figure_top calculation: a short 1-line title
    gives the figure ~0.45 in more vertical space than a 2-line title.
    Bobby's 2026-05-04 ask: "for slide 8 the title was shorter and then
    the figure is a flat rectangle and you should be able to stretch
    that figure up to encroach in the unused empty space."
    """
    n_lines = estimate_title_lines(title, box_width_in, font_size_pt=font_size_pt)
    line_height_in = (font_size_pt * _HEADING_LINE_HEIGHT_FACTOR) / 72.0
    needed = n_lines * line_height_in + 0.10  # small buffer
    return max(min_height_in, needed)


__all__ = [
    "Emu",
    "Inches",
    "Pt",
    "add_picture_fit",
    "apply_font",
    "ensure_blank_layout",
    "estimate_title_box_height",
    "estimate_title_lines",
    "sizes",
]
