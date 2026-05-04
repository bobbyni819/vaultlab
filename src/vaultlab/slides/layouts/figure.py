"""Figure-slide layouts — single figure variants.

Lifted from ``bobby_slides._layout`` (bobby-tools, 2026-04). Provides:

- :func:`add_figure_slide` — figure-with-bullets (default workhorse).
- :func:`add_figure_only_slide` — full-width hero figure, no bullets.
- :func:`add_figure_above_bullets_slide` — vertical-split: figure on top,
  bullets full-width below.
- :func:`add_two_figure_compare_slide` — side-by-side comparison.
- :func:`add_quote_slide` — large centered quote (handy transition slide).

Each layout returns the slide object so callers can attach annotations or
animations afterward.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable

from vaultlab.slides.layouts._helpers import (
    Inches,
    Pt,
    add_picture_fit,
    apply_font,
    ensure_blank_layout,
    sizes,
)


def add_figure_slide(
    pres: Any,
    image_path: str | Path,
    title: str = "",
    caption: str = "",
    bullets: Iterable[str] | None = None,
    citation_source: str = "",
) -> Any:
    """Add a slide with a single large figure, title, caption, and optional bullets.

    Layout:
      - Title at top (heading size).
      - Figure centered, taking most of the slide.
      - Caption below figure (12pt italic).
      - Optional bullets to the right of figure (body size).
      - Citation source as 9pt footer.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title

        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    img_path = Path(image_path)
    bullets_list = list(bullets) if bullets else []
    has_bullets = len(bullets_list) > 0

    fig_top_in = 1.1
    cap_height_in = 0.6 if caption else 0.0  # 12pt × ~2 lines + buffer
    cit_height_in = 0.4 if citation_source else 0.0
    cap_gap = 0.05 if caption else 0.0
    fig_height_in = sh_in - fig_top_in - cap_gap - cap_height_in - cit_height_in
    if fig_height_in < 2.5:
        fig_height_in = 2.5  # never let figure get crushed

    fig_left = Inches(0.3)
    fig_top = Inches(fig_top_in)
    fig_width = Inches(sw_in * 0.60) if has_bullets else Inches(sw_in - 0.6)
    fig_height = Inches(fig_height_in)

    if not has_bullets:
        fig_left = Inches((sw_in - fig_width / 914400) / 2)

    add_picture_fit(slide, str(img_path), fig_left, fig_top, fig_width, fig_height)

    if caption:
        cap_top = Inches(fig_top_in + fig_height_in + cap_gap)
        cap_left = fig_left if has_bullets else Inches(0.5)
        cap_width = fig_width if has_bullets else Inches(sw_in - 1.0)
        cx = slide.shapes.add_textbox(cap_left, cap_top, cap_width, Inches(cap_height_in))
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=12, pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True

    if has_bullets:
        fig_right_in = (fig_left / 914400) + (fig_width / 914400)
        slide_right_margin = 0.3
        avail_left = fig_right_in + 0.4
        avail_right = sw_in - slide_right_margin
        bul_width_in = max(2.5, avail_right - avail_left)
        bul_left = Inches(avail_left)

        bul_top_in = fig_top_in + 0.3
        bul_height_in = fig_height_in - 0.3
        bx = slide.shapes.add_textbox(
            bul_left, Inches(bul_top_in),
            Inches(bul_width_in), Inches(bul_height_in),
        )
        tf = bx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets_list):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = Pt(8)
                p.space_after = Pt(4)
            except Exception:
                pass
        apply_font(tf, size=sizes_d["body"], pres=pres)

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.4),
            Inches(sw_in - 0.6), Inches(0.3),
        )
        cit.text_frame.text = citation_source
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_figure_only_slide(
    pres: Any,
    image_path: str | Path,
    title: str = "",
    caption: str = "",
    citation_source: str = "",
) -> Any:
    """Full-width centered figure with title above, no bullets.

    Use for hero figures or when you want the image to dominate completely.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title

        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    cap_height_in = 0.6 if caption else 0.0  # 12pt × ~2 lines + buffer
    cit_height_in = 0.4 if citation_source else 0.0
    cap_gap = 0.1 if caption else 0.0
    fig_top_in = 1.2
    fig_height_in = sh_in - fig_top_in - cap_gap - cap_height_in - cit_height_in - 0.1

    fig_width_in = sw_in - 1.0
    fig_left = Inches(0.5)
    fig_top = Inches(fig_top_in)

    img_path = Path(image_path)
    add_picture_fit(
        slide, str(img_path), fig_left, fig_top,
        Inches(fig_width_in), Inches(fig_height_in),
    )

    if caption:
        cap_top = Inches(fig_top_in + fig_height_in + cap_gap)
        cx = slide.shapes.add_textbox(
            fig_left, cap_top, Inches(fig_width_in), Inches(cap_height_in)
        )
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=12, pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True
            try:
                from pptx.enum.text import PP_ALIGN
                para.alignment = PP_ALIGN.CENTER
            except Exception:
                pass

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.4),
            Inches(sw_in - 0.6), Inches(0.3),
        )
        cit.text_frame.text = citation_source
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_figure_above_bullets_slide(
    pres: Any,
    image_path: str | Path,
    title: str = "",
    bullets: Iterable[str] | None = None,
    caption: str = "",
    citation_source: str = "",
) -> Any:
    """Vertical-split layout: figure on top half, bullets on bottom half."""
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title

        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    bullets_list = list(bullets) if bullets else []
    has_bullets = len(bullets_list) > 0
    cit_height_in = 0.4 if citation_source else 0.0

    fig_top_in = 1.1
    fig_height_in = sh_in * 0.50
    fig_width_in = sw_in - 1.0
    fig_left = Inches(0.5)

    img_path = Path(image_path)
    add_picture_fit(
        slide, str(img_path), fig_left, Inches(fig_top_in),
        Inches(fig_width_in), Inches(fig_height_in),
    )

    if caption:
        cap_top = Inches(fig_top_in + fig_height_in + 0.05)
        cx = slide.shapes.add_textbox(
            fig_left, cap_top, Inches(fig_width_in), Inches(0.35)
        )
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=12, pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True

    if has_bullets:
        bul_top_in = fig_top_in + fig_height_in + (0.45 if caption else 0.1)
        bul_height_in = sh_in - bul_top_in - cit_height_in - 0.1
        bx = slide.shapes.add_textbox(
            Inches(0.5), Inches(bul_top_in),
            Inches(sw_in - 1.0), Inches(bul_height_in),
        )
        tf = bx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets_list):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = Pt(6)
            except Exception:
                pass
        apply_font(tf, size=20, pres=pres)

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.4),
            Inches(sw_in - 0.6), Inches(0.3),
        )
        cit.text_frame.text = citation_source
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_two_figure_compare_slide(
    pres: Any,
    left_image: str | Path,
    right_image: str | Path,
    title: str = "",
    left_label: str = "",
    right_label: str = "",
    left_caption: str = "",
    right_caption: str = "",
    citation_source: str = "",
) -> Any:
    """Side-by-side comparison of two figures with optional labels and captions."""
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title

        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    cit_height_in = 0.4 if citation_source else 0.0
    fig_top_in = 1.5
    fig_height_in = sh_in - fig_top_in - 1.2 - cit_height_in
    half_width_in = (sw_in - 1.5) / 2

    for i, (img, label, caption) in enumerate([
        (left_image, left_label, left_caption),
        (right_image, right_label, right_caption),
    ]):
        x_left = 0.5 + i * (half_width_in + 0.5)
        if label:
            lbl = slide.shapes.add_textbox(
                Inches(x_left), Inches(fig_top_in - 0.5),
                Inches(half_width_in), Inches(0.4),
            )
            lbl.text_frame.text = label
            apply_font(lbl.text_frame, size=20, bold=True, pres=pres)

        add_picture_fit(
            slide, str(img),
            Inches(x_left), Inches(fig_top_in),
            Inches(half_width_in), Inches(fig_height_in),
        )

        if caption:
            cx = slide.shapes.add_textbox(
                Inches(x_left), Inches(fig_top_in + fig_height_in + 0.1),
                Inches(half_width_in), Inches(0.6),
            )
            cx.text_frame.text = caption
            cx.text_frame.word_wrap = True
            apply_font(cx.text_frame, size=12, pres=pres)
            for para in cx.text_frame.paragraphs:
                for run in para.runs:
                    run.font.italic = True

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.4),
            Inches(sw_in - 0.6), Inches(0.3),
        )
        cit.text_frame.text = citation_source
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_figure_with_side_caption_slide(
    pres: Any,
    image_path: str | Path,
    title: str = "",
    caption: str = "",
    citation_source: str = "",
    bullets: Iterable[str] | None = None,
) -> Any:
    """Wide figure on left, caption + citation + optional bullets on right.

    Use when a figure is wide-aspect (>1.4) and there's room to recover
    by putting the caption + citation in the right-side gutter instead
    of below. Lets the figure use the full slide HEIGHT (under the title)
    rather than being squeezed by the bottom caption row.

    Layout:
      - Title at top (28pt, full width) — height 1.2 in
      - Figure on left, full available height, ~70% of width
      - Right-side text column: caption (12pt italic) + bullets (20pt)
        + citation_source (9pt) at the bottom of the column
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title

        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    # Figure occupies left 62%, near-full height under title.
    # Tightened 2026-05-04: fig_top 1.6→1.4, bot_margin 0.3→0.15 to
    # squeeze every inch of vertical space for the figure (since square
    # figures here are height-bound).
    fig_top_in = 1.4
    fig_height_in = sh_in - fig_top_in - 0.15
    fig_width_in = sw_in * 0.62
    fig_left = Inches(0.4)

    img_path = Path(image_path)
    add_picture_fit(
        slide, str(img_path), fig_left, Inches(fig_top_in),
        Inches(fig_width_in), Inches(fig_height_in),
    )

    # Right-side column: caption (top, biggest) / bullets (middle) /
    # citation (bottom). Bobby's 2026-05-04 ask: "still put the biggest
    # caption on top, put the small figure caption like below, then put
    # the citation on the very bottom."
    col_left_in = 0.4 + fig_width_in + 0.3
    col_width_in = sw_in - col_left_in - 0.4

    bullets_list = list(bullets) if bullets else []
    has_bullets = len(bullets_list) > 0

    # Caption at top of right column — sized to actually-needed height
    # (1 line at 12pt × ~80 chars typical). Was 1.2 in (waste); now 0.7 in.
    cap_top_in = fig_top_in
    cap_height_in = 0.7 if caption else 0.0
    if caption:
        cx = slide.shapes.add_textbox(
            Inches(col_left_in), Inches(cap_top_in),
            Inches(col_width_in), Inches(cap_height_in),
        )
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=12, pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True

    if has_bullets:
        bul_top_in = cap_top_in + cap_height_in + 0.15
        bul_height_in = sh_in - bul_top_in - 0.55  # leave room for citation
        bx = slide.shapes.add_textbox(
            Inches(col_left_in), Inches(bul_top_in),
            Inches(col_width_in), Inches(bul_height_in),
        )
        tf = bx.text_frame
        tf.word_wrap = True
        from pptx.util import Pt as _Pt
        for i, b in enumerate(bullets_list):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = _Pt(8)
                p.space_after = _Pt(4)
            except Exception:
                pass
        apply_font(tf, size=20, pres=pres)

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(col_left_in), Inches(sh_in - 0.45),
            Inches(col_width_in), Inches(0.35),
        )
        cit.text_frame.text = citation_source
        cit.text_frame.word_wrap = True
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_figure_top_caption_br_slide(
    pres: Any,
    image_path: str | Path,
    title: str = "",
    caption: str = "",
    citation_source: str = "",
    bullets: Iterable[str] | None = None,
) -> Any:
    """Wide-flat figure on top + caption/citation in bottom-right corner.

    For figures with aspect ratio 1.5–3.0 (flat horizontal rectangles).
    Layout maximizes figure width by putting caption + citation in the
    bottom-right corner — leaving the rest of the bottom strip empty,
    or filled with bullets on the left.

    Layout:
      - Title at top (28pt, full width) — height 1.2 in
      - Figure: full slide width, takes upper ~70% of remaining height
      - Bottom strip:
          - Bullets (if any): bottom-left, ~50% width
          - Caption + citation: bottom-right, ~45% width

    The figure can be much larger (full slide width × ~70% height) than
    the default figure_slide layout (60% width × 75% height) because the
    bottom-right caption doesn't compete with the figure for horizontal
    space.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(1.2)
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    bullets_list = list(bullets) if bullets else []
    has_bullets = len(bullets_list) > 0

    # Figure occupies full slide width × upper portion of available height.
    # Bottom strip height adapts to bullet count. Tightened 2026-05-04
    # to give the figure more vertical space (Bobby: "shift them down a
    # bit still and make that rectangle a bit larger").
    fig_top_in = 1.4  # was 1.6
    n_bullets = len(bullets_list)
    if n_bullets >= 4:
        bottom_strip_height_in = 2.0  # was 2.4
    elif n_bullets >= 2:
        bottom_strip_height_in = 1.7  # was 2.0
    else:
        bottom_strip_height_in = 1.3  # was 1.5
    fig_height_in = sh_in - fig_top_in - bottom_strip_height_in - 0.15
    fig_width_in = sw_in - 0.6
    fig_left = Inches(0.3)

    img_path = Path(image_path)
    add_picture_fit(
        slide, str(img_path), fig_left, Inches(fig_top_in),
        Inches(fig_width_in), Inches(fig_height_in),
    )

    bottom_top_in = fig_top_in + fig_height_in + 0.15

    # Bullets on bottom-left (when present). Reduced paragraph spacing
    # so 4 bullets fit in the tightened 2.0-in bottom strip.
    if has_bullets:
        bul_width_in = sw_in * 0.50
        bx = slide.shapes.add_textbox(
            Inches(0.5), Inches(bottom_top_in),
            Inches(bul_width_in), Inches(bottom_strip_height_in - 0.05),
        )
        tf = bx.text_frame
        tf.word_wrap = True
        from pptx.util import Pt as _Pt
        for i, b in enumerate(bullets_list):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = _Pt(2)
                p.space_after = _Pt(1)
            except Exception:
                pass
        apply_font(tf, size=18, pres=pres)

    # Caption + citation in bottom-right, tightly stacked. Caption at top
    # of bottom-right column gets only ~0.7 in (was full strip height —
    # huge waste). Citation snug at the very bottom of the slide.
    br_left_in = sw_in * 0.55
    br_width_in = sw_in - br_left_in - 0.3

    if caption:
        cx = slide.shapes.add_textbox(
            Inches(br_left_in), Inches(bottom_top_in),
            Inches(br_width_in), Inches(0.7),
        )
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=12, pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(br_left_in), Inches(sh_in - 0.45),
            Inches(br_width_in), Inches(0.35),
        )
        cit.text_frame.text = citation_source
        cit.text_frame.word_wrap = True
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_quote_slide(
    pres: Any,
    quote: str,
    attribution: str = "",
) -> Any:
    """Big centered quote with optional attribution. Useful for transitions."""
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400

    quote_top_in = sh_in * 0.30
    qx = slide.shapes.add_textbox(
        Inches(1.0), Inches(quote_top_in),
        Inches(sw_in - 2.0), Inches(sh_in * 0.40),
    )
    qx.text_frame.text = f'“{quote}”'
    qx.text_frame.word_wrap = True
    apply_font(qx.text_frame, size=36, bold=False, pres=pres)
    for para in qx.text_frame.paragraphs:
        for run in para.runs:
            run.font.italic = True
        try:
            from pptx.enum.text import PP_ALIGN
            para.alignment = PP_ALIGN.CENTER
        except Exception:
            pass

    if attribution:
        ax = slide.shapes.add_textbox(
            Inches(1.0), Inches(quote_top_in + sh_in * 0.40 + 0.3),
            Inches(sw_in - 2.0), Inches(0.6),
        )
        ax.text_frame.text = f"— {attribution}"
        apply_font(ax.text_frame, size=20, pres=pres)
        for para in ax.text_frame.paragraphs:
            try:
                from pptx.enum.text import PP_ALIGN
                para.alignment = PP_ALIGN.CENTER
            except Exception:
                pass

    return slide


__all__ = [
    "add_figure_above_bullets_slide",
    "add_figure_only_slide",
    "add_figure_slide",
    "add_figure_top_caption_br_slide",
    "add_figure_with_side_caption_slide",
    "add_quote_slide",
    "add_two_figure_compare_slide",
]
