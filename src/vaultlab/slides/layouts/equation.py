"""Equation slide layout — center a single equation prominently.

Designed for math-heavy lab meetings where a single equation is the
focus of the slide. The equation is rendered as plain text at a large
font size; LaTeX-rendering to image is out of scope for this primitive
(callers who need pixel-perfect TeX should render their own image and
fall back to :func:`add_figure_only_slide`).

Honors the vaultlab slide hard rules: Roboto font, min sizes
heading 28 / body 24 / caption 18, descriptive sentence titles.
"""

from __future__ import annotations

from typing import Any

from vaultlab.slides.layouts._helpers import (
    Inches,
    apply_font,
    ensure_blank_layout,
    sizes,
)


def add_equation_slide(
    pres: Any,
    equation: str,
    title: str = "",
    caption: str = "",
) -> Any:
    """Add a slide with a centered equation, descriptive title, and caption.

    Layout (top → bottom):
        - Title (heading size, bold) at top.
        - Equation centered in the middle, large font (44pt).
        - Optional caption (caption size, italic) below the equation.

    Args:
        pres: python-pptx Presentation.
        equation: Equation text (plain string; LaTeX-rendered images can be
            passed to :func:`add_figure_only_slide` instead).
        title: Slide title (descriptive sentence preferred).
        caption: Optional caption below the equation.

    Returns:
        The slide object.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    # Title at top.
    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(0.9)
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    # Equation block — vertically centered between title and caption.
    cap_height_in = 0.8 if caption else 0.0
    eq_top_in = 1.5
    eq_height_in = sh_in - eq_top_in - cap_height_in - 0.6
    if eq_height_in < 2.0:
        eq_height_in = 2.0

    eq = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(eq_top_in),
        Inches(sw_in - 1.0),
        Inches(eq_height_in),
    )
    eq.text_frame.text = equation
    eq.text_frame.word_wrap = True
    apply_font(eq.text_frame, size=44, bold=False, pres=pres)
    try:
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

        eq.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for para in eq.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
    except Exception:
        pass

    # Caption below.
    if caption:
        cap_top_in = eq_top_in + eq_height_in + 0.15
        cx = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(cap_top_in),
            Inches(sw_in - 2.0),
            Inches(cap_height_in - 0.1),
        )
        cx.text_frame.text = caption
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=sizes_d["caption"], pres=pres)
        for para in cx.text_frame.paragraphs:
            for run in para.runs:
                run.font.italic = True
            try:
                from pptx.enum.text import PP_ALIGN

                para.alignment = PP_ALIGN.CENTER
            except Exception:
                pass

    return slide


__all__ = ["add_equation_slide"]
