"""Comparison-table slide layout — two columns side-by-side with bullets.

Use case: Approach A vs Approach B, Before/After, or any pairwise
trade-off discussion in journal clubs. An optional centered "key
insight" callout sits between the columns at the bottom.

Honors the vaultlab hard rules: Roboto, min sizes 28/24/18, no overlap,
descriptive sentence titles.
"""

from __future__ import annotations

from collections.abc import Iterable
from typing import Any

from vaultlab.slides.layouts._helpers import (
    Inches,
    Pt,
    apply_font,
    ensure_blank_layout,
    sizes,
)


def add_comparison_table_slide(
    pres: Any,
    left_header: str,
    right_header: str,
    left_bullets: Iterable[str],
    right_bullets: Iterable[str],
    title: str = "",
    key_insight: str = "",
) -> Any:
    """Add a two-column comparison slide.

    Layout:
        - Title at top (heading size, bold).
        - Two equal-width columns, each with a header + bullet list.
        - Headers in body-size bold; bullets in caption-size.
        - Optional centered "key insight" callout at the bottom.

    Args:
        pres: python-pptx Presentation.
        left_header: Title of the left column (e.g. "Approach A", "Before").
        right_header: Title of the right column.
        left_bullets: Bullet content for the left column.
        right_bullets: Bullet content for the right column.
        title: Overall slide title (descriptive sentence preferred).
        key_insight: Optional callout text shown at the bottom of the slide.

    Returns:
        The slide object.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(sw_in - 1.0), Inches(0.9)
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)

    callout_height_in = 0.9 if key_insight else 0.0
    callout_gap_in = 0.2 if key_insight else 0.0

    # Column layout: equal widths, centered gap.
    col_top_in = 1.4
    col_bottom_in = sh_in - callout_height_in - callout_gap_in - 0.3
    col_height_in = col_bottom_in - col_top_in
    gap_in = 0.4
    side_margin_in = 0.5
    col_width_in = (sw_in - 2 * side_margin_in - gap_in) / 2

    header_height_in = 0.6

    left_lst = list(left_bullets)
    right_lst = list(right_bullets)

    # --- Left column ---
    lx_left = side_margin_in
    hl = slide.shapes.add_textbox(
        Inches(lx_left),
        Inches(col_top_in),
        Inches(col_width_in),
        Inches(header_height_in),
    )
    hl.text_frame.text = left_header
    apply_font(hl.text_frame, size=sizes_d["body"], bold=True, pres=pres)

    if left_lst:
        bl = slide.shapes.add_textbox(
            Inches(lx_left),
            Inches(col_top_in + header_height_in + 0.05),
            Inches(col_width_in),
            Inches(col_height_in - header_height_in - 0.05),
        )
        tf = bl.text_frame
        tf.word_wrap = True
        for i, b in enumerate(left_lst):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = Pt(6)
                p.space_after = Pt(4)
            except Exception:
                pass
        apply_font(tf, size=sizes_d["caption"], pres=pres)

    # --- Right column ---
    rx_left = side_margin_in + col_width_in + gap_in
    hr = slide.shapes.add_textbox(
        Inches(rx_left),
        Inches(col_top_in),
        Inches(col_width_in),
        Inches(header_height_in),
    )
    hr.text_frame.text = right_header
    apply_font(hr.text_frame, size=sizes_d["body"], bold=True, pres=pres)

    if right_lst:
        br = slide.shapes.add_textbox(
            Inches(rx_left),
            Inches(col_top_in + header_height_in + 0.05),
            Inches(col_width_in),
            Inches(col_height_in - header_height_in - 0.05),
        )
        tf = br.text_frame
        tf.word_wrap = True
        for i, b in enumerate(right_lst):
            text = f"•  {b}"
            if i == 0:
                tf.text = text
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = text
            try:
                p.space_before = Pt(6)
                p.space_after = Pt(4)
            except Exception:
                pass
        apply_font(tf, size=sizes_d["caption"], pres=pres)

    # --- Key insight callout ---
    if key_insight:
        cal_top_in = col_bottom_in + callout_gap_in
        cx = slide.shapes.add_textbox(
            Inches(side_margin_in + 0.5),
            Inches(cal_top_in),
            Inches(sw_in - 2 * side_margin_in - 1.0),
            Inches(callout_height_in - 0.1),
        )
        cx.text_frame.text = f"Key insight:  {key_insight}"
        cx.text_frame.word_wrap = True
        apply_font(cx.text_frame, size=sizes_d["body"], bold=True, pres=pres)
        try:
            from pptx.enum.text import PP_ALIGN

            for para in cx.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
        except Exception:
            pass

    return slide


__all__ = ["add_comparison_table_slide"]
