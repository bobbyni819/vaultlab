"""Analogy slide — explain a hard scientific concept via a familiar one.

Bobby's 2026-05-04 ask: not every slide needs a publication figure. For
hard scientific concepts, an analogy slide that maps a familiar idea to
the technical idea (with optional small icon images on each side) is
often more communicative than a 6-panel figure.

Two variants:

- ``add_analogy_slide``: side-by-side layout. Left = familiar concept
  (e.g., "Like a postal sorting room"). Right = scientific concept
  (e.g., "T-cell receptor + MHC + peptide presentation"). Optional
  small icon image per side. Connecting "is like" arrow in the middle.

- ``add_analogy_stacked_slide``: vertical-stack variant for taller
  concept blocks. Familiar on top, scientific mapping below.

Use these instead of figure_only when:

- The concept doesn't have a single canonical publication figure
- The concept is best explained by mapping to a non-scientific everyday thing
- The audience needs intuition, not data
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from vaultlab.slides.layouts._helpers import (
    Inches,
    add_picture_fit,
    apply_font,
    ensure_blank_layout,
    estimate_title_box_height,
    sizes,
)


def add_analogy_slide(
    pres: Any,
    title: str = "",
    *,
    familiar_label: str = "",        # e.g., "Postal sorting room"
    familiar_body: str = "",         # e.g., "Letters arrive..., sorters..."
    familiar_image: str | Path | None = None,
    scientific_label: str = "",      # e.g., "T-cell antigen presentation"
    scientific_body: str = "",       # e.g., "Peptides loaded onto MHC..."
    scientific_image: str | Path | None = None,
    arrow_text: str = "is like",     # connector
    citation_source: str = "",
) -> Any:
    """Side-by-side analogy slide: familiar (left) ↔ scientific (right).

    Layout::

        ─────────────────────────────────────
        |              Title                 |
        ├──────────────┬──────────────────────┤
        | FAMILIAR     |  is like  | SCIENCE  |
        | label        |    →      | label    |
        | [icon]       |           | [icon]   |
        | body         |           | body     |
        ├──────────────┴───────────┴──────────┤
        |             citation                |
        ─────────────────────────────────────
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        title_h_in = estimate_title_box_height(title, sw_in - 1.0)
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(sw_in - 1.0), Inches(title_h_in),
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)
    else:
        title_h_in = 0.0

    body_top_in = 0.3 + title_h_in + 0.2
    body_height_in = sh_in - body_top_in - 0.7  # leave room for citation

    # Three-column layout: left (familiar) / middle (arrow) / right (scientific)
    arrow_w_in = 1.4
    side_w_in = (sw_in - 1.0 - arrow_w_in) / 2  # 0.5 margin each side + middle
    left_x = Inches(0.5)
    arrow_x = Inches(0.5 + side_w_in)
    right_x = Inches(0.5 + side_w_in + arrow_w_in)

    def _zone(x_inch_emu, label, body, image):
        """Build a zone with label (top), optional image, body (bottom)."""
        zone_y_in = body_top_in
        # Label at top — 24pt bold
        if label:
            lbl = slide.shapes.add_textbox(
                x_inch_emu, Inches(zone_y_in),
                Inches(side_w_in), Inches(0.7),
            )
            lbl.text_frame.text = label
            lbl.text_frame.word_wrap = True
            apply_font(lbl.text_frame, size=sizes_d["body"], bold=True, pres=pres)
            zone_y_in += 0.8
        # Image (icon) — small, centered horizontally in zone
        if image and Path(image).exists():
            img_h_in = 2.2
            add_picture_fit(
                slide, str(image),
                x_inch_emu, Inches(zone_y_in),
                Inches(side_w_in), Inches(img_h_in),
            )
            zone_y_in += img_h_in + 0.15
        # Body text — fills remaining space
        if body:
            remaining = body_top_in + body_height_in - zone_y_in
            if remaining < 1.0:
                remaining = 1.0
            bx = slide.shapes.add_textbox(
                x_inch_emu, Inches(zone_y_in),
                Inches(side_w_in), Inches(remaining),
            )
            bx.text_frame.text = body
            bx.text_frame.word_wrap = True
            apply_font(bx.text_frame, size=20, pres=pres)

    _zone(left_x, familiar_label, familiar_body, familiar_image)
    _zone(right_x, scientific_label, scientific_body, scientific_image)

    # Middle column — arrow + connector text
    arrow_y_in = body_top_in + body_height_in / 2 - 0.6
    ax = slide.shapes.add_textbox(
        arrow_x, Inches(arrow_y_in),
        Inches(arrow_w_in), Inches(1.2),
    )
    ax.text_frame.text = f"{arrow_text}\n→"
    ax.text_frame.word_wrap = True
    apply_font(ax.text_frame, size=28, bold=True, pres=pres)
    try:
        from pptx.enum.text import PP_ALIGN
        for para in ax.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
    except Exception:
        pass

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.45),
            Inches(sw_in - 0.6), Inches(0.35),
        )
        cit.text_frame.text = citation_source
        cit.text_frame.word_wrap = True
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


def add_analogy_stacked_slide(
    pres: Any,
    title: str = "",
    *,
    familiar_label: str = "",
    familiar_body: str = "",
    familiar_image: str | Path | None = None,
    scientific_label: str = "",
    scientific_body: str = "",
    scientific_image: str | Path | None = None,
    citation_source: str = "",
) -> Any:
    """Stacked analogy variant: familiar on top half, scientific on bottom half.

    Use when each concept needs more vertical room than the side-by-side
    variant allows.
    """
    slide = pres.slides.add_slide(ensure_blank_layout(pres))
    sw_in = pres.slide_width / 914400
    sh_in = pres.slide_height / 914400
    sizes_d = sizes()

    if title:
        title_h_in = estimate_title_box_height(title, sw_in - 1.0)
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(sw_in - 1.0), Inches(title_h_in),
        )
        tx.text_frame.text = title
        tx.text_frame.word_wrap = True
        apply_font(tx.text_frame, size=sizes_d["heading"], bold=True, pres=pres)
    else:
        title_h_in = 0.0

    body_top_in = 0.3 + title_h_in + 0.15
    body_height_in = sh_in - body_top_in - 0.6
    half_h_in = (body_height_in - 0.4) / 2  # 0.4 gap between zones

    def _zone(y_inch, label, body, image):
        if label:
            lbl = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_inch),
                Inches(sw_in - 1.0), Inches(0.5),
            )
            lbl.text_frame.text = label
            apply_font(lbl.text_frame, size=sizes_d["body"], bold=True, pres=pres)
            y_inch += 0.55
        # Two-column inside zone: image left, body right (if image)
        if image and Path(image).exists():
            img_w_in = sw_in * 0.30
            add_picture_fit(
                slide, str(image),
                Inches(0.5), Inches(y_inch),
                Inches(img_w_in), Inches(half_h_in - 0.5),
            )
            body_left = 0.5 + img_w_in + 0.3
            body_w = sw_in - body_left - 0.5
        else:
            body_left = 0.5
            body_w = sw_in - 1.0
        if body:
            bx = slide.shapes.add_textbox(
                Inches(body_left), Inches(y_inch),
                Inches(body_w), Inches(half_h_in - 0.5),
            )
            bx.text_frame.text = body
            bx.text_frame.word_wrap = True
            apply_font(bx.text_frame, size=18, pres=pres)

    _zone(body_top_in, familiar_label, familiar_body, familiar_image)
    _zone(body_top_in + half_h_in + 0.4, scientific_label, scientific_body, scientific_image)

    if citation_source:
        cit = slide.shapes.add_textbox(
            Inches(0.3), Inches(sh_in - 0.45),
            Inches(sw_in - 0.6), Inches(0.35),
        )
        cit.text_frame.text = citation_source
        cit.text_frame.word_wrap = True
        apply_font(cit.text_frame, size=9, pres=pres)

    return slide


__all__ = ["add_analogy_slide", "add_analogy_stacked_slide"]
