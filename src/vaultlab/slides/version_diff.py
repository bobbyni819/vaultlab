"""Slide-version diff between two ``.pptx`` files.

Compares two rendered decks slide-by-slide and reports which slides are
added, removed, modified, or unchanged. Modification records carry
field-level change tuples so callers can show "shape X went from A to B."

Matching strategy (best of three, in order of preference):

1. **Title-stable match.** A slide's recovered title is hashed and used
   as the matching key. Decks that ship with descriptive sentence
   titles (the lab convention) get a stable match even when slides are
   reordered.
2. **Body-fingerprint fallback.** When titles are missing or duplicated,
   the slide's normalized body text fingerprint is used as a tiebreaker.
3. **Position fallback.** Anything still unmatched is matched by
   zero-indexed position — the safe default for "did the third slide
   change?" workflows.

Field-level changes are computed shape-by-shape. Each text shape's
concatenated text is compared. Picture shapes record an
``"image-changed"`` marker when their MD5 differs between decks.

Public entrypoint: :func:`diff_decks` returning :class:`DeckDiff`.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

__all__ = [
    "DeckDiff",
    "SlideDiff",
    "diff_decks",
]


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass
class SlideDiff:
    """Diff record for a single slide.

    Attributes
    ----------
    slide_index_a : int | None
        Zero-indexed position in deck A. ``None`` when the slide was
        added in deck B and has no counterpart in A.
    slide_index_b : int | None
        Zero-indexed position in deck B. ``None`` when the slide was
        removed from B and existed only in A.
    change_kind : str
        ``"added"`` / ``"removed"`` / ``"modified"`` / ``"unchanged"``.
    field_changes : list[tuple[str, str, str]]
        Field-level change tuples for modified slides: ``(shape_id,
        before, after)``. Empty for added / removed / unchanged.
    """

    slide_index_a: int | None
    slide_index_b: int | None
    change_kind: str
    field_changes: list[tuple[str, str, str]] = field(default_factory=list)


@dataclass
class DeckDiff:
    """Diff between two rendered decks."""

    deck_a: Path
    deck_b: Path
    slides: list[SlideDiff] = field(default_factory=list)

    @property
    def n_added(self) -> int:
        return sum(1 for s in self.slides if s.change_kind == "added")

    @property
    def n_removed(self) -> int:
        return sum(1 for s in self.slides if s.change_kind == "removed")

    @property
    def n_modified(self) -> int:
        return sum(1 for s in self.slides if s.change_kind == "modified")

    @property
    def n_unchanged(self) -> int:
        return sum(1 for s in self.slides if s.change_kind == "unchanged")

    def summary_lines(self) -> list[str]:
        return [
            f"Deck diff: {self.deck_a.name} → {self.deck_b.name}",
            (
                f"  {self.n_added} added, {self.n_removed} removed, "
                f"{self.n_modified} modified, {self.n_unchanged} unchanged."
            ),
        ]


# ---------------------------------------------------------------------------
# Public entrypoint
# ---------------------------------------------------------------------------


def diff_decks(pptx_a: Path | str, pptx_b: Path | str) -> DeckDiff:
    """Compute slide-level diff between two rendered ``.pptx`` decks.

    Parameters
    ----------
    pptx_a, pptx_b
        Paths to the two decks to compare.

    Returns
    -------
    DeckDiff
        Slide-level diff records.

    Raises
    ------
    FileNotFoundError
        If either path doesn't exist.
    """
    pptx_a = Path(pptx_a)
    pptx_b = Path(pptx_b)
    if not pptx_a.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_a}")
    if not pptx_b.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_b}")

    slides_a = _read_slides(pptx_a)
    slides_b = _read_slides(pptx_b)

    diff = DeckDiff(deck_a=pptx_a, deck_b=pptx_b)
    pairs, unmatched_a, unmatched_b = _match_slides(slides_a, slides_b)

    for idx_a, idx_b in pairs:
        sa = slides_a[idx_a]
        sb = slides_b[idx_b]
        field_changes = _slide_field_changes(sa, sb)
        kind = "unchanged" if not field_changes else "modified"
        diff.slides.append(
            SlideDiff(
                slide_index_a=idx_a,
                slide_index_b=idx_b,
                change_kind=kind,
                field_changes=field_changes,
            )
        )

    for idx_a in unmatched_a:
        diff.slides.append(
            SlideDiff(slide_index_a=idx_a, slide_index_b=None, change_kind="removed")
        )

    for idx_b in unmatched_b:
        diff.slides.append(
            SlideDiff(slide_index_a=None, slide_index_b=idx_b, change_kind="added")
        )

    # Stable sort: prefer reading by B's index when present, else A's.
    def _sort_key(s: SlideDiff) -> tuple[int, int]:
        if s.slide_index_b is not None:
            return (0, s.slide_index_b)
        if s.slide_index_a is not None:
            return (1, s.slide_index_a)
        return (2, 0)

    diff.slides.sort(key=_sort_key)
    return diff


# ---------------------------------------------------------------------------
# Slide reading / fingerprinting
# ---------------------------------------------------------------------------


def _read_slides(pptx_path: Path) -> list[dict[str, Any]]:
    """Return a list of slide-summary dicts.

    Each dict has:
        - ``index`` (int): zero-indexed deck position
        - ``title`` (str | None): topmost recovered text line
        - ``body`` (str): newline-joined non-title text
        - ``shapes`` (list[dict]): shape-level summaries for diffing
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover — gated at install time
        raise RuntimeError(
            "python-pptx is required for diff_decks. Install with "
            '`pip install -e ".[slides]"` or `pip install python-pptx`.'
        ) from exc

    prs = Presentation(str(pptx_path))
    out: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        shapes_info: list[dict[str, Any]] = []
        title: str | None = None
        best_top: int | None = None
        body_chunks: list[str] = []

        # First pass: gather shape summaries and pick the topmost text as title.
        for shape_idx, shape in enumerate(slide.shapes):
            shape_summary: dict[str, Any] = {
                "shape_idx": shape_idx,
                "shape_id": _shape_label(shape, shape_idx),
                "kind": _shape_kind(shape),
                "text": "",
                "image_hash": None,
            }
            if getattr(shape, "has_text_frame", False):
                try:
                    text = (shape.text_frame.text or "").strip()
                except Exception:  # pragma: no cover — defensive
                    text = ""
                shape_summary["text"] = text
                if text:
                    try:
                        top = int(shape.top)
                    except (AttributeError, TypeError, ValueError):
                        top = None
                    if top is not None and (best_top is None or top < best_top):
                        best_top = top
                        first_line = (
                            text.splitlines()[0].strip() if text.splitlines() else text
                        )
                        title = first_line
                    body_chunks.append(text)
            elif shape_summary["kind"] == "picture":
                shape_summary["image_hash"] = _picture_hash(shape)
            shapes_info.append(shape_summary)

        out.append(
            {
                "index": idx,
                "title": title,
                "body": "\n".join(body_chunks),
                "shapes": shapes_info,
            }
        )
    return out


def _shape_kind(shape: Any) -> str:
    """Return a coarse shape kind: 'text' / 'picture' / 'other'."""
    if getattr(shape, "has_text_frame", False):
        return "text"
    if getattr(shape, "shape_type", None) == 13:  # 13 = PICTURE
        return "picture"
    return "other"


def _shape_label(shape: Any, fallback_idx: int) -> str:
    """Return a human-readable stable-ish label for a shape."""
    try:
        name = shape.name
    except Exception:  # pragma: no cover — defensive
        name = None
    if name:
        return f"shape[{fallback_idx}]:{name}"
    return f"shape[{fallback_idx}]"


def _picture_hash(shape: Any) -> str | None:
    """Hash a picture shape's image bytes; ``None`` when blob is unreadable."""
    image = getattr(shape, "image", None)
    if image is None:
        return None
    try:
        blob = image.blob
    except Exception:  # pragma: no cover — defensive
        return None
    if not blob:
        return None
    return hashlib.md5(blob, usedforsecurity=False).hexdigest()


# ---------------------------------------------------------------------------
# Matching
# ---------------------------------------------------------------------------


_NORMALIZE_WS = re.compile(r"\s+")


def _normalize(text: str) -> str:
    return _NORMALIZE_WS.sub(" ", text.strip().lower())


def _title_key(slide: dict[str, Any]) -> str | None:
    title = slide.get("title")
    if not title:
        return None
    return _normalize(title)


def _body_fingerprint(slide: dict[str, Any]) -> str:
    body = slide.get("body") or ""
    return hashlib.md5(_normalize(body).encode("utf-8"), usedforsecurity=False).hexdigest()


def _match_slides(
    slides_a: list[dict[str, Any]],
    slides_b: list[dict[str, Any]],
) -> tuple[list[tuple[int, int]], list[int], list[int]]:
    """Pair slides A→B by title, then body fingerprint, then position.

    Returns ``(pairs, unmatched_a, unmatched_b)`` where ``pairs`` is a
    list of ``(idx_a, idx_b)`` tuples.
    """
    pairs: list[tuple[int, int]] = []
    used_a: set[int] = set()
    used_b: set[int] = set()

    # Pass 1 — title-key match (skip duplicate titles to avoid ambiguity).
    title_to_a: dict[str, list[int]] = {}
    for idx_a, slide in enumerate(slides_a):
        key = _title_key(slide)
        if key:
            title_to_a.setdefault(key, []).append(idx_a)
    title_to_b: dict[str, list[int]] = {}
    for idx_b, slide in enumerate(slides_b):
        key = _title_key(slide)
        if key:
            title_to_b.setdefault(key, []).append(idx_b)

    for key, a_indices in title_to_a.items():
        b_indices = title_to_b.get(key, [])
        # Only safe to match 1:1 — duplicates fall through to later passes.
        if len(a_indices) == 1 and len(b_indices) == 1:
            ia, ib = a_indices[0], b_indices[0]
            if ia not in used_a and ib not in used_b:
                pairs.append((ia, ib))
                used_a.add(ia)
                used_b.add(ib)

    # Pass 2 — body-fingerprint match for remaining slides.
    for idx_a, sa in enumerate(slides_a):
        if idx_a in used_a:
            continue
        fa = _body_fingerprint(sa)
        if not sa.get("body"):
            continue
        for idx_b, sb in enumerate(slides_b):
            if idx_b in used_b:
                continue
            if _body_fingerprint(sb) == fa:
                pairs.append((idx_a, idx_b))
                used_a.add(idx_a)
                used_b.add(idx_b)
                break

    # Pass 3 — positional match for the rest, only when both indices still
    # exist on each side (i.e. when the decks have a same-position slide
    # that wasn't matched by title or body).
    n_a, n_b = len(slides_a), len(slides_b)
    for pos in range(max(n_a, n_b)):
        if pos < n_a and pos < n_b and pos not in used_a and pos not in used_b:
            pairs.append((pos, pos))
            used_a.add(pos)
            used_b.add(pos)

    unmatched_a = [i for i in range(n_a) if i not in used_a]
    unmatched_b = [i for i in range(n_b) if i not in used_b]
    return pairs, unmatched_a, unmatched_b


# ---------------------------------------------------------------------------
# Field-level diff
# ---------------------------------------------------------------------------


def _slide_field_changes(
    slide_a: dict[str, Any], slide_b: dict[str, Any]
) -> list[tuple[str, str, str]]:
    """Return ``(shape_id, before, after)`` tuples for shape-level changes."""
    changes: list[tuple[str, str, str]] = []
    shapes_a = slide_a.get("shapes", [])
    shapes_b = slide_b.get("shapes", [])

    n = max(len(shapes_a), len(shapes_b))
    for i in range(n):
        sa = shapes_a[i] if i < len(shapes_a) else None
        sb = shapes_b[i] if i < len(shapes_b) else None
        if sa is None and sb is not None:
            changes.append((sb["shape_id"], "(missing)", _shape_blurb(sb)))
            continue
        if sb is None and sa is not None:
            changes.append((sa["shape_id"], _shape_blurb(sa), "(missing)"))
            continue
        if sa is None or sb is None:
            continue

        # Both present — compare the loaded fields.
        if sa["kind"] != sb["kind"]:
            changes.append(
                (sa["shape_id"], f"kind={sa['kind']}", f"kind={sb['kind']}")
            )
            continue
        if sa["kind"] == "text":
            ta = (sa.get("text") or "").strip()
            tb = (sb.get("text") or "").strip()
            if ta != tb:
                changes.append((sa["shape_id"], ta, tb))
        elif sa["kind"] == "picture":
            ha = sa.get("image_hash")
            hb = sb.get("image_hash")
            if ha != hb:
                changes.append(
                    (sa["shape_id"], f"image-md5={ha or 'none'}", f"image-md5={hb or 'none'}")
                )
    return changes


def _shape_blurb(shape: dict[str, Any]) -> str:
    """Compact summary of a shape (used when one deck is missing the shape)."""
    if shape["kind"] == "text":
        text = (shape.get("text") or "").strip()
        if not text:
            return "text=(empty)"
        head = text.splitlines()[0][:80]
        return f"text={head!r}"
    if shape["kind"] == "picture":
        h = shape.get("image_hash") or "none"
        return f"picture md5={h}"
    return f"{shape['kind']}"
