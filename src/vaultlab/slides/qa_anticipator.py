"""Q&A anticipator — surface likely audience questions for a deck.

Given a rendered ``.pptx``, return a ranked list of questions the
audience is likely to raise. Two operating modes:

* **LLM mode** — when a ``runner_callback`` is provided. The deck is
  serialized into a compact prompt and the callback is invoked once;
  its JSON-shaped return is parsed into
  :class:`AnticipatedQuestion` records. If the callback raises or
  returns unparseable text, the function silently falls back to
  heuristic mode so callers don't have to babysit reliability.

* **Heuristic mode** — keyword templates keyed off slide titles and
  body text. Cheap, deterministic, no network. Covers the four
  highest-value question categories from observing Hickey-lab journal
  clubs:

  1. **Statistical claims** — numeric / p-value / sample-size hints
     (``n=42``, ``p<0.05``, ``%``) → "How was X calculated?"
  2. **Comparisons** — ``vs`` / ``compared to`` / ``versus`` → "Why
     this comparison, not an alternative?"
  3. **Future work** — ``future``, ``next steps``, ``plan to`` → "What's
     the timeline?"
  4. **Limitations** — ``limit``, ``caveat``, ``sample size`` →
     "What's your plan to address X?"

The two modes are intentionally interchangeable. The heuristic floor
makes the function shippable in tests + CI; the LLM ceiling lets
callers upgrade quality when they want richer questions.

Public entrypoint: :func:`anticipate_qa`.
"""

from __future__ import annotations

import json
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

logger = logging.getLogger(__name__)

__all__ = [
    "AnticipatedQuestion",
    "anticipate_qa",
]


# ---------------------------------------------------------------------------
# Data class
# ---------------------------------------------------------------------------


@dataclass
class AnticipatedQuestion:
    """A single anticipated audience question.

    Attributes
    ----------
    question : str
        The full question text as it would be asked.
    anchor_slide_index : int
        Zero-indexed deck position of the slide the question is keyed off.
    why_likely : str
        Short justification — what triggered the question.
    confidence : float
        Heuristic confidence in [0, 1]. Heuristic-mode questions hover
        around 0.5-0.7; LLM-mode questions can carry whatever value the
        runner provides.
    """

    question: str
    anchor_slide_index: int
    why_likely: str
    confidence: float = 0.5


# ---------------------------------------------------------------------------
# Public entrypoint
# ---------------------------------------------------------------------------


def anticipate_qa(
    pptx_path: Path | str,
    *,
    runner_callback: Callable[[str], str] | None = None,
    n_questions: int = 10,
) -> list[AnticipatedQuestion]:
    """Return a ranked list of likely audience questions for the deck.

    Parameters
    ----------
    pptx_path
        Path to the rendered ``.pptx``.
    runner_callback
        Optional ``(prompt: str) -> str`` callable for LLM mode. When
        present, the deck's text is bundled into a prompt and the
        callback is invoked once. Failures fall back to heuristic mode.
    n_questions
        Cap on the number of returned questions.

    Returns
    -------
    list[AnticipatedQuestion]
        Up to ``n_questions`` likely questions, highest-confidence first.

    Raises
    ------
    FileNotFoundError
        If ``pptx_path`` doesn't exist.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_path}")

    slides = _read_slide_text(pptx_path)

    if runner_callback is not None:
        try:
            prompt = _build_llm_prompt(slides, n_questions=n_questions)
            raw = runner_callback(prompt)
            parsed = _parse_llm_response(raw, max_slide_index=len(slides) - 1)
            if parsed:
                return parsed[:n_questions]
            logger.info("qa_anticipator: LLM mode returned no parseable questions; falling back")
        except Exception:  # pragma: no cover — defensive: any runner failure falls back
            logger.exception("qa_anticipator: runner_callback failed; falling back to heuristics")

    return _heuristic_questions(slides, n_questions=n_questions)


# ---------------------------------------------------------------------------
# Slide reading
# ---------------------------------------------------------------------------


def _read_slide_text(pptx_path: Path) -> list[dict[str, Any]]:
    """Return ``[{'index', 'title', 'body', 'all_text'}, ...]`` for the deck."""
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover — gated at install time
        raise RuntimeError(
            "python-pptx is required for anticipate_qa. Install with "
            '`pip install -e ".[slides]"` or `pip install python-pptx`.'
        ) from exc

    prs = Presentation(str(pptx_path))
    out: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        chunks: list[str] = []
        title: str | None = None
        best_top: int | None = None
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                text = (shape.text_frame.text or "").strip()
            except Exception:  # pragma: no cover — defensive
                continue
            if not text:
                continue
            # Track topmost shape's first line as the title.
            try:
                top = int(shape.top)
            except (AttributeError, TypeError, ValueError):
                top = None
            if top is not None and (best_top is None or top < best_top):
                best_top = top
                first_line = text.splitlines()[0].strip() if text.splitlines() else text
                title = first_line
            chunks.append(text)

        all_text = "\n".join(chunks)
        body_lines = [
            line.strip()
            for line in all_text.splitlines()
            if line.strip() and line.strip() != (title or "")
        ]
        out.append(
            {
                "index": idx,
                "title": title,
                "body": "\n".join(body_lines),
                "all_text": all_text,
            }
        )
    return out


# ---------------------------------------------------------------------------
# Heuristic mode
# ---------------------------------------------------------------------------


_STATS_RE = re.compile(
    r"(p\s*[<>=]\s*\d|n\s*=\s*\d|\d+\s*%|±\s*\d|\bSE\b|\bSD\b|\bCI\b)",
    re.IGNORECASE,
)
_COMPARISON_RE = re.compile(
    r"(\b(vs|versus|compared\s+to|compared\s+with|relative\s+to)\b)",
    re.IGNORECASE,
)
_FUTURE_RE = re.compile(
    r"(\bfuture\s+work\b|\bnext\s+steps?\b|\bplan\s+to\b|\bwill\s+extend\b|"
    r"\bin\s+\d+\s*(weeks?|months?|years?))",
    re.IGNORECASE,
)
_LIMIT_RE = re.compile(
    r"(\blimitations?\b|\bcaveats?\b|\bsample\s+size\b|\blimited\b|"
    r"\breplicates?\b|\bsmall\s+(n|sample))",
    re.IGNORECASE,
)


def _heuristic_questions(
    slides: list[dict[str, Any]], *, n_questions: int
) -> list[AnticipatedQuestion]:
    """Generate heuristic Q&A questions from slide titles + bodies."""
    questions: list[AnticipatedQuestion] = []

    for slide in slides:
        idx = slide["index"]
        title = (slide["title"] or "").strip()
        body = slide["body"] or ""
        haystack = f"{title}\n{body}"

        # Stats trigger fires on numeric / p-value / sample-size hints.
        stats_match = _STATS_RE.search(haystack)
        if stats_match:
            token = stats_match.group(0)
            questions.append(
                AnticipatedQuestion(
                    question=(
                        f"How did you calculate {token!r}, and what was the "
                        "underlying statistical model?"
                    ),
                    anchor_slide_index=idx,
                    why_likely=(
                        f"Slide contains a statistical claim ({token!r}). Reviewers "
                        "typically probe the methodology behind reported numbers."
                    ),
                    confidence=0.7,
                )
            )

        # Comparisons trigger when 'vs' / 'compared to' appears.
        cmp_match = _COMPARISON_RE.search(haystack)
        if cmp_match:
            questions.append(
                AnticipatedQuestion(
                    question=(
                        "Why this comparison, and how would the result change "
                        "against an alternative baseline?"
                    ),
                    anchor_slide_index=idx,
                    why_likely=(
                        f"Slide stages a comparison ({cmp_match.group(0)!r}); "
                        "audiences often question the choice of baseline."
                    ),
                    confidence=0.6,
                )
            )

        # Future work → timeline question.
        future_match = _FUTURE_RE.search(haystack)
        if future_match:
            questions.append(
                AnticipatedQuestion(
                    question=(
                        "What's the timeline for the next steps, and what's the "
                        "first milestone?"
                    ),
                    anchor_slide_index=idx,
                    why_likely=(
                        "Slide flags future work — audiences want a concrete schedule "
                        "and the first deliverable."
                    ),
                    confidence=0.6,
                )
            )

        # Limitations → mitigation plan question.
        limit_match = _LIMIT_RE.search(haystack)
        if limit_match:
            questions.append(
                AnticipatedQuestion(
                    question=(
                        "What's your plan to address this limitation in the next "
                        "iteration of the work?"
                    ),
                    anchor_slide_index=idx,
                    why_likely=(
                        "Slide names a limitation; reviewers ask how it will be "
                        "mitigated or controlled for."
                    ),
                    confidence=0.65,
                )
            )

    # Stable sort by confidence (desc), then by slide index (asc) for determinism.
    questions.sort(key=lambda q: (-q.confidence, q.anchor_slide_index))
    return questions[:n_questions]


# ---------------------------------------------------------------------------
# LLM mode
# ---------------------------------------------------------------------------


_LLM_PROMPT_TEMPLATE = """\
You are reviewing a research presentation deck. Anticipate the {n} most
likely audience questions. Anchor each question to the slide that triggers
it.

Return STRICT JSON: an array of objects with keys
    "question" (str),
    "anchor_slide_index" (int, 0-indexed),
    "why_likely" (str),
    "confidence" (float in [0, 1]).

Deck contents (slide index → title → body):
{deck}
"""


def _build_llm_prompt(slides: list[dict[str, Any]], *, n_questions: int) -> str:
    """Serialize the deck into a compact LLM prompt."""
    lines: list[str] = []
    for slide in slides:
        title = slide["title"] or "(no title)"
        body = slide["body"] or "(no body)"
        # Keep the per-slide payload bounded so big decks don't explode the prompt.
        body_trimmed = body[:400]
        lines.append(f"[{slide['index']}] {title}\n    body: {body_trimmed}")
    return _LLM_PROMPT_TEMPLATE.format(n=n_questions, deck="\n".join(lines))


def _parse_llm_response(raw: str, *, max_slide_index: int) -> list[AnticipatedQuestion]:
    """Parse a JSON-shaped LLM response into :class:`AnticipatedQuestion` records.

    Tolerates an LLM that wraps the JSON in prose / fences by grabbing
    the first ``[ ... ]`` array via a non-greedy regex.
    """
    if not raw:
        return []

    candidate = raw.strip()
    # Strip ```json fences if present.
    fence = re.search(r"```(?:json)?\s*(.*?)```", candidate, re.DOTALL | re.IGNORECASE)
    if fence:
        candidate = fence.group(1).strip()

    # Find the first top-level JSON array, if the LLM padded around it.
    array_match = re.search(r"\[\s*(?:\{.*\})\s*(?:,\s*\{.*\}\s*)*\]", candidate, re.DOTALL)
    if array_match:
        candidate = array_match.group(0)

    try:
        data = json.loads(candidate)
    except (json.JSONDecodeError, ValueError):
        return []
    if not isinstance(data, list):
        return []

    out: list[AnticipatedQuestion] = []
    for item in data:
        if not isinstance(item, dict):
            continue
        question = str(item.get("question", "")).strip()
        if not question:
            continue
        try:
            anchor = int(item.get("anchor_slide_index", 0))
        except (TypeError, ValueError):
            continue
        if anchor < 0 or anchor > max_slide_index:
            # Don't trust LLM-fabricated indices outside the deck.
            anchor = max(0, min(anchor, max_slide_index))
        why = str(item.get("why_likely", "")).strip() or "(no rationale provided)"
        try:
            confidence = float(item.get("confidence", 0.5))
        except (TypeError, ValueError):
            confidence = 0.5
        confidence = max(0.0, min(1.0, confidence))
        out.append(
            AnticipatedQuestion(
                question=question,
                anchor_slide_index=anchor,
                why_likely=why,
                confidence=confidence,
            )
        )
    return out
